"""
FileKnowledgeBase — extracted from core/langgraph_agent.py Phase 1.
Includes module-level helpers, constants, and the FileKnowledgeBase class.
"""
from __future__ import annotations
import hashlib, io, os, re, sys, time, json, uuid, gc, math, struct, threading, subprocess
import operator, queue, warnings, contextlib, traceback, copy, itertools
from pathlib import Path
from datetime import datetime, timezone
from contextlib import contextmanager
from dataclasses import dataclass, field
from typing import TypedDict, Literal, List, Dict, Any, Optional, Callable, Annotated, Tuple, Iterator, Union, Sequence

import chromadb
import requests

from database.personal_info_db import PersonalInfoDB
from config import settings
from config.prompts import (
    CLASSIFY_PROMPT, SUMMARY_PROMPT, IMAGE_SUMMARY_PROMPT, IMAGE_OCR_PROMPT,
    get_prompt, normalize_prompt_language,
)
from core.retrieval import (
    KeywordIndexManager,
    KeywordIndexRecord,
    PathScopeMatcher,
    build_cjk_latin_aliases,
    build_lookup_blob,
    compute_lookup_overlap_score,
    ensure_path_scope_matcher,
)
from utils.logger import get_logger
from utils.pdf_utils import HAS_PDF_TEXT, HAS_PYPDFIUM2, extract_pdf_text, render_pdf_pages_to_png
logger = get_logger()

OCR_MAX_TOKENS_DEFAULT = int(os.getenv("OCR_MAX_TOKENS", "2200"))


AUDIO_EXTENSIONS = {'.mp3', '.wav', '.flac', '.aac', '.ogg', '.m4a', '.wma', '.aiff', '.ape'}

VIDEO_EXTENSIONS = {'.mp4', '.m4v', '.avi', '.mov', '.mkv', '.webm', '.flv', '.wmv', '.ts'}

IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.svg', '.ico', '.tiff', '.heic'}

# Canonical category storage uses English keys for indexing/search.
_CATEGORY_ALIASES_TO_EN: Dict[str, str] = {
    # resume
    "简历": "resume", "履历": "resume", "个人简历": "resume",
    "cv": "resume", "resume": "resume", "resumes": "resume",
    # report
    "报告": "report", "分析报告": "report", "研究报告": "report", "报告书": "report",
    "report": "report", "reports": "report",
    # contract
    "合同": "contract", "协议": "contract", "合约": "contract", "协议书": "contract",
    "contract": "contract", "contracts": "contract", "agreement": "contract", "agreements": "contract",
    # note
    "笔记": "note", "备忘": "note", "备忘录": "note",
    "note": "note", "notes": "note", "memo": "note", "memos": "note",
    # manual
    "手册": "manual", "教程": "manual", "指南": "manual", "说明书": "manual", "操作手册": "manual",
    "manual": "manual", "manuals": "manual", "guide": "manual", "guides": "manual", "tutorial": "manual", "tutorials": "manual",
    # paper
    "论文": "paper", "学术论文": "paper", "研究论文": "paper", "期刊": "paper",
    "paper": "paper", "papers": "paper", "article": "paper", "articles": "paper", "journal": "paper",
    # presentation
    "演示": "presentation", "演示文稿": "presentation", "幻灯片": "presentation", "课件": "presentation", "ppt": "presentation",
    "presentation": "presentation", "presentations": "presentation", "slides": "presentation", "deck": "presentation",
    # data
    "数据": "data", "数据集": "data", "表格": "data", "统计表": "data", "工作表": "data", "电子表格": "data", "数据表": "data",
    "dataset": "data", "datasets": "data", "data": "data", "table": "data", "tables": "data",
    "worksheet": "data", "worksheets": "data", "workbook": "data", "workbooks": "data",
    "spreadsheet": "data", "spreadsheets": "data",
    "excel": "data", "xlsx": "data", "xls": "data", "csv": "data", "tsv": "data",
    "numbers": "data",  # Apple Numbers app
    # email
    "邮件": "email", "电子邮件": "email", "eml": "email",
    "email": "email", "emails": "email", "mail": "email", "mails": "email",
    # image
    "照片": "image", "图片": "image", "图像": "image", "相片": "image", "截图": "image",
    "image": "image", "images": "image", "photo": "image", "photos": "image", "picture": "image", "pictures": "image",
    # media
    "音视频": "audio/video", "视频/音频": "audio/video", "影音": "audio/video",
    "audio/video": "audio/video", "audio video": "audio/video", "media": "audio/video",
    "视频": "video", "录像": "video", "影片": "video", "短片": "video",
    "video": "video", "videos": "video", "movie": "video", "movies": "video", "clip": "video", "clips": "video",
    "mp4": "video", "mov": "video", "mkv": "video", "avi": "video", "webm": "video", "m4v": "video",
    "音频": "audio", "录音": "audio", "声音": "audio", "音乐": "audio", "歌曲": "audio",
    "audio": "audio", "audios": "audio", "podcast": "audio", "podcasts": "audio",
    "recording": "audio", "recordings": "audio", "song": "audio", "songs": "audio", "music": "audio",
    "mp3": "audio", "wav": "audio", "m4a": "audio", "flac": "audio", "aac": "audio", "ogg": "audio",
    # book
    "书籍": "book", "图书": "book", "电子书": "book",
    "book": "book", "books": "book", "ebook": "book", "ebooks": "book",
    # code
    "代码": "code", "源码": "code", "脚本": "code", "程序": "code",
    "code": "code", "codes": "code", "source code": "code", "script": "code", "scripts": "code",
    # invoice
    "账单": "invoice", "发票": "invoice", "收据": "invoice", "付款单": "invoice",
    "invoice": "invoice", "invoices": "invoice", "bill": "invoice", "bills": "invoice", "receipt": "invoice", "receipts": "invoice",
    # quotation
    "报价": "quotation", "报价单": "quotation", "报价表": "quotation", "报价书": "quotation", "报价方案": "quotation",
    "quote": "quotation", "quotes": "quotation", "quotation": "quotation", "quotations": "quotation", "proposal": "quotation", "proposals": "quotation",
    # document
    "文档": "document", "资料": "document", "文稿": "document",
    "document": "document", "documents": "document", "doc": "document", "docs": "document",
    "all": "all", "全部": "all", "所有": "all", "文件": "all", "文档全部": "all",
    "其他": "other", "未知": "other", "other": "other", "unknown": "other",
}

_CATEGORY_CANONICAL_TO_ZH: Dict[str, str] = {
    "resume": "简历",
    "report": "报告",
    "contract": "合同",
    "note": "笔记",
    "manual": "手册",
    "paper": "论文",
    "presentation": "演示",
    "data": "数据",
    "email": "邮件",
    "image": "图片",
    "audio": "音频",
    "video": "视频",
    "audio/video": "音视频",
    "book": "书籍",
    "code": "代码",
    "invoice": "发票",
    "quotation": "报价单",
    "document": "文档",
    "all": "全部",
    "other": "其他",
}

_CATEGORY_KEYWORD_HINTS: List[Tuple[str, Tuple[str, ...]]] = [
    ("resume", ("简历", "履历", "resume", "cv")),
    ("report", ("报告", "分析报告", "report")),
    ("contract", ("合同", "协议", "合约", "agreement", "contract")),
    ("note", ("笔记", "备忘", "memo", "note")),
    ("manual", ("手册", "指南", "教程", "说明书", "manual", "guide", "tutorial")),
    ("paper", ("论文", "期刊", "paper", "article", "journal")),
    ("presentation", ("演示", "演示文稿", "幻灯", "ppt", "slides", "presentation", "deck")),
    ("data", ("数据", "表格", "工作表", "电子表格", "dataset", "data", "table",
              "worksheet", "workbook", "spreadsheet", "excel", "csv", "tsv", "numbers")),
    ("email", ("邮件", "电子邮件", "email", "mail")),
    ("image", ("图片", "照片", "截图", "image", "photo", "picture")),
    ("audio/video", ("音视频", "影音", "media")),
    ("video", ("视频", "录像", "影片", "video", "videos", "movie", "clip", "mp4", "mov", "mkv", "avi", "webm", "m4v")),
    ("audio", ("音频", "录音", "声音", "音乐", "歌曲", "audio", "audios", "recording", "song", "music", "mp3", "wav", "m4a", "flac", "aac", "ogg")),
    ("book", ("书籍", "图书", "电子书", "book", "ebook")),
    ("code", ("代码", "源码", "脚本", "code", "script", "source")),
    ("invoice", ("发票", "账单", "收据", "invoice", "bill", "receipt")),
    ("quotation", ("报价", "报价单", "报价书", "quotation", "quote", "proposal")),
    ("document", ("文档", "资料", "文件", "document", "doc")),
    ("all", ("全部", "所有", "all")),
    ("other", ("其他", "未知", "other", "unknown")),
]


def _normalize_category_en(raw: str, default: str = "other") -> str:
    v = str(raw or "").strip().lower()
    if not v:
        return default
    direct = _CATEGORY_ALIASES_TO_EN.get(v)
    if direct:
        return direct
    for canonical, hints in _CATEGORY_KEYWORD_HINTS:
        for h in hints:
            if h.isascii() and len(h) <= 4:
                if re.search(r'(?<![a-z])' + re.escape(h) + r'(?![a-z])', v):
                    return canonical
            else:
                if h in v:
                    return canonical
    return v


def _localize_category_label(category_en: str, language: str) -> str:
    lang = str(language or "en").strip().lower()
    if lang.startswith("zh"):
        return _CATEGORY_CANONICAL_TO_ZH.get(category_en, category_en)
    return category_en


class _ThreadSafeCollectionProxy:
    """Serialize Chroma collection I/O to avoid native rust backend read/write races."""

    _LOCKED_METHODS = {
        "count",
        "get",
        "add",
        "delete",
        "query",
        "upsert",
        "update",
        "peek",
        "modify",
    }

    def __init__(self, collection: Any, io_lock: threading.RLock):
        self._collection = collection
        self._io_lock = io_lock

    def __getattr__(self, name: str) -> Any:
        attr = getattr(self._collection, name)
        if name in self._LOCKED_METHODS and callable(attr):
            def _wrapped(*args, **kwargs):
                with self._io_lock:
                    return attr(*args, **kwargs)
            return _wrapped
        return attr


class _QueryCachePrewarmCancelled(Exception):
    """Internal control-flow exception for cooperative prewarm cancellation."""


# ── English Category Inference ────────────────────────────────────────────────
# Rule-based safety net: when LLM intent pass doesn't emit a category for
# English queries, these trigger-word lists map query content to the correct
# file-type filter.  Ordered from most-specific to least-specific.
_EN_CATEGORY_TRIGGERS: List[Tuple[str, List[str]]] = [
    ("video", [
        "video", "videos", "movie", "movies", "film", "clip", "clips",
        "footage", "reel", "mp4", "mov", "mkv", "avi", "webm", "m4v",
    ]),
    ("audio", [
        "audio", "sound", "recording", "music", "song", "track", "mp3",
        "wav", "flac", "aac", "m4a", "ogg", "wma", "aiff",
        "podcast", "voice memo", "noise", "melody", "beat", "rhythm",
    ]),
    ("image", [
        "image", "images", "photo", "photos", "picture", "pictures",
        "screenshot", "screenshots", "diagram", "chart",
        "graphic", "png", "jpg", "jpeg", "gif", "bmp", "tiff", "heic",
        "illustration", "figure", "drawing", "sketch", "mockup", "thumbnail",
        "wiring diagram", "schematic", "blueprint",
    ]),
    # data/spreadsheet -- must come BEFORE document to avoid spreadsheet->document conflict
    ("data", [
        "worksheet", "worksheets", "workbook", "workbooks",
        "spreadsheet", "spreadsheets", "excel", "xls", "xlsx",
        "csv", "tsv", "numbers file", "numbers document",
    ]),
    # Specific doc subtypes -- BEFORE generic document so they get specific category not 'document'
    ("resume", [
        "resume", "resumes", " cv ", "curriculum vitae", "my cv",
    ]),
    ("invoice", [
        "invoice", "invoices", "receipt", "receipts",
    ]),
    ("contract", [
        "contract", "contracts", "agreement", "agreements",
    ]),
    ("presentation", [
        "presentation", "presentations", "slides", "slide deck",
    ]),
    # generic 'document' for remaining doc-like signals
    ("document", [
        "bill", "bills", "thesis", "dissertation",
        "strategy", "gtm", "roadmap", "plan", "plans",
        "research", "findings", "insights",
        "analysis", "report", "reports", "paper", "papers",
    ]),
    # 'manual' only for explicit product/technical manuals
    ("manual", [
        "product manual", "user guide", "data sheet", "datasheet",
        "installation guide", "technical specification", "spec sheet",
    ]),
]

# Words that look like triggers but should NOT infer category (anti-patterns)
_EN_CATEGORY_ANTITRIGGERS: set = {
    "manager", "management", "managed", "managing",  # "manager" sounds like "manual" to LLMs
}


def _infer_category_from_english_query(query: str) -> str:
    """
    Infer a file-type category from English query signal words.

    Returns the canonical category string (e.g. "video", "audio", "image",
    "document") or an empty string if no signal is detected.

    Design notes:
    - Pure Python, zero LLM calls – executes in microseconds.
    - Used as a post-LLM safety net in _normalize_intent_to_internal_en.
    - Does NOT override a category already set by the LLM.
    - Deliberately conservative: only fires on clear, unambiguous signals.
    • Anti-trigger guard: words like 'manager' must not trigger 'manual'.
    """
    if not query:
        return ""
    q_lower = query.lower()
    q = f" {q_lower} "          # pad so " cv " matches "my cv"
    for category, triggers in _EN_CATEGORY_TRIGGERS:
        for trigger in triggers:
            # Use word-boundary-aware matching via space padding
            t = trigger if " " in trigger else f" {trigger} "
            if t in q:
                # Anti-trigger guard: don't fire if an anti-pattern overwrites
                # e.g. "manual" trigger blocked when "manager" present
                blocked = any(f" {anti} " in q for anti in _EN_CATEGORY_ANTITRIGGERS)
                if blocked and category == "manual":
                    continue
                return category
    return ""
# ─────────────────────────────────────────────────────────────────────────────


# ── BM25 Hybrid Retrieval Utilities ──────────────────────────────────────────
# These functions implement a lightweight in-memory BM25 scoring layer and
# Reciprocal Rank Fusion (RRF) to blend lexical and semantic retrieval signals.

try:
    import jieba as _jieba
    # Suppress jieba's noisy initialization logging
    import logging as _logging
    _logging.getLogger("jieba").setLevel(_logging.WARNING)
    _HAS_JIEBA = True
except ImportError:
    _HAS_JIEBA = False

try:
    from pypinyin import Style as _PinyinStyle, lazy_pinyin as _lazy_pinyin
    _HAS_PYPINYIN = True
except ImportError:
    _HAS_PYPINYIN = False
    _PinyinStyle = None
    _lazy_pinyin = None

try:
    from rank_bm25 import BM25Okapi as _BM25Okapi
    _HAS_BM25 = True
except ImportError:
    _HAS_BM25 = False

import re as _re_bm25_mod

# Bilingual stopwords for BM25 tokenization
_BM25_STOPWORDS = {
    # English
    "a", "an", "the", "is", "are", "was", "were", "be", "been", "being",
    "have", "has", "had", "do", "does", "did", "will", "would", "could",
    "should", "may", "might", "shall", "must", "can", "need", "ought",
    "of", "in", "on", "at", "to", "for", "with", "by", "from", "up",
    "down", "or", "and", "but", "not", "no", "so", "if", "as", "it",
    "its", "this", "that", "these", "those", "he", "she", "we", "they",
    "i", "you", "me", "him", "her", "us", "them", "my", "your", "his",
    "our", "their", "what", "which", "who", "whom", "where", "when",
    "why", "how", "all", "any", "both", "each", "few", "more", "most",
    "some", "such", "than", "then", "there", "here", "about", "into",
    "find", "show", "get", "give", "tell", "search", "look", "want",
    "need", "help", "make", "let", "see", "go", "come", "take", "use",
    "please", "just", "very", "really", "also",
    # Chinese
    "的", "了", "吗", "呢", "吧", "啊", "哦", "请", "帮", "找", "给",
    "看", "想", "需要", "希望", "可以", "我", "你", "他", "她", "它",
    "我们", "你们", "他们", "这", "那", "这个", "那个", "哪个",
    "什么", "怎么", "如何", "哪里", "哪些", "有没有", "是否",
    "文件", "文档", "资料", "内容", "信息", "搜索", "寻找", "找出",
    "相关", "有关", "关于", "或者", "以及", "并且", "帮我", "一下",
    "看看", "看一下", "有",
}


def _bm25_english_variants(token: str) -> List[str]:
    """Return conservative English word-form variants for lexical matching."""
    t = str(token or "").strip().lower()
    if not (t.isascii() and t.isalpha()):
        return []
    if len(t) <= 3:
        return []

    variants: List[str] = []

    def _add(value: str) -> None:
        if len(value) <= 1:
            return
        if value == t or value in _BM25_STOPWORDS:
            return
        if value not in variants:
            variants.append(value)

    if len(t) > 4 and t.endswith("ies"):
        _add(t[:-3] + "y")
    if len(t) > 4 and t.endswith("ves"):
        _add(t[:-3] + "f")
        _add(t[:-3] + "fe")
    if len(t) > 4 and t.endswith("men"):
        _add(t[:-3] + "man")
    if len(t) > 4 and t.endswith("es") and (
        t[:-2].endswith(("s", "x", "z")) or t[:-2].endswith(("ch", "sh"))
    ):
        _add(t[:-2])
    if t.endswith("s") and not t.endswith(("ss", "us", "is")):
        _add(t[:-1])

    return variants


def _bm25_pinyin_variants(token: str) -> List[str]:
    """Return conservative pinyin forms for CJK tokens.

    Index summaries are often English-first and may contain transliterated
    Chinese entity names. Adding pinyin variants lets Chinese queries bridge to
    those summaries without encoding any dataset-specific names.
    """
    t = str(token or "").strip().lower()
    if not t or not _HAS_PYPINYIN or _lazy_pinyin is None:
        return []
    cjk_chars = [ch for ch in t if "\u4e00" <= ch <= "\u9fff"]
    if len(cjk_chars) < 2:
        return []
    if t in _BM25_STOPWORDS or _CATEGORY_ALIASES_TO_EN.get(t):
        return []

    try:
        syllables = [
            str(item or "").strip().lower()
            for item in _lazy_pinyin(t, style=_PinyinStyle.NORMAL, errors="ignore")
            if str(item or "").strip()
        ]
    except Exception:
        return []
    if not syllables:
        return []

    variants: List[str] = []

    def _add(value: str) -> None:
        value = re.sub(r"[^a-z0-9]+", "", str(value or "").lower())
        if len(value) < 2 or value == t or value in _BM25_STOPWORDS:
            return
        if value not in variants:
            variants.append(value)

    joined = "".join(syllables)
    _add(joined)
    for idx in range(len(syllables) - 1):
        _add(syllables[idx] + syllables[idx + 1])

    return variants


def _tokenize_for_bm25(text: str) -> List[str]:
    if not text:
        return []
    text = str(text).strip()

    # Phase 1: Split into raw tokens
    raw_tokens: List[str] = []
    if _HAS_JIEBA:
        raw_tokens = list(_jieba.cut_for_search(text))
    else:
        # Fallback: split on non-alphanum boundaries + CJK char-by-char bigrams
        raw_tokens = _re_bm25_mod.findall(r'[a-zA-Z0-9]+|[\u4e00-\u9fff]+', text)

    # Phase 2: Normalize, filter stopwords, and collect
    tokens: List[str] = []
    for t in raw_tokens:
        t = t.strip().lower()
        if len(t) < 2 and not ('\u4e00' <= t <= '\u9fff'):
            continue
        if t in _BM25_STOPWORDS:
            continue
        tokens.append(t)

    # English file/content search commonly mixes singular filenames or tags with
    # plural user phrasing ("dogs", "photos"). Expand query/index tokens
    # conservatively so BM25 can bridge that gap without stemming CJK terms.
    morph_tokens: List[str] = []
    seen_token_forms = set(tokens)
    for t in tokens:
        for variant in _bm25_english_variants(t):
            if variant not in seen_token_forms:
                morph_tokens.append(variant)
                seen_token_forms.add(variant)
    tokens.extend(morph_tokens)

    # Bridge Chinese user queries to English-first index summaries that contain
    # pinyin entity names, e.g. organization names extracted from invoices.
    pinyin_tokens: List[str] = []
    seen_pinyin_forms = set(tokens)
    for t in tokens:
        for variant in _bm25_pinyin_variants(t):
            if variant not in seen_pinyin_forms:
                pinyin_tokens.append(variant)
                seen_pinyin_forms.add(variant)
    tokens.extend(pinyin_tokens)

    # Phase 3: Category translation expansion
    # If a token maps to a known category, also inject the cross-language equivalent.
    expansion_tokens: List[str] = []
    seen_expansions = set(tokens)
    for t in tokens:
        canon = _CATEGORY_ALIASES_TO_EN.get(t)
        if canon:
            # English → inject Chinese
            zh = _CATEGORY_CANONICAL_TO_ZH.get(canon, "")
            if zh and zh not in seen_expansions:
                expansion_tokens.append(zh)
                seen_expansions.add(zh)
            # Chinese → inject English
            if canon not in seen_expansions:
                expansion_tokens.append(canon)
                seen_expansions.add(canon)
    tokens.extend(expansion_tokens)

    return tokens


_QUERY_SURFACE_STOPWORDS = {
    "find", "search", "show", "list", "display", "look", "locate", "open",
    "get", "give", "tell", "browse", "scan", "help", "please", "me", "my",
    "our", "your", "the", "a", "an", "this", "that", "these", "those", "of",
    "in", "inside", "under", "within", "from", "for", "with", "about", "to",
    "on", "at", "by", "all", "any", "some", "related", "relevant",
    "file", "files", "document", "documents", "doc", "docs",
    "folder", "folders", "directory", "directories", "dir",
    "image", "images", "photo", "photos", "picture", "pictures",
    "video", "videos", "audio", "recording", "recordings",
}


def _build_bilingual_query_surfaces(*texts: str, max_surfaces: int = 6) -> List[str]:
    """Return compact English and CJK query surfaces for multi-route retrieval."""
    surfaces: List[str] = []
    seen: set[str] = set()

    def _add(value: str) -> None:
        cleaned = re.sub(r"\s+", " ", str(value or "").strip())
        if len(cleaned) < 2:
            return
        key = cleaned.casefold()
        if key in seen:
            return
        seen.add(key)
        surfaces.append(cleaned)

    raw_surface_added = False
    for raw in texts:
        text = str(raw or "").strip()
        if not text:
            continue
        if not raw_surface_added:
            _add(text)
            raw_surface_added = True

        cjk_terms: List[str] = []
        cjk_seen: set[str] = set()
        for seg in re.findall(r"[\u4e00-\u9fff]{2,}", text):
            if seg in cjk_seen:
                continue
            cjk_seen.add(seg)
            cjk_terms.append(seg)
            if len(cjk_terms) >= 8:
                break
        if cjk_terms:
            _add(" ".join(cjk_terms))

        latin_terms: List[str] = []
        latin_seen: set[str] = set()
        for token in re.findall(r"[A-Za-z0-9][A-Za-z0-9._-]*", text):
            lowered = token.casefold().strip("._-")
            if len(lowered) < 2 or lowered.isdigit() or lowered in _QUERY_SURFACE_STOPWORDS:
                continue
            if lowered in latin_seen:
                continue
            latin_seen.add(lowered)
            latin_terms.append(token)
            if len(latin_terms) >= 8:
                break
        if latin_terms:
            _add(" ".join(latin_terms))

    return surfaces[: max(1, int(max_surfaces or 1))]


def _build_bm25_and_score(
    all_metas: List[Dict[str, Any]],
    all_ids: List[str],
    query_tokens: List[str],
    allowed_paths: Optional[List[str]] = None,
    category_filter: str = "",
) -> List[Tuple[str, str, float]]:
    """Build an in-memory BM25 index from ChromaDB metadata and score against query.

    Each unique file path becomes one "document" in BM25, composed of:
      file_name + doc_summary + en_tags

    Returns: list of (chroma_id, file_path, bm25_score), sorted by score descending.
    Only returns entries with score > 0.
    """
    if not _HAS_BM25 or not query_tokens:
        return []

    # Deduplicate by file_path (take first chunk's metadata per file)
    seen_paths: Dict[str, Tuple[str, Dict[str, Any]]] = {}
    for mid, meta in zip(all_ids, all_metas):
        fp = meta.get("file_path", "") or ""
        if not fp or fp in seen_paths:
            continue

        # Path allowlist check
        if allowed_paths is not None:
            matched = False
            for ap in allowed_paths:
                if fp.startswith(ap) or os.path.dirname(fp).startswith(ap):
                    matched = True
                    break
            if not matched:
                continue

        # Category filter
        if category_filter:
            doc_cat = _normalize_category_en(meta.get("doc_category", "other"))
            if doc_cat != category_filter:
                continue

        seen_paths[fp] = (mid, meta)

    if not seen_paths:
        return []

    # Build corpus
    corpus_fps: List[str] = []
    corpus_ids: List[str] = []
    tokenized_corpus: List[List[str]] = []

    for fp, (mid, meta) in seen_paths.items():
        fname = meta.get("file_name", "") or ""
        summary = meta.get("doc_summary", "") or ""
        en_tags = meta.get("en_tags", "") or ""
        lookup_aliases = meta.get("lookup_aliases", "") or ""
        table_schema_hint = meta.get("table_schema_hint", "") or ""
        file_name_en = meta.get("file_name_en", "") or ""
        folder_name_en = meta.get("folder_name_en", "") or ""
        # Fix 6 (P2): include path segment words so queries like "Q4 marketing brief"
        # can match files buried deep in folder hierarchies (e.g. /work/Q4_marketing_brief.pdf).
        import re as _bm25_path_re
        _path_seg_tokens = " ".join(_bm25_path_re.findall(r'[A-Za-z0-9一-鿿]+', fp or ""))
        # Compose a "document" from all searchable metadata fields
        # file_name_en is the LLM-translated English name for non-English files
        folder_base = os.path.basename(os.path.dirname(fp or "")) if fp else ""
        lookup_blob = build_lookup_blob(
            fname,
            file_name_en,
            folder_name_en,
            folder_base,
            fp,
        )
        doc_text = f"{fname} {file_name_en} {folder_name_en} {summary} {table_schema_hint} {en_tags} {lookup_aliases} {_path_seg_tokens} {lookup_blob}"
        doc_tokens = _tokenize_for_bm25(doc_text)
        if not doc_tokens:
            doc_tokens = [""]  # BM25Okapi needs non-empty lists
        tokenized_corpus.append(doc_tokens)
        corpus_fps.append(fp)
        corpus_ids.append(mid)

    if not tokenized_corpus:
        return []

    # Build BM25 and score
    bm25 = _BM25Okapi(tokenized_corpus)
    scores = bm25.get_scores(query_tokens)

    results: List[Tuple[str, str, float]] = []
    for i, score in enumerate(scores):
        if score > 0:
            results.append((corpus_ids[i], corpus_fps[i], float(score)))

    results.sort(key=lambda x: x[2], reverse=True)
    return results


def _rrf_fuse(
    vector_ranked: List[Tuple[str, float]],
    bm25_ranked: List[Tuple[str, float]],
    k: int = 60,
) -> List[Tuple[str, float]]:
    """Reciprocal Rank Fusion (RRF) of two ranked lists.

    Each input is a list of (file_path, score) sorted by relevance (best first).
    Returns fused (file_path, rrf_score) sorted descending.

    Formula: rrf_score(d) = Σ  1 / (k + rank_i(d))
    where k=60 is the standard RRF constant.
    """
    fused: Dict[str, float] = {}

    for rank, (fp, _score) in enumerate(vector_ranked, start=1):
        fused[fp] = fused.get(fp, 0.0) + 1.0 / (k + rank)

    for rank, (fp, _score) in enumerate(bm25_ranked, start=1):
        fused[fp] = fused.get(fp, 0.0) + 1.0 / (k + rank)

    result = sorted(fused.items(), key=lambda x: x[1], reverse=True)
    return result
# ─────────────────────────────────────────────────────────────────────────────


def _env_truthy(name: str) -> bool:
    v = (os.getenv(name) or "").strip().lower()
    return v in {"1", "true", "yes", "on", "y"}


def _env_int_first(names: List[str], default: int) -> int:
    for name in names:
        raw = (os.getenv(name) or "").strip()
        if not raw:
            continue
        try:
            return int(raw)
        except Exception:
            continue
    return default


def _indexing_cooperative_yield() -> None:
    try:
        sec = float(os.getenv("FILEAGENT_INDEX_YIELD_SEC", "0.02") or "0.02")
    except Exception:
        sec = 0.02
    if sec > 0:
        time.sleep(min(float(sec), 0.5))

BCE_QUERY_PREFIX = "Represent this sentence for searching relevant passages: "
BCE_DOC_PREFIX = "Represent this document for retrieval: "


from tools import get_all_tools, get_tool, IntentRegistry
from tools.document_tools import (
    count_documents, 
    count_documents_files,
    search_documents,
    summarize_topics,
    set_kb_instance,
    get_kb_instance,
)

# LlamaIndex is kept as an opt-in fallback for rare formats, but it must never
# override text extracted by the built-in parsers. Enable only when explicitly
# needed with FILEAGENT_ENABLE_LLAMA_INDEX_FALLBACK=1.
ENABLE_LLAMA_INDEX_FALLBACK = _env_truthy("FILEAGENT_ENABLE_LLAMA_INDEX_FALLBACK")
SimpleDirectoryReader = None
SentenceSplitter = None
HAS_LLAMA_INDEX = False
if ENABLE_LLAMA_INDEX_FALLBACK:
    try:
        from llama_index.core import SimpleDirectoryReader
        from llama_index.core.node_parser import SentenceSplitter
        HAS_LLAMA_INDEX = True
        logger.info("llama_index fallback enabled for formats not handled by built-in parsers")
    except ImportError:
        logger.info("llama_index fallback enabled but not installed; using built-in parsers only")

class AgentState(TypedDict):
    question: str
    
    query_type: Literal["stats", "search", "chat"]
    category_filter: Optional[str]
    keyword: Optional[str]
    
    search_results: List[Dict[str, Any]]
    reranked_results: List[Dict[str, Any]]
    db_stats: Dict[str, Any]
    
    final_answer: str
    source_files: List[Dict[str, Any]]



def _chunk_with_no_singleton_tail(
    items: List[Any],
    indices: Optional[List[Any]] = None,
    batch_size: int = 10,
) -> Union[List[List[Any]], Tuple[List[List[Any]], List[List[Any]]]]:
    """
    Chunk a list while avoiding a trailing singleton batch when possible.

    Examples:
    - 21 items @ 10 -> [10, 11]
    - 11 items @ 10 -> [11]
    - 31 items @ 10 -> [10, 10, 11]
    """
    batch_size = max(1, int(batch_size or 1))
    batches = [items[i:i + batch_size] for i in range(0, len(items), batch_size)]
    index_batches: Optional[List[List[Any]]] = None
    if indices is not None:
        index_batches = [indices[i:i + batch_size] for i in range(0, len(indices), batch_size)]

    if len(batches) >= 2 and len(batches[-1]) == 1:
        batches[-2].extend(batches[-1])
        batches.pop()
        if index_batches is not None:
            index_batches[-2].extend(index_batches[-1])
            index_batches.pop()

    if index_batches is None:
        return batches
    return batches, index_batches

class FileKnowledgeBase:
    
    def __init__(self, db_path: str = settings.DB_PATH):
        self.db_path = db_path
        self.reranker = None
        self._embed_ctx_local = threading.local()
        self._embedding_error_log_path: Optional[str] = None
        self._index_details_log_path: Optional[str] = None
        self._collection_io_lock = threading.RLock()
        self._reranker_lock = threading.RLock()
        self._write_heavy_mode_lock = threading.Lock()
        self._write_heavy_mode_depth = 0
        self._last_collection_write_ts: float = 0.0
        self._local_llm_available = None
        self._startup_index_prefill_observation_logged = False
        self._last_persist_ts = 0.0
        self._init_db()
        self.personal_info_db = PersonalInfoDB(self.db_path)
        self._init_embedding()
        self._init_reranker()
        self._ignore_cache_path = os.path.join(
            os.path.dirname(os.path.abspath(self.db_path)),
            "index_ignore_cache.json",
        )
        self._ignore_cache_entries: Dict[str, Dict[str, Any]] = {}
        self._ignore_cache_dirty = False
        self._ignore_cache_new_since_save = 0
        self._ignore_cache_last_save_time = time.time()
        self._load_index_ignore_cache()
        self._failed_file_cache_path = os.path.join(
            os.path.dirname(os.path.abspath(self.db_path)),
            "index_failed_files.json",
        )
        self._failed_file_cache_entries: Dict[str, Dict[str, Any]] = {}
        self._failed_file_cache_dirty = False
        self._failed_file_cache_new_since_save = 0
        self._failed_file_cache_last_save_time = time.time()
        self._load_failed_file_cache()
        self._session_ignored_paths = set()
        self._indexed_paths_log_last_sig = ""
        self._indexed_paths_log_last_ts = 0.0
        self._folder_index_lock = threading.Lock()
        self._folder_index_dirty = False
        self._folder_index_auto_rebuild_attempted = False
        self._folder_index_pruned_once = False

        self._meta_cache_lock = threading.Lock()
        self._meta_cache_cond = threading.Condition(self._meta_cache_lock)
        self._meta_cache_version: int = 0
        self._meta_cache_db_count: int = -1
        self._meta_cache_ts: float = 0.0
        self._meta_cache_data: Optional[Dict[str, Any]] = None      # {ids, metadatas, documents}
        self._meta_cache_ttl: float = float(os.getenv("FILEAGENT_META_CACHE_TTL", "120"))  # seconds
        self._meta_cache_building = False
        self._keyword_backend_name = str(
            os.getenv("FILEAGENT_KEYWORD_BACKEND", "bm25") or "bm25"
        ).strip().lower()
        self._keyword_local_scope_max_docs = max(
            1,
            int(os.getenv("FILEAGENT_KEYWORD_LOCAL_SCOPE_MAX_DOCS", "400") or 400),
        )
        self._keyword_rebuild_delay_sec = max(
            0.0,
            float(os.getenv("FILEAGENT_KEYWORD_REBUILD_DELAY_SEC", "2.0") or 2.0),
        )
        self._keyword_index_manager = KeywordIndexManager(
            backend_name=self._keyword_backend_name,
            sidecar_path=self._resolve_keyword_index_sidecar_path(),
            build_records_fn=self._build_keyword_index_records,
            current_chunk_count_fn=lambda: int(self.collection.count() or 0),
            path_allow_fn=self._is_path_allowed,
            local_scope_max_docs=self._keyword_local_scope_max_docs,
            rebuild_delay_sec=self._keyword_rebuild_delay_sec,
        )
        self._query_cache_prewarm_mode: str = str(
            os.getenv("FILEAGENT_QUERY_CACHE_PREWARM", "background") or "background"
        ).strip().lower()
        self._query_cache_prewarm_initial_delay_sec: float = max(
            0.0,
            float(os.getenv("FILEAGENT_QUERY_CACHE_PREWARM_INITIAL_DELAY_SEC", "3.0") or 3.0),
        )
        self._query_cache_prewarm_min_idle_sec: float = max(
            0.0,
            float(os.getenv("FILEAGENT_QUERY_CACHE_PREWARM_MIN_IDLE_SEC", "5.0") or 5.0),
        )
        self._query_cache_prewarm_lock = threading.Lock()
        self._query_cache_prewarm_cancel = threading.Event()
        self._query_cache_prewarm_done = False
        self._query_cache_prewarm_error: Optional[str] = None
        self._query_cache_prewarm_thread: Optional[threading.Thread] = None
        self._query_cache_prewarm_shutdown = threading.Event()

        try:
            self._embedding_error_log_path = self._resolve_embedding_error_log_path()
            logger.info(f"Embedding error log path: {self._embedding_error_log_path}")
        except Exception:
            self._embedding_error_log_path = None
        try:
            self._index_details_log_path = self._resolve_index_details_log_path()
        except Exception:
            self._index_details_log_path = None

    def _resolve_embedding_error_log_path(self) -> str:
        data_dir = (os.getenv("FILEAGENT_DATA_DIR") or "").strip()
        if data_dir.startswith("~"):
            data_dir = os.path.expanduser(data_dir)
        if data_dir:
            base_dir = os.path.abspath(data_dir)
        else:
            base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
        logs_dir = os.path.join(base_dir, "logs")
        try:
            os.makedirs(logs_dir, exist_ok=True)
        except Exception:
            pass
        return os.path.join(logs_dir, "embedding_errors.log")

    def _resolve_index_details_log_path(self) -> str:
        data_dir = (os.getenv("FILEAGENT_DATA_DIR") or "").strip()
        if data_dir.startswith("~"):
            data_dir = os.path.expanduser(data_dir)
        if data_dir:
            base_dir = os.path.abspath(data_dir)
        else:
            base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
        logs_dir = os.path.join(base_dir, "logs")
        try:
            os.makedirs(logs_dir, exist_ok=True)
        except Exception:
            pass
        return os.path.join(logs_dir, "index_details.log")

    def _resolve_keyword_index_sidecar_path(self) -> str:
        override_name = str(
            os.getenv("FILEAGENT_KEYWORD_SIDECAR_FILE", "fileagent_keyword_index.pkl") or "fileagent_keyword_index.pkl"
        ).strip() or "fileagent_keyword_index.pkl"
        db_dir = os.path.abspath(os.path.expanduser(self.db_path))
        try:
            os.makedirs(db_dir, exist_ok=True)
        except Exception:
            pass
        return os.path.join(db_dir, override_name)

    def _build_keyword_index_records(self) -> List[KeywordIndexRecord]:
        cached = self._get_cached_metadata(include_documents=True)
        all_ids = cached.get("ids") or []
        all_metas = cached.get("metadatas") or []
        all_docs = cached.get("documents") or []

        media_chunk_types = {
            "media_summary",
            "media_audio_summary",
            "media_visual_summary",
            "interval_summary",
            "interval_visual",
            "asr_transcript",
            "asr_segment",
            "keyframe",
        }
        media_exts = {
            ".mp4", ".mov", ".mkv", ".avi", ".webm", ".m4v",
            ".mp3", ".wav", ".m4a", ".flac", ".aac", ".ogg",
        }
        max_content_chars = max(
            1000,
            int(os.getenv("FILEAGENT_KEYWORD_CONTENT_MAX_CHARS_PER_FILE", "12000") or 12000),
        )
        max_preview_chars = max(
            400,
            int(os.getenv("FILEAGENT_KEYWORD_CONTENT_PREVIEW_CHARS", "1800") or 1800),
        )

        def _trim_text(value: Any, limit: int = 900) -> str:
            text = re.sub(r"\s+", " ", str(value or "")).strip()
            if len(text) > limit:
                return text[:limit].rstrip()
            return text

        def _append_content(bucket: Dict[str, Any], text: str, *, cap: int = max_content_chars) -> None:
            clean = _trim_text(text)
            if not clean:
                return
            used = int(bucket.get("content_len") or 0)
            if used >= cap:
                return
            remaining = max(0, cap - used)
            if len(clean) > remaining:
                clean = clean[:remaining].rstrip()
            if not clean:
                return
            bucket.setdefault("content_parts", []).append(clean)
            bucket["content_len"] = used + len(clean) + 1

        text_content_cap = max(
            1200,
            int(os.getenv("FILEAGENT_KEYWORD_TEXT_CONTENT_MAX_CHARS_PER_FILE", "3200") or 3200),
        )
        tabular_content_cap = max(
            text_content_cap,
            int(os.getenv("FILEAGENT_KEYWORD_TABULAR_CONTENT_MAX_CHARS_PER_FILE", "8000") or 8000),
        )
        tabular_exts = {".csv", ".tsv", ".xlsx", ".xls", ".numbers"}
        textual_exts = {
            ".pdf", ".doc", ".docx", ".txt", ".md", ".rtf", ".odt", ".pages",
            ".ppt", ".pptx", ".html", ".htm",
        }
        textish_categories = {
            "data", "document", "report", "resume", "paper", "invoice",
            "manual", "note", "contract", "book",
        }

        file_buckets: Dict[str, Dict[str, Any]] = {}
        for mid, meta, doc in itertools.zip_longest(all_ids, all_metas, all_docs, fillvalue=""):
            meta = meta or {}
            fp = str(meta.get("file_path", "") or "").strip()
            if not fp:
                continue
            is_filename_lookup = self._is_filename_lookup_metadata(meta)
            bucket = file_buckets.setdefault(
                fp,
                {
                    "id": str(mid or ""),
                    "meta": dict(meta),
                    "content_parts": [],
                    "content_len": 0,
                },
            )
            if is_filename_lookup or not bucket.get("meta"):
                bucket["id"] = str(mid or "")
                bucket["meta"] = dict(meta)

            chunk_type = str(meta.get("chunk_type") or "").strip().lower()
            file_ext_for_media = str(meta.get("file_extension") or os.path.splitext(fp)[1]).strip().lower()
            media_type = str(meta.get("media_type") or "").strip().lower()
            category = self._meta_category_family(meta)
            is_media_chunk = (
                chunk_type in media_chunk_types
                or media_type in {"audio", "video"}
                or file_ext_for_media in media_exts
                or category in {"audio", "video", "audio/video"}
            )
            if is_media_chunk and not is_filename_lookup:
                _append_content(bucket, meta.get("keyframe_description", ""))
                _append_content(bucket, meta.get("keyframe_ocr_text", ""))
                _append_content(bucket, meta.get("media_visual_summary", ""))
                _append_content(bucket, meta.get("media_audio_summary", ""))
                _append_content(bucket, meta.get("media_summary", ""))
                _append_content(bucket, doc)
                continue

            if is_filename_lookup:
                continue

            is_tabular_chunk = (
                file_ext_for_media in tabular_exts
                or category == "data"
                or chunk_type in {"table", "spreadsheet", "worksheet", "sheet", "csv", "tsv", "excel"}
            )
            is_textish_chunk = is_tabular_chunk or file_ext_for_media in textual_exts or category in textish_categories
            if is_textish_chunk:
                _append_content(
                    bucket,
                    doc,
                    cap=tabular_content_cap if is_tabular_chunk else text_content_cap,
                )

        records: List[KeywordIndexRecord] = []
        for fp, bucket in file_buckets.items():
            mid = str(bucket.get("id") or "")
            meta = dict(bucket.get("meta") or {})

            fname = str(meta.get("file_name", "") or "").strip()
            summary = str(meta.get("doc_summary", "") or "").strip()
            en_tags = str(meta.get("en_tags", "") or "").strip()
            lookup_aliases = str(meta.get("lookup_aliases", "") or "").strip()
            table_schema_hint = str(meta.get("table_schema_hint", "") or "").strip()
            file_name_en = str(meta.get("file_name_en", "") or "").strip()
            folder_name_en = str(meta.get("folder_name_en", "") or "").strip()
            folder_base = os.path.basename(os.path.dirname(fp or "")) if fp else ""
            lookup_blob = build_lookup_blob(
                fname,
                file_name_en,
                folder_name_en,
                folder_base,
                fp,
            )
            latin_aliases = build_cjk_latin_aliases(fname, file_name_en, folder_name_en, folder_base, fp)
            content_index_text = _trim_text(" ".join(bucket.get("content_parts") or []), limit=max_content_chars)
            content_preview = _trim_text(content_index_text, limit=max_preview_chars)
            file_ext = str(meta.get("file_extension", "") or "").strip().lower()
            if not file_ext and fp:
                file_ext = os.path.splitext(fp)[1].lower()
            ext_label = file_ext.lstrip(".")
            doc_text = (
                f"{fname} {file_name_en} {folder_name_en} {summary} "
                f"{table_schema_hint} {en_tags} {lookup_aliases} {lookup_blob} {latin_aliases} "
                f"{content_index_text} "
                f"extension {file_ext} {ext_label} type {ext_label} filetype {ext_label}"
            )
            tokens = _tokenize_for_bm25(doc_text) or [""]

            records.append(KeywordIndexRecord(
                chroma_id=str(mid or ""),
                file_path=fp,
                normalized_path=os.path.normcase(os.path.normpath(os.path.abspath(os.path.expanduser(fp)))),
                file_name=fname or os.path.basename(fp),
                category=str(self._meta_category_family(meta) or ""),
                file_extension=file_ext,
                tokens=tokens,
                content_preview=content_preview,
            ))

        return records

    @staticmethod
    def _filename_lookup_doc_id(file_path: str) -> str:
        norm = os.path.normcase(
            os.path.normpath(os.path.abspath(os.path.expanduser(str(file_path or ""))))
        )
        digest = hashlib.sha256(norm.encode("utf-8")).hexdigest()[:48]
        return f"filename_lookup_{digest}"

    @staticmethod
    def _is_filename_lookup_metadata(meta: Optional[Dict[str, Any]]) -> bool:
        if not meta:
            return False
        return bool(meta.get("filename_lookup_index")) or str(meta.get("chunk_type") or "") == "filename_lookup"

    @staticmethod
    def _sanitize_metadata_for_chroma(meta: Dict[str, Any]) -> Dict[str, Any]:
        cleaned: Dict[str, Any] = {}
        for key, value in (meta or {}).items():
            if value is None:
                continue
            k = str(key)
            if isinstance(value, bool):
                cleaned[k] = value
            elif isinstance(value, str):
                cleaned[k] = value
            elif isinstance(value, int):
                cleaned[k] = int(value)
            elif isinstance(value, float):
                cleaned[k] = float(value) if math.isfinite(float(value)) else str(value)
            else:
                try:
                    cleaned[k] = json.dumps(value, ensure_ascii=False)[:1500]
                except Exception:
                    cleaned[k] = str(value)[:1500]
        return cleaned

    def _build_filename_lookup_record(
        self,
        file_path: str,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> Optional[Tuple[str, str, Dict[str, Any]]]:
        """Build one DB-backed virtual record that makes a file name vector-searchable."""
        base = dict(metadata or {})
        fp = str(base.get("file_path") or file_path or "").strip()
        if not fp:
            return None

        file_name = str(base.get("file_name") or os.path.basename(fp)).strip()
        file_name_no_ext, ext_from_name = os.path.splitext(file_name)
        file_ext = str(base.get("file_extension") or ext_from_name or os.path.splitext(fp)[1]).strip().lower()
        parent_folder = str(base.get("parent_folder") or os.path.basename(os.path.dirname(fp))).strip()
        file_name_en = str(base.get("file_name_en") or file_name_no_ext).strip()
        folder_name_en = str(base.get("folder_name_en") or parent_folder).strip()
        doc_summary = str(base.get("doc_summary") or "").strip()
        if not doc_summary:
            doc_summary = f"Indexed file: {file_name}"
        table_schema_hint = str(base.get("table_schema_hint") or "").strip()
        en_tags = str(base.get("en_tags") or "").strip()
        if not en_tags:
            en_tags = self._build_en_tags(
                file_name,
                file_ext,
                doc_summary=" ".join(x for x in [doc_summary, table_schema_hint] if x),
                parent_folder=parent_folder,
                folder_name_en=folder_name_en,
            )

        lookup_seed = dict(base)
        lookup_seed.update(
            {
                "file_path": fp,
                "file_name": file_name,
                "file_name_no_ext": file_name_no_ext,
                "file_name_en": file_name_en,
                "file_extension": file_ext,
                "parent_folder": parent_folder,
                "folder_name_en": folder_name_en,
                "doc_summary": doc_summary,
                "en_tags": en_tags,
                "table_schema_hint": table_schema_hint,
            }
        )
        lookup_aliases = str(base.get("lookup_aliases") or "").strip()
        rebuilt_aliases = self._build_lookup_aliases(fp, lookup_seed, doc_summary, table_schema_hint)
        if rebuilt_aliases:
            lookup_aliases = rebuilt_aliases

        category_family = self._meta_category_family(base)
        category_leaf = self._meta_category_leaf(base)
        source_role = self._meta_doc_role(base)
        lookup_blob = build_lookup_blob(
            file_name,
            file_name_no_ext,
            file_name_en,
            parent_folder,
            folder_name_en,
            fp,
            doc_summary,
            en_tags,
            table_schema_hint,
            lookup_aliases,
            max_terms=128,
        )
        latin_aliases = build_cjk_latin_aliases(
            file_name,
            file_name_no_ext,
            file_name_en,
            parent_folder,
            folder_name_en,
            fp,
            max_terms=64,
        )
        if latin_aliases:
            lookup_aliases = f"{lookup_aliases} {latin_aliases}".strip()

        doc_parts = [
            "Filename lookup record",
            f"File name: {file_name}",
            f"Stem: {file_name_no_ext}",
            f"English name: {file_name_en}",
            f"Folder: {' '.join(x for x in [parent_folder, folder_name_en] if x)}",
            f"Extension: {file_ext}",
            f"Path: {fp}",
            f"Category: {category_family} {category_leaf}",
            f"Summary: {doc_summary}",
        ]
        if table_schema_hint:
            doc_parts.append(f"Schema hint: {table_schema_hint}")
        if en_tags:
            doc_parts.append(f"Keywords: {en_tags}")
        if lookup_aliases:
            doc_parts.append(f"Aliases: {lookup_aliases}")
        if lookup_blob:
            doc_parts.append(f"Lookup terms: {lookup_blob}")
        doc_text = "\n".join(part for part in doc_parts if str(part or "").strip())[:1600]

        lookup_meta = dict(base)
        lookup_meta.update(
            {
                "file_path": fp,
                "file_name": file_name,
                "file_name_no_ext": file_name_no_ext,
                "file_name_en": file_name_en,
                "file_extension": file_ext,
                "file_type": file_ext.lstrip("."),
                "lookup_type": f"extension:{file_ext.lstrip('.')}" if file_ext else "",
                "parent_folder": parent_folder,
                "folder_name_en": folder_name_en,
                "doc_category": category_family,
                "doc_category_family": category_family,
                "doc_category_leaf": category_leaf,
                "doc_category_raw": str(base.get("doc_category_raw") or category_leaf or category_family),
                "doc_role": "reference",
                "source_doc_role": source_role,
                "doc_summary": doc_summary,
                "en_tags": en_tags,
                "table_schema_hint": table_schema_hint,
                "lookup_aliases": lookup_aliases,
                "chunk_type": "filename_lookup",
                "filename_lookup_index": True,
                "filename_lookup_version": 1,
                "indexed_virtual_record": True,
            }
        )
        return self._filename_lookup_doc_id(fp), doc_text, self._sanitize_metadata_for_chroma(lookup_meta)

    def _upsert_filename_lookup_for_file(
        self,
        file_path: str,
        metadata: Optional[Dict[str, Any]] = None,
        *,
        persist: bool = False,
    ) -> bool:
        record = self._build_filename_lookup_record(file_path, metadata)
        if not record:
            return False
        doc_id, doc_text, lookup_meta = record
        try:
            embeddings = self._embed_texts_for_index([doc_text])
            if not embeddings or not any(embeddings[0]):
                self._index_warning(f"filename lookup upsert skipped zero embedding: {file_path}")
                return False
            self.collection.upsert(
                ids=[doc_id],
                embeddings=[embeddings[0]],
                documents=[doc_text],
                metadatas=[lookup_meta],
            )
            self._invalidate_meta_cache()
            if persist:
                self._maybe_persist(force=True, reason="filename_lookup_upsert")
            return True
        except Exception as e:
            self._index_exception(f"filename lookup upsert failed: {file_path}", e)
            return False

    def backfill_filename_lookup_index(
        self,
        *,
        batch_size: int = 128,
        force: bool = False,
        rebuild_keyword_index: bool = True,
    ) -> Dict[str, Any]:
        """
        Backfill one filename lookup vector record per already-indexed file.

        This scans Chroma metadata only. It does not touch the local filesystem and
        does not re-ingest file contents, so an existing DB can be upgraded in place.
        """
        started = time.time()
        stats: Dict[str, Any] = {
            "ok": False,
            "source_chunks": 0,
            "source_files": 0,
            "existing_lookup_files": 0,
            "planned": 0,
            "upserted": 0,
            "skipped": 0,
            "errors": 0,
            "db_path": self.db_path,
        }
        self.enter_write_heavy_mode(reason="filename_lookup_backfill")
        wrote = False
        try:
            cached = self._get_cached_metadata(include_documents=False)
            all_ids = cached.get("ids") or []
            all_metas = cached.get("metadatas") or []
            stats["source_chunks"] = len(all_ids)

            file_metas: Dict[str, Dict[str, Any]] = {}
            existing_lookup_paths: set[str] = set()
            for meta in all_metas:
                meta = dict(meta or {})
                fp = str(meta.get("file_path") or "").strip()
                if not fp:
                    continue
                if self._is_filename_lookup_metadata(meta):
                    existing_lookup_paths.add(fp)
                    continue
                file_metas.setdefault(fp, meta)

            stats["source_files"] = len(file_metas)
            stats["existing_lookup_files"] = len(existing_lookup_paths)

            records: List[Tuple[str, str, Dict[str, Any]]] = []
            for fp, meta in file_metas.items():
                if not force and fp in existing_lookup_paths:
                    stats["skipped"] += 1
                    continue
                record = self._build_filename_lookup_record(fp, meta)
                if record:
                    records.append(record)
            stats["planned"] = len(records)

            try:
                batch_size = max(1, min(int(batch_size or 128), 512))
            except Exception:
                batch_size = 128

            for start in range(0, len(records), batch_size):
                batch = records[start : start + batch_size]
                docs = [row[1] for row in batch]
                try:
                    embeddings = self._embed_texts_for_index(docs)
                except Exception as e:
                    stats["errors"] += len(batch)
                    self._index_exception(
                        f"filename lookup backfill embedding batch failed offset={start}", e
                    )
                    continue

                valid_rows: List[Tuple[str, List[float], str, Dict[str, Any]]] = []
                for (doc_id, doc_text, meta), emb in zip(batch, embeddings):
                    if not emb or not any(emb):
                        stats["errors"] += 1
                        continue
                    valid_rows.append((doc_id, emb, doc_text, meta))
                if not valid_rows:
                    continue

                try:
                    self.collection.upsert(
                        ids=[row[0] for row in valid_rows],
                        embeddings=[row[1] for row in valid_rows],
                        documents=[row[2] for row in valid_rows],
                        metadatas=[row[3] for row in valid_rows],
                    )
                    wrote = True
                    stats["upserted"] += len(valid_rows)
                    logger.info(
                        "[filename_lookup_backfill] upserted batch %s-%s/%s",
                        start + 1,
                        start + len(batch),
                        len(records),
                    )
                except Exception as e:
                    stats["errors"] += len(valid_rows)
                    self._index_exception(
                        f"filename lookup backfill upsert batch failed offset={start}", e
                    )

            if wrote:
                self._invalidate_meta_cache()
                self._maybe_persist(force=True, reason="filename_lookup_backfill")

            if rebuild_keyword_index:
                try:
                    self._get_cached_metadata(include_documents=False)
                    self._keyword_index_manager.warm_start(background=False)
                except Exception as e:
                    stats["keyword_index_error"] = str(e)
                    self._index_exception("filename lookup keyword index rebuild failed", e)

            stats["ok"] = stats["errors"] == 0
            stats["elapsed_s"] = round(time.time() - started, 3)
            logger.info("[filename_lookup_backfill] done stats=%s", stats)
            return stats
        finally:
            self.leave_write_heavy_mode(reason="filename_lookup_backfill")

    def _append_embedding_error_log(self, msg: str) -> None:
        p = self._embedding_error_log_path
        if not p:
            return
        try:
            with open(p, "a", encoding="utf-8") as f:
                f.write(msg.rstrip() + "\n")
        except Exception:
            pass

    def _append_index_details_log(self, msg: str) -> None:
        p = self._index_details_log_path
        if not p:
            return
        try:
            with open(p, "a", encoding="utf-8") as f:
                f.write(msg.rstrip() + "\n")
        except Exception:
            pass

    def _index_log(self, level: str, msg: str) -> None:
        ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        line = f"{ts} - index - {str(level).upper()} - {msg}"
        self._append_index_details_log(line)

    def _index_info(self, msg: str) -> None:
        self._index_log("INFO", msg)

    def _index_warning(self, msg: str) -> None:
        self._index_log("WARNING", msg)

    def _index_error(self, msg: str) -> None:
        self._index_log("ERROR", msg)

    def _maybe_log_indexed_paths_stats(self, total_count: int, unique_count: int) -> None:
        now = time.time()
        sig = f"{int(total_count)}:{int(unique_count)}"
        try:
            interval = float(os.getenv("FILEAGENT_INDEXED_PATHS_LOG_INTERVAL_SEC", "60"))
        except ValueError:
            interval = 60.0
        interval = max(5.0, min(interval, 3600.0))

        if sig != self._indexed_paths_log_last_sig or (now - self._indexed_paths_log_last_ts) >= interval:
            self._index_info(f"数据库中共 {total_count} 条记录，去重后 {unique_count} 个文件")
            self._indexed_paths_log_last_sig = sig
            self._indexed_paths_log_last_ts = now

    def _index_exception(self, msg: str, err: Any) -> None:
        self._index_log("ERROR", f"{msg} | error={err}")

    def _log_index_file_timing(
        self,
        *,
        file_path: str,
        status: str,
        stage: str,
        total_s: float,
        embed_s: float = 0.0,
        llm_s: float = 0.0,
        vl_s: float = 0.0,
        chunks: int = 0,
        doc_category: str = "",
        extra: str = "",
    ) -> None:
        line = (
            f"[INDEX_FILE] status={status} stage={stage} file={file_path} "
            f"total_s={float(total_s):.3f} embedding_s={float(embed_s):.3f} "
            f"llm_s={float(llm_s):.3f} vl_s={float(vl_s):.3f} "
            f"chunks={int(chunks)} category={doc_category or '-'}"
        )
        if extra:
            line += f" {extra}"
        timing_line = (
            f"[INDEX_TIMING] status={status} file={file_path} "
            f"total_s={float(total_s):.3f} embedding_s={float(embed_s):.3f} "
            f"vl_s={float(vl_s):.3f} llm_s={float(llm_s):.3f} "
            f"chunks={int(chunks)} stage={stage}"
        )
        if extra:
            timing_line += f" {extra}"
        if status == "ok":
            self._index_info(line)
            logger.info(timing_line)
        else:
            self._index_error(line)
            logger.error(timing_line)

    @contextmanager
    def _embed_context(self, ctx: str):
        prev = getattr(self._embed_ctx_local, "value", None)
        self._embed_ctx_local.value = str(ctx or "")
        try:
            yield
        finally:
            if prev is None:
                try:
                    delattr(self._embed_ctx_local, "value")
                except Exception:
                    self._embed_ctx_local.value = ""
            else:
                self._embed_ctx_local.value = prev

    def _get_embed_context(self) -> str:
        v = getattr(self._embed_ctx_local, "value", "")
        if not v:
            return "ctx=<unknown>"
        return str(v)

    def _use_bce_instruction(self) -> bool:
        repo = str(getattr(settings, "EMBEDDING_REPO_ID", "") or "").lower()
        model = str(getattr(settings, "EMBEDDING_MODEL", "") or "").lower()
        path = str(getattr(settings, "LOCAL_EMBEDDING_MODEL_PATH", "") or "").lower()
        token = "bce-embedding-base_v1"
        return (token in repo) or (token in model) or (token in path)

    def _format_query_for_embedding(self, text: str) -> str:
        raw = str(text or "")
        try:
            max_q_chars = int(os.getenv("FILEAGENT_EMBED_MAX_QUERY_CHARS", "1600"))
        except Exception:
            max_q_chars = 1600
        max_q_chars = max(256, min(max_q_chars, 12000))
        if len(raw) > max_q_chars:
            raw = raw[:max_q_chars]
        if not raw:
            return raw
        if self._use_bce_instruction():
            return f"{BCE_QUERY_PREFIX}{raw}"
        return raw

    def _format_doc_for_embedding(self, text: str) -> str:
        raw = str(text or "")
        try:
            max_d_chars = int(os.getenv("FILEAGENT_EMBED_MAX_DOC_CHARS", "4000"))
        except Exception:
            max_d_chars = 4000
        max_d_chars = max(256, min(max_d_chars, 32000))
        if len(raw) > max_d_chars:
            raw = raw[:max_d_chars]
        if not raw:
            return raw
        if self._use_bce_instruction():
            return f"{BCE_DOC_PREFIX}{raw}"
        return raw

    def _embed_query_text(self, text: str) -> List[float]:
        return self.embedding_model.get_text_embedding(self._format_query_for_embedding(text))

    def _embed_doc_text(self, text: str) -> List[float]:
        return self.embedding_model.get_text_embedding(self._format_doc_for_embedding(text))

    def _probe_llama_cpp_runtime(
        self,
        *,
        kind: str,
        model_path: str,
        n_ctx: int,
        n_batch: int,
        n_threads: int,
        n_gpu_layers: int,
        use_llama_embedding: bool = False,
        pooling_type_rank: bool = False,
        timeout_sec: int = 90,
    ) -> Tuple[bool, str]:
        """
        Probe risky llama.cpp model initialization in a subprocess.
        This prevents native aborts in libllama from killing the main backend process.
        """
        import platform

        if platform.system() != "Darwin":
            return True, "probe_skipped_non_darwin"
        if str(os.getenv("FILEAGENT_DISABLE_NATIVE_MODEL_PROBE", "")).strip().lower() in {"1", "true", "yes", "on"}:
            return True, "probe_disabled"

        payload = {
            "model_path": model_path,
            "n_ctx": int(n_ctx),
            "n_batch": int(n_batch),
            "n_threads": int(n_threads),
            "n_gpu_layers": int(n_gpu_layers),
            "use_llama_embedding": bool(use_llama_embedding),
            "pooling_type_rank": bool(pooling_type_rank),
        }
        script = r"""
import json, sys
cfg = json.loads(sys.argv[1])
from llama_cpp import Llama
kwargs = {
    "model_path": cfg["model_path"],
    "n_ctx": int(cfg["n_ctx"]),
    "n_batch": int(cfg["n_batch"]),
    "n_threads": int(cfg["n_threads"]),
    "n_gpu_layers": int(cfg["n_gpu_layers"]),
    "verbose": False,
}
if cfg.get("pooling_type_rank"):
    from llama_cpp import LLAMA_POOLING_TYPE_RANK
    kwargs["pooling_type"] = LLAMA_POOLING_TYPE_RANK
if cfg.get("use_llama_embedding"):
    from llama_cpp.llama_embedding import LlamaEmbedding
    kwargs["n_ubatch"] = int(cfg["n_batch"])
    model = LlamaEmbedding(**kwargs)
else:
    kwargs["embeddings"] = True
    model = Llama(**kwargs)
print("ok")
"""
        try:
            try:
                effective_timeout = int(
                    os.getenv("FILEAGENT_NATIVE_MODEL_PROBE_TIMEOUT_SEC", str(timeout_sec))
                    or str(timeout_sec)
                )
            except Exception:
                effective_timeout = int(timeout_sec)
            effective_timeout = max(5, effective_timeout)
            res = subprocess.run(
                [sys.executable, "-c", script, json.dumps(payload, ensure_ascii=False)],
                capture_output=True,
                text=True,
                timeout=effective_timeout,
            )
        except subprocess.TimeoutExpired:
            return False, f"{kind} probe timeout after {effective_timeout}s"
        except Exception as e:
            return False, f"{kind} probe exception: {e}"

        stdout = (res.stdout or "").strip()
        stderr = (res.stderr or "").strip()
        if res.returncode == 0:
            return True, stdout or "ok"
        detail = stderr or stdout or f"returncode={res.returncode}"
        return False, detail[:400]

    def _maybe_persist(self, *, force: bool = False, reason: str = "") -> None:
        if not force and self._in_write_heavy_mode():
            return
        try:
            interval = float(os.getenv("FILEAGENT_CHROMA_PERSIST_INTERVAL_SEC", "10") or "10")
        except Exception:
            interval = 10.0
        if not force and interval <= 0:
            return
        now = time.time()
        if (not force) and ((now - float(self._last_persist_ts or 0.0)) < interval):
            return
        try:
            fn = getattr(self.chroma_client, "persist", None)
            if callable(fn):
                with self._collection_io_lock:
                    logger.info(f"[ChromaPersist] start: reason={reason or 'n/a'} force={force}")
                    fn()
                    logger.info(f"[ChromaPersist] done: reason={reason or 'n/a'} force={force}")
        except Exception:
            pass
        self._last_persist_ts = now

    def close(self) -> None:
        try:
            self._query_cache_prewarm_shutdown.set()
        except Exception:
            pass

        try:
            keyword_index = getattr(self, "_keyword_index_manager", None)
            if keyword_index is not None:
                keyword_index.close()
        except Exception:
            pass

        try:
            self._maybe_persist()
        except Exception:
            pass

        try:
            self.unload_reranker(reason="kb_close")
        except Exception:
            pass

        try:
            emb_model = getattr(self, "embedding_model", None)
            raw_emb = getattr(emb_model, "_model", emb_model)
            if raw_emb is not None and hasattr(raw_emb, "close"):
                raw_emb.close()  # type: ignore[attr-defined]
        except Exception:
            pass
        finally:
            try:
                self.embedding_model = None
            except Exception:
                pass

        try:
            gc.collect()
        except Exception:
            pass

    def unload_reranker(self, reason: str = "") -> bool:
        unloaded = False
        with self._reranker_lock:
            reranker = getattr(self, "reranker", None)
            raw_reranker = getattr(reranker, "_model", reranker)
            if raw_reranker is None:
                self.reranker = None
                return False
            try:
                if hasattr(raw_reranker, "close"):
                    raw_reranker.close()  # type: ignore[attr-defined]
                unloaded = True
            except Exception as e:
                logger.warning(f"[KB] unload_reranker failed reason={reason or 'n/a'} err={e}")
            finally:
                self.reranker = None
        try:
            gc.collect()
            gc.collect()
        except Exception:
            pass
        if unloaded:
            logger.info(f"[KB] reranker unloaded reason={reason or 'n/a'}")
        return unloaded

    def ensure_reranker_ready(self, reason: str = "") -> bool:
        with self._reranker_lock:
            if getattr(self, "reranker", None) is not None:
                return True
            try:
                self._init_reranker()
            except Exception as e:
                logger.warning(f"[KB] ensure_reranker_ready failed reason={reason or 'n/a'} err={e}")
                return False
            ready = getattr(self, "reranker", None) is not None
        logger.info(f"[KB] reranker restore reason={reason or 'n/a'} ready={ready}")
        return ready

    def request_query_cache_prewarm(self, *, background: Optional[bool] = None, reason: str = "") -> None:
        """Request metadata/keyword cache warmup from an explicit lifecycle phase."""
        mode = self._query_cache_prewarm_mode
        if mode in {"", "0", "false", "off", "disabled", "none"}:
            logger.info("[QueryCachePrewarm] disabled")
            return

        run_background = (mode != "sync") if background is None else bool(background)
        self._query_cache_prewarm_done = False
        self._query_cache_prewarm_error = None

        if not run_background:
            self._query_cache_prewarm_cancel.clear()
            self._run_query_cache_prewarm()
            return

        with self._query_cache_prewarm_lock:
            worker = self._query_cache_prewarm_thread
            if worker is not None and worker.is_alive():
                logger.info(
                    f"[QueryCachePrewarm] request coalesced: reason={reason or 'n/a'}"
                )
                return
            self._query_cache_prewarm_cancel.clear()
            worker = threading.Thread(
                target=self._run_query_cache_prewarm,
                name="fileagent-query-cache-prewarm",
                daemon=True,
            )
            self._query_cache_prewarm_thread = worker
            worker.start()

    def _run_query_cache_prewarm(self) -> None:
        start_ts = time.time()
        try:
            if self._query_cache_prewarm_shutdown.is_set() or self._query_cache_prewarm_cancel.is_set():
                return

            if self._query_cache_prewarm_mode == "background" and self._query_cache_prewarm_initial_delay_sec > 0:
                deadline = time.time() + self._query_cache_prewarm_initial_delay_sec
                while time.time() < deadline:
                    if self._query_cache_prewarm_shutdown.is_set() or self._query_cache_prewarm_cancel.is_set():
                        return
                    if self._in_write_heavy_mode():
                        logger.info("[QueryCachePrewarm] deferred during initial delay because write-heavy mode is active")
                        time.sleep(0.25)
                        continue
                    time.sleep(min(0.1, max(0.0, deadline - time.time())))

            if self._query_cache_prewarm_mode == "background":
                while not self._query_cache_prewarm_shutdown.is_set() and not self._query_cache_prewarm_cancel.is_set():
                    if self._in_write_heavy_mode():
                        time.sleep(0.5)
                        continue
                    if self._has_recent_collection_write():
                        time.sleep(0.5)
                        continue
                    break
            else:
                if self._in_write_heavy_mode():
                    logger.info("[QueryCachePrewarm] deferred because write-heavy mode is active")
                    return
                if self._has_recent_collection_write():
                    logger.info(
                        "[QueryCachePrewarm] deferred due to recent collection writes; "
                        "will rebuild lazily on next request"
                    )
                    return

            current_count = int(self.collection.count() or 0)
            if current_count <= 0:
                logger.info("[QueryCachePrewarm] skipped: empty collection")
                return

            logger.info(
                f"[QueryCachePrewarm] start: mode={self._query_cache_prewarm_mode}, "
                f"chunks={current_count}"
            )

            self._get_cached_metadata(
                include_documents=False,
                cancel_check=lambda: bool(
                    self._query_cache_prewarm_shutdown.is_set()
                    or self._query_cache_prewarm_cancel.is_set()
                    or self._in_write_heavy_mode()
                ),
            )
            if self._query_cache_prewarm_shutdown.is_set() or self._query_cache_prewarm_cancel.is_set():
                return
            if self._query_cache_prewarm_mode == "background":
                while not self._query_cache_prewarm_shutdown.is_set() and not self._query_cache_prewarm_cancel.is_set():
                    if self._in_write_heavy_mode():
                        time.sleep(0.5)
                        continue
                    if self._has_recent_collection_write():
                        time.sleep(0.5)
                        continue
                    break
            else:
                if self._in_write_heavy_mode():
                    logger.info(
                        "[QueryCachePrewarm] skipped keyword warm start because write-heavy mode became active"
                    )
                    return
                if self._has_recent_collection_write():
                    logger.info(
                        "[QueryCachePrewarm] skipped keyword warm start due to ongoing collection writes; "
                        "will rebuild lazily on next request"
                    )
                    return
            if self._query_cache_prewarm_cancel.is_set():
                return
            self._keyword_index_manager.warm_start(
                background=False,
                cancel_check=lambda: bool(
                    self._query_cache_prewarm_shutdown.is_set()
                    or self._query_cache_prewarm_cancel.is_set()
                    or self._in_write_heavy_mode()
                ),
            )

            elapsed = time.time() - start_ts
            logger.info(
                f"[QueryCachePrewarm] ready: mode={self._query_cache_prewarm_mode}, "
                f"chunks={current_count}, took={elapsed:.2f}s"
            )
        except _QueryCachePrewarmCancelled:
            logger.info("[QueryCachePrewarm] cancelled before completion")
        except Exception as e:
            self._query_cache_prewarm_error = str(e)
            logger.warning(f"[QueryCachePrewarm] failed: {e}")
        finally:
            self._query_cache_prewarm_done = True
            with self._query_cache_prewarm_lock:
                current = threading.current_thread()
                if self._query_cache_prewarm_thread is current:
                    self._query_cache_prewarm_thread = None
            gc.collect()

    def ensure_query_resources_ready(self, *, reason: str = "") -> bool:
        """Synchronously prepare metadata and keyword resources before file chat/search.

        Indexing writes Chroma chunks first, then retrieval-side caches are derived
        from those chunks. This gate is intentionally synchronous for lifecycle
        transitions where the UI is about to mark indexing as completed.
        """
        if self._query_cache_prewarm_shutdown.is_set():
            return False
        if self._in_write_heavy_mode():
            logger.info(
                "[QueryCachePrewarm] ensure skipped because write-heavy mode is active: "
                f"reason={reason or 'n/a'}"
            )
            return False
        try:
            current_count = int(self.collection.count() or 0)
        except Exception as e:
            logger.warning(f"[QueryCachePrewarm] ensure failed to read collection count: {e}")
            return False
        if current_count <= 0:
            return True

        t0 = time.time()
        self._query_cache_prewarm_done = False
        self._query_cache_prewarm_error = None
        try:
            logger.info(
                f"[QueryCachePrewarm] ensure start: reason={reason or 'n/a'}, "
                f"chunks={current_count}"
            )
            self._get_cached_metadata(include_documents=False)
            keyword_index = getattr(self, "_keyword_index_manager", None)
            if keyword_index is not None:
                keyword_index.warm_start(background=False)
            ready = bool(self.is_keyword_index_ready())
            elapsed = time.time() - t0
            logger.info(
                f"[QueryCachePrewarm] ensure done: reason={reason or 'n/a'}, "
                f"ready={ready}, chunks={current_count}, took={elapsed:.2f}s"
            )
            self._query_cache_prewarm_done = ready
            return ready
        except Exception as e:
            self._query_cache_prewarm_error = str(e)
            logger.warning(f"[QueryCachePrewarm] ensure failed: {e}")
            return False
        finally:
            if self._query_cache_prewarm_error:
                self._query_cache_prewarm_done = True
    
    def _init_db(self):
        use_default_client = False
        custom_init_error: Optional[Exception] = None
        try:
            from chromadb.config import Settings as ChromaSettings  # type: ignore

            frozen = bool(getattr(sys, "frozen", False))
            api_impl = os.getenv("CHROMA_API_IMPL", "").strip()

            import time
            logger.info(f"[_init_db] START: {time.time()}")
            logger.info(
                f"[_init_db] db_path={self.db_path}, api_impl={(api_impl or '<default>')}, frozen={frozen}"
            )

            chroma_settings = ChromaSettings(
                is_persistent=True,
                persist_directory=self.db_path,
                anonymized_telemetry=False,
            )
            if api_impl:
                chroma_settings.chroma_api_impl = api_impl

            self.chroma_client = chromadb.PersistentClient(path=self.db_path, settings=chroma_settings)
            logger.info("[_init_db] PersistentClient initialized with custom settings")
        except Exception as e:
            use_default_client = True
            custom_init_error = e
            logger.warning(
                f"Failed to initialize ChromaDB with custom settings, falling back to default: {e}"
            )

        if use_default_client:
            try:
                self.chroma_client = chromadb.PersistentClient(path=self.db_path)
                logger.info("[_init_db] PersistentClient initialized with default settings")
            except Exception as fallback_e:
                logger.error(f"Failed to initialize ChromaDB with default settings: {fallback_e}")
                if custom_init_error is not None:
                    logger.error(f"Original custom settings init error: {custom_init_error}")
                raise
            
        logger.info(f"[_init_db] Client Created: {time.time()}")
            
        class DummyDefaultEmbeddingFunction:
            def __call__(self, input):
                return [[0.0] * 1024 for _ in range(len(input))]
            @classmethod
            def name(cls):
                return "default"
                
        raw_collection = self.chroma_client.get_or_create_collection(
            name="local_knowledge",
            metadata={"hnsw:space": "cosine"},
            embedding_function=DummyDefaultEmbeddingFunction()
        )
        self.collection = _ThreadSafeCollectionProxy(raw_collection, self._collection_io_lock)
        logger.info(f"[_init_db] Collection Created: {time.time()}, count: {self.collection.count()}")
        raw_folder_collection = self.chroma_client.get_or_create_collection(
            name="fileagent_folder_index",
            metadata={"hnsw:space": "cosine"},
            embedding_function=DummyDefaultEmbeddingFunction(),
        )
        self.folder_collection = _ThreadSafeCollectionProxy(raw_folder_collection, self._collection_io_lock)
        logger.info(
            f"[_init_db] Folder index collection ready, count: {self.folder_collection.count()}"
        )
    
    @staticmethod
    def _gguf_header_ok(path: str, min_size_bytes: int = 50 * 1024 * 1024) -> bool:
        """
        Lightweight pre-check before handing a GGUF file to llama.cpp.
        Returns False if the file is clearly incomplete or corrupt so we can
        log a clear error instead of letting llama.cpp crash with a cryptic message.
        Checks:
          1. File exists and size >= min_size_bytes (default 50 MB)
          2. First 4 bytes are the GGUF magic: 0x47 0x47 0x55 0x46 ('GGUF')
        Does NOT do a full parse — intentionally fast and offline.
        """
        try:
            if not os.path.isfile(path):
                return False
            if os.path.getsize(path) < min_size_bytes:
                return False
            with open(path, "rb") as f:
                return f.read(4) == b"GGUF"
        except Exception:
            return False

    def _init_embedding(self):
        logger.info(f"初始化 Embedding 模型: {settings.EMBEDDING_MODEL}")
        
        try:
            local_model_path = None
            
            if os.path.exists(settings.LOCAL_EMBEDDING_MODEL_PATH):
                local_model_path = settings.LOCAL_EMBEDDING_MODEL_PATH
                logger.info(f"使用项目内置模型: {local_model_path}")
                # Pre-check GGUF integrity before handing to llama.cpp.
                # A truncated file causes llama.cpp to crash with a cryptic error;
                # better to surface a clear message and leave embedding_model = None.
                if local_model_path.endswith(".gguf") and not self._gguf_header_ok(local_model_path):
                    logger.error(
                        f"❌ Embedding GGUF 文件不完整或已损坏（大小不足或 magic header 无效）: "
                        f"{local_model_path}（大小: {os.path.getsize(local_model_path) if os.path.isfile(local_model_path) else 'N/A'} bytes）"
                        f"，请重新下载。"
                    )
                    self.embedding_model = None
                    return
            else:
                cache_dir = os.path.join(os.path.expanduser("~"), ".cache", "huggingface", "hub")
                model_cache_name = "models--" + settings.EMBEDDING_MODEL.replace('/', '--')
                model_cache_path = os.path.join(cache_dir, model_cache_name)
                
                if os.path.exists(model_cache_path):
                    snapshots_path = os.path.join(model_cache_path, "snapshots")
                    if os.path.exists(snapshots_path):
                        snapshots = os.listdir(snapshots_path)
                        if snapshots:
                            local_model_path = os.path.join(snapshots_path, snapshots[0])
                            logger.info(f"使用 HF 缓存: {local_model_path}")

            if local_model_path and local_model_path.endswith('.gguf') and os.path.isfile(local_model_path):
                logger.info(f"使用 GGUF 格式加载 Embedding 模型: {local_model_path}")
                try:
                    try:
                        from llama_cpp.llama_embedding import LlamaEmbedding  # type: ignore
                    except Exception:
                        LlamaEmbedding = None  # type: ignore
                    import multiprocessing
                    n_threads = max(1, multiprocessing.cpu_count() - 2)
                    import platform
                    if platform.system() == "Darwin":
                        n_threads = min(n_threads, 4)
                        
                    # BGE-M3 GGUF supports a much larger training context; 1024 avoids
                    # repeatedly splitting enriched chunks while staying modest on memory.
                    try:
                        _n_ctx = int(os.getenv("FILEAGENT_EMBED_N_CTX", "1024"))
                    except Exception:
                        _n_ctx = 1024
                    _n_ctx = max(256, min(_n_ctx, 8192))
                    try:
                        _n_batch = int(os.getenv("FILEAGENT_EMBED_N_BATCH", str(min(_n_ctx, 1024))))
                    except Exception:
                        _n_batch = min(_n_ctx, 1024)
                    _n_batch = max(32, min(_n_batch, _n_ctx))
                    try:
                        _embed_max_tokens = int(
                            os.getenv("FILEAGENT_EMBED_MAX_TOKENS", str(min(_n_ctx, _n_batch)))
                        )
                    except Exception:
                        _embed_max_tokens = min(_n_ctx, _n_batch)
                    _embed_max_tokens = max(32, min(_embed_max_tokens, _n_ctx, _n_batch))
                    logger.info(
                        f"Embedding GGUF 参数: n_ctx={_n_ctx}, n_batch={_n_batch}, "
                        f"max_tokens={_embed_max_tokens}, "
                        f"index_chunk_chars={self._index_chunk_max_chars()}"
                    )

                    _embed_runtime = str(os.getenv("FILEAGENT_EMBED_RUNTIME", "auto")).strip().lower()
                    if _embed_runtime not in {"auto", "llama", "llama_embedding"}:
                        _embed_runtime = "auto"
                    if _embed_runtime == "llama":
                        logger.warning(
                            "FILEAGENT_EMBED_RUNTIME=llama 已不再支持；当前构建强制使用 LlamaEmbedding。"
                        )
                        _embed_runtime = "llama_embedding"
                    try:
                        _embed_gpu_layers = int(os.getenv("FILEAGENT_EMBED_N_GPU_LAYERS", "-1"))
                    except Exception:
                        _embed_gpu_layers = -1

                    _prefer_llama_embedding = bool(
                        LlamaEmbedding is not None and _embed_runtime in {"auto", "llama_embedding"}
                    )
                    logger.info(
                        f"Embedding runtime 规划: prefer_llama_embedding={_prefer_llama_embedding} "
                        f"gpu_layers={_embed_gpu_layers}"
                    )

                    if LlamaEmbedding is None or not _prefer_llama_embedding:
                        raise RuntimeError("LlamaEmbedding runtime unavailable for GGUF embedding")

                    runtime_name = "LlamaEmbedding"
                    _ok, _detail = self._probe_llama_cpp_runtime(
                        kind="embedding_llama_embedding",
                        model_path=local_model_path,
                        n_ctx=_n_ctx,
                        n_batch=_n_batch,
                        n_threads=n_threads,
                        n_gpu_layers=_embed_gpu_layers,
                        use_llama_embedding=True,
                    )
                    if not _ok:
                        raise RuntimeError(
                            f"LlamaEmbedding 预检失败，跳过主进程初始化以避免崩溃: {_detail}"
                        )
                    llm = LlamaEmbedding(
                        model_path=local_model_path,
                        n_ctx=_n_ctx,
                        n_batch=_n_batch,
                        n_ubatch=_n_batch,
                        n_threads=n_threads,
                        n_gpu_layers=_embed_gpu_layers,
                        verbose=False,
                    )
                    logger.info("✅ Embedding runtime: LlamaEmbedding")
                    class _GGUFEmbeddingAdapter:
                        def __init__(
                            self,
                            model,
                            runtime_name: str = "LlamaEmbedding",
                            vector_dim: int = 1024,
                            max_tokens: int = 512,
                            context_getter: Optional[Callable[[], str]] = None,
                            error_sink: Optional[Callable[[str], None]] = None,
                        ):
                            self._model = model
                            self._runtime_name = str(runtime_name or "LlamaEmbedding")
                            self._dim = max(128, int(vector_dim or 1024))
                            self._max_tokens = max(32, int(max_tokens or 512))
                            self._lock = threading.RLock()
                            self._global_fail_count = 0
                            self._ctx_fail_count: Dict[str, int] = {}
                            self._context_getter = context_getter or (lambda: "ctx=<unknown>")
                            self._error_sink = error_sink
                            self._truncation_log_count: Dict[str, int] = {}
                            self._split_log_count: Dict[str, int] = {}

                        def _get_query_embedding(self, query: str) -> list[float]:
                            return self.embed(query)
                            
                        def _get_text_embedding(self, text: str) -> list[float]:
                            return self.embed(text)
                            
                        async def _aget_query_embedding(self, query: str) -> list[float]:
                            return self.embed(query)
                            
                        async def _aget_text_embedding(self, text: str) -> list[float]:
                            return self.embed(text)
                            
                        def get_agg_embedding_from_queries(self, queries: list[str]) -> list[float]:
                            return self.embed(queries[0]) if queries else []
                            
                        def __call__(self, input: list[str]) -> list[list[float]]:
                            embeddings = [self.embed(text) for text in input]
                            return embeddings
                            
                        @classmethod
                        def name(cls):
                            return "gguf_embedding"

                        def _log_failure(self, stage: str, err: Any, txt: str) -> None:
                            self._global_fail_count += 1
                            try:
                                ctx = str(self._context_getter() or "ctx=<unknown>")
                            except Exception:
                                ctx = "ctx=<unknown>"
                            key = f"{ctx}|{stage}"
                            c = int(self._ctx_fail_count.get(key, 0)) + 1
                            self._ctx_fail_count[key] = c
                            if c <= 3 or c % 50 == 0:
                                sample = (txt or "").replace("\n", " ")[:100]
                                msg = (
                                    f"[Embedding] {stage} failed | {ctx} "
                                    f"| ctx_fail_count={c} global_fail_count={self._global_fail_count} "
                                    f"| text_len={len(txt or '')} sample={sample!r} | error={err}"
                                )
                                logger.error(msg)
                                try:
                                    if self._error_sink:
                                        self._error_sink(msg)
                                except Exception:
                                    pass

                        def _log_truncation(self, original_text: str, original_tokens: int, final_tokens: int) -> None:
                            try:
                                ctx = str(self._context_getter() or "ctx=<unknown>")
                            except Exception:
                                ctx = "ctx=<unknown>"
                            c = int(self._truncation_log_count.get(ctx, 0)) + 1
                            self._truncation_log_count[ctx] = c
                            if c <= 3 or c % 20 == 0:
                                sample = (original_text or "").replace("\n", " ")[:100]
                                logger.warning(
                                    f"[Embedding] truncate overlong text | {ctx} "
                                    f"| token_count={original_tokens} kept_tokens={final_tokens} "
                                    f"| text_len={len(original_text or '')} sample={sample!r}"
                                )

                        def _log_split(
                            self,
                            original_text: str,
                            original_tokens: int,
                            part_count: int,
                            budget_tokens: int,
                            overlap_tokens: int,
                            reserve_tokens: int,
                        ) -> None:
                            try:
                                ctx = str(self._context_getter() or "ctx=<unknown>")
                            except Exception:
                                ctx = "ctx=<unknown>"
                            c = int(self._split_log_count.get(ctx, 0)) + 1
                            self._split_log_count[ctx] = c
                            if c <= 3 or c % 20 == 0:
                                sample = (original_text or "").replace("\n", " ")[:100]
                                logger.info(
                                    f"[Embedding] split overlong text | {ctx} "
                                    f"| token_count={original_tokens} parts={part_count} "
                                    f"budget_tokens={budget_tokens} overlap_tokens={overlap_tokens} "
                                    f"reserve_tokens={reserve_tokens} text_len={len(original_text or '')} "
                                    f"sample={sample!r}"
                                )

                        def _tokenize_text(self, txt: str, *, add_bos: bool) -> list[int]:
                            model = self._model
                            text_bytes = (txt or "").encode("utf-8")
                            return model.tokenize(text_bytes, add_bos=add_bos, special=True) or []

                        def _detokenize_text(self, tokens: list[int]) -> str:
                            if not tokens:
                                return ""
                            detok = getattr(self._model, "detokenize", None)
                            if callable(detok):
                                try:
                                    detok_bytes = detok(tokens)
                                    if isinstance(detok_bytes, (bytes, bytearray)):
                                        return bytes(detok_bytes).decode("utf-8", errors="ignore")
                                except Exception:
                                    pass
                            return ""

                        def split_text_for_embedding(
                            self,
                            text: str,
                            *,
                            reserve_text: str = "",
                            overlap_tokens: Optional[int] = None,
                            max_tokens: Optional[int] = None,
                        ) -> List[str]:
                            txt = text if isinstance(text, str) else str(text or "")
                            txt = txt.strip()
                            if not txt:
                                return []
                            target_tokens = max(32, min(int(max_tokens or self._max_tokens), self._max_tokens))
                            with self._lock:
                                reserve_tokens = len(self._tokenize_text(reserve_text, add_bos=False)) if reserve_text else 0
                                budget_tokens = max(32, target_tokens - reserve_tokens - 1)
                                tokens = self._tokenize_text(txt, add_bos=False)
                                total_tokens = len(tokens)
                                if total_tokens <= budget_tokens:
                                    return [txt]
                                if overlap_tokens is None:
                                    overlap_tokens = max(16, min(64, budget_tokens // 8))
                                overlap_tokens = max(0, min(int(overlap_tokens), max(0, budget_tokens // 3)))
                                step_tokens = max(1, budget_tokens - overlap_tokens)
                                parts: List[str] = []
                                start = 0
                                while start < total_tokens:
                                    end = min(total_tokens, start + budget_tokens)
                                    part = self._detokenize_text(tokens[start:end]).strip()
                                    if not part:
                                        ratio_start = start / max(1, total_tokens)
                                        ratio_end = end / max(1, total_tokens)
                                        ch_start = max(0, min(len(txt), int(len(txt) * ratio_start)))
                                        ch_end = max(ch_start + 1, min(len(txt), int(len(txt) * ratio_end)))
                                        part = txt[ch_start:ch_end].strip()
                                    if part and (not parts or part != parts[-1]):
                                        parts.append(part)
                                    if end >= total_tokens:
                                        break
                                    start += step_tokens
                            if parts:
                                self._log_split(
                                    txt,
                                    total_tokens,
                                    len(parts),
                                    budget_tokens,
                                    int(overlap_tokens or 0),
                                    reserve_tokens,
                                )
                                return parts
                            return [txt]

                        def _tokenize_with_limit(self, txt: str) -> tuple[list[int], str, bool]:
                            model = self._model
                            tokens = self._tokenize_text(txt, add_bos=True)
                            if not tokens:
                                return [], txt, False
                            original_tokens = len(tokens)
                            if original_tokens <= self._max_tokens:
                                return tokens, txt, False
                            tokens = tokens[: self._max_tokens]
                            safe_txt = txt
                            detok_txt = self._detokenize_text(tokens)
                            if detok_txt:
                                safe_txt = detok_txt
                            if safe_txt == txt:
                                approx_chars = max(
                                    64,
                                    int(len(txt or "") * (self._max_tokens / max(1, original_tokens))) - 8,
                                )
                                safe_txt = (txt or "")[:approx_chars]
                                for _ in range(6):
                                    probe_tokens = self._tokenize_text(safe_txt, add_bos=True)
                                    if not probe_tokens:
                                        break
                                    if len(probe_tokens) <= self._max_tokens:
                                        tokens = probe_tokens
                                        break
                                    approx_chars = max(
                                        32,
                                        int(len(safe_txt) * (self._max_tokens / max(1, len(probe_tokens)))) - 4,
                                    )
                                    safe_txt = safe_txt[:approx_chars]
                            self._log_truncation(txt, original_tokens, len(tokens))
                            return tokens, safe_txt, True

                        def _extract_vector(self, res: Any) -> Optional[list[float]]:
                            # Common structured return: {"data":[{"embedding":[...]}]}
                            try:
                                if isinstance(res, dict):
                                    data = res.get("data")
                                    if isinstance(data, list) and data:
                                        emb = data[0].get("embedding")
                                        if isinstance(emb, list) and emb:
                                            return [float(x) for x in emb]
                                if isinstance(res, list) and res:
                                    first = res[0]
                                    if isinstance(first, list):
                                        return [float(x) for x in first]
                                    if isinstance(first, (int, float)):
                                        return [float(x) for x in res]
                            except Exception:
                                return None
                            return None

                        def _embed_via_batch_get_one(self, tokens: list[int]) -> Optional[list[float]]:
                            """
                            Use llama_batch_get_one() + llama_decode() directly from the C API.
                            This keeps the primary embedding path on the low-level stable
                            batch/decode API instead of higher-level helper wrappers.
                            """
                            try:
                                import llama_cpp as _llc
                                model = self._model
                                if not tokens:
                                    return None
                                n_tokens = len(tokens)
                                n_embd = _llc.llama_n_embd(model.model)
                                if n_embd <= 0:
                                    return None
                                # Clear KV cache for a clean forward pass
                                try:
                                    _llc.llama_kv_cache_clear(model.ctx)
                                except Exception:
                                    pass
                                # Build a single-sequence batch directly.
                                c_tokens = (_llc.llama_token * n_tokens)(*tokens)
                                batch = _llc.llama_batch_get_one(c_tokens, n_tokens)
                                ret = _llc.llama_decode(model.ctx, batch)
                                if ret != 0:
                                    return None
                                # Extract embeddings — try pooled APIs first, then fallbacks
                                emb_ptr = None
                                for _fn_name, _args in (
                                    ("llama_get_embeddings_seq", (model.ctx, 0)),
                                    ("llama_get_embeddings_ith", (model.ctx, n_tokens - 1)),
                                    ("llama_get_embeddings", (model.ctx,)),
                                ):
                                    try:
                                        _fn = getattr(_llc, _fn_name, None)
                                        if _fn is None:
                                            continue
                                        emb_ptr = _fn(*_args)
                                        if emb_ptr:
                                            break
                                    except Exception:
                                        continue
                                if not emb_ptr:
                                    return None
                                raw = [float(emb_ptr[i]) for i in range(n_embd)]
                                # L2 normalise (BGE/bge-m3 embedding models expect unit vectors)
                                norm = (sum(x * x for x in raw) ** 0.5) or 1.0
                                return [x / norm for x in raw]
                            except Exception:
                                return None

                        def embed(self, text: str) -> list[float]:
                            txt = text if isinstance(text, str) else str(text or "")
                            # Primary: low-level C API path.
                            with self._lock:
                                tokens, safe_txt, was_truncated = self._tokenize_with_limit(txt)
                                vec = self._embed_via_batch_get_one(tokens)
                            if vec:
                                return vec
                            # Fallback: high-level embed() API (may fail on some versions)
                            try:
                                with self._lock:
                                    res = self._model.embed(safe_txt if was_truncated else txt)
                                vec = self._extract_vector(res)
                                if vec:
                                    return vec
                                self._log_failure(
                                    "embed",
                                    "empty_or_invalid_vector",
                                    safe_txt if was_truncated else txt,
                                )
                            except Exception as e:
                                self._log_failure("embed", e, safe_txt if was_truncated else txt)
                            return [0.0] * self._dim
                            
                        def get_text_embedding(self, text: str) -> list[float]:
                            return self.embed(text)
                    
                    try:
                        _dim = int(llm.n_embd())
                    except Exception:
                        _dim = 1024
                    self.embedding_model = _GGUFEmbeddingAdapter(
                        llm,
                        runtime_name=runtime_name,
                        vector_dim=_dim,
                        max_tokens=_embed_max_tokens,
                        context_getter=self._get_embed_context,
                        error_sink=self._append_embedding_error_log,
                    )
                    logger.info(
                        f"[Embedding] adapter ready: runtime={runtime_name} dim={_dim}"
                    )
                except ImportError:
                    logger.info(f"llama_cpp_python 未安装，无法加载 GGUF embedding 模型")
                    self.embedding_model = None
            elif local_model_path and os.path.exists(local_model_path):
                if ENABLE_LLAMA_INDEX_FALLBACK:
                    try:
                        from llama_index.embeddings.huggingface import HuggingFaceEmbedding
                        self.embedding_model = HuggingFaceEmbedding(
                            model_name=local_model_path,
                            trust_remote_code=True
                        )
                    except ImportError:
                        logger.info("HuggingFaceEmbedding 未安装")
                        self.embedding_model = None
                else:
                    logger.info("LlamaIndex HuggingFaceEmbedding disabled; using packaged GGUF embedding path only")
                    self.embedding_model = None
            else:
                logger.error(f"找不到本地 Embedding 模型文件: {settings.LOCAL_EMBEDDING_MODEL_PATH}，请检查是否已正确下载")
                self.embedding_model = None
            
            if self.embedding_model is not None:
                test_vec = self.embedding_model.get_text_embedding("测试")
                logger.info(f"✅ Embedding 加载成功，维度: {len(test_vec)}")
                if self._use_bce_instruction():
                    logger.info("✅ Embedding 指令模板: BCE query/doc 双前缀已启用")
            else:
                logger.error(f"❌ Embedding 缺失")
            
        except Exception as e:
            logger.error(f"❌ Embedding 加载失败: {e}")
            # Do NOT re-raise: a missing embedding only disables indexing, it must not
            # crash the entire FileKnowledgeBase / FileAgent startup.  Search over
            # already-indexed documents still works; user can re-index after fixing.
            self.embedding_model = None
    
    def _init_reranker(self):
        if os.getenv("FILEAGENT_DISABLE_RERANKER", "false").lower() in {"1", "true", "yes", "on"}:
            logger.info(f"Reranker disabled by FILEAGENT_DISABLE_RERANKER.")
            self.reranker = None
            return

        local_gguf = settings.LOCAL_RERANKER_MODEL_PATH
        if isinstance(local_gguf, str) and local_gguf.lower().endswith(".gguf") and os.path.isfile(local_gguf):
            logger.info(f"使用项目内置 Reranker(GGUF): {local_gguf}")
            try:
                import llama_cpp  # type: ignore
                from llama_cpp import LLAMA_POOLING_TYPE_RANK  # type: ignore
                try:
                    from llama_cpp.llama_embedding import LlamaEmbedding, NORM_MODE_NONE  # type: ignore
                except Exception:
                    LlamaEmbedding = None  # type: ignore
                    NORM_MODE_NONE = -1  # type: ignore
                import multiprocessing
                n_threads = max(1, multiprocessing.cpu_count() - 2)
                import platform
                if platform.system() == "Darwin":
                    n_threads = min(n_threads, 4)

                _gguf_bn = os.path.basename(local_gguf).lower()
                _rerank_n_ctx = 2048 if "qwen3-reranker" in _gguf_bn else 512
                _rerank_n_batch = min(512, _rerank_n_ctx)

                _rerank_runtime = str(os.getenv("FILEAGENT_RERANK_RUNTIME", "auto")).strip().lower()
                if _rerank_runtime not in {"auto", "llama", "llama_embedding"}:
                    _rerank_runtime = "auto"
                if _rerank_runtime == "llama":
                    logger.warning(
                        "FILEAGENT_RERANK_RUNTIME=llama 已不再支持；当前构建强制使用 LlamaEmbedding。"
                    )
                    _rerank_runtime = "llama_embedding"
                try:
                    _rerank_gpu_layers = int(os.getenv("FILEAGENT_RERANK_N_GPU_LAYERS", "-1"))
                except Exception:
                    _rerank_gpu_layers = -1

                _prefer_rerank_llama_embedding = bool(
                    LlamaEmbedding is not None and _rerank_runtime in {"auto", "llama_embedding"}
                )
                logger.info(
                    f"Reranker runtime 规划: prefer_llama_embedding={_prefer_rerank_llama_embedding} "
                    f"gpu_layers={_rerank_gpu_layers}"
                )

                if LlamaEmbedding is None or not _prefer_rerank_llama_embedding:
                    raise RuntimeError("LlamaEmbedding runtime unavailable for GGUF reranker")

                _ok, _detail = self._probe_llama_cpp_runtime(
                    kind="reranker_llama_embedding",
                    model_path=local_gguf,
                    n_ctx=_rerank_n_ctx,
                    n_batch=_rerank_n_batch,
                    n_threads=n_threads,
                    n_gpu_layers=_rerank_gpu_layers,
                    use_llama_embedding=True,
                    pooling_type_rank=True,
                )
                if not _ok:
                    raise RuntimeError(
                        f"LlamaEmbedding Reranker 预检失败，跳过主进程初始化以避免崩溃: {_detail}"
                    )

                llm = LlamaEmbedding(
                    model_path=local_gguf,
                    pooling_type=LLAMA_POOLING_TYPE_RANK,
                    n_ctx=_rerank_n_ctx,
                    n_batch=_rerank_n_batch,
                    n_ubatch=_rerank_n_batch,
                    n_threads=n_threads,
                    n_gpu_layers=_rerank_gpu_layers,
                    verbose=False,
                )
                logger.info("✅ Reranker runtime: LlamaEmbedding")

                class _GGUFRerankerAdapter:
                    def __init__(self, model):
                        self._model = model
                        self._lock = threading.RLock()
                        self._batch_warned = False
                        self._rerank_template_checked = False
                        self._rerank_template: Optional[str] = None

                    @staticmethod
                    def _coerce_score_value(value: Any) -> float:
                        if isinstance(value, (int, float)):
                            return float(value)
                        if hasattr(value, "item"):
                            try:
                                return float(value.item())
                            except Exception:
                                pass
                        if isinstance(value, dict):
                            score_val = value.get("score")
                            return _GGUFRerankerAdapter._coerce_score_value(score_val)
                        if isinstance(value, (list, tuple)):
                            if not value:
                                return 0.0
                            # Some runtimes surface [yes_logit, no_logit] or [score].
                            return _GGUFRerankerAdapter._coerce_score_value(value[0])
                        raise TypeError(f"Unsupported reranker score value: {type(value).__name__}")

                    @staticmethod
                    def _normalize_pairs(pairs: Any) -> List[Tuple[str, str]]:
                        if not pairs:
                            return []
                        if (
                            isinstance(pairs, (list, tuple))
                            and len(pairs) >= 2
                            and not isinstance(pairs[0], (list, tuple))
                        ):
                            return [(str(pairs[0] or ""), str(pairs[1] or ""))]
                        normalized: List[Tuple[str, str]] = []
                        for item in pairs:
                            if not isinstance(item, (list, tuple)) or len(item) < 2:
                                continue
                            normalized.append((str(item[0] or ""), str(item[1] or "")))
                        return normalized

                    def _mark_batch_warning(self, err: Exception) -> None:
                        if self._batch_warned:
                            return
                        self._batch_warned = True
                        logger.warning(
                            f"Reranker batch embed path failed, falling back to stable serial embed scoring: {err}"
                        )

                    def _get_rerank_template(self) -> Optional[str]:
                        if self._rerank_template_checked:
                            return self._rerank_template
                        self._rerank_template_checked = True
                        try:
                            tpl = llama_cpp.llama_model_chat_template(self._model._model.model, b"rerank")
                            if tpl:
                                self._rerank_template = tpl.decode("utf-8")
                        except Exception:
                            self._rerank_template = None
                        return self._rerank_template

                    def _build_rank_batch_inputs(self, pairs: List[Tuple[str, str]]) -> List[List[int]]:
                        rerank_template = self._get_rerank_template()
                        batch_inputs: List[List[int]] = []
                        query_cache: Dict[str, List[int]] = {}
                        eos_id = self._model.token_eos()
                        sep_id = self._model.token_sep() if self._model.token_sep() != -1 else eos_id

                        for query, document in pairs:
                            if rerank_template:
                                prompt = rerank_template.replace("{query}", query).replace("{document}", document)
                                tokens = self._model.tokenize(
                                    prompt.encode("utf-8"),
                                    add_bos=False,
                                    special=True,
                                )
                                batch_inputs.append(tokens)
                                continue

                            q_tokens = query_cache.get(query)
                            if q_tokens is None:
                                q_tokens = self._model.tokenize(query.encode("utf-8"), add_bos=True, special=True)
                                if q_tokens and q_tokens[-1] == eos_id:
                                    q_tokens = q_tokens[:-1]
                                query_cache[query] = q_tokens

                            d_tokens = self._model.tokenize(document.encode("utf-8"), add_bos=False, special=True)
                            full_seq = list(q_tokens) + [sep_id] + d_tokens
                            if not full_seq or full_seq[-1] != eos_id:
                                full_seq.append(eos_id)
                            batch_inputs.append(full_seq)

                        return batch_inputs

                    def _extract_scores(self, raw_results: Any, expected_len: int) -> List[float]:
                        if expected_len <= 0:
                            return []
                        if isinstance(raw_results, (int, float)):
                            return [float(raw_results)] + [0.0] * (expected_len - 1)
                        if not isinstance(raw_results, list):
                            return [0.0] * expected_len

                        if raw_results and isinstance(raw_results[0], dict):
                            parsed_scores = [0.0] * expected_len
                            for i, item in enumerate(raw_results):
                                if isinstance(item, dict) and "corpus_id" in item and "score" in item:
                                    idx = int(item["corpus_id"])
                                    if 0 <= idx < expected_len:
                                        parsed_scores[idx] = self._coerce_score_value(item["score"])
                                elif i < expected_len:
                                    parsed_scores[i] = self._coerce_score_value(item)
                            return parsed_scores

                        parsed = [self._coerce_score_value(item) for item in raw_results[:expected_len]]
                        if len(parsed) < expected_len:
                            parsed.extend([0.0] * (expected_len - len(parsed)))
                        return parsed

                    def _compute_embed_scores(self, pairs: List[Tuple[str, str]]) -> List[float]:
                        batch_inputs = self._build_rank_batch_inputs(pairs)
                        with self._lock:
                            raw_results = self._model.embed(batch_inputs, normalize=NORM_MODE_NONE)
                        return self._extract_scores(raw_results, len(pairs))

                    def compute_score(self, pairs):
                        normalized_pairs = self._normalize_pairs(pairs)
                        if not normalized_pairs:
                            return []
                        try:
                            return self._compute_embed_scores(normalized_pairs)
                        except Exception as e:
                            self._mark_batch_warning(e)

                        scores: List[float] = []
                        for pair in normalized_pairs:
                            try:
                                scores.extend(self._compute_embed_scores([pair]))
                            except Exception:
                                scores.append(0.0)
                        return scores

                self.reranker = _GGUFRerankerAdapter(llm)
                logger.info(f"✅ GGUF Reranker 初始化成功")
                return
            except Exception as e:
                logger.error(f"⚠️ GGUF Reranker init failed. Disabled. err={e}")
                self.reranker = None
                return

        self.reranker = None
        return
    
    
    def _get_local_llm_client(self):
        from services.inproc_openai_client import get_inproc_openai_client

        return get_inproc_openai_client()

    def _get_configured_index_quantization_file(self, model_id: str) -> str:
        mid = (model_id or "").strip()
        if not mid:
            return ""
        try:
            from services.preference_manager import PreferenceManager
            import config.settings as agent_settings

            base_dir = getattr(
                agent_settings,
                "BASE_DIR",
                os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
            )
            pm = PreferenceManager(base_dir)
            return (pm.get_selected_quantization_file(mid) or "").strip()
        except Exception:
            return ""

    def _is_local_llm_hot_for_indexing(self, model_id: str) -> bool:
        mid = (model_id or "").strip()
        if not mid:
            return False

        def _int_env(name: str, default: int) -> int:
            try:
                raw = (os.getenv(name) or "").strip()
                if not raw:
                    return int(default)
                return int(raw)
            except Exception:
                return int(default)

        try:
            llm_mgr = getattr(self, "_llm_manager", None)
            if llm_mgr is None:
                from services.local_llm import get_local_llm_manager

                llm_mgr = get_local_llm_manager()

            resolved = llm_mgr.resolve_target_model(
                mid,
                preferred_quantization_file=self._get_configured_index_quantization_file(mid) or None,
            )
            if not resolved:
                return False

            model_config, model_path, mmproj_path = resolved
            resolved_model_id = str(model_config.get("id") or mid).strip()
            expected_n_ctx = 5120
            expected_n_batch = 512
            if mmproj_path and os.path.exists(mmproj_path):
                expected_n_ctx = min(expected_n_ctx, _int_env("FILEAGENT_VL_N_CTX", 5120))
                expected_n_batch = min(expected_n_batch, _int_env("FILEAGENT_VL_N_BATCH", 512))
            expected_n_ubatch = min(
                expected_n_batch,
                _int_env("FILEAGENT_LLM_N_UBATCH", expected_n_batch),
            )

            current_n_ctx = int(getattr(llm_mgr, "current_n_ctx", 0) or 0)
            current_n_batch = int(getattr(llm_mgr, "current_n_batch", 0) or 0)
            current_n_ubatch = int(getattr(llm_mgr, "current_n_ubatch", 0) or 0)

            return bool(
                getattr(llm_mgr, "_llama", None) is not None
                and str(getattr(llm_mgr, "current_model_id", None) or "").strip() == resolved_model_id
                and str(getattr(llm_mgr, "current_model_path", None) or "") == str(model_path or "")
                and getattr(llm_mgr, "current_mmproj_path", None) == mmproj_path
                and current_n_ctx >= expected_n_ctx
                and current_n_batch >= expected_n_batch
                and current_n_ubatch >= expected_n_ubatch
            )
        except Exception:
            return False

    def _maybe_log_startup_index_prefill_observation(self, file_path: str, *, use_smart_indexing: bool) -> None:
        if not use_smart_indexing or getattr(self, "_startup_index_prefill_observation_logged", False):
            return

        idx_model = self._get_configured_index_model_id()
        if not idx_model:
            return

        try:
            llm_mgr = getattr(self, "_llm_manager", None)
            if llm_mgr is None:
                from services.local_llm import get_local_llm_manager

                llm_mgr = get_local_llm_manager()

            getter = getattr(llm_mgr, "get_startup_index_prefill_status", None)
            if not callable(getter):
                return

            status = dict(getter() or {})
            warmup_state = str(status.get("state") or "unknown").strip() or "unknown"
            reason = " ".join(str(status.get("reason") or "").split())[:240]
            ready_now = self._is_local_llm_hot_for_indexing(idx_model)
            current_model_id = str(getattr(llm_mgr, "current_model_id", None) or "").strip()
            current_model_path = str(getattr(llm_mgr, "current_model_path", None) or "")
            target_model_id = str(status.get("target_model_id") or idx_model).strip()
            target_model_path = str(status.get("target_model_path") or "")

            target_match = True
            if target_model_id and current_model_id and target_model_id != current_model_id:
                target_match = False
            if target_model_path and current_model_path and target_model_path != current_model_path:
                target_match = False

            prefill_hit = bool(warmup_state == "done" and ready_now and target_match)
            observer = getattr(llm_mgr, "record_startup_index_prefill_observation", None)
            snapshot: Dict[str, Any]
            if callable(observer):
                snapshot = dict(
                    observer(
                        hit=prefill_hit,
                        file_path=os.path.abspath(file_path),
                    )
                    or {}
                )
            else:
                snapshot = {
                    "observations_total": 1,
                    "hits_total": 1 if prefill_hit else 0,
                    "hit_rate": 1.0 if prefill_hit else 0.0,
                }

            warmup_elapsed_ms = int(status.get("elapsed_ms", 0) or 0)
            completed_at = float(status.get("completed_at", 0.0) or 0.0)
            warmup_age_ms = int(max(0.0, time.time() - completed_at) * 1000) if completed_at > 0 else -1

            self._index_info(
                f"[StartupWarmup] first_smart_index file={os.path.abspath(file_path)} "
                f"warmup_state={warmup_state} ready_now={str(bool(ready_now)).lower()} "
                f"prefill_hit={str(bool(prefill_hit)).lower()} target_match={str(bool(target_match)).lower()} "
                f"observations={int(snapshot.get('observations_total', 0) or 0)} "
                f"hits={int(snapshot.get('hits_total', 0) or 0)} "
                f"hit_rate={float(snapshot.get('hit_rate', 0.0) or 0.0):.3f} "
                f"warmup_elapsed_ms={warmup_elapsed_ms} warmup_age_ms={warmup_age_ms} "
                f"current_model={current_model_id or '<none>'} target_model={target_model_id or '<none>'} "
                f"reason={reason or 'n/a'}"
            )
            self._startup_index_prefill_observation_logged = True
        except Exception as e:
            logger.warning(f"[StartupWarmup] observation failed: {e}")
    
    def _test_local_llm_connection(self) -> bool:
        if getattr(self, "_local_llm_available", None) is True:
            return True
            
        import time
        last_fail = getattr(self, "_local_llm_last_fail", 0)
        if hasattr(self, "_local_llm_available") and self._local_llm_available is False:
            if time.time() - last_fail < 30:
                return False
                
        self._index_info(f"正在测试本地 LLM 连接...")
        try:
            idx_model = self._require_configured_index_model_id("测试本地 LLM 连接")
            if not idx_model:
                self._local_llm_available = False
                self._local_llm_last_fail = time.time()
                return False
            if self._is_local_llm_hot_for_indexing(idx_model):
                self._index_info(f"✅ 本地 LLM 已就绪，跳过额外探测")
                self._local_llm_available = True
                return True
            client = self._get_local_llm_client()
            response = client.chat.completions.create(
                model=idx_model,
                messages=[{"role": "user", "content": "你好"}],
                max_tokens=10,
                temperature=0.0,
                stream=False,
            )
            if response.choices and response.choices[0].message.content:
                self._index_info(f"✅ 本地 LLM 连接正常")
                self._local_llm_available = True
                return True
            else:
                self._index_error(f"⚠️ 本地 LLM 返回空响应")
                self._local_llm_available = False
                self._local_llm_last_fail = time.time()
                return False
        except Exception as e:
            self._index_exception("⚠️ 本地 LLM 不可用，可能会在后台继续尝试准备", e)
            self._local_llm_available = False
            self._local_llm_last_fail = time.time()
            return False

    def _append_model_prompt_suffix(self, prompt: str, model_id: Optional[str] = None) -> str:
        base = str(prompt or "")
        try:
            from services.local_llm import get_local_llm_manager
            llm_mgr = self._llm_manager or get_local_llm_manager()
            mid = model_id or getattr(llm_mgr, "current_model_id", None)
            if not mid:
                return base
            cfg = llm_mgr.get_target_model_config(mid) or {}
            suffix = str(cfg.get("intent_prompt_suffix") or "").strip()
            if not suffix:
                return base
            if base.rstrip().endswith(suffix):
                return base
            return f"{base}\n{suffix}"
        except Exception:
            return base
    
    def _looks_like_readme(self, file_name: str, file_ext: str) -> bool:
        n = (file_name or "").lower()
        ext = (file_ext or "").lower()
        if ext not in (".md", ".txt", ".rst"):
            return False
        return any(k in n for k in ["readme", "read.me", "quickstart", "install", "changelog", "license"])

    def _looks_like_manual_doc(self, file_name: str, content: str = "", file_ext: str = "") -> bool:
        name = (file_name or "").lower()
        text = (content or "").lower()
        ext = (file_ext or "").lower()
        if self._looks_like_readme(file_name, file_ext):
            return True
        if ext not in (".md", ".txt", ".rst", ".pdf", ".doc", ".docx"):
            return False
        name_markers = (
            "faq", "datasheet", "data sheet", "spec", "specification",
            "guide", "manual", "tutorial", "howto", "how-to", "install",
            "diagram", "architecture", "change", "changes", "changelog",
            "说明", "指南", "手册", "教程", "架构图", "修改说明", "接口文档", "快速开始",
        )
        text_markers = (
            "quick start", "getting started", "installation", "usage", "requirements",
            "setup", "步骤", "安装", "使用说明", "操作步骤", "注意事项",
        )
        if any(marker in name for marker in name_markers):
            return True
        if ext in (".md", ".txt", ".rst") and any(marker in text for marker in text_markers):
            return True
        return False

    def _looks_like_form_document(self, file_name: str, content: str = "", file_ext: str = "") -> bool:
        name = (file_name or "").lower()
        text = (content or "")
        text_lower = text.lower()
        ext = (file_ext or "").lower()
        if ext and ext not in (".md", ".txt", ".rst", ".pdf", ".doc", ".docx"):
            return False

        name_markers = (
            "form", "application", "registration", "申报表", "申请表", "报名表",
            "登记表", "审批表", "信息表", "认定表", "申请书",
        )
        field_markers = (
            "姓名", "性别", "出生年月", "身份证号码", "身份证号", "电话", "联系电话",
            "申请人", "签名", "所在院部", "专业", "学历", "申报职业", "申报级别",
            "评价机构", "发证日期", "name", "gender", "date of birth",
            "id number", "phone", "signature", "applicant",
        )
        if any(marker in name for marker in name_markers):
            return True
        hits = sum(1 for marker in field_markers if marker in text or marker in text_lower)
        return hits >= 4

    def _looks_like_data_doc(self, file_name: str, content: str = "", file_ext: str = "") -> bool:
        name = (file_name or "").lower()
        text = (content or "").lower()
        ext = (file_ext or "").lower()
        if ext in (".csv", ".tsv", ".xlsx", ".xls", ".numbers", ".json", ".jsonl", ".sql"):
            return True
        name_markers = (
            "dataset", "datasets", "training_files", "validation_files", "schema",
            "inventory", "manifest", "metrics", "benchmark", "filelist",
            "数据集", "清单", "列表", "表格",
        )
        text_markers = ("rows", "columns", "csv", "jsonl", "tsv", "数据字段", "字段说明")
        return any(marker in name for marker in name_markers) or any(marker in text for marker in text_markers)

    def _looks_like_chat_transcript(self, content: str, file_name: str = "", file_ext: str = "") -> bool:
        text = (content or "")
        text_lower = text.lower()
        ext = (file_ext or "").lower()
        if ext and ext not in (".md", ".txt", ".rst", ".pdf", ".doc", ".docx"):
            return False

        user_hits = text.count("用户:") + text_lower.count("user:")
        assistant_hits = (
            text.count("Claude:")
            + text.count("助手:")
            + text.count("Assistant:")
            + text.count("ChatGPT:")
            + text_lower.count("claude:")
            + text_lower.count("assistant:")
            + text_lower.count("chatgpt:")
        )
        if user_hits >= 2 and assistant_hits >= 1:
            return True
        return "上面是我们之前的对话内容" in text

    def _looks_like_report_doc(self, file_name: str, content: str = "", file_ext: str = "") -> bool:
        blob = " ".join([(file_name or "").lower(), (content or "").lower()[:1200]])
        markers = (
            "analysis", "brief", "roadmap", "report", "proposal", "plan",
            "architecture", "summary", "review", "overview", "comparison", "update",
            "分析", "方案", "规划", "路线图", "报告", "总结", "架构", "评估", "解读", "对比", "更新",
        )
        return any(marker in blob for marker in markers)

    def _looks_like_paper(
        self,
        content: str,
        page_count: Optional[int],
        file_name: str = "",
        file_ext: str = "",
    ) -> bool:
        name = (file_name or "").lower()
        ext = (file_ext or "").lower()
        t = (content or "").lower()
        arxiv_like_name = bool(re.search(r"\b\d{4}\.\d{4,5}(?:v\d+)?\b", name))
        if self._looks_like_form_document(file_name, content, file_ext=file_ext):
            return False
        if self._looks_like_chat_transcript(content, file_name=file_name, file_ext=file_ext):
            return False
        strong_report_name_markers = (
            "report", "analysis", "summary", "brief", "roadmap", "proposal", "plan",
            "报告", "分析", "方案", "规划", "路线图", "总结",
        )
        if not arxiv_like_name and any(marker in name for marker in strong_report_name_markers):
            return False
        if ext == ".pdf" and arxiv_like_name:
            return True
        hits = [
            "abstract", "keywords", "references", "arxiv", "doi",
            "摘要", "关键词", "参考文献", "致谢", "引言", "方法", "实验", "结论",
        ]
        if page_count is not None and page_count <= 40 and any(k in t for k in hits):
            return True
        score = sum(1 for k in hits if k in t)
        return score >= 3

    def _looks_like_book(self, file_name: str, content: str, file_ext: str = "", page_count: Optional[int] = None) -> bool:
        name = (file_name or "").lower()
        text = (content or "").lower()
        ext = (file_ext or "").lower()

        if self._looks_like_readme(file_name, file_ext):
            return False
        if self._looks_like_manual_doc(file_name, content, file_ext=file_ext):
            return False
        if self._looks_like_form_document(file_name, content, file_ext=file_ext):
            return False
        if self._looks_like_chat_transcript(content, file_name=file_name, file_ext=file_ext):
            return False
        if self._looks_like_report_doc(file_name, content, file_ext=file_ext):
            return False

        if ext in (".epub", ".mobi", ".azw3"):
            return True

        if self._looks_like_paper(content, page_count, file_name=file_name, file_ext=file_ext):
            return False

        name_hits = ["电子书", "pdf电子书", "kindle", "雅书", "豆瓣"]
        if any(k in name for k in name_hits):
            if page_count is None or page_count >= 40:
                return True

        book_struct_hits = ["isbn", "出版社", "版权页", "目录", "前言", "序言", "译者", "第1章", "第一章", "chapter 1"]
        if any(k in text for k in book_struct_hits):
            if page_count is None or page_count >= 60:
                return True

        if "《" in (file_name or "") and "》" in (file_name or "") and "简历" not in (file_name or ""):
            if page_count is None or page_count >= 60:
                return True

        if ext == ".pdf" and page_count is not None and page_count >= 80:
            return True

        return False

    @staticmethod
    def _sanitize_llm_category(raw_category: str) -> str:
        c = str(raw_category or "").strip()
        if not c:
            return "other"
        c = c.splitlines()[0].strip()
        for prefix in ("类别：", "类别:", "分类：", "分类:", "Category:", "Category：", "Type:", "Type："):
            if c.startswith(prefix):
                c = c[len(prefix):].strip()
        c = c.strip("。，；.,;:：'\"[]【】()（）*` ")
        if not c:
            return "other"
        if len(c) > 15:
            return "other"
        mapped = _normalize_category_en(c, default="other")
        if mapped == "all":
            return "other"
        return mapped

    @staticmethod
    def _extract_json_object(raw_text: str) -> Dict[str, Any]:
        text = str(raw_text or "").strip()
        if not text:
            return {}
        try:
            parsed = json.loads(text)
            return parsed if isinstance(parsed, dict) else {}
        except Exception:
            pass
        match = re.search(r"\{[\s\S]*\}", text)
        if not match:
            match_text = text
        else:
            match_text = match.group(0)
        try:
            parsed = json.loads(match_text)
            return parsed if isinstance(parsed, dict) else {}
        except Exception:
            pass

        # Best-effort salvage for truncated JSON-like responses. This keeps
        # usable taxonomy/summary data instead of forcing retries on near-valid
        # model output.
        partial: Dict[str, Any] = {}

        def _json_string_field(name: str) -> str:
            m = re.search(rf'"{name}"\s*:\s*"((?:\\.|[^"\\])*)"', match_text, flags=re.DOTALL)
            if not m:
                return ""
            try:
                return json.loads(f'"{m.group(1)}"')
            except Exception:
                return m.group(1).replace('\\"', '"').strip()

        for key in ("family", "leaf_category", "role", "summary", "file_name_en"):
            value = _json_string_field(key)
            if value:
                partial[key] = value

        m_conf = re.search(r'"confidence"\s*:\s*([0-9]+(?:\.[0-9]+)?)', match_text)
        if m_conf:
            try:
                partial["confidence"] = float(m_conf.group(1))
            except Exception:
                pass

        if "extracts" not in partial:
            partial["extracts"] = []
        return partial

    @staticmethod
    def _sanitize_taxonomy_leaf(raw_leaf: Any, fallback: str = "other") -> str:
        leaf = str(raw_leaf or "").strip().lower()
        if not leaf:
            return fallback
        leaf = leaf.splitlines()[0].strip()
        leaf = re.sub(r"[^a-z0-9_\-/ ]+", "_", leaf)
        leaf = re.sub(r"[\s\-]+", "_", leaf).strip("_")
        if not leaf:
            return fallback
        if len(leaf) > 48:
            leaf = leaf[:48].strip("_")
        return leaf or fallback

    @staticmethod
    def _sanitize_doc_role(raw_role: Any, fallback: str = "other") -> str:
        role = str(raw_role or "").strip().lower()
        role_map = {
            "primary": "primary_source",
            "primary_source": "primary_source",
            "source": "primary_source",
            "summary": "summary",
            "summarization": "summary",
            "explainer": "explainer",
            "explanation": "explainer",
            "analysis": "analysis",
            "analytical": "analysis",
            "transcript": "transcript",
            "chat_transcript": "transcript",
            "ocr_result": "ocr_result",
            "ocr": "ocr_result",
            "generated_doc": "generated_doc",
            "generated": "generated_doc",
            "reference": "reference",
            "other": "other",
        }
        return role_map.get(role, fallback)

    def _infer_doc_role(
        self,
        *,
        content: str = "",
        file_name: str = "",
        file_ext: str = "",
        family: str = "other",
        leaf_category: str = "",
    ) -> str:
        name = (file_name or "").lower()
        leaf = (leaf_category or "").lower()
        text = (content or "").lower()

        if "ocr" in name or leaf in {"ocr_result", "ocr_extract"}:
            return "ocr_result"
        if self._looks_like_chat_transcript(content, file_name=file_name, file_ext=file_ext):
            return "transcript"
        if any(tok in name for tok in ("summary", "摘要", "总结", "概述")) or leaf.endswith("_summary"):
            return "summary"
        if any(tok in name for tok in ("explainer", "解释", "解读", "含义")) or leaf.endswith("_explainer"):
            return "explainer"
        if any(tok in name for tok in ("analysis", "分析", "评估", "review", "comparison", "对比")) or leaf.endswith("_analysis"):
            return "analysis"
        if any(tok in text[:400] for tok in ("generated by", "ai generated", "由ai生成", "整理如下")):
            return "generated_doc"
        if family in {"paper", "contract", "resume", "manual", "book", "invoice", "quotation", "presentation", "data"}:
            return "primary_source"
        if family in {"document", "report", "note", "email"}:
            return "reference"
        return "other"

    def _classify_document_taxonomy(
        self,
        content: str,
        *,
        file_name: str = "",
        file_ext: str = "",
        page_count: Optional[int] = None,
    ) -> Dict[str, Any]:
        """
        Return a layered taxonomy:
          - family: stable retrieval family
          - leaf_category: dynamic fine-grained label
          - role: semantic role such as primary_source / summary / analysis
          - confidence: model confidence if available
        """
        fallback_family = self._guard_category_by_extension(
            file_ext,
            self._classify_document(content, file_name=file_name, file_ext=file_ext, page_count=page_count),
        )
        fallback_leaf = self._sanitize_taxonomy_leaf(fallback_family, fallback=fallback_family)
        fallback_role = self._infer_doc_role(
            content=content,
            file_name=file_name,
            file_ext=file_ext,
            family=fallback_family,
            leaf_category=fallback_leaf,
        )

        prompt = get_prompt("CLASSIFY_TAXONOMY_PROMPT", "en").format(
            file_name=file_name,
            file_ext=file_ext,
            page_count=("" if page_count is None else str(page_count)),
            content=(content or "")[:2000],
        )
        idx_model = self._require_configured_index_model_id("taxonomy 分类")
        if not idx_model:
            parsed = {}
        else:
            prompt = self._append_model_prompt_suffix(prompt, idx_model)

            try:
                client = self._get_local_llm_client()
                response = client.chat.completions.create(
                    model=idx_model,
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=180,
                    temperature=0.0,
                    stream=False,
                )
                parsed = self._extract_json_object(response.choices[0].message.content)
            except Exception as e:
                self._index_exception("taxonomy 分类失败", e)
                parsed = {}

        family = self._guard_category_by_extension(
            file_ext,
            _normalize_category_en(parsed.get("family", fallback_family), default=fallback_family),
        )
        if family in {"", "all", "unknown"}:
            family = fallback_family

        leaf_category = self._sanitize_taxonomy_leaf(parsed.get("leaf_category"), fallback=fallback_leaf)
        role = self._sanitize_doc_role(
            parsed.get("role"),
            fallback=self._infer_doc_role(
                content=content,
                file_name=file_name,
                file_ext=file_ext,
                family=family,
                leaf_category=leaf_category,
            ),
        )

        try:
            confidence = float(parsed.get("confidence", 0.0) or 0.0)
        except Exception:
            confidence = 0.0
        confidence = min(max(confidence, 0.0), 1.0)

        # Keep family stable, but let leaf_category grow organically.
        if not leaf_category:
            leaf_category = fallback_leaf

        return {
            "family": family,
            "leaf_category": leaf_category,
            "role": role,
            "confidence": confidence,
        }

    _EXT_CATEGORY_ALLOW: Dict[str, frozenset] = {
        ".pdf":  frozenset({"document", "report", "paper", "book", "manual", "contract", "invoice",
                             "resume", "form", "presentation", "other"}),
        ".docx": frozenset({"document", "report", "paper", "book", "manual", "contract", "letter",
                             "resume", "form", "presentation", "other"}),
        ".doc":  frozenset({"document", "report", "paper", "book", "manual", "contract", "letter",
                             "resume", "form", "presentation", "other"}),
        ".pptx": frozenset({"presentation", "document", "report", "other"}),
        ".ppt":  frozenset({"presentation", "document", "report", "other"}),
        ".md":   frozenset({"manual", "document", "report", "paper", "other"}),
        ".txt":  frozenset({"manual", "document", "other"}),
        ".rst":  frozenset({"manual", "document", "other"}),
        ".jpg":  frozenset({"image", "photo", "other"}),
        ".jpeg": frozenset({"image", "photo", "other"}),
        ".png":  frozenset({"image", "photo", "other"}),
        ".gif":  frozenset({"image", "other"}),
        ".bmp":  frozenset({"image", "other"}),
        ".webp": frozenset({"image", "other"}),
        ".heic": frozenset({"image", "photo", "other"}),
        ".tiff": frozenset({"image", "other"}),
        ".csv":  frozenset({"data", "other"}),
        ".xlsx": frozenset({"data", "other"}),
        ".xls":  frozenset({"data", "other"}),
    }
    _EXT_CATEGORY_DEFAULT: Dict[str, str] = {
        ".pdf": "document",  ".docx": "document", ".doc": "document",
        ".pptx": "presentation", ".ppt": "presentation",
        ".md": "manual",   ".txt": "manual",   ".rst": "manual",
        ".jpg": "image",   ".jpeg": "image",   ".png": "image",
        ".gif": "image",   ".bmp": "image",    ".webp": "image",
        ".heic": "image",  ".tiff": "image",
        ".csv": "data",    ".xlsx": "data",    ".xls": "data",
    }

    def _guard_category_by_extension(self, ext: str, category: str) -> str:
        ext = (ext or "").lower()
        if not ext:
            return category

        allow = self._EXT_CATEGORY_ALLOW.get(ext)
        if allow is None:
            return category

        if category in allow:
            return category


        _DOC_EXTS = {".pdf", ".docx", ".doc", ".md", ".txt", ".rst", ".pptx", ".ppt"}
        _MEDIA_ONLY_CATS = {"image", "photo", "audio/video", "audio", "video"}

        if ext in _DOC_EXTS and category in _MEDIA_ONLY_CATS:
            default = self._EXT_CATEGORY_DEFAULT.get(ext, "other")
            logger.warning(
                f"[category_guard] ext={ext} LLM_category='{category}' 跨媒体类型越界 → 回退为 '{default}'"
            )
            return default

        _IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".heic", ".tiff"}
        _DOC_ONLY_CATS = {"document", "report", "paper", "manual", "contract", "resume", "book",
                          "presentation", "invoice", "code", "data", "email", "note", "quotation"}

        if ext in _IMAGE_EXTS and category in _DOC_ONLY_CATS:
            default = self._EXT_CATEGORY_DEFAULT.get(ext, "other")
            logger.warning(
                f"[category_guard] ext={ext} LLM_category='{category}' 跨媒体类型越界 → 回退为 '{default}'"
            )
            return default

        _DATA_EXTS = {".csv", ".tsv", ".xlsx", ".xls", ".numbers", ".json", ".jsonl"}
        _NON_DATA_CATS = {
            "document", "report", "paper", "manual", "contract", "resume",
            "book", "presentation", "invoice", "quotation", "email",
            "image", "photo", "audio/video", "audio", "video", "code",
        }
        if ext in _DATA_EXTS and category in _NON_DATA_CATS:
            default = self._EXT_CATEGORY_DEFAULT.get(ext, "data")
            logger.warning(
                f"[category_guard] ext={ext} LLM_category='{category}' 偏离结构化数据家族 → 回退为 '{default}'"
            )
            return default

        logger.info(f"[category_guard] ext={ext} 接受 LLM 自定义类别 '{category}'（开放式分类）")
        return category

    def _infer_doc_category_family(
        self,
        raw_category: str,
        *,
        content: str = "",
        file_name: str = "",
        file_ext: str = "",
        page_count: Optional[int] = None,
    ) -> str:
        """
        Derive a stable retrieval family from an open-ended raw category.

        `doc_category` should be stable enough for filtering/retrieval, while we
        still preserve the raw/open label separately for analysis and future
        refinement.
        """
        canonical_categories = {
            "resume", "report", "contract", "note", "manual", "paper",
            "presentation", "data", "email", "image", "audio", "video", "audio/video",
            "book", "code", "invoice", "quotation", "document", "other",
        }
        raw = str(raw_category or "").strip()
        normalized = _normalize_category_en(raw, default="")
        ext = (file_ext or "").lower()
        blob = " ".join([raw.lower(), (file_name or "").lower()])

        if ext in {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".heic", ".tiff", ".svg"}:
            return "image"
        if ext in {".mp3", ".wav", ".flac", ".aac", ".ogg", ".m4a", ".wma", ".aiff", ".ape"}:
            return "audio"
        if ext in {".mp4", ".m4v", ".avi", ".mov", ".mkv", ".webm", ".flv", ".wmv", ".ts"}:
            return "video"
        if ext in {".csv", ".tsv", ".xlsx", ".xls", ".numbers", ".json", ".jsonl"}:
            return "data"
        if ext in {".ppt", ".pptx", ".key"}:
            return "presentation"
        if ext in {".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".cpp", ".c", ".h", ".hpp", ".go", ".rs", ".swift", ".kt", ".sql"}:
            return "code"

        if self._looks_like_manual_doc(file_name, content, file_ext=file_ext):
            return "manual"
        if self._looks_like_form_document(file_name, content, file_ext=file_ext):
            return "document"
        if self._looks_like_data_doc(file_name, content, file_ext=file_ext):
            return "data"
        if self._looks_like_chat_transcript(content, file_name=file_name, file_ext=file_ext):
            if self._looks_like_report_doc(file_name, content, file_ext=file_ext):
                return "report"
            return "note"
        if self._looks_like_paper(content, page_count, file_name=file_name, file_ext=file_ext):
            return "paper"
        if self._looks_like_book(file_name, content, file_ext=file_ext, page_count=page_count):
            return "book"
        if self._looks_like_report_doc(file_name, content, file_ext=file_ext):
            return "report"

        manual_markers = (
            "readme", "faq", "datasheet", "data sheet", "spec", "specification",
            "guide", "manual", "tutorial", "howto", "how-to", "install",
            "diagram", "architecture", "technical", "技术", "架构", "说明", "指南",
        )
        data_markers = (
            "json", "schema", "table", "dataset", "graph", "chart", "log",
            "trace", "csv", "统计", "数据", "表格",
        )
        report_markers = (
            "analysis", "brief", "roadmap", "report", "proposal", "plan",
            "总结", "分析", "方案", "规划", "路线图", "报告",
        )
        note_markers = ("memo", "note", "conversation", "备忘", "笔记", "对话")

        if any(marker in blob for marker in manual_markers):
            return "manual"
        if any(marker in blob for marker in data_markers):
            return "data"
        if any(marker in blob for marker in report_markers):
            return "report"
        if any(marker in blob for marker in note_markers):
            return "note"
        if normalized in canonical_categories:
            return normalized

        if ext in {".md", ".txt", ".rst", ".pdf", ".doc", ".docx"}:
            return "document"
        return "other"

    @staticmethod
    def _meta_category_family(meta: Dict[str, Any]) -> str:
        media_type = str(meta.get("media_type") or "").strip().lower()
        raw = meta.get("doc_category_family") or meta.get("doc_category") or "other"
        normalized = _normalize_category_en(raw, default="other")
        if normalized == "audio/video" and media_type in {"audio", "video"}:
            return media_type
        return normalized

    @staticmethod
    def _meta_category_leaf(meta: Dict[str, Any]) -> str:
        media_type = str(meta.get("media_type") or "").strip().lower()
        leaf = str(
            meta.get("doc_category_leaf")
            or meta.get("doc_category_raw")
            or meta.get("doc_category_family")
            or meta.get("doc_category")
            or "other"
        ).strip()
        if _normalize_category_en(leaf, default="other") == "audio/video" and media_type in {"audio", "video"}:
            return media_type
        return leaf or "other"

    @staticmethod
    def _meta_doc_role(meta: Dict[str, Any]) -> str:
        role = str(meta.get("doc_role") or "").strip().lower()
        if role:
            return role
        return "primary_source"

    def _role_score_multiplier(self, *, target_category: str = "", meta: Optional[Dict[str, Any]] = None) -> float:
        """
        Light role-aware reranking layer.

        Fixes the modeling mismatch where a user searches for a family (for
        example papers), but summary/explainer documents about those files rank
        above the original source documents.
        """
        if not meta:
            return 1.0
        target = _normalize_category_en(target_category or "", default="")
        role = self._meta_doc_role(meta)
        leaf = self._meta_category_leaf(meta)

        if target == "paper":
            if role == "primary_source":
                return 1.12
            if role in {"summary", "explainer", "analysis", "generated_doc"}:
                return 0.72
            if leaf.endswith("_summary") or leaf.endswith("_explainer") or leaf.endswith("_analysis"):
                return 0.72
            if role in {"transcript", "ocr_result"}:
                return 0.8

        if target in {"report", "document"}:
            if role == "primary_source":
                return 1.05
            if role in {"summary", "explainer"}:
                return 0.85

        return 1.0

    def _classify_document(self, content: str, file_name: str = "", file_ext: str = "", page_count: Optional[int] = None) -> str:
        ext = (file_ext or "").lower()
        name_lower = (file_name or "").lower()
        
        if self._looks_like_manual_doc(file_name, content, file_ext=file_ext):
            return "manual"
        if ext in (".ppt", ".pptx", ".key"):
            return "presentation"
        if ext in (".csv", ".xls", ".xlsx", ".sql", ".tsv"):
            return "data"
        if self._looks_like_form_document(file_name, content, file_ext=file_ext):
            return "document"
        if self._looks_like_data_doc(file_name, content, file_ext=file_ext):
            return "data"
        if self._looks_like_chat_transcript(content, file_name=file_name, file_ext=file_ext):
            if self._looks_like_report_doc(file_name, content, file_ext=file_ext):
                return "report"
            return "note"
            
        if "简历" in name_lower or "cv" in name_lower or "resume" in name_lower:
            return "resume"
        if "合同" in name_lower or "协议" in name_lower:
            return "contract"
        # ------------------------------------
        
        try:
            if self._looks_like_book(file_name, content, file_ext=file_ext, page_count=page_count):
                return "book"
            if self._looks_like_paper(content, page_count, file_name=file_name, file_ext=file_ext):
                return "paper"
            if self._looks_like_report_doc(file_name, content, file_ext=file_ext):
                return "report"

            client = self._get_local_llm_client()
            prompt = get_prompt("CLASSIFY_PROMPT", "en").format(
                file_name=file_name,
                file_ext=file_ext,
                page_count=("" if page_count is None else str(page_count)),
                content=content[:2000],
            )
            idx_model = self._require_configured_index_model_id("文档分类")
            if not idx_model:
                return "other"
            prompt = self._append_model_prompt_suffix(prompt, idx_model)

            response = client.chat.completions.create(
                model=idx_model,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=20,
                temperature=0.0,
                stream=False,
            )
            
            category = self._sanitize_llm_category(response.choices[0].message.content)
            
            if category == "resume" and self._looks_like_book(file_name, content, file_ext=file_ext, page_count=page_count):
                return "book"
            if category == "book":
                if self._looks_like_manual_doc(file_name, content, file_ext=file_ext):
                    return "manual"
                if self._looks_like_form_document(file_name, content, file_ext=file_ext):
                    return "document"
                if self._looks_like_data_doc(file_name, content, file_ext=file_ext):
                    return "data"
                if self._looks_like_chat_transcript(content, file_name=file_name, file_ext=file_ext):
                    if self._looks_like_report_doc(file_name, content, file_ext=file_ext):
                        return "report"
                    return "note"
                if self._looks_like_report_doc(file_name, content, file_ext=file_ext):
                    return "report"
                if self._looks_like_paper(content, page_count, file_name=file_name, file_ext=file_ext):
                    return "paper"
            if category == "paper":
                if self._looks_like_form_document(file_name, content, file_ext=file_ext):
                    return "document"
                if self._looks_like_chat_transcript(content, file_name=file_name, file_ext=file_ext):
                    if self._looks_like_report_doc(file_name, content, file_ext=file_ext):
                        return "report"
                    return "note"
                if self._looks_like_manual_doc(file_name, content, file_ext=file_ext):
                    return "manual"
            if category == "manual":
                if self._looks_like_paper(content, page_count, file_name=file_name, file_ext=file_ext):
                    return "paper"
                if self._looks_like_data_doc(file_name, content, file_ext=file_ext):
                    return "data"
                if (
                    self._looks_like_report_doc(file_name, content, file_ext=file_ext)
                    and not self._looks_like_manual_doc(file_name, content, file_ext=file_ext)
                ):
                    return "report"
            if category == "report":
                if self._looks_like_paper(content, page_count, file_name=file_name, file_ext=file_ext):
                    return "paper"
                if self._looks_like_form_document(file_name, content, file_ext=file_ext):
                    return "document"
            if category == "other":
                if self._looks_like_data_doc(file_name, content, file_ext=file_ext):
                    return "data"
                if self._looks_like_report_doc(file_name, content, file_ext=file_ext):
                    return "report"
            return category
        except Exception as e:
            self._index_exception("分类失败", e)
            return "other"

    def _doc_summary_utc_now(self) -> str:
        return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")

    def _snapshot_loaded_llm_model_id(self) -> str:
        try:
            from services.local_llm import get_local_llm_manager

            mid = getattr(get_local_llm_manager(), "current_model_id", None) or ""
            mid = str(mid).strip()
            if mid:
                return mid
        except Exception:
            pass
        try:
            from services.preference_manager import PreferenceManager
            import config.settings as agent_settings

            base_dir = getattr(
                agent_settings,
                "BASE_DIR",
                os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
            )
            pm = PreferenceManager(base_dir)
            mid = (pm.get_selected_index_model_id() or "").strip()
            return mid
        except Exception:
            return ""

    def _release_local_llm_before_media_embedding(self, reason: str = "") -> None:
        try:
            from services.local_llm import get_local_llm_manager

            llm_mgr = get_local_llm_manager()
            current_id = str(getattr(llm_mgr, "current_model_id", None) or "").strip()
            if not current_id:
                return
            logger.info(
                f"[MediaExpert] releasing local LLM before media embedding: "
                f"model_id={current_id} reason={reason or 'n/a'}"
            )
            llm_mgr.unload_model()
            try:
                time.sleep(0.2)
            except Exception:
                pass
        except Exception as e:
            logger.warning(f"[MediaExpert] failed to release local LLM before media embedding: {e}")

    def _get_configured_index_model_id(self) -> str:
        try:
            from services.preference_manager import PreferenceManager
            import config.settings as agent_settings
            base_dir = getattr(
                agent_settings,
                "BASE_DIR",
                os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
            )
            pm = PreferenceManager(base_dir)
            mid = (pm.get_selected_index_model_id() or "").strip()
            if mid:
                return mid
        except Exception:
            pass
        return ""

    def _require_configured_index_model_id(self, stage: str = "index_llm") -> str:
        """Return the Add Sources index model id, or skip index-side LLM/VL work."""
        mid = self._get_configured_index_model_id()
        if not mid:
            self._index_warning(
                f"{stage}: 未配置 Add Sources 索引模型，跳过本次索引侧 LLM/VL 调用，避免串用 Chat 模型"
            )
        return mid

    def _set_doc_summary_provenance(self, metadata: Dict[str, Any], model_id: str) -> None:
        mid = (model_id or "").strip() or "none"
        metadata["doc_summary_model_id"] = mid[:240]
        metadata["doc_summary_saved_at"] = self._doc_summary_utc_now()

    _NON_EN_RATIO_THRESHOLD = 0.20
    _SUMMARY_MAX_RETRIES    = 3
    _INDEX_EN_STOPWORDS = frozenset({
        "a", "an", "and", "as", "at", "by", "for", "from", "in", "is", "it",
        "of", "on", "or", "that", "the", "this", "to", "with", "image",
        "icon", "logo", "photo", "chart", "table", "document", "screenshot",
        "shows", "contains",
    })
    _INDEX_NON_EN_STOPWORDS = frozenset({
        # Spanish
        "aqui", "aquí", "claro", "con", "de", "del", "el", "en", "es", "esta",
        "este", "imagen", "la", "las", "los", "para", "por", "que", "resumen",
        "se", "un", "una", "y",
        # French
        "avec", "bien", "bonjour", "cette", "ceci", "de", "des", "est", "image",
        "la", "le", "les", "pour", "sur", "une", "voici",
        # Portuguese
        "com", "da", "das", "de", "do", "dos", "esta", "imagem", "para", "que",
        "resumo", "uma",
    })

    @staticmethod
    def _non_english_ratio(text: str) -> float:
        alpha_chars = [c for c in text if c.isalpha()]
        if not alpha_chars:
            return 0.0
        non_en = sum(1 for c in alpha_chars if ord(c) > 127)
        return non_en / len(alpha_chars)

    @staticmethod
    def _label_needs_translation(text: str) -> bool:
        raw = str(text or "").strip()
        if not raw:
            return False
        return any(ch.isalpha() and ord(ch) > 127 for ch in raw)

    def _generate_summary(self, content: str, file_name: str = "") -> Tuple[str, str]:
        import re as _re

        stripped_content = content.strip()
        if len(stripped_content) < 30:
            text = self._ensure_english_index_text(stripped_content.replace("\n", " "), max_len=600)
            return text, "passthrough_no_llm"

        try:
            client = self._get_local_llm_client()

            model_id = self._require_configured_index_model_id("文档摘要")
            if not model_id:
                return "", "missing_index_model"

            # All models use the bilingual SUMMARY_PROMPT_INDEXING which:
            # - Writes primarily in English (for English query retrieval)
            # - Preserves Chinese proper nouns (person/company/product names in original)
            # - Extracts years of experience for resumes and quantitative metrics for reports
            # This is optimal for Chinese-content + English-query cross-lingual retrieval.
            prompt_name = "SUMMARY_PROMPT_INDEXING"

            prompt = get_prompt(prompt_name, "en").format(content=content[:1500])
            
            prompt = self._append_model_prompt_suffix(prompt, model_id)

            self._index_info(f"[_generate_summary] model={model_id} prompt_len={len(prompt)}")

            summary = ""
            used_model = "unknown"

            for attempt in range(1, self._SUMMARY_MAX_RETRIES + 1):
                temp = 0.0 if attempt == 1 else round(0.1 * (attempt - 1), 1)

                response = client.chat.completions.create(
                    model=model_id,
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=600,
                    temperature=temp,
                    stream=False,
                )

                raw = (response.choices[0].message.content or "").strip()
                raw = _re.sub(r"^(?:\*\*)?Summary:(?:\*\*)?\s*", "", raw, flags=_re.IGNORECASE).strip()

                if not raw:
                    self._index_warning(
                        f"[_generate_summary] attempt={attempt} LLM returned empty, retrying..."
                    )
                    continue

                ratio = self._non_english_ratio(raw)
                if ratio <= self._NON_EN_RATIO_THRESHOLD:
                    summary = raw[:600]
                    used_model = self._snapshot_loaded_llm_model_id() or "unknown"
                    if attempt > 1:
                        self._index_info(
                            f"[_generate_summary] quality OK after {attempt} attempts "
                            f"(non_en_ratio={ratio:.2%})"
                        )
                    break
                else:
                    self._index_warning(
                        f"[_generate_summary] attempt={attempt}/{self._SUMMARY_MAX_RETRIES} "
                        f"non_en_ratio={ratio:.2%} > {self._NON_EN_RATIO_THRESHOLD:.0%} threshold, "
                        f"retrying (temp={temp})... preview: {raw[:80]!r}"
                    )

            else:
                self._index_warning(
                    f"[_generate_summary] FAILED after {self._SUMMARY_MAX_RETRIES} retries — "
                    f"non-English ratio consistently exceeded {self._NON_EN_RATIO_THRESHOLD:.0%}. "
                    f"model={model_id} content_preview={content[:60]!r}. "
                    f"Summary will NOT be indexed."
                )
                return "", "error_non_english"

            return summary, used_model

        except Exception as e:
            self._index_exception("摘要生成失败", e)
            return "", "error"

    def _generate_summary_and_extract(self, content: str, file_name: str = "") -> Tuple[str, str, List[Dict[str, Any]], str]:
        import re as _re
        import json

        stripped_content = content.strip()
        if len(stripped_content) < 30:
            text = self._ensure_english_index_text(stripped_content.replace("\n", " "), max_len=600)
            return text, "", [], "passthrough_no_llm"

        try:
            client = self._get_local_llm_client()
            model_id = self._require_configured_index_model_id("文档摘要与信息提取")
            if not model_id:
                return "", "", [], "missing_index_model"

            prompt_name = "SUMMARY_AND_EXTRACT_PROMPT"
            prompt = get_prompt(prompt_name, "en").format(content=content[:1500])
            
            prompt = self._append_model_prompt_suffix(prompt, model_id)

            self._index_info(f"[_generate_summary_and_extract] model={model_id} prompt_len={len(prompt)}")

            summary = ""
            extracts = []
            used_model = "unknown"

            for attempt in range(1, self._SUMMARY_MAX_RETRIES + 1):
                temp = 0.0 if attempt == 1 else round(0.1 * (attempt - 1), 1)

                response = client.chat.completions.create(
                    model=model_id,
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=1000,
                    temperature=temp,
                    stream=False,
                )

                raw = (response.choices[0].message.content or "").strip()
                
                if raw.startswith("```"):
                    raw = _re.sub(r"^```(?:json)?\s*", "", raw)
                    raw = _re.sub(r"\s*```$", "", raw).strip()
                
                try:
                    parsed = json.loads(raw)
                    raw_summary = parsed.get("summary", "")
                    raw_extracts = parsed.get("extracts", [])
                    summary_str = str(raw_summary).strip()
                except Exception as parse_e:
                    self._index_warning(f"[_generate_summary_and_extract] JSON parse failed, retrying... {parse_e}")
                    continue
                
                if not summary_str:
                    self._index_warning(
                        f"[_generate_summary_and_extract] attempt={attempt} LLM returned empty summary, retrying..."
                    )
                    continue

                ratio = self._non_english_ratio(summary_str)
                if ratio <= self._NON_EN_RATIO_THRESHOLD:
                    summary = summary_str[:600]
                    if isinstance(raw_extracts, list):
                        extracts = [e for e in raw_extracts if isinstance(e, dict) and 'content' in e and 'type' in e]
                    
                    if extracts:
                        contact_subtypes = {"phone", "email", "address", "social_media", "contact"}
                        found_cnt_types = set(e["type"].lower() for e in extracts if isinstance(e.get("type"), str) and e["type"].lower() in contact_subtypes)
                        if found_cnt_types:
                            summary += f" [Contact Information (联系方式): {', '.join(found_cnt_types)}]"
                    
                    used_model = self._snapshot_loaded_llm_model_id() or "unknown"
                    if attempt > 1:
                        self._index_info(f"[_generate_summary_and_extract] quality OK after {attempt} attempts")
                        
                    if extracts:
                        self._index_info(f"[_generate_summary_and_extract] extracted {len(extracts)} items.")
                    break
                else:
                    self._index_warning(
                        f"[_generate_summary_and_extract] attempt={attempt}/{self._SUMMARY_MAX_RETRIES} "
                        f"non_en_ratio={ratio:.2%} > threshold, retrying... preview: {summary_str[:80]!r}"
                    )

            else:
                self._index_warning(
                    f"[_generate_summary_and_extract] FAILED after {self._SUMMARY_MAX_RETRIES} retries. "
                    f"model={model_id}. Summary will NOT be indexed."
                )
                return "", "", [], "error_format_or_english"

            return summary, "", extracts, used_model

        except Exception as e:
            self._index_exception("摘要与信息提取失败", e)
            return "", "", [], "error"

    def _classify_and_summarize_unified(
        self,
        content: str,
        *,
        file_name: str = "",
        file_ext: str = "",
        page_count: Optional[int] = None,
    ) -> Dict[str, Any]:
        """Single LLM pass for text-file taxonomy, summary, and extracts."""
        import re as _re

        stripped_content = (content or "").strip()
        fallback_family = self._guard_category_by_extension(
            file_ext,
            self._infer_doc_category_family(
                stripped_content,
                file_name=file_name,
                file_ext=file_ext,
                page_count=page_count,
            ),
        )
        fallback_leaf = self._sanitize_taxonomy_leaf(fallback_family, fallback=fallback_family)
        fallback_role = self._infer_doc_role(
            content=stripped_content,
            file_name=file_name,
            file_ext=file_ext,
            family=fallback_family,
            leaf_category=fallback_leaf,
        )

        def _fallback_result(summary: str = "", model_id: str = "unified_fallback") -> Dict[str, Any]:
            base_summary = str(summary or "").strip()
            if not base_summary:
                name_hint = os.path.splitext(file_name or "")[0].strip() or "document"
                base_summary = f"{fallback_family.title()} file {name_hint}. Extracted text available for retrieval."
            return {
                "taxonomy": {
                    "family": fallback_family,
                    "leaf_category": fallback_leaf,
                    "role": fallback_role,
                    "confidence": 0.0,
                },
                "summary": base_summary[:700],
                "file_name_en": "",
                "extracts": [],
                "model": model_id,
            }

        if len(stripped_content) < 30:
            summary = self._ensure_english_index_text(stripped_content.replace("\n", " "), max_len=600)
            return _fallback_result(summary=summary, model_id="passthrough_no_llm")

        model_id = self._require_configured_index_model_id("unified 文档分类摘要")
        if not model_id:
            return _fallback_result(model_id="missing_index_model")

        ext_lower = (file_ext or "").lower()
        try:
            ultra_short_doc_chars = int(os.getenv("FILEAGENT_UNIFIED_ULTRA_SHORT_DOC_CHARS", "1200") or 1200)
        except Exception:
            ultra_short_doc_chars = 1200
        ultra_short_doc_chars = max(400, min(ultra_short_doc_chars, 2400))
        try:
            short_doc_chars = int(os.getenv("FILEAGENT_UNIFIED_SHORT_DOC_CHARS", "3200") or 3200)
        except Exception:
            short_doc_chars = 3200
        short_doc_chars = max(ultra_short_doc_chars + 200, min(short_doc_chars, 6000))
        try:
            medium_doc_chars = int(os.getenv("FILEAGENT_UNIFIED_MEDIUM_DOC_CHARS", "6200") or 6200)
        except Exception:
            medium_doc_chars = 6200
        medium_doc_chars = max(short_doc_chars + 200, min(medium_doc_chars, 12000))
        ultra_short_doc_exts = {
            ".md", ".txt", ".rst",
        }
        short_doc_exts = {
            ".md", ".txt", ".rst", ".pdf", ".doc", ".docx", ".ppt", ".pptx", ".key",
        }
        tabular_doc_exts = {
            ".csv", ".tsv", ".xlsx", ".xls", ".numbers", ".json", ".jsonl",
        }
        tabular_doc_mode = ext_lower in tabular_doc_exts
        tabular_large_doc_mode = (
            tabular_doc_mode
            and len(stripped_content) > medium_doc_chars
        )
        ultra_short_doc_mode = (
            ext_lower in ultra_short_doc_exts
            and len(stripped_content) <= ultra_short_doc_chars
            and (page_count is None or int(page_count or 0) <= 8)
        )
        short_doc_mode = (
            not ultra_short_doc_mode
            and not tabular_doc_mode
            and
            ext_lower in short_doc_exts
            and len(stripped_content) <= short_doc_chars
            and (page_count is None or int(page_count or 0) <= 24)
        )
        medium_doc_mode = (
            not ultra_short_doc_mode
            and not tabular_doc_mode
            and
            ext_lower in short_doc_exts
            and not short_doc_mode
            and len(stripped_content) <= medium_doc_chars
            and (page_count is None or int(page_count or 0) <= 48)
        )

        if ultra_short_doc_mode:
            default_content_chars = 900
            content_chars = _env_int_first(
                [
                    "FILEAGENT_UNIFIED_ULTRA_SHORT_CONTENT_CHARS",
                    "FILEAGENT_UNIFIED_INDEX_CONTENT_CHARS",
                ],
                default_content_chars,
            )
        elif tabular_large_doc_mode:
            default_content_chars = 1800
            content_chars = _env_int_first(
                [
                    "FILEAGENT_UNIFIED_TABULAR_LARGE_CONTENT_CHARS",
                    "FILEAGENT_UNIFIED_TABULAR_CONTENT_CHARS",
                    "FILEAGENT_UNIFIED_INDEX_CONTENT_CHARS",
                ],
                default_content_chars,
            )
        elif tabular_doc_mode:
            default_content_chars = 1200
            content_chars = _env_int_first(
                [
                    "FILEAGENT_UNIFIED_TABULAR_CONTENT_CHARS",
                    "FILEAGENT_UNIFIED_INDEX_CONTENT_CHARS",
                ],
                default_content_chars,
            )
        elif short_doc_mode:
            default_content_chars = 1500
            content_chars = _env_int_first(
                [
                    "FILEAGENT_UNIFIED_SHORT_CONTENT_CHARS",
                    "FILEAGENT_UNIFIED_INDEX_CONTENT_CHARS",
                ],
                default_content_chars,
            )
        elif medium_doc_mode:
            default_content_chars = 1800
            content_chars = _env_int_first(
                [
                    "FILEAGENT_UNIFIED_MEDIUM_CONTENT_CHARS",
                    "FILEAGENT_UNIFIED_INDEX_CONTENT_CHARS",
                ],
                default_content_chars,
            )
        else:
            default_content_chars = 2000
            content_chars = _env_int_first(
                [
                    "FILEAGENT_UNIFIED_FULL_CONTENT_CHARS",
                    "FILEAGENT_UNIFIED_INDEX_CONTENT_CHARS",
                ],
                default_content_chars,
            )
        if ultra_short_doc_mode:
            min_content_chars = 500
        elif tabular_large_doc_mode:
            min_content_chars = 1000
        elif tabular_doc_mode:
            min_content_chars = 700
        elif short_doc_mode:
            min_content_chars = 900
        else:
            min_content_chars = 1200 if medium_doc_mode else 1200
        content_chars = max(min_content_chars, min(content_chars, 5000))
        prompt_content = stripped_content[:content_chars]

        if ultra_short_doc_mode:
            default_max_tokens = 360
            response_max_tokens = _env_int_first(
                [
                    "FILEAGENT_UNIFIED_ULTRA_SHORT_MAX_TOKENS",
                    "FILEAGENT_UNIFIED_INDEX_MAX_TOKENS",
                ],
                default_max_tokens,
            )
        elif tabular_large_doc_mode:
            default_max_tokens = 560
            response_max_tokens = _env_int_first(
                [
                    "FILEAGENT_UNIFIED_TABULAR_LARGE_MAX_TOKENS",
                    "FILEAGENT_UNIFIED_TABULAR_MAX_TOKENS",
                    "FILEAGENT_UNIFIED_INDEX_MAX_TOKENS",
                ],
                default_max_tokens,
            )
        elif tabular_doc_mode:
            default_max_tokens = 420
            response_max_tokens = _env_int_first(
                [
                    "FILEAGENT_UNIFIED_TABULAR_MAX_TOKENS",
                    "FILEAGENT_UNIFIED_INDEX_MAX_TOKENS",
                ],
                default_max_tokens,
            )
        elif short_doc_mode:
            default_max_tokens = 460
            response_max_tokens = _env_int_first(
                [
                    "FILEAGENT_UNIFIED_SHORT_MAX_TOKENS",
                    "FILEAGENT_UNIFIED_INDEX_MAX_TOKENS",
                ],
                default_max_tokens,
            )
        elif medium_doc_mode:
            default_max_tokens = 680
            response_max_tokens = _env_int_first(
                [
                    "FILEAGENT_UNIFIED_MEDIUM_MAX_TOKENS",
                    "FILEAGENT_UNIFIED_INDEX_MAX_TOKENS",
                ],
                default_max_tokens,
            )
        else:
            default_max_tokens = 960
            response_max_tokens = _env_int_first(
                [
                    "FILEAGENT_UNIFIED_FULL_MAX_TOKENS",
                    "FILEAGENT_UNIFIED_INDEX_MAX_TOKENS",
                ],
                default_max_tokens,
            )
        response_max_tokens = max(192, min(response_max_tokens, 1200))

        if ultra_short_doc_mode:
            prompt = f"""You are building compact retrieval metadata for a short local file.
Return ONE strict JSON object. No markdown fences or extra text.

Use these hints unless the content clearly contradicts them:
- family_hint: {fallback_family}
- role_hint: {fallback_role}

Rules:
- family MUST be one of ["resume", "report", "contract", "note", "manual", "paper", "presentation", "data", "email", "image", "audio", "video", "book", "code", "invoice", "quotation", "document", "other"].
- leaf_category should be a short snake_case label.
- role MUST be one of ["primary_source", "summary", "explainer", "analysis", "transcript", "ocr_result", "generated_doc", "reference", "other"].
- confidence must be a float between 0 and 1.
- summary MUST be one English-first paragraph, no markdown, <= 80 words.
- Include the main subject plus key names, identifiers, numbers, or contact clues a user may search.
- Preserve Chinese proper nouns when useful.
- extracts should be [] unless obvious personal/contact data is present. Keep at most 4 extracts.

File:
- name: {file_name}
- ext: {file_ext}
- pages: {"" if page_count is None else page_count}

Content excerpt:
{prompt_content}

JSON:
{{"family":"document","leaf_category":"short_note","role":"primary_source","confidence":0.9,"summary":"...","extracts":[]}}
"""
        elif tabular_doc_mode:
            if tabular_large_doc_mode:
                prompt = f"""You are building retrieval metadata for a large tabular data file.
Return ONE strict JSON object. No markdown fences or extra text.

Use these hints unless the content clearly contradicts them:
- family_hint: data
- role_hint: primary_source

Rules:
- family MUST be "data" unless the content clearly shows another family.
- leaf_category should be a short snake_case label such as "recruiting_sheet", "metrics_table", "inventory_export", "experiment_results".
- role MUST be one of ["primary_source", "summary", "explainer", "analysis", "transcript", "ocr_result", "generated_doc", "reference", "other"].
- confidence must be a float between 0 and 1.
- summary MUST be one English-first paragraph, no markdown, <= 120 words.
- The excerpt may be partial. Prioritize sheet names, visible header fields, repeated measures, entities, dates, locations, IDs, and natural search keywords.
- Preserve Chinese proper nouns when useful.
- extracts should usually be [] unless clear personal/contact data is present. Keep at most 4 extracts.

File:
- name: {file_name}
- ext: {file_ext}
- pages: {"" if page_count is None else page_count}

Content excerpt:
{prompt_content}

JSON:
{{"family":"data","leaf_category":"metrics_table","role":"primary_source","confidence":0.9,"summary":"...","extracts":[]}}
"""
            else:
                prompt = f"""You are building retrieval metadata for a tabular data file.
Return ONE strict JSON object. No markdown fences or extra text.

Use these hints unless the content clearly contradicts them:
- family_hint: data
- role_hint: primary_source

Rules:
- family MUST be "data" unless the content clearly shows another family.
- leaf_category should be a short snake_case label such as "recruiting_sheet", "metrics_table", "inventory_export", "experiment_results".
- role MUST be one of ["primary_source", "summary", "explainer", "analysis", "transcript", "ocr_result", "generated_doc", "reference", "other"].
- confidence must be a float between 0 and 1.
- summary MUST be one English-first paragraph, no markdown, <= 100 words.
- Focus on dataset topic, sheets/columns, entities, dates, measures, IDs, job terms, locations, or search keywords visible in the excerpt.
- Preserve Chinese proper nouns when useful.
- extracts should usually be [] unless clear personal/contact data is present. Keep at most 4 extracts.

File:
- name: {file_name}
- ext: {file_ext}
- pages: {"" if page_count is None else page_count}

Content excerpt:
{prompt_content}

JSON:
{{"family":"data","leaf_category":"metrics_table","role":"primary_source","confidence":0.9,"summary":"...","extracts":[]}}
"""
        elif short_doc_mode:
            prompt = f"""You are building fast retrieval metadata for a local semantic file search system.
Return ONE strict JSON object. Do not use markdown fences or extra text.

Use these hints unless the content clearly contradicts them:
- family_hint: {fallback_family}
- role_hint: {fallback_role}

Rules:
- family MUST be one of ["resume", "report", "contract", "note", "manual", "paper", "presentation", "data", "email", "image", "audio", "video", "book", "code", "invoice", "quotation", "document", "other"].
- leaf_category should be a short snake_case label such as "architecture_plan", "wiring_diagram", "strategic_report", "research_presentation".
- role MUST be one of ["primary_source", "summary", "explainer", "analysis", "transcript", "ocr_result", "generated_doc", "reference", "other"].
- confidence must be a float between 0 and 1.
- summary MUST be exactly one paragraph, English-first, no markdown, <= 110 words.
- In summary include the main subject, 1-3 key entities or numbers, and natural retrieval terms a user would search.
- Preserve Chinese proper nouns when useful.
- extracts should be [] unless obvious high-value personal/contact data is present. Keep at most 4 extracts.

File:
- name: {file_name}
- ext: {file_ext}
- pages: {"" if page_count is None else page_count}

Content excerpt:
{prompt_content}

JSON:
{{"family":"manual","leaf_category":"faq","role":"reference","confidence":0.9,"summary":"...","extracts":[]}}
"""
        elif medium_doc_mode:
            prompt = f"""You are building retrieval metadata for a medium-length local document.
Return ONE strict JSON object. No markdown fences or extra narration.

Use these hints unless the content clearly contradicts them:
- family_hint: {fallback_family}
- role_hint: {fallback_role}

Rules:
- family MUST be one of ["resume", "report", "contract", "note", "manual", "paper", "presentation", "data", "email", "image", "audio", "video", "book", "code", "invoice", "quotation", "document", "other"].
- leaf_category should be a short snake_case label.
- role MUST be one of ["primary_source", "summary", "explainer", "analysis", "transcript", "ocr_result", "generated_doc", "reference", "other"].
- confidence must be a float between 0 and 1.
- summary MUST be one paragraph, English-first, no markdown, <= 140 words.
- Summary should include: main topic, 2-4 important entities/numbers, and practical retrieval phrases a user may search.
- For reports/plans include outcomes, KPIs, milestones, or decisions when present.
- For manuals/technical docs include product/system name and covered components/topics.
- Preserve Chinese proper nouns when useful.
- extracts should be [] unless obvious personal/contact data appears. Keep at most 5 extracts.

File:
- name: {file_name}
- ext: {file_ext}
- pages: {"" if page_count is None else page_count}

Content excerpt:
{prompt_content}

JSON:
{{"family":"report","leaf_category":"analysis_report","role":"analysis","confidence":0.9,"summary":"...","extracts":[]}}
"""
        else:
            prompt = f"""You are building index metadata for a local semantic file retrieval system.
Return ONE strict JSON object. Do not use markdown fences or extra text.

Tasks in this single response:
1. Classify the file into taxonomy metadata.
2. Write a high-recall retrieval summary.
3. Extract high-value personal/sensitive information when present.

Taxonomy rules:
- family MUST be one of ["resume", "report", "contract", "note", "manual", "paper", "presentation", "data", "email", "image", "audio", "video", "book", "code", "invoice", "quotation", "document", "other"].
- leaf_category should be a short reusable label such as "faq", "architecture_plan", "research_paper", "configuration_file", "meeting_notes".
- role MUST be one of ["primary_source", "summary", "explainer", "analysis", "transcript", "ocr_result", "generated_doc", "reference", "other"].
- confidence must be a float between 0 and 1.

Summary rules:
- Output EXACTLY ONE paragraph in the "summary" field. Write primarily in English.
- Preserve Chinese proper nouns in original Chinese and add Pinyin or explicit English names when available.
- Include English retrieval keywords a user would type to find this file.
- For reports/analysis include key findings and numbers if present.
- For manuals/technical docs include product/system name and covered topics.
- If content is empty or meaningless, summary must be "Empty or unreadable document."

Extracts:
- extracts is an array of objects with owner, type, description, content.
- If no sensitive/high-value personal info exists, use [].
- Keep only the highest-value 6 extracts.

File info:
- Name: {file_name}
- Extension: {file_ext}
- PDF pages: {"" if page_count is None else page_count}

Content:
{prompt_content}

Output JSON shape:
{{"family":"manual","leaf_category":"faq","role":"reference","confidence":0.9,"summary":"...","extracts":[]}}
"""
        prompt = self._append_model_prompt_suffix(prompt, model_id)
        if ultra_short_doc_mode:
            prompt_mode = "ultra_short_doc"
        elif tabular_large_doc_mode:
            prompt_mode = "tabular_doc_large"
        elif tabular_doc_mode:
            prompt_mode = "tabular_doc"
        elif short_doc_mode:
            prompt_mode = "short_doc"
        elif medium_doc_mode:
            prompt_mode = "medium_doc"
        else:
            prompt_mode = "full"
        self._index_info(
            f"[_classify_and_summarize_unified] mode={prompt_mode} model={model_id} "
            f"prompt_len={len(prompt)} content_chars={len(prompt_content)} max_tokens={response_max_tokens}"
        )

        client = self._get_local_llm_client()
        last_raw = ""
        for attempt in range(1, self._SUMMARY_MAX_RETRIES + 1):
            temp = 0.0 if attempt == 1 else round(0.1 * (attempt - 1), 1)
            try:
                response = client.chat.completions.create(
                    model=model_id,
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=response_max_tokens,
                    temperature=temp,
                    stream=False,
                )
                raw = (response.choices[0].message.content or "").strip()
                last_raw = raw
                if raw.startswith("```"):
                    raw = _re.sub(r"^```(?:json)?\s*", "", raw)
                    raw = _re.sub(r"\s*```$", "", raw).strip()
                parsed = self._extract_json_object(raw)
            except Exception as e:
                self._index_exception(f"unified 分类摘要失败 attempt={attempt}", e)
                parsed = {}

            summary_str = str(parsed.get("summary") or "").strip()
            if not parsed or not summary_str:
                self._index_warning(
                    f"[_classify_and_summarize_unified] attempt={attempt} missing JSON/summary, retrying..."
                )
                continue
            ratio = self._non_english_ratio(summary_str)
            if ratio > self._NON_EN_RATIO_THRESHOLD:
                self._index_warning(
                    f"[_classify_and_summarize_unified] attempt={attempt}/{self._SUMMARY_MAX_RETRIES} "
                    f"non_en_ratio={ratio:.2%} > threshold, retrying... preview: {summary_str[:80]!r}"
                )
                continue

            family = self._guard_category_by_extension(
                file_ext,
                _normalize_category_en(parsed.get("family", fallback_family), default=fallback_family),
            )
            if family in {"", "all", "unknown"}:
                family = fallback_family
            leaf_category = self._sanitize_taxonomy_leaf(parsed.get("leaf_category"), fallback=fallback_leaf)
            inferred_role = self._infer_doc_role(
                content=stripped_content,
                file_name=file_name,
                file_ext=file_ext,
                family=family,
                leaf_category=leaf_category,
            )
            role = self._sanitize_doc_role(
                parsed.get("role"),
                fallback=inferred_role,
            )
            if inferred_role == "primary_source" and role in {
                "summary", "explainer", "analysis", "generated_doc", "reference"
            }:
                role = "primary_source"
            try:
                confidence = float(parsed.get("confidence", 0.0) or 0.0)
            except Exception:
                confidence = 0.0
            confidence = min(max(confidence, 0.0), 1.0)

            raw_extracts = parsed.get("extracts", [])
            extracts: List[Dict[str, Any]] = []
            if isinstance(raw_extracts, list):
                extracts = [
                    e for e in raw_extracts
                    if isinstance(e, dict) and "content" in e and "type" in e
                ]
            if extracts:
                contact_subtypes = {"phone", "email", "address", "social_media", "contact"}
                found_cnt_types = {
                    str(e.get("type", "")).lower()
                    for e in extracts
                    if str(e.get("type", "")).lower() in contact_subtypes
                }
                if found_cnt_types:
                    summary_str += f" [Contact Information (联系方式): {', '.join(sorted(found_cnt_types))}]"

            if attempt > 1:
                self._index_info(f"[_classify_and_summarize_unified] quality OK after {attempt} attempts")
            if extracts:
                self._index_info(f"[_classify_and_summarize_unified] extracted {len(extracts)} items.")
            return {
                "taxonomy": {
                    "family": family,
                    "leaf_category": leaf_category,
                    "role": role,
                    "confidence": confidence,
                },
                "summary": summary_str[:700],
                "file_name_en": "",
                "extracts": extracts,
                "model": self._snapshot_loaded_llm_model_id() or model_id,
            }

        self._index_warning(
            f"[_classify_and_summarize_unified] FAILED after {self._SUMMARY_MAX_RETRIES} retries. "
            f"model={model_id} raw_preview={last_raw[:120]!r}. Falling back to heuristic taxonomy."
        )
        return _fallback_result(model_id="error_format_or_english")

    def _store_personal_info(self, file_path: str, extracts: List[Dict[str, Any]]) -> None:
        if not extracts or not hasattr(self, 'personal_info_db'):
            return
        try:
            file_name = os.path.basename(file_path)
            for record in extracts:
                record["source_file"] = file_path
                record["source_file_name"] = file_name
            self.personal_info_db.upsert_batch(extracts)
        except Exception as e:
            logger.error(f"[_store_personal_info] failed for {os.path.basename(file_path)}: {e}")

    
    def _generate_image_summary(self, image_path: str) -> str:
        import base64
        from io import BytesIO
        
        MAX_IMAGE_SIZE = 1024
        MIN_IMAGE_SIZE = 4

        try:
            sz = os.path.getsize(image_path)
            self._index_info(
                f"图片摘要：读取/编码 {os.path.basename(image_path)}（约 {sz / 1024 / 1024:.2f} MB），"
                f"随后 VL 推理可能需数十秒～数分钟，期间日志会较少"
            )
        except Exception:
            self._index_info(f"图片摘要：处理 {os.path.basename(image_path)} …")
        
        try:
            from PIL import Image
            HAS_PIL = True
        except ImportError:
            HAS_PIL = False
        
        try:
            ext = os.path.splitext(image_path)[1].lower()
            
            if HAS_PIL:
                with self._open_image_for_vl(image_path) as img:
                    original_size = img.size
                    
                    if img.width < MIN_IMAGE_SIZE or img.height < MIN_IMAGE_SIZE:
                        self._index_warning(
                            f"[_generate_image_summary] 图片尺寸过小 ({img.width}x{img.height})，"
                            f"跳过 VL 推理避免崩溃：{os.path.basename(image_path)}"
                        )
                        return f"Image file: {os.path.basename(image_path)}"

                    if img.mode in ('RGBA', 'LA', 'P'):
                        img = img.convert('RGB')
                    elif img.mode != 'RGB':
                        img = img.convert('RGB')
                    
                    width, height = img.size
                    if width > MAX_IMAGE_SIZE or height > MAX_IMAGE_SIZE:
                        ratio = min(MAX_IMAGE_SIZE / width, MAX_IMAGE_SIZE / height)
                        new_size = (int(width * ratio), int(height * ratio))
                        img = img.resize(new_size, Image.Resampling.LANCZOS)
                        self._index_info(f"缩放 {original_size} -> {new_size}")
                    
                    buffer = BytesIO()
                    img.save(buffer, format='JPEG', quality=85)
                    image_data = base64.b64encode(buffer.getvalue()).decode("utf-8")
                    mime_type = 'image/jpeg'
            else:
                with open(image_path, "rb") as f:
                    image_data = base64.b64encode(f.read()).decode("utf-8")
                
                mime_types = {
                    '.jpg': 'image/jpeg',
                    '.jpeg': 'image/jpeg',
                    '.png': 'image/png',
                    '.gif': 'image/gif',
                }
                mime_type = mime_types.get(ext, 'image/jpeg')
            
            client = self._get_local_llm_client()
            idx_model = self._require_configured_index_model_id("图片摘要")
            if not idx_model:
                return ""
            self._index_info(
                f"图片摘要：开始 VL 推理（encode 完成），model={idx_model}，请稍候…"
            )

            base_prompt = get_prompt("IMAGE_SUMMARY_PROMPT", "en")
            summary = ""
            for attempt in range(1, 3):
                prompt = base_prompt
                if attempt > 1:
                    prompt += (
                        "\n\nIMPORTANT: Your previous response was invalid for indexing. "
                        "Return exactly one plain English paragraph describing the image. "
                        "Do not emit any tool-call syntax, XML tags, JSON, Markdown, labels, "
                        "or special tokens."
                    )
                response = client.chat.completions.create(
                    model=idx_model,
                    messages=[{
                        "role": "user",
                        "content": [
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:{mime_type};base64,{image_data}"}
                            },
                            {
                                "type": "text",
                                "text": self._append_model_prompt_suffix(prompt, idx_model)
                            }
                        ]
                    }],
                    max_tokens=500,
                    temperature=0.0,
                    stream=False,
                )

                summary = self._clean_image_summary(response.choices[0].message.content or "")
                if not self._looks_like_invalid_image_summary(summary):
                    break
                self._index_warning(
                    f"图片摘要输出疑似无效，attempt={attempt}/2，preview={summary[:120]!r}"
                )
                summary = ""

            if not summary:
                self._index_warning(
                    f"图片摘要重试后仍无有效输出，降级为文件名兜底：{os.path.basename(image_path)}"
                )
                return f"Image file: {os.path.basename(image_path)}"

            self._index_info(f"生成成功，长度: {len(summary)} 字符")
            return self._ensure_english_index_text(summary, max_len=600, file_name=os.path.basename(image_path))
        except Exception as e:
            self._index_exception("图片摘要失败", e)
            return ""

    def _open_image_for_vl(self, image_path: str):
        """Open image data for VL, using macOS conversion for HEIC/HEIF when Pillow lacks support."""
        from PIL import Image

        try:
            img = Image.open(image_path)
            img.load()
            return img
        except Exception as pil_exc:
            ext = os.path.splitext(image_path)[1].lower()
            if sys.platform == "darwin" and ext in {".heic", ".heif"}:
                import tempfile

                tmp_path = ""
                try:
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                        tmp_path = tmp.name
                    res = subprocess.run(
                        ["/usr/bin/sips", "-s", "format", "png", image_path, "--out", tmp_path],
                        capture_output=True,
                        text=True,
                        timeout=30,
                    )
                    if res.returncode == 0 and os.path.exists(tmp_path) and os.path.getsize(tmp_path) > 0:
                        img = Image.open(tmp_path)
                        img.load()
                        self._index_info(
                            f"HEIC/HEIF 图片已通过 sips 转换用于 VL: {os.path.basename(image_path)}"
                        )
                        return img
                    err = (res.stderr or res.stdout or "").strip()
                    self._index_warning(
                        f"HEIC/HEIF sips 转换失败: {os.path.basename(image_path)} | {err[:300]}"
                    )
                except Exception as conv_exc:
                    self._index_exception("HEIC/HEIF 图片转换失败", conv_exc)
                finally:
                    if tmp_path:
                        try:
                            os.unlink(tmp_path)
                        except Exception:
                            pass
            raise pil_exc

    def _encode_image_for_vl(
        self,
        image_path: str,
        *,
        max_image_size: int = 1024,
        min_image_size: int = 4,
    ) -> str:
        import base64
        from io import BytesIO

        try:
            try:
                from PIL import Image

                with self._open_image_for_vl(image_path) as img:
                    if img.width < min_image_size or img.height < min_image_size:
                        return ""
                    if img.mode not in ("RGB",):
                        img = img.convert("RGB")
                    w, h = img.size
                    if w > max_image_size or h > max_image_size:
                        ratio = min(max_image_size / w, max_image_size / h)
                        img = img.resize((int(w * ratio), int(h * ratio)), Image.Resampling.LANCZOS)
                    buf = BytesIO()
                    img.save(buf, format="JPEG", quality=85)
                    image_data = base64.b64encode(buf.getvalue()).decode("utf-8")
            except ImportError:
                with open(image_path, "rb") as f:
                    image_data = base64.b64encode(f.read()).decode("utf-8")
            return f"data:image/jpeg;base64,{image_data}"
        except Exception as exc:
            self._index_exception("视频帧编码失败", exc)
            return ""

    @staticmethod
    def _clean_video_frame_summary(summary: str) -> str:
        cleaned = str(summary or "").strip()
        if not cleaned:
            return ""
        cleaned = re.sub(r"<think>.*?</think>", "", cleaned, flags=re.DOTALL).strip()
        cleaned = re.sub(r"<\|thinking\|>.*?<\|/thinking\|>", "", cleaned, flags=re.DOTALL).strip()
        import re as _re_p

        preamble_pat = _re_p.compile(
            r'^(?:'
            r'based on (?:the |this )?(?:image(?: provided)?|provided image|the frame)\s*[,;]\s*'
            r'|based on (?:the )?previous frame(?:[^,;]{0,40})?[,;]\s*'
            r'|in this (?:video )?frame\s*[,;:]\s*'
            r'|here is a (?:detailed )?description[^:]{0,60}:\s*'
            r'|output frame description:\s*'
            r'|frame description:\s*'
            r')+',
            _re_p.IGNORECASE,
        )
        cleaned = preamble_pat.sub("", cleaned).strip()
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
        if cleaned and cleaned[0].islower():
            cleaned = cleaned[0].upper() + cleaned[1:]
        return cleaned

    @staticmethod
    def _clean_image_summary(summary: str) -> str:
        cleaned = str(summary or "").strip()
        if not cleaned:
            return ""
        cleaned = re.sub(r"<think>.*?</think>", "", cleaned, flags=re.DOTALL).strip()
        cleaned = re.sub(r"<\|thinking\|>.*?<\|/thinking\|>", "", cleaned, flags=re.DOTALL).strip()
        cleaned = re.sub(r"^```(?:text|markdown)?\s*", "", cleaned, flags=re.IGNORECASE).strip()
        cleaned = re.sub(r"\s*```$", "", cleaned).strip()
        cleaned = re.sub(
            r"^\s*<\|tool_call\|?>.*?(?:\n+|$)",
            "",
            cleaned,
            flags=re.IGNORECASE | re.DOTALL,
        ).strip()
        cleaned = re.sub(
            r"^\s*(?:call|function_call)\s*:\s*[^\n]+(?:\n+|$)",
            "",
            cleaned,
            flags=re.IGNORECASE,
        ).strip()
        cleaned = re.sub(
            r"^(?:"
            r"(?:sure|certainly|of course|here(?:'s| is)|below is|let me|i can|i will)\b[^.!?\n]{0,140}[.!?]\s*"
            r"|(?:claro|con gusto|por supuesto|aqui|aquí|bien sûr|bien sur|voici|bonjour)\b[^.!?\n]{0,160}[.!?]\s*"
            r")+",
            "",
            cleaned,
            flags=re.IGNORECASE,
        ).strip()
        cleaned = re.sub(r"^\s*[-*]\s*", "", cleaned)
        cleaned = re.sub(
            r"^\s*\*{0,2}(?:output image summary|image summary|summary|description)\s*:\*{0,2}\s*",
            "",
            cleaned,
            flags=re.IGNORECASE,
        ).strip()
        cleaned = re.sub(
            r"^(?:"
            r"based on (?:the |this |the provided |provided )?image\s*[,;:]\s*"
            r"|in (?:the |this )image\s*[,;:]\s*"
            r"|(?:the |this )image (?:shows|depicts|displays|features|contains)\s+"
            r"|this is an image of\s+"
            r")+",
            "",
            cleaned,
            flags=re.IGNORECASE,
        ).strip()
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
        if cleaned and cleaned[0].islower():
            cleaned = cleaned[0].upper() + cleaned[1:]
        return cleaned

    @classmethod
    def _looks_like_invalid_image_summary(cls, summary: str) -> bool:
        text = str(summary or "").strip()
        if not text:
            return True
        lower = text.lower()
        if any(marker in lower for marker in ("<|tool_call", "call:google/", "function_call:", "tool output")):
            return True
        if lower.startswith(("call:", "function_call:", "tool_call:")):
            return True
        if len(text) > 120 and lower.count("image_and_text") >= 2:
            return True
        return False

    @classmethod
    def _should_rewrite_index_text_to_english(cls, text: str) -> bool:
        raw = str(text or "").strip()
        if not raw:
            return False
        if any("\u4e00" <= ch <= "\u9fff" for ch in raw):
            return True
        if cls._non_english_ratio(raw) > cls._NON_EN_RATIO_THRESHOLD:
            return True

        lower = raw.lower()
        if lower.startswith((
            "claro",
            "con gusto",
            "por supuesto",
            "aqui",
            "aquí",
            "bien sûr",
            "bien sur",
            "voici",
            "bonjour",
        )):
            return True

        tokens = re.findall(r"[a-zA-ZÀ-ÿ']+", lower)
        if len(tokens) < 4:
            return False
        en_hits = sum(1 for tok in tokens if tok in cls._INDEX_EN_STOPWORDS)
        non_en_hits = sum(1 for tok in tokens if tok in cls._INDEX_NON_EN_STOPWORDS)
        return non_en_hits >= 3 and non_en_hits >= en_hits + 2

    @classmethod
    def _extract_labeled_frame_summaries(
        cls,
        raw_text: str,
        frame_labels: Sequence[str],
    ) -> Dict[str, str]:
        labels = [str(label or "").strip() for label in frame_labels if str(label or "").strip()]
        if not labels:
            return {}
        text = str(raw_text or "").strip()
        if not text:
            return {}

        parsed_entries: Dict[str, str] = {}
        try:
            parsed = cls._extract_json_object(text)
            frames = parsed.get("frames") if isinstance(parsed, dict) else None
            if isinstance(frames, list):
                for item in frames:
                    if not isinstance(item, dict):
                        continue
                    label = str(item.get("label") or item.get("frame") or item.get("id") or "").strip()
                    summary = cls._clean_video_frame_summary(str(item.get("summary") or item.get("description") or ""))
                    if label in labels and summary:
                        parsed_entries[label] = summary
            if parsed_entries:
                return parsed_entries
        except Exception:
            pass

        for label in labels:
            next_label_pat = r"(?:[A-Z]\d+|" + "|".join(re.escape(v) for v in labels) + r")"
            pattern = re.compile(
                rf"(?:^|\n)\s*(?:Frame\s+)?{re.escape(label)}\s*[:\-]\s*(.*?)(?=(?:\n\s*(?:Frame\s+)?{next_label_pat}\s*[:\-])|\Z)",
                re.IGNORECASE | re.DOTALL,
            )
            match = pattern.search(text)
            if not match:
                continue
            summary = cls._clean_video_frame_summary(match.group(1))
            if summary:
                parsed_entries[label] = summary
        if parsed_entries:
            return parsed_entries

        lines = [line.strip() for line in text.splitlines() if line.strip()]
        if len(lines) == len(labels):
            fallback: Dict[str, str] = {}
            for label, line in zip(labels, lines):
                summary = cls._clean_video_frame_summary(re.sub(r"^(?:Frame\s+)?[A-Z]\d+\s*[:\-]\s*", "", line, flags=re.IGNORECASE))
                if summary:
                    fallback[label] = summary
            if len(fallback) == len(labels):
                return fallback
        return {}

    def _generate_video_frame_batch_summaries(
        self,
        frame_items: Sequence[Tuple[str, str]],
        *,
        prev_description: str = "",
    ) -> Dict[str, str]:
        """Describe a small ordered batch of frames in one VL call.

        Returns a mapping of frame label -> cleaned summary. Missing labels are
        intentionally omitted so callers can fall back to single-frame analysis.
        """
        labeled_uris: List[Tuple[str, str]] = []
        for label, image_path in frame_items:
            clean_label = str(label or "").strip()
            if not clean_label:
                continue
            data_uri = self._encode_image_for_vl(image_path)
            if data_uri:
                labeled_uris.append((clean_label, data_uri))
        if not labeled_uris:
            return {}

        frame_schema = ", ".join(f'"{label}"' for label, _ in labeled_uris)
        prompt_lines = [
            "You are analyzing a short ordered sequence of video frames from the same interval.",
            "For each labeled frame, describe that frame specifically in 1-3 sentences.",
            "Focus on what is visible, screen text/UI, objects, actions, and how it differs from nearby frames when relevant.",
            "Write every summary in English. If visible text is Chinese or another language, translate it into English and preserve the original text only when useful for retrieval.",
            "Do not merge multiple frames into one description.",
            "Return JSON only in this exact shape: {\"frames\":[{\"label\":\"F1\",\"summary\":\"...\"}]}",
            f"Include one object for every label in this batch: {frame_schema}.",
        ]
        if prev_description:
            prompt_lines.append(
                "The previous frame summary before this batch was: "
                f"{prev_description[:400]}"
            )
        prompt_lines.append("Keep each summary factual and concise.")

        idx_model = self._require_configured_index_model_id("视频帧批量摘要")
        if not idx_model:
            return {}
        prompt_text = self._append_model_prompt_suffix("\n".join(prompt_lines), idx_model)
        content: List[Dict[str, Any]] = [{"type": "text", "text": prompt_text}]
        for label, data_uri in labeled_uris:
            content.append({"type": "text", "text": f"Frame {label}"})
            content.append({"type": "image_url", "image_url": {"url": data_uri}})

        try:
            client = self._get_local_llm_client()
            response = client.chat.completions.create(
                model=idx_model,
                messages=[{"role": "user", "content": content}],
                max_tokens=2200,
                temperature=0.0,
                stream=False,
            )
            raw = str(response.choices[0].message.content or "").strip()
            parsed = self._extract_labeled_frame_summaries(
                raw,
                [label for label, _ in labeled_uris],
            )
            return {
                label: self._ensure_english_index_text(
                    summary,
                    max_len=360,
                    file_name="video frame",
                )
                for label, summary in parsed.items()
                if summary
            }
        except Exception as exc:
            self._index_exception("视频帧批量摘要失败", exc)
            return {}

    def _generate_video_frame_summary(self, image_path: str, prev_description: str = "") -> str:
        """VL summary for a single video keyframe."""
        try:
            batch_result = self._generate_video_frame_batch_summaries(
                [("F1", image_path)],
                prev_description=prev_description,
            )
            return self._clean_video_frame_summary(batch_result.get("F1") or "")
        except Exception as e:
            self._index_exception("视频帧摘要失败", e)
            return ""

    def _generate_video_frame_ocr(self, image_path: str) -> str:
        """
        Lightweight OCR for selected keyframes.
        Reuses the existing VL OCR path instead of introducing a separate OCR engine,
        so package size stays flat and only incremental inference time is added.
        """
        try:
            with open(image_path, "rb") as f:
                image_bytes = f.read()
            ext = os.path.splitext(image_path)[1].lower()
            mime_type = "image/png" if ext == ".png" else "image/jpeg"
            max_tokens = int(os.getenv("MEDIA_FRAME_OCR_MAX_TOKENS", str(max(160, OCR_MAX_TOKENS_DEFAULT // 8))))
            ocr = self._qwen_vl_image_text(image_bytes, mime_type, IMAGE_OCR_PROMPT, max_tokens=max_tokens)
            ocr = re.sub(r"<think>.*?</think>", "", (ocr or ""), flags=re.DOTALL).strip()
            ocr = re.sub(r"\s+", " ", ocr)
            max_chars = int(os.getenv("MEDIA_FRAME_OCR_MAX_CHARS", "400"))
            return ocr[:max(80, max_chars)]
        except Exception as e:
            self._index_exception("视频帧 OCR 失败", e)
            return ""

    def _ensure_english_index_text(self, text: str, max_len: int = 200, file_name: str = "") -> str:
        raw = str(text or "").strip()
        if not raw:
            return ""
        raw = raw[: max(32, max_len)]
        if not self._should_rewrite_index_text_to_english(raw):
            return raw
        try:
            client = self._get_local_llm_client()
            prompt = (
                "Translate the following text into concise natural English for retrieval indexing.\n"
                "- Keep names, numbers, and technical terms unchanged when needed.\n"
                "- Remove greetings, wrappers, labels, and assistant-style preambles.\n"
                "- Output a single plain English paragraph only, no extra explanation.\n\n"
            )
            
            prompt += f"{raw}"
            idx_model = self._require_configured_index_model_id("索引文本翻译")
            if not idx_model:
                return raw[:max_len]
            prompt = self._append_model_prompt_suffix(prompt, idx_model)
            resp = client.chat.completions.create(
                model=idx_model,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=min(600, max(64, max_len)),
                temperature=0.0,
                stream=False,
            )
            en_text = (resp.choices[0].message.content or "").strip()
            return (en_text or raw)[:max_len]
        except Exception:
            return raw[:max_len]

    def _qwen_vl_image_text(self, image_bytes: bytes, mime_type: str, prompt: str, max_tokens: int = 900) -> str:
        import base64
        from io import BytesIO
        
        MAX_IMAGE_SIZE = 1024
        MIN_IMAGE_SIZE = 4
        try:
            from PIL import Image
            with Image.open(BytesIO(image_bytes)) as img:
                if img.width < MIN_IMAGE_SIZE or img.height < MIN_IMAGE_SIZE:
                    self._index_warning(
                        f"[_qwen_vl_image_text] 图片尺寸过小 ({img.width}x{img.height})，跳过 VL 推理"
                    )
                    return ""

                if img.mode in ('RGBA', 'LA', 'P') or img.mode != 'RGB':
                    img = img.convert('RGB')
                
                width, height = img.size
                if width > MAX_IMAGE_SIZE or height > MAX_IMAGE_SIZE:
                    ratio = min(MAX_IMAGE_SIZE / width, MAX_IMAGE_SIZE / height)
                    new_size = (int(width * ratio), int(height * ratio))
                    img = img.resize(new_size, Image.Resampling.LANCZOS)
                
                buffer = BytesIO()
                img.save(buffer, format='JPEG', quality=85)
                image_bytes = buffer.getvalue()
                mime_type = 'image/jpeg'
        except Exception as e:
            self._index_exception("Image resize failed", e)

        try:
            client = self._get_local_llm_client()
            image_data = base64.b64encode(image_bytes).decode("utf-8")
            
            idx_model = self._require_configured_index_model_id("VL OCR")
            if not idx_model:
                return ""
            prompt_with_suffix = self._append_model_prompt_suffix(prompt, idx_model)
            
            resp = client.chat.completions.create(
                model=idx_model,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_data}"}},
                        {"type": "text", "text": prompt_with_suffix},
                    ],
                }],
                max_tokens=max_tokens,
                temperature=0.0,
                stream=False,
            )
            return (resp.choices[0].message.content or "").strip()
        except Exception as e:
            self._index_exception("VL OCR 失败", e)
            return ""

    def _pdf_pages_to_ocr_chunks(self, pdf_path: str, max_pages: int = 3, dpi: int = 72) -> List[str]:
        if not HAS_PYPDFIUM2:
            return []
        chunks: List[str] = []
        try:
            for page_number, img_bytes in render_pdf_pages_to_png(pdf_path, max_pages=max_pages, dpi=dpi):
                max_tokens = int(os.getenv("PDF_OCR_MAX_TOKENS", str(OCR_MAX_TOKENS_DEFAULT)))
                ocr = self._qwen_vl_image_text(img_bytes, "image/png", IMAGE_OCR_PROMPT, max_tokens=max_tokens)
                ocr = (ocr or "").strip()
                if ocr:
                    chunks.append(f"[PDF OCR 第{page_number}页]\n{ocr[:5000]}")
        except Exception as e:
            self._index_exception("PDF OCR 失败", e)
        return chunks

    def _docx_images_to_ocr_chunks(self, docx_path: str, max_images: int = 3) -> List[str]:
        import zipfile
        chunks: List[str] = []
        try:
            with zipfile.ZipFile(docx_path) as z:
                media = [n for n in z.namelist() if n.startswith("word/media/") and not n.endswith("/")]
                for idx, name in enumerate(sorted(media)[: int(max_images)], 1):
                    data = z.read(name)
                    ext = os.path.splitext(name)[1].lower()
                    mime = "image/png" if ext == ".png" else "image/jpeg"
                    max_tokens = int(os.getenv("DOCX_OCR_MAX_TOKENS", str(max(1200, OCR_MAX_TOKENS_DEFAULT // 2))))
                    ocr = self._qwen_vl_image_text(data, mime, IMAGE_OCR_PROMPT, max_tokens=max_tokens)
                    if ocr:
                        chunks.append(f"[DOCX OCR 图片{idx}]\n{ocr.strip()[:5000]}")
        except Exception as e:
            self._index_exception("DOCX OCR 失败", e)
        return chunks

    def _pptx_images_to_ocr_chunks(self, pptx_path: str, max_images: int = 3) -> List[str]:
        import zipfile
        chunks: List[str] = []
        try:
            with zipfile.ZipFile(pptx_path) as z:
                media = [n for n in z.namelist() if n.startswith("ppt/media/") and not n.endswith("/")]
                for idx, name in enumerate(sorted(media)[: int(max_images)], 1):
                    data = z.read(name)
                    ext = os.path.splitext(name)[1].lower()
                    mime = "image/png" if ext == ".png" else "image/jpeg"
                    max_tokens = int(os.getenv("PPTX_OCR_MAX_TOKENS", str(max(1200, OCR_MAX_TOKENS_DEFAULT // 2))))
                    ocr = self._qwen_vl_image_text(data, mime, IMAGE_OCR_PROMPT, max_tokens=max_tokens)
                    if ocr:
                        chunks.append(f"[PPTX OCR 图片{idx}]\n{ocr.strip()[:5000]}")
        except Exception as e:
            self._index_exception("PPTX OCR 失败", e)
        return chunks

    @staticmethod
    def _clean_extracted_container_text(text: str) -> str:
        """Normalize text extracted from XML/HTML/binary containers."""
        if not text:
            return ""
        try:
            import html
            text = html.unescape(text)
        except Exception:
            pass
        text = text.replace("\x00", " ")
        text = re.sub(r"[\r\t]+", " ", text)
        text = re.sub(r"[ \f\v]+", " ", text)
        text = re.sub(r"\n\s*\n\s*\n+", "\n\n", text)
        lines = [ln.strip() for ln in text.splitlines()]
        return "\n".join(ln for ln in lines if ln).strip()

    @classmethod
    def _xmlish_bytes_to_text(cls, data: bytes) -> str:
        if not data:
            return ""
        raw = data.decode("utf-8", errors="ignore")
        raw = re.sub(r"<(?:text:line-break|br|w:br|a:br)\b[^>]*/?>", "\n", raw, flags=re.I)
        raw = re.sub(r"</(?:w:p|a:p|p|text:p|text:h|table:table-row|tr|li|h[1-6])>", "\n", raw, flags=re.I)
        raw = re.sub(r"<[^>]+>", " ", raw)
        return cls._clean_extracted_container_text(raw)

    def _extract_zip_text_members(
        self,
        file_path: str,
        *,
        ext: str,
        max_chars: int = 2_000_000,
    ) -> str:
        """Extract useful text from zip-based document containers."""
        import zipfile

        def _member_priority(name: str) -> Tuple[int, str]:
            low = name.lower()
            if low in {"word/document.xml", "content.xml"}:
                return (0, low)
            if low.startswith("ppt/slides/slide") and low.endswith(".xml"):
                return (1, low)
            if low.endswith((".xhtml", ".html", ".htm")):
                return (2, low)
            if low.endswith(".xml"):
                return (3, low)
            if low.endswith(".txt"):
                return (4, low)
            if low.endswith(".pdf") and ("quicklook/" in low or "preview" in low):
                return (5, low)
            return (9, low)

        def _wanted(name: str) -> bool:
            low = name.lower()
            if low.endswith("/"):
                return False
            if low.startswith(("__macosx/", "docprops/", "_rels/")):
                return False
            if low.endswith(".pdf") and ("quicklook/" in low or "preview" in low):
                return True
            if ext in {".doc", ".docx"}:
                return low.startswith("word/") and low.endswith(".xml")
            if ext in {".ppt", ".pptx", ".key"}:
                return (
                    (low.startswith("ppt/slides/") and low.endswith(".xml"))
                    or (low.startswith("ppt/notesSlides/") and low.endswith(".xml"))
                    or low.endswith(".html")
                    or low.endswith(".xhtml")
                    or (low.endswith(".pdf") and ("quicklook/" in low or "preview" in low))
                )
            if ext in {".odt", ".ods", ".odp"}:
                return low == "content.xml" or low.endswith((".html", ".xhtml", ".txt"))
            if ext == ".epub":
                return low.endswith((".xhtml", ".html", ".htm", ".xml", ".txt"))
            if ext in {".pages", ".numbers"}:
                return (
                    low.endswith((".xml", ".html", ".xhtml", ".txt"))
                    or (low.endswith(".pdf") and ("quicklook/" in low or "preview" in low))
                )
            return low.endswith((".xml", ".html", ".xhtml", ".txt"))

        parts: List[str] = []
        try:
            with zipfile.ZipFile(file_path) as z:
                names = sorted([n for n in z.namelist() if _wanted(n)], key=_member_priority)
                for name in names:
                    if _joined_len := (sum(len(p) for p in parts) + max(0, len(parts) - 1) * 2):
                        if _joined_len >= max_chars:
                            break
                    low = name.lower()
                    try:
                        data = z.read(name)
                    except Exception:
                        continue
                    if low.endswith(".pdf") and HAS_PDF_TEXT:
                        try:
                            pdf_text, _ = extract_pdf_text(data=data, max_chars=max_chars)
                            text = self._clean_extracted_container_text(pdf_text)
                        except Exception:
                            text = ""
                    elif low.endswith(".txt"):
                        text = self._clean_extracted_container_text(data.decode("utf-8", errors="ignore"))
                    else:
                        text = self._xmlish_bytes_to_text(data)
                    if text:
                        parts.append(text[: max(0, max_chars - sum(len(p) for p in parts))])
        except zipfile.BadZipFile:
            return ""
        except Exception as e:
            self._index_exception(f"Zip container text extraction failed: {os.path.basename(file_path)}", e)
            return ""
        return self._clean_extracted_container_text("\n\n".join(parts)[:max_chars])

    def _extract_binary_strings_text(self, file_path: str, *, max_bytes: int = 32_000_000, max_chars: int = 1_000_000) -> str:
        """Best-effort fallback for old binary Office / MOBI style containers."""
        try:
            with open(file_path, "rb") as f:
                data = f.read(max_bytes)
        except Exception:
            return ""

        def _useful(chunk: str) -> bool:
            s = chunk.strip()
            if len(s) < 4:
                return False
            if not re.search(r"[A-Za-z0-9\u4e00-\u9fff]", s):
                return False
            alnum = len(re.findall(r"[A-Za-z0-9\u4e00-\u9fff]", s))
            return (alnum / max(1, len(s))) >= 0.20

        parts: List[str] = []
        seen: set[str] = set()
        for enc in ("utf-16le", "utf-8", "latin-1"):
            try:
                decoded = data.decode(enc, errors="ignore")
            except Exception:
                continue
            decoded = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]+", "\n", decoded)
            for chunk in re.findall(r"[^\n]{4,}", decoded):
                chunk = self._clean_extracted_container_text(chunk)
                if not _useful(chunk):
                    continue
                key = chunk[:160]
                if key in seen:
                    continue
                seen.add(key)
                parts.append(chunk)
                if sum(len(p) for p in parts) >= max_chars:
                    break
            if sum(len(p) for p in parts) >= max_chars:
                break
        return self._clean_extracted_container_text("\n".join(parts)[:max_chars])

    def _extract_legacy_office_text(self, file_path: str, *, ext: str) -> str:
        zip_text = self._extract_zip_text_members(file_path, ext=ext)
        if zip_text:
            return zip_text

        if sys.platform == "darwin" and ext in {".doc", ".rtf"}:
            try:
                res = subprocess.run(
                    ["/usr/bin/textutil", "-convert", "txt", "-stdout", file_path],
                    capture_output=True,
                    text=True,
                    timeout=20,
                )
                text = self._clean_extracted_container_text((res.stdout or "").strip())
                if text and "Error reading" not in text:
                    return text
            except Exception:
                pass

        return self._extract_binary_strings_text(file_path)

    def _matches_category_keyword(
        self,
        *,
        keyword: str,
        original_query: Optional[str],
        file_name: str,
        file_name_no_ext: str,
        doc_summary: str,
        file_path: str,
    ) -> Tuple[bool, str]:
        kw_lower = str(keyword or "").lower().strip()
        if not kw_lower:
            return True, ""

        text_to_search = f"{file_name} {file_name_no_ext} {doc_summary} {file_path}".lower()
        oq_lower = str(original_query or "").lower()
        import re

        oq_cjk = "".join(re.findall(r"[\u4e00-\u9fff]+", oq_lower))
        synonym_hit = ""

        _SYNONYM_MAP: Dict[str, set[str]] = {
            "dog": {"puppy", "puppies", "canine", "doggy", "doggie", "hound", "pooch"},
            "cat": {"kitten", "kittens", "kitty", "feline", "kitty cat"},
            "car": {"vehicle", "automobile", "sedan", "suv", "truck", "van"},
            "bike": {"bicycle", "cycling", "cyclist", "motorcycle", "motorbike"},
            "baby": {"infant", "toddler", "newborn", "child", "kid"},
            "house": {"home", "building", "residence", "property", "apartment"},
            "food": {"meal", "dish", "cuisine", "recipe", "snack", "beverage"},
            "flower": {"blossom", "bloom", "floral", "petal", "bouquet"},
            "tree": {"forest", "woodland", "foliage", "plant", "vegetation"},
            "bird": {"avian", "sparrow", "eagle", "parrot", "flock"},
            "fish": {"aquatic", "seafood", "marine", "underwater"},
            "people": {"person", "human", "individual", "man", "woman", "girl", "boy"},
            "狗": {"小狗", "幼犬", "犬", "狗狗"},
            "猫": {"小猫", "猫咪", "幼猫", "猫猫"},
        }

        kw_parts = kw_lower.split() if " " in kw_lower else [kw_lower]
        fname_lower = str(file_name or "").lower()
        fbase = os.path.splitext(fname_lower)[0]
        doc_summary_lower = str(doc_summary or "").lower()
        file_path_lower = str(file_path or "").lower()

        for part in kw_parts:
            part_match = False
            if part in fname_lower or part in doc_summary_lower or part in file_path_lower:
                part_match = True
            elif len(fbase) >= 2 and fbase in part:
                part_match = True
            elif part in _SYNONYM_MAP:
                for syn in _SYNONYM_MAP[part]:
                    if syn in doc_summary_lower:
                        part_match = True
                        synonym_hit = syn
                        break
            if not part_match:
                break
        else:
            return True, synonym_hit

        if oq_lower and oq_lower in text_to_search:
            return True, synonym_hit
        if oq_cjk and len(oq_cjk) >= 2 and oq_cjk in text_to_search:
            return True, synonym_hit
        return False, ""

    def count_by_category(
        self,
        category: str = None,
        keyword: str = None,
        file_extensions: List[str] = None,
        allowed_paths: Optional[List[str]] = None,
        folder: str = None,
        original_query: Optional[str] = None,
        compatible_category_match: bool = False,
    ) -> Dict[str, Any]:
        try:
            from core.retrieval.category_engine import get_compatible_categories

            target_category = _normalize_category_en(category or "", default="")
            if target_category in {"", "all", "unknown"}:
                target_category = ""
            compatible_categories = get_compatible_categories(target_category)
            inventory_category = ""
            if target_category and not compatible_category_match and target_category != "audio/video":
                inventory_category = target_category

            inventory = self.indexed_file_inventory(
                query=keyword or "",
                allowed_paths=allowed_paths,
                category_filter=inventory_category,
                file_extensions=file_extensions,
                limit=0,
                hydrate=True,
            )
            if not inventory.get("ready"):
                logger.warning(
                    "[count_by_category] keyword inventory not ready; refusing query-time metadata scan"
                )
                return {"count": 0, "files": [], "raw_count": 0, "index_ready": False}

            raw_count = int(inventory.get("raw_count") or 0)
            file_dict: Dict[str, Dict[str, Any]] = {}
            media_categories = {"audio/video", "audio", "video"}

            for item in inventory.get("files") or []:
                meta = dict(item.get("metadata") or {})
                file_path = str(item.get("file_path") or meta.get("file_path") or "").strip()
                file_name = str(item.get("file_name") or meta.get("file_name") or os.path.basename(file_path)).strip()
                doc_summary = str(item.get("doc_summary") or meta.get("doc_summary") or "")
                doc_category = _normalize_category_en(item.get("doc_category") or self._meta_category_family(meta), default="other")
                file_key = file_path or file_name
                if not file_key:
                    continue

                if target_category:
                    if target_category == "audio/video":
                        if doc_category not in media_categories:
                            continue
                    elif compatible_category_match:
                        if doc_category not in compatible_categories:
                            continue
                    elif doc_category != target_category:
                        continue

                if folder:
                    parent_dir = os.path.dirname(file_path).lower()
                    folder_blob = " ".join(
                        [
                            parent_dir,
                            os.path.basename(parent_dir),
                            str(meta.get("folder_name_en") or ""),
                            str(meta.get("lookup_aliases") or ""),
                        ]
                    ).lower()
                    if folder.lower() not in folder_blob:
                        continue

                file_dict[file_key] = {
                    "file_name": file_name,
                    "file_path": file_path,
                    "doc_summary": doc_summary,
                    "doc_category": doc_category,
                    "doc_category_leaf": self._meta_category_leaf(meta),
                    "doc_role": self._meta_doc_role(meta),
                    "file_name_en": str(meta.get("file_name_en") or "").strip(),
                    "folder_name_en": str(meta.get("folder_name_en") or "").strip(),
                    "lookup_aliases": str(meta.get("lookup_aliases") or "").strip(),
                    "table_schema_hint": str(meta.get("table_schema_hint") or "").strip(),
                    "en_tags": str(meta.get("en_tags") or "").strip(),
                    "metadata": meta,
                    "hit_chunks": int(item.get("hit_chunks") or 1),
                    "_first_chunk_id": str(item.get("_first_chunk_id") or ""),
                    "_matched_synonym": "",
                }

            unique_count = len(file_dict)
            logger.info(
                "[count_by_category] indexed inventory filtered files=%d raw_records=%d category=%s ext=%s keyword=%r",
                unique_count,
                raw_count,
                target_category or "all",
                file_extensions,
                keyword,
            )
            return {
                "count": unique_count,
                "files": list(file_dict.values()),
                "raw_count": raw_count
            }
        except Exception as e:
            logger.error(f"统计失败: {e}")
            return {"count": 0, "files": [], "raw_count": 0}
    
    def count_all_categories(self, allowed_paths: Optional[List[str]] = None) -> Dict[str, int]:
        try:
            from core.retrieval.category_engine import normalize_meta_category_name, persist_category_registry

            category_files: Dict[str, set] = {}
            inventory = self.indexed_file_inventory(
                allowed_paths=allowed_paths,
                limit=0,
                hydrate=True,
            )
            if not inventory.get("ready"):
                logger.warning(
                    "[count_all_categories] keyword inventory not ready; refusing query-time metadata scan"
                )
                return {}
            for item in inventory.get("files") or []:
                meta = dict(item.get("metadata") or {})
                category = normalize_meta_category_name(meta or {"doc_category": item.get("doc_category", "other")})
                file_path = str(item.get("file_path") or meta.get("file_path") or "")
                file_name = str(item.get("file_name") or meta.get("file_name") or "")
                file_key = file_path or file_name
                if category not in category_files:
                    category_files[category] = set()
                if file_key:
                    category_files[category].add(file_key)

            category_counts = {cat: len(files) for cat, files in category_files.items()}
            if allowed_paths is None:
                try:
                    persist_category_registry(category_counts, source="count_all_categories")
                except Exception:
                    pass
            return category_counts
        except Exception as e:
            logger.error(f"统计所有分类失败: {e}")
            return {}

    def _extract_anchor_terms(self, text: str, max_terms: int = 6) -> List[str]:
        import re

        src = str(text or "").strip()
        if not src:
            return []
        raw_terms: List[str] = []
        raw_terms.extend(re.findall(r"[\u4e00-\u9fff]{2,6}", src))
        raw_terms.extend(re.findall(r"[A-Za-z][A-Za-z0-9._-]{2,31}", src))
        stop = {
            "什么", "怎么", "如何", "一下", "这个", "那个", "主要", "内容", "协作", "模式",
            "文件", "查找", "搜索", "找到", "获取", "打开", "显示", "帮我", "帮助",
            "what", "how", "does", "is", "are", "the", "this", "that", "main", "content",
            "collaboration", "pattern", "between",
            "find", "show", "get", "give", "help", "look", "search", "open", "see",
            "tell", "make", "take", "use", "list", "fetch", "retrieve", "display",
            "want", "need", "can", "will", "please", "let", "try", "check",
            "my", "me", "you", "your", "our", "we", "they", "them", "it", "its",
            "for", "and", "with", "from", "about", "into", "some", "all", "any",
        }
        out: List[str] = []
        seen = set()
        for term in raw_terms:
            t = str(term).strip()
            if not t:
                continue
            t = re.sub(r"^[的与和在把将对请问呢吗啊呀了]+", "", t)
            t = re.sub(r"[的与和呢吗啊呀了]+$", "", t)
            if t.endswith("是") and len(t) > 2:
                t = t[:-1]
            if len(t) < 2 or len(t) > 32:
                continue
            tl = t.lower()
            if tl in stop or tl in seen:
                continue
            seen.add(tl)
            out.append(t)
            if len(out) >= max_terms:
                break
        return out

    def _build_en_tags(
        self,
        file_name: str,
        ext: str,
        doc_summary: str = "",
        parent_folder: str = "",
        folder_name_en: str = "",
    ) -> str:
        """
        Generate English semantic tags for a file at index time.

        These tags are stored as `en_tags` metadata in ChromaDB and allow
        English queries to match Chinese-named files via the lexical fallback.

        Strategy (four cascading levels):
        1. File-format type keywords (e.g. ".wav" → "audio recording sound WAV")
        2. English words already in the filename
        3. Technical identifiers / model numbers (e.g. "MODEL-123", "ESP32")
        4. English words from the doc_summary (if already generated)
        """
        import re as _re

        tags: List[str] = []

        # Level 1: Format → fixed English type keywords
        _FORMAT_TAGS: Dict[str, str] = {
            ".mp3":  "audio music song MP3",
            ".wav":  "audio recording sound WAV",
            ".flac": "audio lossless FLAC",
            ".aac":  "audio AAC",
            ".ogg":  "audio OGG",
            ".m4a":  "audio M4A",
            ".wma":  "audio WMA",
            ".aiff": "audio AIFF",
            ".mp4":  "video MP4",
            ".mov":  "video MOV",
            ".avi":  "video AVI",
            ".mkv":  "video MKV",
            ".pdf":  "document PDF",
            ".docx": "document Word",
            ".doc":  "document Word",
            ".xlsx": "spreadsheet Excel",
            ".xls":  "spreadsheet Excel",
            ".pptx": "presentation slides PowerPoint",
            ".ppt":  "presentation slides PowerPoint",
            ".png":  "image PNG",
            ".jpg":  "image photo JPEG",
            ".jpeg": "image photo JPEG",
            ".heic": "image photo HEIC",
            ".gif":  "image GIF",
            ".svg":  "image vector SVG",
        }
        type_tag = _FORMAT_TAGS.get(ext.lower(), "")
        if type_tag:
            tags.append(type_tag)

        # Level 2: Extract English words from filename (≥2 chars, alpha/numeric)
        name_no_ext = os.path.splitext(file_name)[0]
        en_words = _re.findall(r'[A-Za-z]{2,}', name_no_ext)
        tags.extend(en_words)

        # Level 3: Technical identifiers / product numbers (e.g. "MODEL-123", "ESP32")
        tech_ids = _re.findall(r'[A-Z][A-Za-z]*\d+[A-Za-z0-9]*|\d+[A-Za-z]{2,}', name_no_ext)
        tags.extend(tech_ids)

        # Level 4: Pull English words from doc_summary (if already present)
        if doc_summary:
            summary_en = _re.findall(r'[A-Za-z]{3,}', doc_summary)
            # Only add high-value words (skip very common English words)
            _SKIP = {"the", "and", "for", "are", "was", "this", "that", "with",
                     "from", "its", "not", "has", "file", "document", "content"}
            tags.extend(w for w in summary_en if w.lower() not in _SKIP)

        # Level 5: Include folder semantics so queries can match by directory intent.
        for folder_text in (parent_folder, folder_name_en):
            if not folder_text:
                continue
            tags.extend(_re.findall(r'[A-Za-z]{2,}', folder_text))
            tags.extend(_re.findall(r'[A-Z][A-Za-z]*\d+[A-Za-z0-9]*|\d+[A-Za-z]{2,}', folder_text))

        # Deduplicate while preserving order
        seen: set = set()
        unique_tags: List[str] = []
        for t in tags:
            tl = t.lower()
            if tl not in seen and len(t) >= 2:
                seen.add(tl)
                unique_tags.append(t)

        result = " ".join(unique_tags)
        return result[:500]          # cap at 500 chars to avoid ChromaDB metadata bloat

    def _build_lookup_aliases(
        self,
        file_path: str,
        metadata: Dict[str, Any],
        *extra_texts: str,
    ) -> str:
        """Build a compact, reusable alias blob for lexical/identifier retrieval."""
        try:
            file_name = str(metadata.get("file_name") or "").strip()
            file_name_no_ext = str(metadata.get("file_name_no_ext") or "").strip()
            file_name_en = str(metadata.get("file_name_en") or "").strip()
            parent_folder = str(metadata.get("parent_folder") or "").strip()
            folder_name_en = str(metadata.get("folder_name_en") or "").strip()
            doc_summary = str(metadata.get("doc_summary") or "").strip()
            en_tags = str(metadata.get("en_tags") or "").strip()
            table_schema_hint = str(metadata.get("table_schema_hint") or "").strip()
            latin_aliases = build_cjk_latin_aliases(
                file_name,
                file_name_no_ext,
                file_name_en,
                parent_folder,
                folder_name_en,
                file_path,
                max_terms=64,
            )
            blob = build_lookup_blob(
                file_name,
                file_name_no_ext,
                file_name_en,
                parent_folder,
                folder_name_en,
                file_path,
                doc_summary,
                en_tags,
                table_schema_hint,
                latin_aliases,
                *extra_texts,
                max_terms=96,
            )
            return blob[:700]
        except Exception:
            return ""

    @staticmethod
    def _truncate_index_anchor(text: str, max_chars: int) -> str:
        raw = re.sub(r"\s+", " ", str(text or "").strip())
        if not raw:
            return ""
        if len(raw) <= max_chars:
            return raw
        cut = raw[:max_chars].rstrip()
        for sep in ("。", ".", "！", "!", "？", "?", "；", ";", "，", ",", "、", " "):
            idx = cut.rfind(sep)
            if idx >= max(24, max_chars // 2):
                cut = cut[: idx + 1].rstrip()
                break
        return cut.strip()

    def _build_compact_embedding_context(
        self,
        file_path: str,
        metadata: Dict[str, Any],
        *,
        ext: str = "",
        max_chars: Optional[int] = None,
    ) -> str:
        """
        Build a short, high-signal prefix for chunk embedding.

        The full summary / aliases remain in metadata, but repeating them in
        every chunk makes token-aware splitting explode into tiny chunks. This
        context keeps only the most useful anchors for retrieval.
        """
        try:
            limit = int(
                max_chars
                if max_chars is not None
                else os.getenv("FILEAGENT_EMBED_CONTEXT_MAX_CHARS", "320")
            )
        except Exception:
            limit = 320
        limit = max(160, min(limit, 768))

        file_name = str(metadata.get("file_name") or os.path.basename(file_path) or "").strip()
        file_name_no_ext = str(metadata.get("file_name_no_ext") or "").strip()
        if not file_name_no_ext and file_name:
            file_name_no_ext = os.path.splitext(file_name)[0].strip()
        file_name_en = str(metadata.get("file_name_en") or "").strip()
        parent_folder = str(metadata.get("parent_folder") or "").strip()
        folder_name_en = str(metadata.get("folder_name_en") or "").strip()
        doc_summary = str(metadata.get("doc_summary") or "").strip()
        table_schema_hint = str(metadata.get("table_schema_hint") or "").strip()
        en_tags = str(metadata.get("en_tags") or "").strip()

        # Compact alias surface: keep just the most useful file/folder anchors.
        alias_seed = build_lookup_blob(
            file_name,
            file_name_no_ext,
            file_name_en,
            parent_folder,
            folder_name_en,
            max_terms=24,
        )
        if alias_seed:
            pinyin_seed = build_cjk_latin_aliases(
                file_name,
                file_name_no_ext,
                file_name_en,
                parent_folder,
                folder_name_en,
                max_terms=8,
            )
            if pinyin_seed:
                alias_seed = f"{alias_seed} {pinyin_seed}".strip()

        summary_anchor = doc_summary or table_schema_hint
        summary_anchor = self._truncate_index_anchor(summary_anchor, 180)

        tag_terms: List[str] = []
        for term in re.findall(r"[A-Za-z0-9]{2,}|[\u4e00-\u9fff]{2,}", en_tags):
            cleaned = term.strip()
            if cleaned and cleaned.lower() not in {"document", "file", "data", "content"}:
                tag_terms.append(cleaned)
            if len(tag_terms) >= 10:
                break
        tag_anchor = " ".join(tag_terms)

        parts: List[str] = []
        if file_name or file_name_en:
            file_part = " ".join(
                p for p in [
                    self._truncate_index_anchor(file_name, 48),
                    self._truncate_index_anchor(file_name_en, 48),
                ] if p
            ).strip()
            if file_part:
                parts.append(file_part)
        if parent_folder or folder_name_en:
            folder_part = " ".join(
                p for p in [
                    self._truncate_index_anchor(parent_folder, 36),
                    self._truncate_index_anchor(folder_name_en, 36),
                ] if p
            ).strip()
            if folder_part:
                parts.append(folder_part)
        if ext.lower() in {".csv", ".tsv", ".xlsx", ".xls", ".numbers"} and table_schema_hint:
            parts.append(self._truncate_index_anchor(table_schema_hint, 140))
        elif summary_anchor:
            parts.append(summary_anchor)
        if tag_anchor:
            parts.append(tag_anchor)
        if alias_seed:
            parts.append(alias_seed)

        context = " | ".join(part for part in parts if part).strip()
        if not context:
            context = self._truncate_index_anchor(file_name or file_path, 96)
        return context[:limit]

    def _build_tabular_index_hint(self, text: str, file_name: str, ext: str) -> str:
        """Create a deterministic schema-style hint for CSV/XLSX/Numbers files."""
        import re as _re

        raw = str(text or "").strip()
        if not raw:
            return ""

        lines = [line.strip() for line in raw.splitlines() if line.strip()]
        if not lines:
            return ""

        sheet_names: List[str] = []
        headers: List[str] = []
        approx_rows = 0

        for line in lines[:120]:
            if line.startswith("[Sheet: ") and line.endswith("]"):
                sheet_name = line[len("[Sheet: ") : -1].strip()
                if sheet_name and sheet_name not in sheet_names:
                    sheet_names.append(sheet_name)
                continue
            if "|" in line and not headers:
                cells = [c.strip() for c in line.split("|")]
                filtered = [c for c in cells if c and c != "---"]
                if filtered:
                    headers = filtered[:12]
                    continue
            if not line.startswith("[Sheet:") and "---" not in line:
                approx_rows += 1

        parts: List[str] = [f"Tabular data file {file_name}"]
        ext_label = {
            ".csv": "CSV",
            ".tsv": "TSV",
            ".xlsx": "Excel workbook",
            ".xls": "Excel workbook",
            ".numbers": "Numbers workbook",
        }.get(ext.lower(), "table file")
        parts.append(f"Format: {ext_label}")
        if sheet_names:
            parts.append(f"Sheets: {', '.join(sheet_names[:4])}")
        if headers:
            parts.append(f"Columns: {', '.join(headers[:12])}")
        if approx_rows > 0:
            parts.append(f"Approx rows: {max(approx_rows - 1, 0)}")

        # Keep a few stable identifier-ish tokens from the content so
        # filename-like / code-like lookup queries have more anchors.
        token_pool = []
        for token in headers + sheet_names:
            token_pool.extend(_re.findall(r"[A-Za-z0-9_./-]+|[\u4e00-\u9fff]+", token))
        if token_pool:
            dedup = []
            seen = set()
            for token in token_pool:
                t = str(token or "").strip()
                if not t:
                    continue
                tl = t.lower()
                if tl in seen:
                    continue
                seen.add(tl)
                dedup.append(t)
                if len(dedup) >= 16:
                    break
            if dedup:
                parts.append(f"Schema keywords: {', '.join(dedup)}")

        return ". ".join(parts)[:700]

    def _default_tabular_summary(self, file_name: str, ext: str) -> str:
        """Return a stable fallback summary for tabular files when LLM summary is empty."""
        ext_label = {
            ".csv": "CSV",
            ".tsv": "TSV",
            ".xlsx": "Excel workbook",
            ".xls": "Excel workbook",
            ".numbers": "Numbers workbook",
        }.get(str(ext or "").lower(), "table file")
        return f"Tabular data file {file_name}. Format: {ext_label}"[:700]

    _FOLDER_INDEX_SKIP_BASENAMES = frozenset(
        {
            "",
            ".",
            "..",
            "downloads",
            "documents",
            "desktop",
            "桌面",
            "下载",
            "文档",
            "users",
            "library",
        }
    )

    def _folder_index_skip_basename(self, name: str) -> bool:
        n = (name or "").strip().lower()
        if not n or n in (".", ".."):
            return True
        return n in self._FOLDER_INDEX_SKIP_BASENAMES

    def _folder_chain_should_ignore_disk_file(
        self,
        file_path: str,
        *,
        relax_rules: bool = False,
    ) -> bool:
        if not relax_rules:
            return self._should_ignore_file(file_path)

        abs_path = os.path.abspath(file_path)
        if abs_path in self._session_ignored_paths:
            return True

        file_name = os.path.basename(file_path)
        lower_name = file_name.lower()
        if lower_name in {".ds_store"}:
            return True
        if lower_name.startswith(".~") or lower_name.startswith("~$") or lower_name.startswith("._"):
            return True
        if lower_name.endswith((".swp", ".swo", ".tmp", ".bak")):
            return True
        return False

    def _folder_chain_disk_search_roots(
        self,
        allowed_paths: Optional[List[str] | PathScopeMatcher],
    ) -> List[str]:
        import config.settings as settings

        matcher = ensure_path_scope_matcher(allowed_paths)
        roots: List[str] = []
        raw_candidates: List[str] = []
        if not matcher.allow_all and matcher.roots:
            raw_candidates.extend(matcher.roots)
        else:
            if getattr(settings, "USE_WHITELIST_MODE", False):
                raw_candidates.extend(sorted(getattr(settings, "INCLUDE_PATHS", set()) or []))
            if not raw_candidates:
                raw_candidates.extend(
                    [
                        os.path.join(settings.HOME_DIR, "Documents"),
                        os.path.join(settings.HOME_DIR, "Downloads"),
                        os.path.join(settings.HOME_DIR, "Desktop"),
                    ]
                )

        seen: set[str] = set()
        for raw in raw_candidates:
            try:
                norm = os.path.normpath(os.path.abspath(os.path.expanduser(str(raw or "").strip())))
            except Exception:
                continue
            if not norm:
                continue
            if os.path.isfile(norm):
                norm = os.path.dirname(norm)
            if not os.path.isdir(norm) or norm in seen:
                continue
            seen.add(norm)
            roots.append(norm)
        return roots

    def discover_exact_folder_literal_roots_from_disk(
        self,
        needles: List[str],
        allowed_paths: Optional[List[str] | PathScopeMatcher] = None,
        *,
        max_roots: int = 16,
    ) -> List[str]:
        if not needles:
            return []

        import config.settings as settings

        matcher = ensure_path_scope_matcher(allowed_paths)
        search_roots = self._folder_chain_disk_search_roots(matcher)
        if not search_roots:
            return []

        try:
            max_dirs = max(256, int(os.getenv("FOLDER_ROOT_DISK_FALLBACK_MAX_DIRS", "15000") or 15000))
        except (TypeError, ValueError):
            max_dirs = 15000
        try:
            max_ms = max(250, int(os.getenv("FOLDER_ROOT_DISK_FALLBACK_MAX_MS", "2000") or 2000))
        except (TypeError, ValueError):
            max_ms = 2000

        deadline = time.time() + (max_ms / 1000.0)
        skip_dir_names = {
            "node_modules",
            "__pycache__",
            ".git",
            ".svn",
            ".hg",
            ".vscode",
            ".cursor",
            ".idea",
            "site-packages",
            "venv",
            "env",
            ".env",
        }

        found: List[str] = []
        seen_found: set[str] = set()
        seen_dirs = 0

        def _allow_exact_skip_root(norm_path: str) -> bool:
            try:
                bn = os.path.basename(str(norm_path or "").rstrip(os.sep))
            except Exception:
                return False
            if not bn or not self._folder_index_skip_basename(bn):
                return False
            return any(
                self._folder_literal_needle_matches_name_exact(nd, folder_name=bn)
                for nd in needles
            )

        for base in search_roots:
            try:
                base_norm = os.path.normpath(os.path.abspath(os.path.expanduser(base)))
            except Exception:
                continue
            if not os.path.isdir(base_norm):
                continue

            for root, subdirs, _files in os.walk(base_norm):
                seen_dirs += 1
                if seen_dirs > max_dirs or time.time() >= deadline:
                    logger.info(
                        "[folder_chain] disk root discovery budget reached: searched=%d found=%d",
                        seen_dirs,
                        len(found),
                    )
                    return found

                try:
                    norm_root = os.path.normpath(root)
                except Exception:
                    norm_root = root

                if matcher.allow_all:
                    try:
                        rel_path = os.path.relpath(norm_root, settings.HOME_DIR)
                    except Exception:
                        rel_path = norm_root
                    top_dir = rel_path.split(os.sep)[0] if os.sep in rel_path else rel_path
                    if top_dir in getattr(settings, "IGNORE_TOP_LEVEL_DIRS", set()):
                        subdirs[:] = []
                        continue
                elif not matcher.allows_folder(norm_root):
                    subdirs[:] = []
                    continue

                try:
                    base_name = os.path.basename(norm_root.rstrip(os.sep))
                except Exception:
                    base_name = ""
                if (
                    base_name
                    and any(
                        self._folder_literal_needle_matches_name_exact(nd, folder_name=base_name)
                        for nd in needles
                    )
                    and (
                        self._folder_expand_root_is_meaningful(norm_root)
                        or _allow_exact_skip_root(norm_root)
                    )
                    and norm_root not in seen_found
                ):
                    seen_found.add(norm_root)
                    found.append(norm_root)
                    if len(found) >= max_roots:
                        return found

                filtered_subdirs: List[str] = []
                for name in sorted(subdirs):
                    if name in skip_dir_names:
                        continue
                    filtered_subdirs.append(name)
                subdirs[:] = filtered_subdirs
        return found

    def _folder_index_doc_id(self, folder_path: str) -> str:
        norm = os.path.normpath(os.path.abspath(os.path.expanduser(str(folder_path))))
        h = hashlib.sha256(norm.encode("utf-8")).hexdigest()[:48]
        return f"fd_{h}"

    def _iter_parent_dirs_for_folder_index(self, file_path: str) -> List[Tuple[str, str]]:
        try:
            abs_fp = os.path.abspath(os.path.expanduser(str(file_path)))
        except Exception:
            return []
        out: List[Tuple[str, str]] = []
        current_dir = os.path.dirname(abs_fp)
        while current_dir and current_dir not in ("/", "\\"):
            fname = os.path.basename(current_dir)
            if fname and not self._folder_index_skip_basename(fname):
                norm = os.path.normpath(current_dir)
                out.append((fname, norm))
            parent = os.path.dirname(current_dir)
            if parent == current_dir:
                break
            current_dir = parent
        return out

    @staticmethod
    def _looks_like_translation_noise_token(token: str) -> bool:
        cleaned = str(token or "").strip().strip("_-~. ")
        if not cleaned:
            return True
        if re.fullmatch(r"[A-Fa-f0-9]{8,}", cleaned):
            return True
        if re.fullmatch(r"[A-Za-z0-9]{12,}", cleaned):
            digit_count = sum(ch.isdigit() for ch in cleaned)
            alpha_count = sum(ch.isalpha() for ch in cleaned)
            if digit_count >= 3 and alpha_count >= 3:
                return True
        return False

    def _prepare_translation_source(self, text: str, *, kind: str = "file") -> str:
        """
        Strip noisy hashes / ids from mixed-language file or folder labels before
        sending them to the translator.
        """
        raw = str(text or "").strip()
        if not raw:
            return ""
        source = raw
        if kind == "file":
            source = os.path.splitext(source)[0].strip()
        source = re.sub(r"\s+", " ", source)

        tokens: List[str] = []
        for part in re.split(r"[\\/|_]+", source):
            chunk = str(part or "").strip().strip("~`'\"[](){}<>")
            if not chunk:
                continue
            if self._looks_like_translation_noise_token(chunk):
                continue
            cjk_runs = re.findall(r"[\u4e00-\u9fff]{1,24}", chunk)
            if cjk_runs:
                tokens.extend(cjk_runs[:4])
                latin_tail = re.findall(r"[A-Za-z][A-Za-z0-9.+-]{1,20}", chunk)
                for tail in latin_tail[:2]:
                    if not self._looks_like_translation_noise_token(tail):
                        tokens.append(tail)
                continue
            normalized = re.sub(r"\s+", " ", chunk).strip()
            if normalized and not self._looks_like_translation_noise_token(normalized):
                tokens.append(normalized[:28])

        if not tokens:
            cjk_only = " ".join(re.findall(r"[\u4e00-\u9fff]{1,24}", source)[:4]).strip()
            if cjk_only:
                return cjk_only
            return source[:80]

        deduped: List[str] = []
        seen: set[str] = set()
        for token in tokens:
            key = token.lower()
            if key in seen:
                continue
            seen.add(key)
            deduped.append(token)
            if len(deduped) >= 8:
                break
        limit = 48 if kind == "folder" else 80
        return " ".join(deduped)[:limit].strip()

    def _translate_file_name_to_en(self, fname: str) -> str:
        if not self._label_needs_translation(fname):
            return ""

        if hasattr(self, "_file_translation_cache"):
            if fname in self._file_translation_cache:
                return self._file_translation_cache[fname]
        else:
            self._file_translation_cache = {}

        try:
            res = self.collection.get(where={"file_name": fname}, limit=1, include=["metadatas"])
            if res and res.get("metadatas") and len(res["metadatas"]) > 0:
                en_name = self._sanitize_translation_label(res["metadatas"][0].get("file_name_en"))
                if en_name:
                    self._file_translation_cache[fname] = en_name
                    return en_name
        except Exception:
            pass

        try:
            translation_source = self._prepare_translation_source(fname, kind="file") or fname
            sys_prompt = (
                "You translate mixed Chinese file labels into concise English retrieval anchors. "
                "Keep meaningful English tokens when useful, ignore random hashes / ids / suffix noise, "
                "and output ONLY the English label."
            )
            idx_model = self._require_configured_index_model_id("文件名翻译")
            if not idx_model:
                return ""
            sys_prompt = self._append_model_prompt_suffix(sys_prompt, idx_model)
            
            client = self._get_local_llm_client()
            res = client.chat.completions.create(
                model=idx_model,
                messages=[
                    {"role": "system", "content": sys_prompt},
                    {"role": "user", "content": translation_source},
                ],
                max_tokens=64,
                temperature=0.0,
                stream=False,
            )
            text = self._sanitize_translation_label(res.choices[0].message.content)
            self._file_translation_cache[fname] = text
            return text
        except Exception as e:
            logger.warning(f"File name translation failed for {fname}: {e}")
            return ""

    def _resolve_file_name_en_for_metadata(
        self,
        file_name: str,
        *,
        existing_translation: str = "",
    ) -> str:
        """
        Finalize file_name_en for metadata writes.

        File names now use the dedicated translator only. For non-English
        filenames, avoid relying on summary-side JSON fields that can drift.
        """
        if not file_name:
            return ""
        file_name_no_ext = os.path.splitext(file_name)[0]
        if not self._label_needs_translation(file_name_no_ext):
            return file_name_no_ext

        cached = self._sanitize_translation_label(existing_translation)
        if cached:
            return cached

        direct = self._sanitize_translation_label(self._translate_file_name_to_en(file_name))
        if direct:
            return direct
        return file_name_no_ext

    # Hardcoded macOS / system folder name aliases.
    # These are well-known names with canonical English equivalents; using a table
    # instead of LLM guarantees deterministic, correct translations every time.
    _MACOS_FOLDER_ALIASES: Dict[str, str] = {
        "未命名文件夹": "Unnamed Folder",
        "桌面": "Desktop",
        "下载": "Downloads",
        "文稿": "Documents",
        "音乐": "Music",
        "图片": "Pictures",
        "影片": "Movies",
        "公共": "Public",
        "应用程序": "Applications",
        "资源库": "Library",
        "废纸篓": "Trash",
        "个人收藏": "Favorites",
        "最近使用": "Recents",
        "网络": "Network",
        "系统": "System",
        "工作": "Work",
        "学习": "Study",
        "项目": "Projects",
        "备份": "Backup",
        "归档": "Archive",
        "临时": "Temp",
        "新建文件夹": "New Folder",
        "untitled folder": "Untitled Folder",
    }

    def _translate_folder_name(self, fname: str) -> str:
        if not self._label_needs_translation(fname):
            return ""
        
        alias = self.__class__._MACOS_FOLDER_ALIASES.get(fname) or self.__class__._MACOS_FOLDER_ALIASES.get(fname.strip())
        if alias:
            if hasattr(self, "_folder_translation_cache"):
                self._folder_translation_cache[fname] = alias
            else:
                self._folder_translation_cache = {fname: alias}
            return alias

        if hasattr(self, "_folder_translation_cache"):
            if fname in self._folder_translation_cache:
                return self._folder_translation_cache[fname]
        else:
            self._folder_translation_cache = {}

        try:
            fc = getattr(self, "folder_collection", None)
            if fc is not None:
                res = fc.get(where={"folder_name": fname}, limit=1, include=["metadatas"])
                if res and res.get("metadatas") and len(res["metadatas"]) > 0:
                    en_name = self._sanitize_translation_label(res["metadatas"][0].get("folder_name_en"))
                    if en_name:
                        self._folder_translation_cache[fname] = en_name
                        return en_name
        except Exception:
            pass

        try:
            translation_source = self._prepare_translation_source(fname, kind="folder") or fname
            sys_prompt = (
                "Translate the given Chinese folder label into concise English. "
                "Ignore random hashes, ids, or meaningless suffix noise. "
                "Output ONLY the English folder label."
            )
            idx_model = self._require_configured_index_model_id("文件夹名翻译")
            if not idx_model:
                return ""
            sys_prompt = self._append_model_prompt_suffix(sys_prompt, idx_model)
            
            client = self._get_local_llm_client()
            res = client.chat.completions.create(
                model=idx_model,
                messages=[
                    {"role": "system", "content": sys_prompt},
                    {"role": "user", "content": translation_source},
                ],
                max_tokens=64,
                temperature=0.0,
                stream=False,
            )
            text = self._sanitize_translation_label(res.choices[0].message.content)
            self._folder_translation_cache[fname] = text
            return text
        except Exception as e:
            logger.warning(f"Folder translation failed for {fname}: {e}")
            return ""

    def _sanitize_translation_label(self, text: Any) -> str:
        """Keep only short label-like translations and drop leaked reasoning text."""
        cleaned = str(text or "").strip()
        if not cleaned:
            return ""
        try:
            from services.inproc_openai_client import _strip_think_blocks

            cleaned = _strip_think_blocks(cleaned).strip()
        except Exception:
            cleaned = re.sub(r"<think>.*?(</think>|$)", "", cleaned, flags=re.DOTALL).strip()
        if not cleaned:
            return ""
        lines = [ln.strip().strip("\"'`") for ln in cleaned.splitlines() if ln.strip()]
        if not lines:
            return ""
        cleaned = lines[0]
        cleaned = re.sub(r"^(translation|english|translated)\s*[:：-]\s*", "", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\s+", " ", cleaned).strip().strip("\"'`")
        lowered = cleaned.lower()
        bad_markers = (
            "the user wants",
            "translate the chinese",
            "concise english phrase",
            "let me break this down",
            "output just the english",
            "translation of",
            "here is the translation",
        )
        if any(marker in lowered for marker in bad_markers):
            return ""
        if len(cleaned) > 140:
            return ""
        return cleaned

    def _upsert_folder_index_for_file(self, file_path: str) -> None:
        if self._in_write_heavy_mode():
            self._folder_index_dirty = True
            return
        fc = getattr(self, "folder_collection", None)
        if fc is None:
            return
        em = getattr(self, "embedding_model", None)
        if em is None:
            return
        try:
            pairs = self._iter_parent_dirs_for_folder_index(file_path)
            if not pairs:
                return
            doc_ids_to_fetch = [self._folder_index_doc_id(fp) for _, fp in pairs]
            old_translations = {}
            try:
                existing = fc.get(ids=doc_ids_to_fetch, include=["metadatas"])
                for meta in (existing.get("metadatas") or []):
                    if meta:
                        fp = meta.get("folder_path")
                        en = self._sanitize_translation_label(meta.get("folder_name_en"))
                        if fp and en:
                            old_translations[fp] = en
            except Exception:
                pass

            ids: List[str] = []
            embeddings: List[List[float]] = []
            documents: List[str] = []
            metadatas: List[Dict[str, Any]] = []
            for fname, fpath in pairs:
                doc_id = self._folder_index_doc_id(fpath)
                fname_en = old_translations.get(fpath, "")
                if not fname_en and self._label_needs_translation(fname):
                    fname_en = self._translate_folder_name(fname)
                
                doc_text = f"{fname} {fname_en}".strip()
                with self._embed_context(f"folder_index_upsert name={fname}"):
                    emb = self._embed_query_text(doc_text)
                ids.append(doc_id)
                embeddings.append(emb)
                documents.append(doc_text)
                metadatas.append(
                    {
                        "folder_path": fpath,
                        "folder_name": fname,
                        "folder_name_en": fname_en,
                        "kind": "folder_index",
                    }
                )
            with self._folder_index_lock:
                fc.upsert(ids=ids, embeddings=embeddings, documents=documents, metadatas=metadatas)
        except Exception as e:
            self._index_warning(f"folder index upsert 跳过: {file_path} | {e}")

    def rebuild_folder_index_if_dirty(self) -> int:
        if not self._folder_index_dirty:
            return 0
        self._folder_index_dirty = False
        return int(self.rebuild_folder_index() or 0)

    def rebuild_folder_index(self) -> int:
        fc = getattr(self, "folder_collection", None)
        if fc is None:
            return 0
        em = getattr(self, "embedding_model", None)
        if em is None:
            return 0
        folder_paths: Dict[str, str] = {}
        try:
            total = int(self.collection.count())
            page_size = 2000
            offset = 0
            while offset < total:
                batch = self.collection.get(
                    include=["metadatas"],
                    limit=page_size,
                    offset=offset,
                )
                metas = batch.get("metadatas") or []
                if not metas:
                    break
                for m in metas:
                    fp = m.get("file_path")
                    if not fp:
                        continue
                    for fname, fpath in self._iter_parent_dirs_for_folder_index(str(fp)):
                        folder_paths[fpath] = fname
                offset += len(metas)
        except Exception as e:
            logger.error(f"rebuild_folder_index 扫描失败: {e}")
            return 0
        if not folder_paths:
            with self._folder_index_lock:
                try:
                    existing = fc.get(include=[])
                    old_ids = existing.get("ids") or []
                    if old_ids:
                        for i in range(0, len(old_ids), 500):
                            fc.delete(ids=old_ids[i : i + 500])
                except Exception as e:
                    logger.warning(f"rebuild_folder_index 清空旧索引: {e}")
            try:
                self._upsert_filename_lookup_for_file(file_path, metadata, persist=False)
                self._maybe_persist()
            except Exception:
                pass
            logger.info("[folder_index] 重建完成：无已索引文件，文件夹索引已清空")
            return 0
        items = list(folder_paths.items())
        batch_sz = 48
        n_written = 0
        with self._folder_index_lock:
            old_translations = {}
            try:
                existing = fc.get(include=["metadatas"])
                old_ids = existing.get("ids") or []
                for meta in (existing.get("metadatas") or []):
                    if meta:
                        fp = meta.get("folder_path")
                        en = self._sanitize_translation_label(meta.get("folder_name_en"))
                        if fp and en:
                            old_translations[fp] = en
                if old_ids:
                    for i in range(0, len(old_ids), 500):
                        fc.delete(ids=old_ids[i : i + 500])
            except Exception as e:
                logger.warning(f"rebuild_folder_index 清空旧索引: {e}")
            for i in range(0, len(items), batch_sz):
                chunk = items[i : i + batch_sz]
                ids: List[str] = []
                embeddings: List[List[float]] = []
                documents: List[str] = []
                metadatas: List[Dict[str, Any]] = []
                for fpath, fname in chunk:
                    doc_id = self._folder_index_doc_id(fpath)
                    fname_en = old_translations.get(fpath, "")
                    if not fname_en and self._label_needs_translation(fname):
                        fname_en = self._translate_folder_name(fname)
                    
                    doc_text = f"{fname} {fname_en}".strip()
                    with self._embed_context(f"folder_index_rebuild name={fname}"):
                        emb = self._embed_query_text(doc_text)
                    ids.append(doc_id)
                    embeddings.append(emb)
                    documents.append(doc_text)
                    metadatas.append(
                        {
                            "folder_path": fpath,
                            "folder_name": fname,
                            "folder_name_en": fname_en,
                            "kind": "folder_index",
                        }
                    )
                fc.upsert(ids=ids, embeddings=embeddings, documents=documents, metadatas=metadatas)
                n_written += len(chunk)
        try:
            self._maybe_persist()
        except Exception:
            pass
        logger.info(f"[folder_index] 重建完成，文件夹条目数={n_written}")
        return n_written

    def prune_folder_index_orphans(self) -> int:
        fc = getattr(self, "folder_collection", None)
        if fc is None:
            return 0
        try:
            if fc.count() == 0:
                return 0
        except Exception:
            return 0
        valid: set = set()
        try:
            total = int(self.collection.count())
            page_size = 2000
            offset = 0
            while offset < total:
                batch = self.collection.get(
                    include=["metadatas"],
                    limit=page_size,
                    offset=offset,
                )
                metas = batch.get("metadatas") or []
                if not metas:
                    break
                for m in metas:
                    fp = m.get("file_path")
                    if not fp:
                        continue
                    for _, fpath in self._iter_parent_dirs_for_folder_index(str(fp)):
                        valid.add(fpath)
                offset += len(metas)
        except Exception as e:
            logger.warning(f"prune_folder_index_orphans 扫描文件失败: {e}")
            return 0
        to_delete: List[str] = []
        try:
            batch = fc.get(include=["metadatas"])
            ids = batch.get("ids") or []
            metas = batch.get("metadatas") or []
            for doc_id, meta in zip(ids, metas):
                p = str((meta or {}).get("folder_path") or "")
                if p and p not in valid:
                    to_delete.append(doc_id)
        except Exception as e:
            logger.warning(f"prune_folder_index_orphans 读取 folder 集合失败: {e}")
            return 0
        if not to_delete:
            return 0
        with self._folder_index_lock:
            for i in range(0, len(to_delete), 500):
                fc.delete(ids=to_delete[i : i + 500])
        try:
            self._maybe_persist()
        except Exception:
            pass
        logger.info(f"[folder_index] 已移除 {len(to_delete)} 条无已索引文件支撑的文件夹条目")
        return len(to_delete)

    def _maybe_auto_rebuild_folder_index(self) -> None:
        if str(os.getenv("FILEAGENT_SKIP_FOLDER_INDEX_AUTO_REBUILD", "") or "").strip().lower() in (
            "1",
            "true",
            "yes",
            "on",
        ):
            return
        fc = getattr(self, "folder_collection", None)
        if fc is None:
            return
        if self._folder_index_auto_rebuild_attempted:
            return
        try:
            if fc.count() > 0:
                self._folder_index_auto_rebuild_attempted = True
                return
            if self.collection.count() == 0:
                self._folder_index_auto_rebuild_attempted = True
                return
        except Exception:
            return
        self._folder_index_auto_rebuild_attempted = True
        self.rebuild_folder_index()

    @staticmethod
    def _is_path_allowed(
        file_path: str,
        allowed_paths: Optional[List[str] | PathScopeMatcher],
    ) -> bool:
        return ensure_path_scope_matcher(allowed_paths).allows_file(file_path)

    def _folder_allowed(
        self,
        folder_path: str,
        allowed_paths: Optional[List[str] | PathScopeMatcher],
    ) -> bool:
        return ensure_path_scope_matcher(allowed_paths).allows_folder(folder_path)

    def _folder_path_literal_needles(self, original_query: Optional[str], query: str) -> List[str]:
        needles: List[str] = []
        oq = str(original_query or "").strip()
        tq = str(query or "").strip()

        def _explicit_folder_label_needles(text: str) -> List[str]:
            raw_text = str(text or "").strip()
            if not raw_text:
                return []
            labels: List[str] = []
            seen_labels: set[str] = set()
            patterns = (
                r"\b(?:in|inside|under|within|from)\s+(?:the\s+)?(?:folder|directory|dir)\s+(?:named\s+|called\s+)?(?P<label>[^,;!?。\n\r]+)",
                r"\b(?:folder|directory|dir)\s+(?:named\s+|called\s+)?(?P<label>[^,;!?。\n\r]+)",
            )
            tail_stop = re.compile(
                r"\b(?:with|containing|that|which|where|whose|about|and|or|then|please|pls)\b",
                re.IGNORECASE,
            )
            generic = {
                "file", "files", "folder", "folders", "directory", "directories", "dir",
                "document", "documents", "doc", "docs", "image", "images", "audio", "video",
                "videos", "all", "the", "my", "a", "an",
            }

            def _add_label(value: str) -> None:
                label = str(value or "").strip(" \"'“”‘’`.,;:!?()[]{}<>")
                if not label:
                    return
                label = tail_stop.split(label, maxsplit=1)[0].strip(" \"'“”‘’`.,;:!?()[]{}<>")
                if not label:
                    return
                label = re.sub(r"\s+", " ", label)
                if len(label) < 2 or len(label) > 120:
                    return
                key = label.casefold()
                if key in generic or key in seen_labels:
                    return
                seen_labels.add(key)
                labels.append(key)
                base = os.path.basename(label.rstrip("/\\"))
                base_key = base.casefold().strip()
                if base_key and base_key != key and len(base_key) >= 2 and base_key not in generic and base_key not in seen_labels:
                    seen_labels.add(base_key)
                    labels.append(base_key)

            for pattern in patterns:
                for match in re.finditer(pattern, raw_text, flags=re.IGNORECASE):
                    _add_label(match.group("label"))
            return labels

        explicit_folder_needles = _explicit_folder_label_needles(f"{oq}\n{tq}")
        if any("\u4e00" <= ch <= "\u9fff" for ch in oq):
            _STOPS = {
                "帮我", "找一下", "找下", "查找", "搜索",
                "给我", "看看", "看下", "看一下",
                "请帮", "我的", "你们", "我们", "他们",
                "一下", "一些", "一个", "所有的", "全部的",
                "这个", "那个", "哪个", "什么", "怎么", "如何", "哪里", "哪些",
                "所有", "全部", "相关", "关于", "有关", "以及", "其他", "其它",
                "进行", "通过", "使用", "利用", "实现", "重要", "主要", "基本", "一般", "常见",
                "寻找", "获取", "下载", "打开", "浏览", "显示", "列出",
                "需要", "可能", "应该", "必须", "文件", "文档", "资料", "内容", "信息",
                "图片", "图像", "照片", "视频", "音频",
                "帮", "找", "查", "搜", "给", "看", "请", "我", "你", "他", "她", "下", "些",
                "的", "了", "吗", "呢", "吧", "啊", "哦", "是", "有", "在", "和", "与", "或",
                "这", "那", "哪",
            }
            cleaned_oq = oq
            for sw in sorted(_STOPS, key=len, reverse=True):
                cleaned_oq = cleaned_oq.replace(sw, " ")

            seen: set = set()

            def _add(s: str) -> None:
                if len(s) < 2:
                    return
                k = s.casefold()
                if k in seen:
                    return
                seen.add(k)
                needles.append(k)

            for seg in re.findall(r"[\u4e00-\u9fff]+", cleaned_oq):
                if len(seg) < 2:
                    continue
                _add(seg)
                upper = min(8, len(seg))
                for win in range(upper, 1, -1):
                    for i in range(0, len(seg) - win + 1):
                        _add(seg[i : i + win])
                        if len(needles) >= 72:
                            break
                    if len(needles) >= 72:
                        break
                if len(needles) >= 72:
                    break
            for label in explicit_folder_needles:
                _add(label)
                if len(needles) >= 96:
                    break
            _ASCII_STOPS = {
                "find", "search", "look", "show", "folder", "directory", "dir",
                "file", "files", "document", "documents", "doc", "docs",
                "image", "images", "photo", "photos", "picture", "pictures",
                "audio", "video", "videos", "recording", "recordings",
                "table", "tables", "spreadsheet", "spreadsheets",
            }
            for token in re.findall(r"[A-Za-z0-9][A-Za-z0-9._-]{2,}", oq):
                lowered = token.casefold()
                if lowered in _ASCII_STOPS or lowered.isdigit():
                    continue
                _add(lowered)
                if len(needles) >= 96:
                    break
            for surface in _build_bilingual_query_surfaces(tq, max_surfaces=4):
                if len(surface) < 5:
                    continue
                if any("\u4e00" <= ch <= "\u9fff" for ch in surface):
                    continue
                if not any(ch.isascii() and ch.isalpha() for ch in surface):
                    continue
                _add(surface)
                if len(needles) >= 112:
                    break
            return needles
        seen_non_cjk: set[str] = set()
        for label in explicit_folder_needles:
            if label and label not in seen_non_cjk:
                needles.append(label)
                seen_non_cjk.add(label)
        for surface in _build_bilingual_query_surfaces(tq, original_query or "", max_surfaces=4):
            key = surface.casefold()
            if len(surface) >= 5 and key not in seen_non_cjk:
                needles.append(key)
                seen_non_cjk.add(key)
        return needles

    @staticmethod
    def _folder_literal_needle_matches_path(needle: str, *, folder_name: str = "", folder_name_en: str = "", path: str = "") -> bool:
        nd = str(needle or "").casefold().strip()
        if not nd:
            return False
        names = [
            str(folder_name or "").casefold().strip(),
            str(folder_name_en or "").casefold().strip(),
        ]
        path_cf = str(path or "").casefold()
        if len(nd) < 4 and re.fullmatch(r"[a-z0-9][a-z0-9._-]*", nd):
            segments = [seg.casefold() for seg in re.split(r"[\\/]+", str(path or "")) if seg]
            return nd in names or nd in segments
        return any(nd in name for name in names if name) or nd in path_cf

    def _folder_expand_root_is_meaningful(self, norm_path: str) -> bool:
        p = str(norm_path or "").strip()
        if not p:
            return False
        try:
            bn = os.path.basename(p.rstrip(os.sep))
        except Exception:
            return False
        if self._folder_index_skip_basename(bn):
            return False
        try:
            home = os.path.normpath(os.path.expanduser("~")).rstrip(os.sep)
            if p.rstrip(os.sep) == home:
                return False
        except Exception:
            pass
        return True

    @staticmethod
    def _folder_literal_needle_matches_name_exact(needle: str, *, folder_name: str = "", folder_name_en: str = "") -> bool:
        nd = str(needle or "").casefold().strip()
        if not nd:
            return False
        return nd in {
            str(folder_name or "").casefold().strip(),
            str(folder_name_en or "").casefold().strip(),
        }

    def _finalize_folder_expand_roots_for_cjk_query(
        self, roots: List[str], path_needles: List[str]
    ) -> List[str]:
        import os
        if str(os.getenv("FOLDER_CHAIN_DISABLE_CJK_ROOT_NARROW", "") or "").strip().lower() in (
            "1",
            "true",
            "yes",
            "on",
        ):
            return roots
            
        if not path_needles:
            return roots
        norm_roots: List[str] = []
        for p in roots or []:
            try:
                np = os.path.normpath(os.path.expanduser(str(p or "").strip()))
            except Exception:
                continue
            if np:
                norm_roots.append(np)

        nds = [n for n in (path_needles or []) if isinstance(n, str) and len(n) >= 2]
        hit: List[str] = []
        if nds:
            for np in norm_roots:
                npc = np.casefold()
                if any(nd in npc for nd in nds):
                    hit.append(np)

        def _root_matches_needles(np: str) -> bool:
            try:
                bn = os.path.basename(str(np or "").rstrip(os.sep))
            except Exception:
                bn = ""
            return any(
                self._folder_literal_needle_matches_path(nd, folder_name=bn, path=np)
                for nd in nds
            )

        def _explicit_skip_root_allowed(np: str) -> bool:
            try:
                bn = os.path.basename(str(np or "").rstrip(os.sep))
            except Exception:
                return False
            if not bn or not self._folder_index_skip_basename(bn):
                return False
            return any(
                self._folder_literal_needle_matches_name_exact(nd, folder_name=bn)
                for nd in nds
            )

        def _shallowest_meaningful(pool: List[str]) -> List[str]:

            good = [
                p for p in pool
                if self._folder_expand_root_is_meaningful(p) or _explicit_skip_root_allowed(p)
            ]
            if not good:
                return []
            good.sort(key=lambda x: (len(x), x.count(os.sep)))
            return [good[0]]

        if not hit:
            return []
        hit = [np for np in hit if _root_matches_needles(np)]
        if not hit:
            return []
        hit_good = _shallowest_meaningful(hit)
        return hit_good if hit_good else []

    def _merge_folder_index_literal_hits(
        self,
        fc: Any,
        path_best_sim: Dict[str, float],
        needles: List[str],
        allowed_paths: Optional[List[str] | PathScopeMatcher],
        *,
        hit_score: float = 0.995,
        literal_hit_paths: Optional[set[str]] = None,
        exact_literal_hit_paths: Optional[set[str]] = None,
    ) -> None:
        if not needles or fc is None:
            return
        page = 3000
        offset = 0
        while True:
            batch = fc.get(include=["metadatas"], limit=page, offset=offset)
            metas = batch.get("metadatas") or []
            if not metas:
                break
            for meta in metas:
                fpath = str(meta.get("folder_path") or "")
                fname = str(meta.get("folder_name") or "")
                fname_en = str(meta.get("folder_name_en") or "")
                if not fpath or not self._folder_allowed(fpath, allowed_paths):
                    continue
                # Include folder_name_en so English queries match Chinese folder names
                blob_cf = f"{fname}\n{fname_en}\n{fpath}".casefold()
                hit = False
                exact_name_hit = False
                for nd in needles:
                    if len(nd) < 2:
                        continue
                    if self._folder_literal_needle_matches_path(
                        nd,
                        folder_name=fname,
                        folder_name_en=fname_en,
                        path=fpath,
                    ):
                        hit = True
                        if self._folder_literal_needle_matches_name_exact(nd, folder_name=fname, folder_name_en=fname_en):
                            exact_name_hit = True
                        break
                if hit:
                    score = 0.9995 if exact_name_hit else float(hit_score)
                    path_best_sim[fpath] = max(path_best_sim.get(fpath, 0.0), score)
                    if literal_hit_paths is not None:
                        literal_hit_paths.add(fpath)
                    if exact_name_hit and exact_literal_hit_paths is not None:
                        exact_literal_hit_paths.add(fpath)
            offset += len(metas)
            if len(metas) < page:
                break

    def _merge_parent_dirs_from_indexed_files_literal(
        self,
        path_best_sim: Dict[str, float],
        needles: List[str],
        allowed_paths: Optional[List[str] | PathScopeMatcher],
        *,
        hit_score: float = 0.99,
        literal_hit_paths: Optional[set[str]] = None,
        exact_literal_hit_paths: Optional[set[str]] = None,
    ) -> None:
        if not needles:
            return
        keyword_index = getattr(self, "_keyword_index_manager", None)
        if keyword_index is None:
            return
        records = keyword_index.list_records(allowed_paths=allowed_paths, require_current=False)  # type: ignore[arg-type]
        for rec in records:
            fp = str(rec.file_path or "")
            if not fp:
                continue
            try:
                cur = os.path.dirname(os.path.abspath(os.path.expanduser(fp)))
            except Exception:
                continue
            seen_dir: set = set()
            while cur and cur not in ("/", "\\"):
                if cur in seen_dir:
                    break
                seen_dir.add(cur)
                try:
                    norm = os.path.normpath(cur)
                except Exception:
                    break
                bn = os.path.basename(norm)
                blob_cf = f"{bn}\n{norm}".casefold()
                hit = False
                exact_name_hit = False
                for nd in needles:
                    if len(nd) < 2:
                        continue
                    if self._folder_literal_needle_matches_path(nd, folder_name=bn, path=norm):
                        hit = True
                        if self._folder_literal_needle_matches_name_exact(nd, folder_name=bn):
                            exact_name_hit = True
                        break
                if not hit or not bn:
                    parent = os.path.dirname(cur)
                    if parent == cur:
                        break
                    cur = parent
                    continue
                skip_name = self._folder_index_skip_basename(bn)
                if skip_name and not exact_name_hit:
                    parent = os.path.dirname(cur)
                    if parent == cur:
                        break
                    cur = parent
                    continue
                score = 0.9995 if exact_name_hit else float(hit_score)
                path_best_sim[norm] = max(path_best_sim.get(norm, 0.0), score)
                if literal_hit_paths is not None:
                    literal_hit_paths.add(norm)
                if exact_name_hit and exact_literal_hit_paths is not None:
                    exact_literal_hit_paths.add(norm)
                parent = os.path.dirname(cur)
                if parent == cur:
                    break
                cur = parent

    def collect_exact_folder_literal_roots(
        self,
        needles: List[str],
        allowed_paths: Optional[List[str] | PathScopeMatcher] = None,
        *,
        max_roots: int = 16,
    ) -> List[str]:
        if not needles:
            return []
        keyword_index = getattr(self, "_keyword_index_manager", None)
        if keyword_index is None:
            return []
        scope_matcher = ensure_path_scope_matcher(allowed_paths)
        roots: List[str] = []
        seen: set[str] = set()
        try:
            records = keyword_index.list_records(allowed_paths=scope_matcher, require_current=False)  # type: ignore[arg-type]
        except Exception:
            return []
        for rec in records:
            fp = str(rec.file_path or "")
            if not fp:
                continue
            try:
                cur = os.path.dirname(os.path.abspath(os.path.expanduser(fp)))
            except Exception:
                continue
            seen_dir: set[str] = set()
            while cur and cur not in ("/", "\\"):
                if cur in seen_dir:
                    break
                seen_dir.add(cur)
                try:
                    norm = os.path.normpath(cur)
                except Exception:
                    break
                bn = os.path.basename(norm)
                if (
                    bn
                    and self._folder_allowed(norm, scope_matcher)
                    and any(
                        len(nd) >= 2 and self._folder_literal_needle_matches_name_exact(nd, folder_name=bn)
                        for nd in needles
                    )
                ):
                    if norm not in seen:
                        seen.add(norm)
                        roots.append(norm)
                        if len(roots) >= max_roots:
                            return roots
                parent = os.path.dirname(cur)
                if parent == cur:
                    break
                cur = parent
        return roots

    def folder_semantic_search(
        self,
        query: str,
        original_query: Optional[str] = None,
        allowed_paths: Optional[List[str]] = None,
    ) -> List[Dict[str, Any]]:
        folder_results: List[Dict[str, Any]] = []
        fc = getattr(self, "folder_collection", None)
        if fc is None:
            return folder_results
        if not self._folder_index_pruned_once:
            self._folder_index_pruned_once = True
            try:
                self.prune_folder_index_orphans()
            except Exception as e:
                logger.warning(f"[folder_index] prune 跳过: {e}")
        self._maybe_auto_rebuild_folder_index()
        try:
            if fc.count() == 0:
                return folder_results
        except Exception:
            return folder_results

        try:
            try:
                vec_dist_max = float(os.getenv("FILEAGENT_FOLDER_VEC_DIST_MAX", "0.65") or "0.65")
            except ValueError:
                vec_dist_max = 0.65
            vec_dist_max = max(0.15, min(vec_dist_max, 0.95))

            matched_folder_paths: set = set()
            q_texts: List[str] = []
            q = str(query or "").strip()
            oq = str(original_query or "").strip()
            if q:
                q_texts.append(q)
            if oq and oq.lower() != q.lower():
                q_texts.append(oq)

            for qtxt in q_texts:
                try:
                    with self._embed_context(f"folder_index_query q={qtxt[:80]!r}"):
                        qemb = self._embed_query_text(qtxt)
                    n_fc = int(fc.count())
                    k = min(max(20, n_fc), 120)
                    res = fc.query(
                        query_embeddings=[qemb],
                        n_results=k,
                        include=["metadatas", "distances"],
                    )
                    dists = (res.get("distances") or [[]])[0]
                    metas = (res.get("metadatas") or [[]])[0]
                    for dist, meta in zip(dists, metas):
                        if meta is None:
                            continue
                        fpath = str(meta.get("folder_path") or "")
                        if not self._folder_allowed(fpath, allowed_paths):
                            continue
                        d = float(dist) if dist is not None else 1.0
                        if d <= vec_dist_max:
                            matched_folder_paths.add(fpath)
                except Exception as e:
                    logger.warning(f"folder_index 向量查询失败: {e}")

            needles: List[str] = []
            for src in (oq, q):
                if not src:
                    continue
                sl = src.lower()
                if sl not in needles:
                    needles.append(sl)
                cjk = "".join(ch for ch in src if "\u4e00" <= ch <= "\u9fff")
                if len(cjk) >= 2 and cjk.lower() not in needles:
                    needles.append(cjk.lower())
                if len(cjk) >= 2:
                    pref = cjk[:2].lower()
                    if pref not in needles:
                        needles.append(pref)

            page = 3000
            offset = 0
            while True:
                batch = fc.get(include=["metadatas"], limit=page, offset=offset)
                metas = batch.get("metadatas") or []
                if not metas:
                    break
                for meta in metas:
                    fpath = str(meta.get("folder_path") or "")
                    fname = str(meta.get("folder_name") or "")
                    if not fpath:
                        continue
                    if not self._folder_allowed(fpath, allowed_paths):
                        continue
                    fl = fname.lower()
                    pl = fpath.lower()
                    hit = False
                    for nd in needles:
                        if not nd:
                            continue
                        if nd in fl or nd in pl:
                            hit = True
                            break
                    if hit:
                        matched_folder_paths.add(fpath)
                offset += len(metas)
                if len(metas) < page:
                    break

            if not matched_folder_paths:
                return folder_results

            inventory_pack = self.indexed_file_inventory(
                allowed_paths=list(matched_folder_paths),
                limit=0,
                hydrate=True,
            )
            inventory_items = list(inventory_pack.get("files") or []) if inventory_pack.get("ready") else []

            existing_paths: set = set()
            for parent_dir in matched_folder_paths:
                folder_name = os.path.basename(parent_dir)
                if self._folder_index_skip_basename(folder_name):
                    continue
                prefix = parent_dir + os.sep
                added_count = 0
                for item in inventory_items:
                    m = dict(item.get("metadata") or {})
                    fp = str(item.get("file_path") or m.get("file_path") or "")
                    if not fp.startswith(prefix):
                        continue
                    if fp in existing_paths:
                        continue
                    raw_summary = item.get("doc_summary") or m.get("doc_summary", "")
                    if raw_summary and str(raw_summary).startswith("{") and '"doc_summary"' in str(
                        raw_summary
                    ):
                        try:
                            summary_dict = json.loads(raw_summary)
                            raw_summary = summary_dict.get("doc_summary", raw_summary)
                        except Exception:
                            pass
                    folder_results.append(
                        {
                            "text": item.get("text") or raw_summary or item.get("file_name") or "",
                            "metadata": m,
                            "distance": 0.0,
                            "file_name": item.get("file_name") or m.get("file_name", ""),
                            "file_path": fp,
                            "doc_summary": raw_summary,
                            "doc_category": item.get("doc_category") or self._meta_category_family(m),
                            "score": 1.0,
                            "is_folder_match": True,
                        }
                    )
                    existing_paths.add(fp)
                    added_count += 1
                    if added_count >= 20:
                        break
        except Exception as e:
            logger.error(f"Semantic folder matching failed: {e}")

        return folder_results

    def collect_folder_index_candidates(
        self,
        query: str,
        original_query: Optional[str] = None,
        allowed_paths: Optional[List[str] | PathScopeMatcher] = None,
    ) -> List[Dict[str, Any]]:
        rows: List[Dict[str, Any]] = []
        path_best_sim: Dict[str, float] = {}
        fc = getattr(self, "folder_collection", None)
        scope_matcher = ensure_path_scope_matcher(allowed_paths)

        if fc is not None:
            if not self._folder_index_pruned_once:
                self._folder_index_pruned_once = True
                try:
                    self.prune_folder_index_orphans()
                except Exception as e:
                    logger.warning(f"[folder_index] prune 跳过: {e}")
            self._maybe_auto_rebuild_folder_index()

        fc_nonempty = False
        if fc is not None:
            try:
                fc_nonempty = int(fc.count() or 0) > 0
            except Exception:
                fc_nonempty = False

        q0 = str(query or "").strip()
        oq0 = str(original_query or "").strip()
        path_needles = self._folder_path_literal_needles(oq0, q0)
        literal_hit_paths: set[str] = set()
        exact_literal_hit_paths: set[str] = set()

        try:
            if fc_nonempty and fc is not None:
                try:
                    vec_dist_max = float(os.getenv("FILEAGENT_FOLDER_VEC_DIST_MAX", "0.65") or "0.65")
                except ValueError:
                    vec_dist_max = 0.65
                vec_dist_max = max(0.15, min(vec_dist_max, 0.95))
                n_fc = int(fc.count())
                k = min(max(20, n_fc), 120) if n_fc > 0 else 0

                q_texts: List[Tuple[str, str]] = []
                seen_cf: set = set()
                for label, raw in (
                    ("retrieval", q0),
                    ("original", oq0),
                    *[
                        (f"surface_{idx}", surface)
                        for idx, surface in enumerate(
                            _build_bilingual_query_surfaces(q0, oq0, max_surfaces=4),
                            1,
                        )
                    ],
                ):
                    t = str(raw or "").strip()
                    if not t:
                        continue
                    ck = t.casefold()
                    if ck in seen_cf:
                        continue
                    seen_cf.add(ck)
                    q_texts.append((label, t))

                for which, qtxt in q_texts:
                    try:
                        with self._embed_context(
                            f"folder_index_candidate/{which} q={qtxt[:80]!r}"
                        ):
                            qemb = self._embed_query_text(qtxt)
                        if k <= 0:
                            continue
                        res = fc.query(
                            query_embeddings=[qemb],
                            n_results=k,
                            include=["metadatas", "distances"],
                        )
                        dists = (res.get("distances") or [[]])[0]
                        metas = (res.get("metadatas") or [[]])[0]
                        for dist, meta in zip(dists, metas):
                            if meta is None:
                                continue
                            fpath = str(meta.get("folder_path") or "")
                            if not self._folder_allowed(fpath, scope_matcher):
                                continue
                            d = float(dist) if dist is not None else 1.0
                            if d > vec_dist_max:
                                continue
                            sim = max(0.0, min(1.0, 1.0 - d))
                            path_best_sim[fpath] = max(path_best_sim.get(fpath, 0.0), sim)
                    except Exception as e:
                        logger.warning(f"folder_index_candidate 向量查询失败 ({which}): {e}")

                if path_needles:
                    self._merge_folder_index_literal_hits(
                        fc,
                        path_best_sim,
                        path_needles,
                        scope_matcher,
                        literal_hit_paths=literal_hit_paths,
                        exact_literal_hit_paths=exact_literal_hit_paths,
                    )

            if path_needles:
                self._merge_parent_dirs_from_indexed_files_literal(
                    path_best_sim,
                    path_needles,
                    scope_matcher,
                    literal_hit_paths=literal_hit_paths,
                    exact_literal_hit_paths=exact_literal_hit_paths,
                )

            if not path_best_sim:
                return []

            for fpath, sim in sorted(path_best_sim.items(), key=lambda x: -x[1]):
                folder_name = os.path.basename(fpath)
                literal_hit = fpath in literal_hit_paths
                if self._folder_index_skip_basename(folder_name) and not literal_hit:
                    continue
                if not self._folder_allowed(fpath, scope_matcher):
                    continue
                line = f"Folder: {folder_name}\nPath: {fpath}"
                row = {
                    "text": line,
                    "file_path": fpath,
                    "file_name": folder_name,
                    "folder_path": fpath,
                    "doc_summary": "",
                    "doc_category": "other",
                    "metadata": {},
                    "distance": max(0.0, 1.0 - float(sim)),
                    "score": float(sim),
                    "_folder_index_row": True,
                }
                if literal_hit:
                    row["_folder_literal_hit"] = True
                    row["_direct_score"] = 95
                if fpath in exact_literal_hit_paths:
                    row["_folder_exact_name_hit"] = True
                rows.append(row)
        except Exception as e:
            logger.error(f"collect_folder_index_candidates failed: {e}")
            return []

        cap = max(10, int(os.getenv("FOLDER_INDEX_CANDIDATE_CAP", "80") or 80))
        return rows[:cap]

    def expand_folder_paths_to_chain_sources(
        self,
        folder_paths: List[str],
        *,
        max_per_folder: int = 200,
        allow_skip_basename_roots: bool = False,
        relax_ignore_rules: bool = False,
    ) -> List[Dict[str, Any]]:
        out: List[Dict[str, Any]] = []
        if not folder_paths:
            return out
        try:
            m_raw = int(max_per_folder)
        except (TypeError, ValueError):
            m_raw = 200
        unlimited = m_raw <= 0
        m = max(1, m_raw) if not unlimited else 0

        inventory_pack = self.indexed_file_inventory(
            allowed_paths=folder_paths,
            limit=0,
            hydrate=True,
        )
        inventory_items = list(inventory_pack.get("files") or []) if inventory_pack.get("ready") else []

        try:
            fallback_cap = max(0, int(os.getenv("FOLDER_CHAIN_FALLBACK_MAX_FILES", "64") or 64))
        except (TypeError, ValueError):
            fallback_cap = 64
        try:
            fallback_max_depth = max(0, int(os.getenv("FOLDER_CHAIN_FALLBACK_MAX_DEPTH", "4") or 4))
        except (TypeError, ValueError):
            fallback_max_depth = 4
        try:
            fallback_max_dirs = max(1, int(os.getenv("FOLDER_CHAIN_FALLBACK_MAX_DIRS", "256") or 256))
        except (TypeError, ValueError):
            fallback_max_dirs = 256

        existing_paths: set = set()
        for parent_dir in folder_paths:
            parent_dir = os.path.normpath(os.path.expanduser(str(parent_dir or "").strip()))
            if not parent_dir:
                continue
            folder_name = os.path.basename(parent_dir)
            if self._folder_index_skip_basename(folder_name) and not allow_skip_basename_roots:
                continue
            prefix = parent_dir + os.sep
            norm_parent = parent_dir.rstrip(os.sep)
            matches: List[Tuple[str, Dict[str, Any], Dict[str, Any]]] = []
            for item in inventory_items:
                meta = dict(item.get("metadata") or {})
                fp = str(item.get("file_path") or meta.get("file_path") or "")
                if not fp.startswith(prefix):
                    continue
                if fp.rstrip(os.sep) == norm_parent:
                    continue
                matches.append((fp, meta, item))
            matches.sort(key=lambda x: x[0])
            added_count = 0
            for fp, meta, item in matches:
                if fp in existing_paths:
                    continue
                if not unlimited and added_count >= m:
                    break
                doc0 = str(item.get("text") or "")
                raw_summary = meta.get("doc_summary", "")
                if raw_summary and str(raw_summary).startswith("{") and '"doc_summary"' in str(raw_summary):
                    try:
                        summary_dict = json.loads(raw_summary)
                        raw_summary = summary_dict.get("doc_summary", raw_summary)
                    except Exception:
                        pass
                try:
                    rel = os.path.relpath(fp, parent_dir)
                except (ValueError, OSError):
                    rel = os.path.basename(fp)
                rel_disp = rel.replace("\\", "/")
                out.append(
                    {
                        "text": doc0,
                        "metadata": meta,
                        "distance": 0.0,
                        "file_name": meta.get("file_name", ""),
                        "file_path": fp,
                        "doc_summary": raw_summary,
                        "doc_category": self._meta_category_family(meta),
                        "score": 1.0,
                        "rerank_score": 4.0,
                        "is_folder_chain_match": True,
                        "folder_chain_root": parent_dir,
                        "folder_chain_relative_path": rel_disp,
                    }
                )
                existing_paths.add(fp)
                added_count += 1

            if fallback_cap <= 0:
                continue
            remaining = fallback_cap if unlimited else max(0, m - added_count)
            if remaining <= 0:
                continue
            disk_fallback_items = self._expand_folder_paths_to_chain_sources_from_disk(
                parent_dir,
                existing_paths=existing_paths,
                max_files=remaining,
                max_depth=fallback_max_depth,
                max_dirs=fallback_max_dirs,
                relax_ignore_rules=relax_ignore_rules,
            )
            if disk_fallback_items:
                out.extend(disk_fallback_items)
                existing_paths.update(
                    str(item.get("file_path") or "").strip()
                    for item in disk_fallback_items
                    if str(item.get("file_path") or "").strip()
                )
        return out

    def _expand_folder_paths_to_chain_sources_from_disk(
        self,
        parent_dir: str,
        *,
        existing_paths: set,
        max_files: int,
        max_depth: int,
        max_dirs: int,
        relax_ignore_rules: bool = False,
    ) -> List[Dict[str, Any]]:
        if max_files <= 0 or not os.path.isdir(parent_dir):
            return []

        import config.settings as settings

        out: List[Dict[str, Any]] = []
        seen_dirs = 0
        parent_dir = os.path.normpath(os.path.expanduser(str(parent_dir or "").strip()))
        if not parent_dir:
            return []

        for root, subdirs, files in os.walk(parent_dir):
            seen_dirs += 1
            if seen_dirs > max_dirs:
                logger.info(
                    "[folder_chain] disk fallback dir budget reached: root=%r max_dirs=%d",
                    parent_dir,
                    max_dirs,
                )
                break

            try:
                rel_dir = os.path.relpath(root, parent_dir)
            except (ValueError, OSError):
                rel_dir = "."
            depth = 0 if rel_dir in {".", ""} else rel_dir.count(os.sep) + 1
            if depth >= max_depth:
                subdirs[:] = []
            else:
                filtered_subdirs: List[str] = []
                for name in sorted(subdirs):
                    if name in settings.IGNORE_PATTERNS or name.startswith(".") or name.endswith(".app"):
                        continue
                    filtered_subdirs.append(name)
                subdirs[:] = filtered_subdirs

            for file_name in sorted(files):
                fp = os.path.abspath(os.path.join(root, file_name))
                if fp in existing_paths or self._folder_chain_should_ignore_disk_file(
                    fp,
                    relax_rules=relax_ignore_rules,
                ):
                    continue
                try:
                    rel = os.path.relpath(fp, parent_dir)
                except (ValueError, OSError):
                    rel = os.path.basename(fp)
                rel_disp = rel.replace("\\", "/")
                ext = os.path.splitext(file_name)[1].lower()
                doc_category = self._infer_doc_category_family(
                    "",
                    file_name=file_name,
                    file_ext=ext,
                )
                meta = {
                    "file_path": fp,
                    "file_name": file_name,
                    "doc_category": doc_category,
                    "doc_category_family": doc_category,
                    "_folder_chain_disk_fallback": True,
                }
                out.append(
                    {
                        "text": rel_disp,
                        "metadata": meta,
                        "distance": 0.0,
                        "file_name": file_name,
                        "file_path": fp,
                        "doc_summary": "",
                        "doc_category": doc_category,
                        "score": 0.98,
                        "rerank_score": 3.8,
                        "is_folder_chain_match": True,
                        "folder_chain_root": parent_dir,
                        "folder_chain_relative_path": rel_disp,
                    }
                )
                if len(out) >= max_files:
                    logger.info(
                        "[folder_chain] disk fallback file budget reached: root=%r max_files=%d",
                        parent_dir,
                        max_files,
                    )
                    return out
        if out:
            logger.info(
                "[folder_chain] disk fallback added %d file(s) under %r",
                len(out),
                parent_dir,
            )
        return out

    def vector_search(self, query: str, n_results: int = 30, allowed_paths: Optional[List[str]] = None, category_filter: Optional[str] = None, keyword: Optional[str] = None, folder: Optional[str] = None, original_query: Optional[str] = None, file_extensions: Optional[List[str]] = None) -> List[Dict[str, Any]]:
        try:
            target_category = _normalize_category_en(category_filter or "", default="")
            if target_category in {"", "all", "unknown"}:
                target_category = ""
            scope_matcher = ensure_path_scope_matcher(allowed_paths)
            clean_kw_str = ""
            if keyword:
                import ast
                try:
                    parsed_kw = ast.literal_eval(keyword)
                    if isinstance(parsed_kw, list):
                        clean_kw_str = " ".join(str(k) for k in parsed_kw)
                    else:
                        clean_kw_str = str(keyword)
                except Exception:
                    clean_kw_str = keyword.replace("[", "").replace("]", "").replace("'", "").replace('"', "").replace(",", " ")
            
            queries_to_embed = _build_bilingual_query_surfaces(
                query,
                original_query or "",
                clean_kw_str,
                max_surfaces=5,
            )
            if not queries_to_embed:
                queries_to_embed = [query or clean_kw_str]
                
            with self._embed_context(f"query_embedding query_len={len(query or '')} keyword_len={len(clean_kw_str or '')}"):
                query_embeddings = [self._embed_query_text(q) for q in queries_to_embed]
            
            search_k = n_results
            # allowed_paths:
            if allowed_paths is not None:
                search_k = max(n_results * 3, 100) 
            
            where_clause = None
            if target_category:
                from core.retrieval.category_engine import get_compatible_categories

                aliases = set(get_compatible_categories(target_category) or {target_category})

                for cat in list(aliases):
                    for k, v in _CATEGORY_ALIASES_TO_EN.items():
                        if v == cat:
                            aliases.add(k)
                    for canonical, hints in _CATEGORY_KEYWORD_HINTS:
                        if canonical == cat:
                            aliases.update(hints)
                aliases = sorted(aliases)
                
                if len(aliases) == 1:
                    where_clause = {"doc_category": aliases[0]}
                elif len(aliases) > 1:
                    where_clause = {"doc_category": {"$in": aliases}}

            query_kwargs = {
                "query_embeddings": query_embeddings,
                "n_results": search_k,
                "include": ["documents", "metadatas", "distances"]
            }
            if where_clause:
                query_kwargs["where"] = where_clause
            
            results = self.collection.query(**query_kwargs)
            
            # ========= BM25 + RRF Hybrid Retrieval =========
            # Replace old substring-match Lexical Fallback with BM25 scoring.
            # 1. Pull all metadata once (reused by BM25)
            # 2. Score with BM25Okapi (jieba-tokenized, bilingual, category-translated)
            # 3. Fuse Vector + BM25 rankings via RRF (done after vector result collection)
            import re as _re_lex

            _bm25_query_blob = build_lookup_blob(query, clean_kw_str, original_query or "")
            _bm25_query_text = f"{query} {clean_kw_str} {original_query or ''} {_bm25_query_blob}".strip()
            _bm25_query_tokens = _tokenize_for_bm25(_bm25_query_text)
            _, _media_bm25_expanded_terms, _wants_media_bm25 = self._media_content_query_terms(
                _bm25_query_text,
                _bm25_query_tokens,
                category_filter=target_category,
                file_extensions=file_extensions,
            )
            if _wants_media_bm25 and _media_bm25_expanded_terms:
                _seen_bm25_tokens = set(_bm25_query_tokens)
                _bm25_query_tokens.extend(
                    term for term in _media_bm25_expanded_terms
                    if term and term not in _seen_bm25_tokens
                )

            _bm25_hits: List[Tuple[str, str, float]] = []
            _bm25_meta_by_path: Dict[str, Dict[str, Any]] = {}
            _bm25_doc_by_id: Dict[str, str] = {}
            try:
                # 🔥 Use cached BM25 index instead of rebuilding from scratch each query
                if _bm25_query_tokens and _HAS_BM25:
                    _bm25_hits = self._get_or_build_bm25(
                        query_tokens=_bm25_query_tokens,
                        allowed_paths=scope_matcher,
                        category_filter=target_category,
                        file_extensions=file_extensions,
                    )
                    _bm25_fetch_ids = list(dict.fromkeys([bid for bid, _, _ in _bm25_hits[:30]]))
                    if _bm25_fetch_ids:
                        _doc_batch = self.collection.get(ids=_bm25_fetch_ids, include=["documents", "metadatas"])
                        for mid, doc, meta in zip(
                            _doc_batch.get("ids") or [],
                            _doc_batch.get("documents") or [],
                            _doc_batch.get("metadatas") or [],
                        ):
                            _mid = str(mid or "")
                            _bm25_doc_by_id[_mid] = doc or ""
                            if isinstance(meta, dict):
                                _fp = str(meta.get("file_path") or "").strip()
                                if _fp:
                                    _bm25_meta_by_path[_fp] = meta

                    logger.info(
                        f"[BM25] scored {len(_bm25_hits)} files with tokens={_bm25_query_tokens[:8]}... "
                        f"top3={[(os.path.basename(fp), f'{sc:.2f}') for _, fp, sc in _bm25_hits[:3]]}"
                    )
                else:
                    logger.info(f"[BM25] skipped (has_bm25={_HAS_BM25}, tokens={len(_bm25_query_tokens)})")
            except Exception as _bm25_err:
                logger.warning(f"[BM25] index build failed: {_bm25_err}")
            # =======================================================
            
            search_results = []
            seen_chunks = set()
            
            for q_idx in range(len(results['documents'])):
                for doc, meta, dist in zip(
                    results['documents'][q_idx],
                    results['metadatas'][q_idx],
                    results['distances'][q_idx]
                ):
                    file_path = meta.get('file_path', '')
                    chunk_id = f"{file_path}_{doc[:50]}"
                    if chunk_id in seen_chunks:
                        continue
                    
                    if not self._is_path_allowed(file_path, scope_matcher):
                        continue
                    
                    if folder:
                        parent_dir = os.path.dirname(file_path).lower()
                        folder_blob = " ".join(
                            [
                                parent_dir,
                                os.path.basename(parent_dir),
                                str(meta.get("folder_name_en") or ""),
                                str(meta.get("lookup_aliases") or ""),
                            ]
                        ).lower()
                        if folder.lower() not in folder_blob:
                            continue

                    ext_matched = False
                    if file_extensions:
                        ext = os.path.splitext(file_path)[1].lower()
                        if not any(ext == e or file_path.lower().endswith(e) for e in file_extensions):
                            continue
                        ext_matched = True

                    if target_category:
                        hit_cat = self._meta_category_family(meta)
                        if ext_matched and target_category in ["document", "image", "video", "audio", "other", "unknown"]:
                            pass
                        elif hit_cat not in aliases:
                            continue

                    seen_chunks.add(chunk_id)
                    similarity = 1 - dist
                    role_multiplier = self._role_score_multiplier(target_category=target_category, meta=meta)
                    similarity *= role_multiplier
                    
                    is_lexical = False
                    if '_is_strong_lexical' in meta:
                        kw = meta['_is_strong_lexical']
                        if kw.lower() in meta.get('file_name', '').lower() or kw.lower() in file_path.lower():
                            is_lexical = True

                    # Inject media anchors to doc_summary for better LLM context
                    doc_summary = meta.get('doc_summary', '')
                    ctype = meta.get('chunk_type', '')
                    if ctype in ("asr_transcript", "asr_segment", "keyframe", "interval_summary", "interval_visual"):
                        has_cjk = any('\u4e00' <= ch <= '\u9fff' for ch in query)
                        
                        def _fmt_time(ts_val):
                            try:
                                m, s = divmod(int(float(ts_val)), 60)
                                h, m = divmod(m, 60)
                                if has_cjk:
                                    if h > 0: return f"约 {h}小时{m}分{s}秒"
                                    elif m > 0: return f"约 {m}分{s}秒"
                                    else: return f"约 {s}秒"
                                else:
                                    if h > 0: return f"~ {h}h {m}m {s}s"
                                    elif m > 0: return f"~ {m}m {s}s"
                                    else: return f"~ {s}s"
                            except:
                                return ""

                        if ctype in ("asr_transcript", "asr_segment") and "asr_start_sec" in meta:
                            t_str = _fmt_time(meta["asr_start_sec"])
                            if t_str: doc_summary = f"🎙️ [{t_str}] {doc_summary}"
                        elif ctype == "keyframe" and "keyframe_time_sec" in meta:
                            t_str = _fmt_time(meta["keyframe_time_sec"])
                            suffix = " 画面" if has_cjk else " visual"
                            if t_str: doc_summary = f"🖼️ [{t_str}{suffix}] {doc_summary}"
                        elif ctype == "interval_summary" and "interval_start_sec" in meta:
                            start_str = _fmt_time(meta["interval_start_sec"])
                            end_str = _fmt_time(meta.get("interval_end_sec", meta["interval_start_sec"]))
                            if start_str and end_str:
                                prefix = "🧭" if has_cjk else "🧭"
                                doc_summary = f"{prefix} [{start_str} - {end_str}] {doc_summary}"
                        elif ctype == "interval_visual" and ("interval_visual_time_sec" in meta or "keyframe_time_sec" in meta):
                            t_raw = meta.get("interval_visual_time_sec", meta.get("keyframe_time_sec"))
                            t_str = _fmt_time(t_raw)
                            suffix = " 区间画面" if has_cjk else " interval visual"
                            if t_str:
                                doc_summary = f"🖼️ [{t_str}{suffix}] {doc_summary}"

                    search_results.append({
                        'text': doc,
                        'metadata': meta,
                        'distance': dist,
                        'file_name': meta.get('file_name', ''),
                        'file_path': file_path,
                        'doc_summary': doc_summary,
                        'doc_category': self._meta_category_family(meta),
                        'doc_category_leaf': self._meta_category_leaf(meta),
                        'doc_role': self._meta_doc_role(meta),
                        'score': similarity,
                        '_is_lexical_hit': is_lexical,
                        '_direct_score': 100 if is_lexical else 0,
                        '_role_multiplier': role_multiplier,
                    })
            # --- RRF Fusion: Merge Vector + BM25 rankings ---
            if _bm25_hits:
                # Build vector ranked list (best score per file, descending)
                _vector_file_best: Dict[str, float] = {}
                for sr in search_results:
                    fp = sr.get('file_path', '')
                    sc = float(sr.get('score', 0.0))
                    if sc > _vector_file_best.get(fp, -1.0):
                        _vector_file_best[fp] = sc

                _vec_ranked_tuples = sorted(_vector_file_best.items(), key=lambda x: x[1], reverse=True)
                _bm25_ranked_tuples = [(fp, sc) for _, fp, sc in _bm25_hits]

                # Fuse via RRF
                _rrf_scores = _rrf_fuse(_vec_ranked_tuples, _bm25_ranked_tuples, k=60)
                _rrf_score_map: Dict[str, float] = {fp: sc for fp, sc in _rrf_scores}
                _bm25_score_map: Dict[str, float] = {fp: sc for _, fp, sc in _bm25_hits}

                # Overwrite scores in search_results with RRF scores
                for sr in search_results:
                    fp = sr.get('file_path', '')
                    rrf_sc = _rrf_score_map.get(fp, 0.0)
                    bm25_sc = _bm25_score_map.get(fp, 0.0)
                    # RRF scores are ~0.01-0.03; scale up for downstream compatibility
                    role_multiplier = float(sr.get('_role_multiplier', 1.0) or 1.0)
                    sr['score'] = rrf_sc * 30.0 * role_multiplier
                    sr['_rrf_score'] = rrf_sc
                    sr['_bm25_score'] = bm25_sc
                    sr['_vector_score'] = _vector_file_best.get(fp, 0.0)
                    if bm25_sc > 0:
                        sr['_is_lexical_hit'] = True
                        sr['_direct_score'] = min(100, int(bm25_sc * 10))

                # Inject BM25-only hits that weren't in vector results
                _existing_paths = {sr.get('file_path') for sr in search_results}
                _bm25_top_score = _bm25_hits[0][2] if _bm25_hits else 1.0
                for _bm25_rank, (_bid, _bfp, _bsc) in enumerate(_bm25_hits[:30], start=1):
                    if _bfp in _existing_paths:
                        continue
                    _bmeta = _bm25_meta_by_path.get(_bfp, {})
                    _bdoc = _bm25_doc_by_id.get(_bid, "")
                    _rrf_sc_bm25_side = 1.0 / (60 + _bm25_rank)
                    _bm25_dominant = _bsc >= _bm25_top_score * 0.5
                    _rrf_sc = _rrf_sc_bm25_side + (1.0 / (60 + 10) if _bm25_dominant else 0.0)
                    _role_multiplier = self._role_score_multiplier(target_category=target_category, meta=_bmeta)
                    _injected_score = max(_rrf_sc * 30.0 * _role_multiplier, 1.1 if _bm25_dominant else 1.0)
                    search_results.append({
                        'text': _bdoc,
                        'metadata': _bmeta,
                        'distance': 0.5,
                        'file_name': _bmeta.get('file_name', os.path.basename(_bfp)),
                        'file_path': _bfp,
                        'doc_summary': _bmeta.get('doc_summary', ''),
                        'doc_category': self._meta_category_family(_bmeta),
                        'doc_category_leaf': self._meta_category_leaf(_bmeta),
                        'doc_role': self._meta_doc_role(_bmeta),
                        'score': _injected_score,
                        '_is_lexical_hit': True,
                        '_direct_score': min(100, int(_bsc * 10)),
                        '_rrf_score': _rrf_sc,
                        '_bm25_score': _bsc,
                        '_vector_score': 0.0,
                        '_role_multiplier': _role_multiplier,
                    })
                    _existing_paths.add(_bfp)

                logger.info(
                    f"[RRF] fused {len(_rrf_scores)} unique files. "
                    f"top5={[(os.path.basename(fp), f'{sc:.4f}') for fp, sc in _rrf_scores[:5]]}"
                )

            search_results_sorted = sorted(search_results, key=lambda x: x['score'], reverse=True)
            search_results = search_results_sorted[:n_results]

            if _bm25_hits:
                _in_results_fps = {sr.get('file_path') for sr in search_results}
                _bm25_rescued = 0
                # Fix 5 (P2): adaptive rescue threshold = max(bm25_score) * 0.3.
                # Replaces hardcoded >20.0, which could miss all hits in small corpora
                # where top BM25 score may be <20 (e.g. short-document collections).
                _all_bm25_scores = [float(sr.get('_bm25_score', 0) or 0) for sr in search_results_sorted[n_results:]]
                # Adaptive: 30% of the top BM25 score in the overflow pool, floor=5.0, hard-cap=20.0
                # This ensures rescue fires proportionally even in small corpora where top score < 20.
                _bm25_rescue_threshold = (max(_all_bm25_scores) * 0.3 if _all_bm25_scores else 20.0)
                _bm25_rescue_threshold = max(5.0, min(_bm25_rescue_threshold, 20.0))
                for sr in search_results_sorted[n_results:]:
                    if sr.get('file_path') in _in_results_fps:
                        continue
                    _sr_bm25 = float(sr.get('_bm25_score', 0) or 0)
                    if _sr_bm25 > _bm25_rescue_threshold:
                        # Boost score > 1.0 so rescue hits survive the final sort vs supplement files (score=1.0)
                        sr['score'] = max(float(sr.get('score', 0) or 0), 1.1)
                        search_results.append(sr)
                        _in_results_fps.add(sr.get('file_path'))
                        _bm25_rescued += 1
                        if _bm25_rescued >= 5:
                            break
                if _bm25_rescued:
                    logger.info(f"[BM25-rescue] 挽救 {_bm25_rescued} 条被截断的 BM25 高分命中 (threshold={_bm25_rescue_threshold:.1f})")
            
            # --- 2. Query-time metadata supplement recall intentionally disabled ---
            # Filename/path/alias surfaces are indexed into the BM25 sidecar at ingest
            # time. General search should stay on indexed vector + indexed lexical
            # retrieval; doing a full metadata scan here can block short entity
            # searches before the semantic route even has a chance to return.

            # Fix 1 (P0): Add _bypass_category_filter=True so these results are NEVER
            #   discarded by the category strict-filter downstream (doc_category="other"
            #   would otherwise be silently dropped).
            # Fix 2 (P0): Expanded type detection to cover address/birthday/social_media/website.
            if hasattr(self, "personal_info_db"):
                query_lower = (query or "").lower()
                orig_lower = (original_query or "").lower()
                combined_q = f"{query_lower} {clean_kw_str} {orig_lower}"

                detected_types = []
                if any(k in combined_q for k in ["电话", "号码", "手机", "phone", "mobile", "tel", "telephone"]):
                    detected_types.append("phone")
                if any(k in combined_q for k in ["邮箱", "邮件", "email", "e-mail", "mail", "联系方式"]):
                    detected_types.append("email")
                address_like = any(
                    k in combined_q
                    for k in [
                        "地址",
                        "住址",
                        "家庭住址",
                        "居住",
                        "居住地",
                        "住所",
                        "住宅",
                        "address",
                        "location",
                        "home address",
                        "residence",
                        "residential address",
                    ]
                ) or bool(
                    re.search(
                        r"\b(?:his|her|their|its)\s+(?:home|residence|residential\s+address)\b|"
                        r"\bwhere.{0,32}\b(?:he|she|they|his|her|their).{0,32}"
                        r"(?:live|lives|reside|resides|home|address|location)\b|"
                        r"(?:他|她|这个人|这位|该候选人|候选人).{0,12}"
                        r"(?:住哪|住在哪里|住址|家庭住址|居住地|地址|家在哪|家在哪里)",
                        combined_q,
                        re.IGNORECASE,
                    )
                )
                if address_like:
                    detected_types.append("address")
                if any(k in combined_q for k in ["生日", "出生日期", "birthday", "date of birth", "dob"]):
                    detected_types.append("birthday")
                if any(k in combined_q for k in ["微信", "wechat", "qq", "linkedin", "twitter", "instagram",
                                                   "社交", "social", "网站", "website", "homepage"]):
                    detected_types.append("social_media")
                if any(k in combined_q for k in ["api key", "apikey", "密钥", "token", "api_key"]):
                    detected_types.append("api_key")
                if any(k in combined_q for k in ["pwd", "密码", "password"]):
                    detected_types.append("password")
                if any(k in combined_q for k in ["身份证", "证件号", "id card", "passport", "护照"]):
                    detected_types.append("id_card")

                if detected_types:
                    try:
                        db_filter_term = clean_kw_str.strip()
                        if not db_filter_term:
                            # Fallback: strip known attribute words from original query, keep the name part
                            import re as _pi_re
                            _attr_strip = _pi_re.compile(
                                r"(的|是多少|是什么|怎么|号码|地址|邮箱|手机|电话|联系方式|生日|密码|\s*'s\s*(phone|email|address|birthday|mobile|contact).*)", _pi_re.IGNORECASE
                            )
                            db_filter_term = _attr_strip.sub("", query).strip()
                        pi_res = self.personal_info_db.search(query=db_filter_term, types=detected_types, limit=10)
                        for p in pi_res:
                            ptype = p.get("info_type", "")
                            pdesc = p.get("description", "")
                            pcontent = p.get("content", "")
                            powner = p.get("owner_name", "")
                            src_file = p.get("source_file", "")
                            if allowed_paths is not None and not self._is_path_allowed(src_file, scope_matcher):
                                continue
                            src_name = p.get("source_file_name", os.path.basename(src_file) if src_file else "")

                            syn_text = f"【系统提取的敏感信息】\n类型: {ptype}\n所属人: {powner}\n描述: {pdesc}\n具体内容: {pcontent}\n来源文件: {src_name}"

                            pi_doc = {
                                "text": syn_text,
                                "metadata": {
                                    "file_name": src_name,
                                    "file_path": src_file,
                                    "doc_category": "other",
                                    "doc_summary": syn_text,
                                },
                                "distance": 0.0,
                                "file_name": src_name,
                                "file_path": src_file,
                                "doc_summary": syn_text,
                                "doc_category": "other",
                                "score": 20.0,
                                "_is_lexical_hit": True,
                                "_direct_score": 100,
                                "_bypass_category_filter": True,  # Fix 1: bypass strict category downstream
                            }
                            search_results.append(pi_doc)
                        if pi_res:
                            logger.info(
                                f"[PersonalInfoDB] Injected {len(pi_res)} for types={detected_types}, "
                                f"filter_term_chars={len(db_filter_term or '')}"
                            )
                    except Exception as e:
                        logger.warning(f"[PersonalInfoDB] Error injecting synthetic results: {e}")

            if _bm25_hits:
                # Build a score-boost map: file_path -> boosted_score for top BM25 hits
                # Only boost when BM25 top score is strong enough (>= 20.0 absolute threshold).
                # This prevents content-only queries (where BM25 scores are uniformly low/scattered)
                # from boosting unrelated files that happen to share a common token.
                _bm25_top_sc = _bm25_hits[0][2] if _bm25_hits else 0.0
                _BM25_BOOST_MIN_SCORE = 20.0  # require at least this absolute BM25 score to trigger boost
                _bm25_boost_map: Dict[str, float] = {}
                _bm25_force_injected = 0
                if _bm25_top_sc >= _BM25_BOOST_MIN_SCORE:
                    for _rank_i, (_fbid, _fbfp, _fbsc) in enumerate(_bm25_hits[:10]):
                        if _fbsc < _bm25_top_sc * 0.5:
                            break
                        # Assign a high score: 2.0 for rank1, decreasing by rank
                        _bm25_boost_map[_fbfp] = 2.0 - (_rank_i * 0.1)

                if _bm25_boost_map:
                    # Step 1: Boost score for BM25 strong hits already in search_results
                    _existing_fps = set()
                    for _sr in search_results:
                        _fp = _sr.get('file_path')
                        if _fp and _fp in _bm25_boost_map:
                            _sr['score'] = max(float(_sr.get('score', 0) or 0), _bm25_boost_map[_fp])
                            _sr['_is_lexical_hit'] = True
                            _sr['_bm25_score'] = max(float(_sr.get('_bm25_score', 0) or 0),
                                                     next(sc for _, fp, sc in _bm25_hits if fp == _fp))
                            _existing_fps.add(_fp)

                    # Step 2: Inject BM25 strong hits NOT yet in search_results
                    for _rank_i, (_fbid, _fbfp, _fbsc) in enumerate(_bm25_hits[:10]):
                        if _fbfp in _existing_fps:
                            continue
                        if _fbsc < _bm25_top_sc * 0.5:
                            break
                        _fbmeta = _bm25_meta_by_path.get(_fbfp, {})
                        _fbdoc = _bm25_doc_by_id.get(_fbid, "")
                        search_results.append({
                            'text': _fbdoc,
                            'metadata': _fbmeta,
                            'distance': 0.5,
                            'file_name': _fbmeta.get('file_name', os.path.basename(_fbfp)),
                            'file_path': _fbfp,
                            'doc_summary': _fbmeta.get('doc_summary', ''),
                            'doc_category': self._meta_category_family(_fbmeta),
                            'score': 2.0 - (_rank_i * 0.1),
                            '_is_lexical_hit': True,
                            '_direct_score': min(100, int(_fbsc * 10)),
                            '_bm25_score': _fbsc,
                            '_vector_score': 0.0,
                        })
                        _existing_fps.add(_fbfp)
                        _bm25_force_injected += 1

                    logger.info(f"[BM25-force] boosted {len(_bm25_boost_map)} BM25 strong hits, injected {_bm25_force_injected} new")

            search_results = sorted(search_results, key=lambda x: x.get("score", 0.0), reverse=True)[:n_results]
            
            # Keep default logs free of user query text and absolute local paths.
            logger.debug("Vector search completed: query_chars=%s category_filter=%s", len(query or ""), target_category or "(none)")
            if not search_results:
                logger.debug("未找到结果！")
            else:
                logger.debug(
                    f"找到 {len(search_results)} 个结果（已应用路径/文件夹/分类等过滤；分数为向量相似度，非最终 rerank）："
                )
                for i, res in enumerate(search_results):
                    logger.debug(
                        f"[{i+1}] 得分: {res.get('score', 0):.4f} | 文件: {os.path.basename(str(res.get('file_path') or ''))}"
                    )
            
            return search_results
        except Exception as e:
            logger.error(f"向量检索失败: {e}")
            return []
    
    def rerank(self, query: str, documents: List[Dict[str, Any]], top_k: int = 10) -> List[Dict[str, Any]]:
        if not documents:
            return []

        with self._reranker_lock:
            reranker = getattr(self, "reranker", None)
            if not reranker:
                try:
                    for d in documents:
                        if "rerank_score" not in d:
                            vec_score = float(d.get("score", 0.0) or 0.0)
                            # vec=0.85 → +1.0，vec=0.75 → 0.0，vec=0.65 → -1.0
                            d["rerank_score"] = (vec_score - 0.75) * 10.0
                    return sorted(documents, key=lambda x: float(x.get("rerank_score", 0.0) or 0.0), reverse=True)[:top_k]
                except Exception:
                    return documents[:top_k]
            return self._rerank_with_loaded_model(query, documents, top_k, reranker)

    def _rerank_with_loaded_model(
        self,
        query: str,
        documents: List[Dict[str, Any]],
        top_k: int,
        reranker: Any,
    ) -> List[Dict[str, Any]]:
        try:
            _rerank_input_cap = max(30, min(top_k, len(documents)))
            # Ensure highest-scoring candidates are included in the rerank window
            documents_sorted = sorted(documents, key=lambda x: float(x.get("score", 0.0) or 0.0), reverse=True)
            docs_to_rerank = documents_sorted[:_rerank_input_cap]

            _rerank_fps = {d.get("file_path") for d in docs_to_rerank}
            _bm25_inject_count = 0
            for doc in documents_sorted:
                if doc.get("file_path") in _rerank_fps:
                    continue
                bm25_sc = float(doc.get("_bm25_score", 0) or 0)
                if bm25_sc > 10.0:  # Only inject strong BM25 hits
                    docs_to_rerank.append(doc)
                    _rerank_fps.add(doc.get("file_path"))
                    _bm25_inject_count += 1
                    if _bm25_inject_count >= 10:
                        break
            if _bm25_inject_count:
                logger.info(
                    f"[rerank] 注入 {_bm25_inject_count} 条 BM25 高分命中到重排窗口"
                )

            pairs_to_compute = []
            compute_indices = []
            import math

            def _content_for_rerank(doc: Dict[str, Any]) -> str:
                t = str(doc.get('text') or "").strip()
                s = str(doc.get('doc_summary') or "").strip()
                if "text extraction unavailable" in t.lower():
                    rescue = self._build_lightweight_placeholder_preview(doc.get("metadata") or doc)
                    if rescue:
                        return rescue
                    if s:
                        return s
                return f"{s}\n{t}" if s else t
            
            for i, doc in enumerate(docs_to_rerank):
                is_folder = doc.get('is_folder_match', False)
                is_folder_index_row = bool(doc.get("_folder_index_row"))
                is_bm25_force = (float(doc.get("_bm25_score", 0) or 0) > 0
                                 and float(doc.get("score", 0) or 0) >= 1.5)
                dist_val = doc.get('distance')
                if dist_val is None:
                    dist_val = 1.0
                
                is_exact = (float(doc.get('score', 0.0) or 0.0) >= 0.95 and float(dist_val) == 0.0)
                
                if is_folder or is_folder_index_row:
                    # Folder-level entries: content too short for cross-encoder, keep bypass
                    sim = float(doc.get("score", 0.0) or 0.0)
                    sim = min(1.0, max(0.0, sim))
                    doc["rerank_score"] = 4.0 + sim * 0.01
                elif is_bm25_force:
                    content_for_rerank = _content_for_rerank(doc)
                    pairs_to_compute.append([query, content_for_rerank[:1000]])
                    compute_indices.append(i)
                    doc["_is_bm25_force_rerank"] = True  # flag for floor logic below
                elif is_exact:
                    # Lexical filename hit: still run through reranker model for content validation
                    # Mark it so we can apply a floor after scoring
                    content_for_rerank = _content_for_rerank(doc)
                    pairs_to_compute.append([query, content_for_rerank[:1000]])
                    compute_indices.append(i)
                    doc["_is_lexical_rerank"] = True  # flag for floor logic later
                else:
                    content_for_rerank = _content_for_rerank(doc)
                    _doc_syn = str(doc.get("_matched_synonym") or "").strip()
                    _query_for_this_doc = f"{query} ({_doc_syn})" if _doc_syn else query
                    pairs_to_compute.append([_query_for_this_doc, content_for_rerank[:1000]])
                    compute_indices.append(i)
            
            if pairs_to_compute:
                batch_size = 10
                batched_pairs, batched_indices = _chunk_with_no_singleton_tail(
                    pairs_to_compute,
                    compute_indices,
                    batch_size=batch_size,
                )
                
                def _score_batch(bp, bidx):
                    try:
                        scores = reranker.compute_score(bp)
                        if not isinstance(scores, list):
                            scores = [scores]
                        for idx, score in zip(bidx, scores):
                            if isinstance(score, list):
                                score = score[0] if score else 0.0
                            docs_to_rerank[idx]['rerank_score'] = float(score) if not math.isnan(float(score)) else -99.0
                    except Exception as model_err:
                        logger.warning(f"Reranker batch fail, falling back: {model_err}")
                        for idx, pair in zip(bidx, bp):
                            try:
                                sc = reranker.compute_score(pair)
                                if isinstance(sc, list):
                                    sc = sc[0] if sc else 0.0
                                docs_to_rerank[idx]['rerank_score'] = float(sc) if not math.isnan(float(sc)) else -99.0
                            except Exception:
                                docs_to_rerank[idx]['rerank_score'] = float(docs_to_rerank[idx].get("score", 0.0) or 0.0)

                # Reranker models often don't support multithreading if they share the same GPU/CPU context.
                # Just execute the batches serially to be safe, since they're already batched.
                for bp, bidx in zip(batched_pairs, batched_indices):
                    _score_batch(bp, bidx)
            
            # Post-rerank floor logic:
            for d in docs_to_rerank:
                if d.get("_is_lexical_rerank"):
                    # Exact filename hit (distance=0.0): mild floor at -0.3 so it survives
                    # downstream RELEVANCE_THRESHOLD (-0.5) unless truly irrelevant.
                    rk = float(d.get("rerank_score", -99.0) or -99.0)
                    d["rerank_score"] = max(rk, -0.3)
                elif d.get("_is_bm25_force_rerank"):
                    pass  # no floor — rerank cross-encoder score is the sole judge

            for d in docs_to_rerank:
                if "rerank_score" not in d or math.isnan(d["rerank_score"]):
                    original_score = float(d.get("score", 0.0) or 0.0)
                    if original_score >= 0.6:
                         d["rerank_score"] = original_score
                    else:
                         d["rerank_score"] = max(-1.0, original_score)
            
            sorted_docs = sorted(docs_to_rerank, key=lambda x: x['rerank_score'], reverse=True)

            if str(os.getenv("FILEAGENT_RERANK_SCORE_DEBUG", "") or "").strip().lower() in (
                "1",
                "true",
                "yes",
                "on",
            ):
                qpv = str(query or "")[:160]
                logger.info(
                    f"[rerank-debug] query_preview={qpv!r} n={len(docs_to_rerank)} "
                    f"RELEVANCE_THRESHOLD={getattr(settings, 'RELEVANCE_THRESHOLD', None)} "
                    "(raw_* = GGUF 模型 logits，非 Python 加分)"
                )
                for j, doc_j in enumerate(docs_to_rerank):
                    fn_j = os.path.basename(str(doc_j.get("file_path") or doc_j.get("file_name") or ""))
                    vec_j = doc_j.get("score")
                    dist_j = doc_j.get("distance")
                    is_folder = doc_j.get("is_folder_match", False)
                    is_fidx = bool(doc_j.get("_folder_index_row"))
                    forced = (
                        is_folder
                        or is_fidx
                        or (float(vec_j or 0) >= 0.95 and float(dist_j or 1) == 0.0)
                    )
                    logger.info(
                        f"[rerank-debug]  #{j} file={fn_j!r} vec_score={vec_j} distance={dist_j} "
                        f"final_rerank_score={doc_j.get('rerank_score')} forced={forced}"
                    )

            return sorted_docs[:top_k]
        except Exception as e:
            logger.error(f"重排序失败: {e}")
            for d in documents:
                if "rerank_score" not in d:
                    original_score = float(d.get("score", 0.0) or 0.0)
                    if original_score >= 0.5:
                        d["rerank_score"] = original_score
                    else:
                        d["rerank_score"] = max(-1.0, original_score)
            
            sorted_docs = sorted(documents, key=lambda x: float(x.get("rerank_score", 0.0) or 0.0), reverse=True)
            
            if str(os.getenv("FILEAGENT_RERANK_SCORE_DEBUG", "") or "").strip().lower() in (
                "1",
                "true",
                "yes",
                "on",
            ):
                qpv = str(query or "")[:160]
                logger.info(
                    f"[rerank-debug] query_preview={qpv!r} n={len(sorted_docs)} "
                    f"RELEVANCE_THRESHOLD={getattr(settings, 'RELEVANCE_THRESHOLD', None)} "
                    "(raw_* = GGUF 模型 logits，非 Python 加分)"
                )
                for j, doc_j in enumerate(sorted_docs):
                    fn_j = os.path.basename(str(doc_j.get("file_path") or doc_j.get("file_name") or ""))
                    vec_j = doc_j.get("score")
                    dist_j = doc_j.get("distance")
                    is_folder = doc_j.get("is_folder_match", False)
                    is_fidx = bool(doc_j.get("_folder_index_row"))
                    forced = (
                        is_folder
                        or is_fidx
                        or (float(vec_j or 0) >= 0.95 and float(dist_j or 1) == 0.0)
                    )
                    logger.info(
                        f"[rerank-debug]  #{j} file={fn_j!r} vec_score={vec_j} distance={dist_j} "
                        f"final_rerank_score={doc_j.get('rerank_score')} forced={forced}"
                    )
            return sorted_docs[:top_k]
    
    def get_file_chunks(self, file_path: str) -> List[Dict[str, Any]]:
        try:
            results = self.collection.get(
                where={"file_path": file_path},
                include=["documents", "metadatas"]
            )
            
            chunks = []
            if results and results['documents']:
                for i, doc in enumerate(results['documents']):
                    metadata = results['metadatas'][i] if results['metadatas'] else {}
                    chunks.append({
                        'text': doc,
                        'file_name': metadata.get('file_name', ''),
                        'file_path': metadata.get('file_path', ''),
                        'doc_summary': metadata.get('doc_summary', ''),
                        'doc_category': metadata.get('doc_category', ''),
                    })
            
            logger.info(f"获取文件 {os.path.basename(file_path)} 的 {len(chunks)} 个文本块")
            return chunks
        except Exception as e:
            logger.error(f"获取文件块失败: {e}")
            return []

    
    def _should_ignore_file(self, file_path: str) -> bool:
        abs_path = os.path.abspath(file_path)
        if abs_path in self._session_ignored_paths:
            return True

        file_name = os.path.basename(file_path)
        _, ext = os.path.splitext(file_name)
        ext = ext.lower()

        lower_name = file_name.lower()
        if lower_name.startswith(".~") or lower_name.startswith("~$") or lower_name.startswith("._"):
            self._index_info(f"忽略文件(临时文件): {file_name}")
            self._session_ignored_paths.add(abs_path)
            return True
        if lower_name.endswith((".swp", ".swo", ".tmp", ".bak")):
            self._index_info(f"忽略文件(临时文件后缀): {file_name}")
            self._session_ignored_paths.add(abs_path)
            return True
        
        import config.settings as settings

        if ext not in settings.ALLOWED_EXTENSIONS:
            self._index_info(f"忽略文件(不支持的后缀 {ext}): {file_name}")
            self._session_ignored_paths.add(abs_path)
            return True
        
        if file_name in settings.EXCLUDE_FILENAMES:
            self._index_info(f"忽略文件(文件名黑名单): {file_name}")
            self._session_ignored_paths.add(abs_path)
            return True
        
        for prefix in settings.EXCLUDE_FILENAME_PREFIXES:
            if file_name.lower().startswith(prefix.lower()):
                self._index_info(f"忽略文件(前缀黑名单): {file_name}")
                self._session_ignored_paths.add(abs_path)
                return True
        
        if getattr(settings, "USE_WHITELIST_MODE", False):
            in_whitelist = False
            for include_path in settings.INCLUDE_PATHS:
                include_path = os.path.expanduser(include_path)
                if abs_path.startswith(include_path) or include_path.startswith(abs_path):
                    in_whitelist = True
                    break
            if not in_whitelist:
                self._index_info(f"忽略文件(不在白名单中): {abs_path}")
                self._session_ignored_paths.add(abs_path)
                return True
        else:
            for pattern in settings.IGNORE_PATTERNS:
                if pattern in abs_path:
                    self._index_info(f"忽略文件(路径黑名单匹配 {pattern}): {abs_path}")
                    self._session_ignored_paths.add(abs_path)
                    return True
            
            for exclude_path in settings.EXCLUDE_PATHS:
                if exclude_path in abs_path:
                    self._index_info(f"忽略文件(排除路径 {exclude_path}): {abs_path}")
                    self._session_ignored_paths.add(abs_path)
                    return True
        
        return False

    # ═══════════════════════════════════════════════════════════════════
    # Metadata + BM25 Cache — avoid per-query full collection scans
    # ═══════════════════════════════════════════════════════════════════

    def enter_write_heavy_mode(self, reason: str = "") -> None:
        with self._write_heavy_mode_lock:
            self._write_heavy_mode_depth += 1
            depth = self._write_heavy_mode_depth
        self.cancel_query_cache_prewarm(wait=True, reason=f"enter_write_heavy_mode:{reason or 'n/a'}")
        self._note_collection_write()
        logger.info(
            f"[CollectionMode] enter write-heavy mode depth={depth} "
            f"reason={reason or 'n/a'}"
        )

    def leave_write_heavy_mode(self, reason: str = "") -> None:
        with self._write_heavy_mode_lock:
            if self._write_heavy_mode_depth > 0:
                self._write_heavy_mode_depth -= 1
            depth = self._write_heavy_mode_depth
        self._note_collection_write()
        logger.info(
            f"[CollectionMode] leave write-heavy mode depth={depth} "
            f"reason={reason or 'n/a'}"
        )

    def _in_write_heavy_mode(self) -> bool:
        with self._write_heavy_mode_lock:
            return self._write_heavy_mode_depth > 0

    def cancel_query_cache_prewarm(self, *, wait: bool = False, timeout_sec: float = 10.0, reason: str = "") -> None:
        self._query_cache_prewarm_cancel.set()
        worker: Optional[threading.Thread] = None
        with self._query_cache_prewarm_lock:
            worker = self._query_cache_prewarm_thread
        if wait and worker is not None and worker.is_alive():
            worker.join(timeout=max(0.0, float(timeout_sec or 0.0)))
            if worker.is_alive():
                logger.warning(
                    f"[QueryCachePrewarm] cancel wait timed out after {timeout_sec:.1f}s "
                    f"reason={reason or 'n/a'}"
                )
        logger.info(f"[QueryCachePrewarm] cancel requested: reason={reason or 'n/a'}")

    def _note_collection_write(self) -> None:
        self._last_collection_write_ts = time.time()

    def _has_recent_collection_write(self, min_idle_sec: Optional[float] = None) -> bool:
        idle_sec = self._query_cache_prewarm_min_idle_sec if min_idle_sec is None else max(0.0, float(min_idle_sec))
        last_write = float(getattr(self, "_last_collection_write_ts", 0.0) or 0.0)
        if last_write <= 0:
            return False
        return (time.time() - last_write) < idle_sec

    def _invalidate_meta_cache(self) -> None:
        """Invalidate metadata + keyword cache. Call after any collection write."""
        self._note_collection_write()
        with self._meta_cache_lock:
            self._meta_cache_version += 1
            self._meta_cache_data = None
            self._meta_cache_db_count = -1
            self._meta_cache_ts = 0.0
            self._meta_cache_cond.notify_all()
            logger.debug(f"[MetaCache] invalidated → version={self._meta_cache_version}")
        keyword_index = getattr(self, "_keyword_index_manager", None)
        if keyword_index is not None:
            try:
                keyword_index.invalidate(reason="collection_write")
            except Exception as e:
                logger.warning(f"[KeywordIndex] invalidate failed: {e}")

    def _get_cached_metadata(
        self,
        include_documents: bool = False,
        cancel_check: Optional[Callable[[], bool]] = None,
    ) -> Dict[str, Any]:
        """Get metadata (and optionally documents) from cache or collection.

        Returns dict with keys: 'ids', 'metadatas', and optionally 'documents'.
        Uses count-based + TTL invalidation to keep cache fresh.
        """
        while True:
            if cancel_check and cancel_check():
                raise _QueryCachePrewarmCancelled()
            now = time.time()
            current_count = int(self.collection.count() or 0)

            with self._meta_cache_lock:
                cache_valid = (
                    self._meta_cache_data is not None
                    and current_count == self._meta_cache_db_count
                    and (now - self._meta_cache_ts) < self._meta_cache_ttl
                )
                if cache_valid and include_documents and "documents" not in self._meta_cache_data:
                    cache_valid = False

                if cache_valid:
                    logger.debug(
                        f"[MetaCache] HIT (version={self._meta_cache_version}, "
                        f"count={current_count}, age={now - self._meta_cache_ts:.1f}s)"
                    )
                    return self._meta_cache_data  # type: ignore

                if self._meta_cache_building:
                    self._meta_cache_cond.wait(timeout=0.1)
                    continue

                build_version = self._meta_cache_version
                self._meta_cache_building = True

            _t0 = time.time()
            include_fields = ["metadatas"]
            if include_documents:
                include_fields.append("documents")

            all_ids: List[str] = []
            all_metas: List[Dict] = []
            all_docs: List[str] = []
            page_size = 5000
            for offset in range(0, max(current_count, 1), page_size):
                if cancel_check and cancel_check():
                    with self._meta_cache_lock:
                        self._meta_cache_building = False
                        self._meta_cache_cond.notify_all()
                    raise _QueryCachePrewarmCancelled()
                batch = self.collection.get(
                    include=include_fields,
                    limit=page_size,
                    offset=offset,
                )
                all_ids.extend(batch.get("ids") or [])
                all_metas.extend(batch.get("metadatas") or [])
                if include_documents:
                    all_docs.extend(batch.get("documents") or [])

            result: Dict[str, Any] = {"ids": all_ids, "metadatas": all_metas}
            if include_documents:
                result["documents"] = all_docs

            with self._meta_cache_lock:
                if self._meta_cache_version != build_version:
                    self._meta_cache_building = False
                    self._meta_cache_cond.notify_all()
                    continue
                if cancel_check and cancel_check():
                    self._meta_cache_building = False
                    self._meta_cache_cond.notify_all()
                    raise _QueryCachePrewarmCancelled()

                self._meta_cache_data = result
                self._meta_cache_db_count = current_count
                self._meta_cache_ts = time.time()
                self._meta_cache_version += 1
                built_version = self._meta_cache_version
                self._meta_cache_building = False
                self._meta_cache_cond.notify_all()

            _elapsed = time.time() - _t0
            logger.info(
                f"[MetaCache] REBUILD: {len(all_ids)} chunks, "
                f"version={built_version}, took={_elapsed:.2f}s"
            )
            return result

    def _get_or_build_bm25(
        self,
        query_tokens: List[str],
        allowed_paths: Optional[List[str] | PathScopeMatcher] = None,
        category_filter: str = "",
        file_extensions: Optional[List[str]] = None,
    ) -> List[Tuple[str, str, float]]:
        """Score files against query using the keyword index manager."""
        if not _HAS_BM25 or not query_tokens:
            return []
        keyword_index = getattr(self, "_keyword_index_manager", None)
        if keyword_index is None:
            return []
        return keyword_index.score(
            query_tokens,
            allowed_paths=allowed_paths,
            category_filter=category_filter,
            file_extensions=file_extensions,
        )

    def is_keyword_index_ready(self) -> bool:
        keyword_index = getattr(self, "_keyword_index_manager", None)
        if keyword_index is None:
            return False
        try:
            return bool(keyword_index.is_ready())
        except Exception:
            return False

    @staticmethod
    def _media_content_query_terms(
        raw_query: str,
        query_tokens: List[str],
        *,
        category_filter: str = "",
        file_extensions: Optional[List[str]] = None,
    ) -> Tuple[List[str], List[str], bool]:
        media_exts = {
            ".mp4", ".mov", ".mkv", ".avi", ".webm", ".m4v",
            ".mp3", ".wav", ".m4a", ".flac", ".aac", ".ogg",
        }
        target_exts = {
            ext if str(ext or "").startswith(".") else f".{ext}"
            for ext in (file_extensions or [])
            if str(ext or "").strip()
        }
        target_exts = {str(ext or "").strip().lower() for ext in target_exts if str(ext or "").strip()}
        target_category = _normalize_category_en(category_filter or "", default="")
        if target_category in {"", "all", "unknown"}:
            target_category = ""

        query_lower = str(raw_query or "").lower()
        wants_media = (
            target_category in {"audio", "video", "audio/video"}
            or bool(target_exts.intersection(media_exts))
            or bool(re.search(r"\b(video|videos|audio|recording|clip|movie|mp4|mov|mkv|avi|webm|m4v)\b", query_lower))
            or bool(re.search(r"(视频|音频|录像|录音|片段)", query_lower))
        )
        if not wants_media:
            return [], [], False

        generic_media_terms = {
            "audio", "video", "videos", "recording", "recordings", "clip", "clips", "movie", "movies",
            "file", "files", "media", "indexed", "find", "search", "show", "locate", "which", "where",
            "saw", "see", "seen", "watch", "watched", "mp4", "mov", "mkv", "avi", "webm", "m4v",
            "mp3", "wav", "m4a", "flac", "aac", "ogg", "audio/video", "type", "source",
            "download", "downloaded", "created", "resolution", "duration", "content",
            "summary", "visual", "frame", "frames", "keyframe", "sound", "sounds",
            "effect", "effects", "sfx",
            "视频", "音频", "音视频", "文件", "片段", "录像", "录音", "音效",
            "找", "寻找", "搜索", "展示", "显示", "有", "的", "那个", "哪个", "哪些",
            "里面", "里", "中", "在", "看到", "看见", "帮我", "一下",
            "含有", "包含", "包括", "带有", "具有", "或", "或者", "以及",
        }
        content_terms = [
            term for term in dict.fromkeys(str(t or "").strip().lower() for t in query_tokens)
            if term and "/" not in term and term not in generic_media_terms and term not in _BM25_STOPWORDS
        ]

        # Do not infer broad semantic neighbors here. Media topic search should
        # be backed by terms actually present in the user's query and indexed
        # evidence; otherwise category candidates can swamp precise searches.
        expanded_terms: List[str] = []

        return content_terms, expanded_terms, True

    @staticmethod
    def _match_terms_in_text(
        terms: List[str],
        text_lower: str,
        text_tokens: set[str],
    ) -> List[str]:
        matched: List[str] = []
        for term in terms:
            term = str(term or "").strip().lower()
            if not term:
                continue
            if any("\u4e00" <= ch <= "\u9fff" for ch in term):
                if len(term) >= 2 and term in text_lower:
                    matched.append(term)
                elif term in text_tokens:
                    matched.append(term)
                continue
            if term in text_tokens or (len(term) >= 4 and term in text_lower):
                matched.append(term)
        return matched

    @staticmethod
    def _media_content_direct_score(
        content_terms: List[str],
        matched_terms: List[str],
        expanded_matches: List[str],
        overlap: int,
        fallback_score: float,
    ) -> int:
        matched_set = set(matched_terms)
        term_set = set(content_terms)
        score = 0
        if matched_terms:
            matched_count = len(matched_set)
            required_count = max(1, len(term_set))
            if required_count >= 3:
                score = 98 if matched_count >= 3 else 94 if matched_count >= 2 else 88
            else:
                score = 98 if matched_count >= required_count else 92
        if expanded_matches:
            score = max(score, min(92, 76 + len(expanded_matches) * 4))
        if overlap > 0 and (matched_terms or expanded_matches):
            score = max(score, 85)
        if score <= 0:
            score = max(50, min(68, int(fallback_score * 5)))
        return int(max(0, min(100, score)))

    @staticmethod
    def _infer_media_extension_filters(raw_query: str) -> Optional[List[str]]:
        query_lower = str(raw_query or "").lower()
        specific_exts: List[str] = []
        ext_aliases = {
            ".mp4": ["mp4"],
            ".mov": ["mov"],
            ".mkv": ["mkv"],
            ".avi": ["avi"],
            ".webm": ["webm"],
            ".m4v": ["m4v"],
            ".mp3": ["mp3"],
            ".wav": ["wav"],
            ".m4a": ["m4a"],
            ".flac": ["flac"],
            ".aac": ["aac"],
            ".ogg": ["ogg"],
        }
        for ext, aliases in ext_aliases.items():
            if any(re.search(rf"\b{re.escape(alias)}\b", query_lower) for alias in aliases):
                specific_exts.append(ext)
        if specific_exts:
            return list(dict.fromkeys(specific_exts))

        has_video_word = bool(
            re.search(r"\b(video|videos|movie|movies|clip|clips)\b", query_lower)
            or re.search(r"(视频|录像|影片|短片)", query_lower)
        )
        has_audio_word = bool(
            re.search(r"\b(audio|audios|sound|sounds|voice|voices|music|song|songs)\b", query_lower)
            or re.search(r"(音频|录音|声音|音乐|歌曲)", query_lower)
        )
        if has_video_word and not has_audio_word:
            return [".mp4", ".mov", ".mkv", ".avi", ".webm", ".m4v"]
        if has_audio_word and not has_video_word:
            return [".mp3", ".wav", ".m4a", ".flac", ".aac", ".ogg"]
        return None

    def _indexed_media_content_fallback_search(
        self,
        raw_query: str,
        query_tokens: List[str],
        *,
        allowed_paths: Optional[List[str] | PathScopeMatcher] = None,
        category_filter: str = "",
        file_extensions: Optional[List[str]] = None,
        limit: int = 30,
    ) -> List[Dict[str, Any]]:
        """Lexically scan indexed media chunk text when vector/BM25 recall is unavailable."""
        media_exts = {
            ".mp4", ".mov", ".mkv", ".avi", ".webm", ".m4v",
            ".mp3", ".wav", ".m4a", ".flac", ".aac", ".ogg",
        }
        target_exts = set(self._normalize_extension_filters(file_extensions) or [])
        target_category = _normalize_category_en(category_filter or "", default="")
        if target_category in {"", "all", "unknown"}:
            target_category = ""

        content_terms, expanded_terms, wants_media = self._media_content_query_terms(
            raw_query,
            query_tokens,
            category_filter=target_category,
            file_extensions=list(target_exts),
        )
        if not wants_media:
            return []

        media_chunk_types = {
            "media_summary",
            "media_audio_summary",
            "media_visual_summary",
            "interval_summary",
            "interval_visual",
            "asr_transcript",
            "asr_segment",
            "keyframe",
        }
        scope_matcher = ensure_path_scope_matcher(allowed_paths)

        compatible_categories = {target_category} if target_category else set()
        if target_category:
            try:
                from core.retrieval.category_engine import get_compatible_categories

                compatible_categories = {
                    str(cat or "").strip().lower()
                    for cat in (get_compatible_categories(target_category) or {target_category})
                    if str(cat or "").strip()
                } or {target_category}
            except Exception:
                compatible_categories = {target_category}

        try:
            cached = self._get_cached_metadata(include_documents=True)
        except Exception as exc:
            logger.warning("[indexed_media_content_fallback] metadata scan failed: %s", exc)
            return []

        rows: Dict[str, Dict[str, Any]] = {}
        for mid, meta, doc in itertools.zip_longest(
            cached.get("ids") or [],
            cached.get("metadatas") or [],
            cached.get("documents") or [],
            fillvalue="",
        ):
            meta = dict(meta or {})
            fp = str(meta.get("file_path") or "").strip()
            if not fp or not scope_matcher.allows_file(fp):
                continue
            file_ext = str(meta.get("file_extension") or os.path.splitext(fp)[1]).strip().lower()
            if target_exts and file_ext not in target_exts:
                continue
            category = self._meta_category_family(meta)
            if target_category and str(category or "").strip().lower() not in compatible_categories:
                continue

            chunk_type = str(meta.get("chunk_type") or "").strip().lower()
            media_type = str(meta.get("media_type") or "").strip().lower()
            if not (
                chunk_type in media_chunk_types
                or media_type in {"audio", "video"}
                or file_ext in media_exts
                or category in {"audio", "video", "audio/video"}
            ):
                continue
            if self._is_filename_lookup_metadata(meta):
                continue

            text_blob = " ".join(
                str(part or "").strip()
                for part in [
                    meta.get("keyframe_description", ""),
                    meta.get("keyframe_ocr_text", ""),
                    meta.get("media_visual_summary", ""),
                    meta.get("media_audio_summary", ""),
                    meta.get("media_summary", ""),
                    doc,
                ]
                if str(part or "").strip()
            )
            text_blob = re.sub(r"\s+", " ", text_blob).strip()
            if not text_blob:
                continue

            text_lower = text_blob.lower()
            text_tokens = set(_tokenize_for_bm25(text_blob))
            matched_terms = self._match_terms_in_text(content_terms, text_lower, text_tokens)
            expanded_matches = self._match_terms_in_text(expanded_terms, text_lower, text_tokens)
            overlap = compute_lookup_overlap_score(raw_query, text_blob)
            inventory_only = not content_terms and not expanded_terms
            if not inventory_only and not matched_terms and not expanded_matches:
                continue

            score = float(len(matched_terms) * 20 + len(expanded_matches) * 6)
            if matched_terms or expanded_matches:
                score += max(0, overlap) * 3
            if inventory_only:
                score = 1.0
            existing = rows.get(fp)
            if existing and float(existing.get("_bm25_score", 0.0) or 0.0) >= score:
                continue

            file_name = str(meta.get("file_name") or os.path.basename(fp)).strip()
            excerpt = text_blob[:1800].rstrip()
            doc_summary = str(meta.get("doc_summary") or "").strip() or excerpt[:240]
            direct_score = (
                70 if inventory_only else self._media_content_direct_score(
                    content_terms,
                    matched_terms,
                    expanded_matches,
                    int(overlap or 0),
                    score,
                )
            )
            rows[fp] = {
                "text": excerpt,
                "metadata": meta,
                "distance": 0.0,
                "file_name": file_name,
                "file_path": fp,
                "file_name_en": meta.get("file_name_en", ""),
                "folder_name_en": meta.get("folder_name_en", ""),
                "doc_summary": doc_summary,
                "doc_category": category or media_type or "audio/video",
                "doc_category_leaf": self._meta_category_leaf(meta),
                "doc_role": self._meta_doc_role(meta),
                "lookup_aliases": meta.get("lookup_aliases", ""),
                "score": max(float(settings.RELEVANCE_THRESHOLD) + 0.2, min(2.0, score / 10.0)),
                "rerank_score": float(settings.RELEVANCE_THRESHOLD) + 0.3,
                "_is_lexical_hit": True,
                "_direct_score": direct_score,
                "_bm25_score": score,
                "_keyword_index_hit": True,
                "_keyword_content_fallback": True,
                "_keyword_index_id": str(mid or ""),
                "_matched_terms": matched_terms[:8],
                "_expanded_terms": expanded_matches[:8],
            }

        ranked = sorted(
            rows.values(),
            key=lambda item: (
                int(item.get("_direct_score", 0) or 0),
                float(item.get("_bm25_score", 0.0) or 0.0),
                str(item.get("file_name") or "").lower(),
            ),
            reverse=True,
        )
        return ranked[: max(1, min(int(limit or 30), 200))]

    def indexed_keyword_search(
        self,
        query: str,
        *,
        allowed_paths: Optional[List[str] | PathScopeMatcher] = None,
        category_filter: str = "",
        file_extensions: Optional[List[str]] = None,
        limit: int = 30,
    ) -> List[Dict[str, Any]]:
        """
        Search the ingest-time lexical index and hydrate only the top Chroma rows.

        This is the query-time replacement for broad metadata scans: file names,
        stems, path segments, aliases, summaries, and schema hints are indexed into
        the BM25 sidecar during ingestion. Search only scores that index, then
        fetches the top ids directly from Chroma.
        """
        raw_query = str(query or "").strip()
        if not raw_query:
            return []

        target_category = _normalize_category_en(category_filter or "", default="")
        if target_category in {"", "all", "unknown"}:
            target_category = ""

        normalized_exts: Optional[List[str]] = None
        if file_extensions:
            normalized_exts = []
            for raw_ext in file_extensions:
                ext = str(raw_ext or "").strip().lower()
                if not ext:
                    continue
                if not ext.startswith("."):
                    ext = f".{ext}"
                normalized_exts.append(ext)
            normalized_exts = list(dict.fromkeys(normalized_exts)) or None
        if normalized_exts is None:
            normalized_exts = self._infer_media_extension_filters(raw_query)

        query_blob = build_lookup_blob(raw_query, max_terms=48)
        tokens = _tokenize_for_bm25(f"{raw_query} {query_blob}".strip())
        if not tokens:
            return []
        media_content_terms, media_expanded_terms, wants_media_content = self._media_content_query_terms(
            raw_query,
            tokens,
            category_filter=target_category,
            file_extensions=normalized_exts,
        )

        try:
            hit_limit = max(1, min(int(limit or 30), 200))
        except Exception:
            hit_limit = 30

        fallback_hits = self._indexed_media_content_fallback_search(
            raw_query,
            tokens,
            allowed_paths=allowed_paths,
            category_filter=target_category,
            file_extensions=normalized_exts,
            limit=hit_limit,
        ) if wants_media_content else []

        hits = self._get_or_build_bm25(
            query_tokens=tokens,
            allowed_paths=allowed_paths,
            category_filter=target_category,
            file_extensions=normalized_exts,
        )[:hit_limit]
        if not hits:
            if fallback_hits:
                logger.info(
                    "[indexed_keyword_search] media content fallback returned %d hit(s) for query=%r",
                    len(fallback_hits),
                    raw_query,
                )
                return fallback_hits
            return []

        ids = list(dict.fromkeys(str(hit_id or "") for hit_id, _, _ in hits if str(hit_id or "").strip()))
        records_by_id: Dict[str, KeywordIndexRecord] = {}
        records_by_path: Dict[str, KeywordIndexRecord] = {}
        try:
            keyword_index = getattr(self, "_keyword_index_manager", None)
            if keyword_index is not None:
                for rec in keyword_index.records_snapshot(require_current=False):
                    if str(rec.chroma_id or "").strip():
                        records_by_id[str(rec.chroma_id or "")] = rec
                    if str(rec.file_path or "").strip():
                        records_by_path[str(rec.file_path or "")] = rec
        except Exception:
            records_by_id = {}
            records_by_path = {}
        hydrated: Dict[str, tuple[str, Dict[str, Any]]] = {}
        if ids:
            try:
                batch = self.collection.get(ids=ids, include=["documents", "metadatas"])
                for mid, doc, meta in zip(
                    batch.get("ids") or [],
                    batch.get("documents") or [],
                    batch.get("metadatas") or [],
                ):
                    hydrated[str(mid or "")] = (str(doc or ""), dict(meta or {}))
            except Exception as exc:
                logger.warning("[indexed_keyword_search] hydrate failed: %s", exc)

        results: List[Dict[str, Any]] = []
        seen_paths: set[str] = set()
        for hit_id, file_path, bm25_score in hits:
            fp = str(file_path or "").strip()
            if not fp or fp in seen_paths:
                continue
            seen_paths.add(fp)
            doc, meta = hydrated.get(str(hit_id or ""), ("", {}))
            rec = records_by_id.get(str(hit_id or "")) or records_by_path.get(fp)
            content_preview = str(getattr(rec, "content_preview", "") or "").strip() if rec else ""
            if not meta:
                meta = {
                    "file_path": fp,
                    "file_name": os.path.basename(fp),
                    "doc_category": target_category or "other",
                    "doc_summary": content_preview[:240] if content_preview else "",
                }
            file_name = str(meta.get("file_name") or os.path.basename(fp))
            summary = str(meta.get("doc_summary") or "")
            score = float(bm25_score or 0.0)
            topic_overlap = compute_lookup_overlap_score(raw_query, content_preview) if content_preview else 0
            if wants_media_content and content_preview:
                preview_lower = content_preview.lower()
                preview_tokens = set(_tokenize_for_bm25(content_preview))
                exact_content_matches = self._match_terms_in_text(media_content_terms, preview_lower, preview_tokens)
                expanded_content_matches = self._match_terms_in_text(media_expanded_terms, preview_lower, preview_tokens)
            else:
                exact_content_matches = []
                expanded_content_matches = []
            if wants_media_content and media_content_terms and not exact_content_matches and not expanded_content_matches:
                continue
            if wants_media_content:
                direct_score = self._media_content_direct_score(
                    media_content_terms,
                    exact_content_matches,
                    expanded_content_matches,
                    int(topic_overlap or 0),
                    score,
                )
            else:
                direct_score = min(100, max(50, int(score * 10), 92 if topic_overlap > 0 else 0))
            results.append(
                {
                    "text": content_preview or doc or summary,
                    "metadata": meta,
                    "distance": 0.5,
                    "file_name": file_name,
                    "file_path": fp,
                    "file_name_en": meta.get("file_name_en", ""),
                    "folder_name_en": meta.get("folder_name_en", ""),
                    "doc_summary": summary or content_preview[:240],
                    "doc_category": self._meta_category_family(meta),
                    "doc_category_leaf": self._meta_category_leaf(meta),
                    "doc_role": self._meta_doc_role(meta),
                    "lookup_aliases": meta.get("lookup_aliases", ""),
                    "score": max(1.0, min(2.0, score / 10.0)),
                    "_is_lexical_hit": True,
                    "_direct_score": direct_score,
                    "_bm25_score": score,
                    "_keyword_index_hit": True,
                    "_keyword_content_hit": bool(exact_content_matches),
                    "_keyword_content_expanded_hit": bool(expanded_content_matches),
                    "_matched_terms": exact_content_matches[:8],
                    "_expanded_terms": expanded_content_matches[:8],
                    "_keyword_index_id": str(hit_id or ""),
                }
            )
        if fallback_hits:
            by_path: Dict[str, Dict[str, Any]] = {
                str(item.get("file_path") or (item.get("metadata") or {}).get("file_path") or ""): item
                for item in results
                if str(item.get("file_path") or (item.get("metadata") or {}).get("file_path") or "").strip()
            }
            for item in fallback_hits:
                fp = str(item.get("file_path") or (item.get("metadata") or {}).get("file_path") or "").strip()
                if not fp:
                    continue
                prev = by_path.get(fp)
                if prev is None:
                    by_path[fp] = item
                    continue
                prev_key = (
                    int(prev.get("_direct_score", 0) or 0),
                    float(prev.get("_bm25_score", 0.0) or 0.0),
                )
                item_key = (
                    int(item.get("_direct_score", 0) or 0),
                    float(item.get("_bm25_score", 0.0) or 0.0),
                )
                if item_key > prev_key:
                    by_path[fp] = item
            results = list(by_path.values())

        results.sort(
            key=lambda item: (
                int(item.get("_direct_score", 0) or 0),
                float(item.get("_bm25_score", 0.0) or 0.0),
                str(item.get("file_name") or "").lower(),
            ),
            reverse=True,
        )
        return results[:hit_limit]

    def indexed_exact_filename_lookup(
        self,
        file_name: str,
        *,
        allowed_paths: Optional[List[str] | PathScopeMatcher] = None,
        file_extensions: Optional[List[str]] = None,
        match_mode: str = "exact_filename",
        limit: int = 6,
        hydrate: bool = True,
        include_documents: bool = True,
    ) -> Dict[str, Any]:
        """Find exact filename/stem matches from the keyword sidecar.

        This is intentionally independent of semantic/BM25 scoring. Explicit
        filename requests should not miss just because a broad category hint
        ("document") differs from an indexed leaf category ("report").
        """
        from core.retrieval.filename_canonicalizer import (
            compact_filename_key,
            filename_stem_key_matches_query,
        )

        target = os.path.basename(str(file_name or "").strip())
        if not target:
            return {"ready": False, "count": 0, "files": [], "raw_count": 0}

        normalized_exts = self._normalize_extension_filters(file_extensions)
        target_exts = {str(ext or "").lower() for ext in (normalized_exts or []) if str(ext or "").strip()}
        target_name_key = compact_filename_key(target)
        target_stem_keys = {
            compact_filename_key(os.path.splitext(target)[0]),
            compact_filename_key(target),
        }
        target_stem_keys.discard("")
        mode = str(match_mode or "exact_filename").strip().lower()
        if mode not in {"exact_filename", "exact_stem"}:
            mode = "exact_filename"

        keyword_index = getattr(self, "_keyword_index_manager", None)
        if keyword_index is None:
            return {"ready": False, "count": 0, "files": [], "raw_count": 0}
        current_records = keyword_index.records_snapshot(require_current=True)
        available_records = current_records or keyword_index.records_snapshot(require_current=False)
        if not available_records:
            return {"ready": False, "count": 0, "files": [], "raw_count": 0}

        scope_matcher = ensure_path_scope_matcher(allowed_paths)
        rows: List[Tuple[KeywordIndexRecord, float]] = []
        seen_paths: set[str] = set()
        for rec in available_records:
            fp = str(rec.file_path or "").strip()
            if not fp or fp in seen_paths:
                continue
            if not scope_matcher.allows_file(rec.normalized_path or fp):
                continue

            rec_ext = str(rec.file_extension or os.path.splitext(fp)[1]).strip().lower()
            if target_exts and rec_ext not in target_exts:
                continue

            candidate_name = os.path.basename(str(rec.file_name or fp).strip())
            candidate_name_key = compact_filename_key(candidate_name)
            candidate_stem_key = compact_filename_key(os.path.splitext(candidate_name)[0])
            if mode == "exact_filename":
                matched = bool(target_name_key and candidate_name_key == target_name_key)
            else:
                matched = bool(
                    candidate_stem_key
                    and any(
                        filename_stem_key_matches_query(candidate_stem_key, target_key)
                        for target_key in target_stem_keys
                    )
                )
            if not matched:
                continue

            seen_paths.add(fp)
            rows.append((rec, 1.0))

        rows.sort(key=lambda row: (row[0].file_path.count(os.sep), row[0].file_name.lower(), row[0].file_path))
        raw_count = len(rows)
        try:
            lim_val = int(limit or 0)
        except Exception:
            lim_val = 6
        selected = rows if lim_val <= 0 else rows[: max(0, lim_val)]

        hydrated: Dict[str, Tuple[str, Dict[str, Any]]] = {}
        if hydrate and selected:
            hydrated = self._hydrate_keyword_records(selected, include_documents=include_documents)

        files = [
            self._inventory_item_from_record(
                rec,
                score,
                hydrated.get(str(rec.chroma_id or "")),
            )
            for rec, score in selected
        ]
        return {
            "ready": True,
            "count": raw_count,
            "files": files,
            "raw_count": raw_count,
            "limited": len(files) < raw_count,
            "stale": not bool(current_records),
        }

    @staticmethod
    def _normalize_extension_filters(file_extensions: Optional[List[str]]) -> Optional[List[str]]:
        if not file_extensions:
            return None
        out: List[str] = []
        for raw_ext in file_extensions:
            ext = str(raw_ext or "").strip().lower()
            if not ext:
                continue
            if not ext.startswith("."):
                ext = f".{ext}"
            if ext not in out:
                out.append(ext)
        return out or None

    def _hydrate_keyword_records(
        self,
        records: List[Tuple[KeywordIndexRecord, float]],
        *,
        include_documents: bool = True,
    ) -> Dict[str, Tuple[str, Dict[str, Any]]]:
        ids = list(dict.fromkeys(str(rec.chroma_id or "") for rec, _ in records if str(rec.chroma_id or "").strip()))
        if not ids:
            return {}
        hydrated: Dict[str, Tuple[str, Dict[str, Any]]] = {}
        include_fields = ["metadatas"]
        if include_documents:
            include_fields.append("documents")
        batch_size = 500
        for start in range(0, len(ids), batch_size):
            batch_ids = ids[start : start + batch_size]
            try:
                batch = self.collection.get(ids=batch_ids, include=include_fields)
                docs = batch.get("documents") or []
                metas = batch.get("metadatas") or []
                mids = batch.get("ids") or []
                if not include_documents:
                    docs = [""] * len(mids)
                for mid, doc, meta in zip(mids, docs, metas):
                    hydrated[str(mid or "")] = (str(doc or ""), dict(meta or {}))
            except Exception as exc:
                logger.warning("[indexed_file_inventory] hydrate batch failed: %s", exc)
        return hydrated

    def _inventory_item_from_record(
        self,
        rec: KeywordIndexRecord,
        score: float,
        hydrated: Optional[Tuple[str, Dict[str, Any]]] = None,
    ) -> Dict[str, Any]:
        doc, meta = hydrated or ("", {})
        meta = dict(meta or {})
        fp = str(meta.get("file_path") or rec.file_path or "").strip()
        file_name = str(meta.get("file_name") or rec.file_name or os.path.basename(fp)).strip()
        doc_summary = str(meta.get("doc_summary") or "").strip()
        if not doc_summary and self._is_filename_lookup_metadata(meta):
            doc_summary = f"Indexed file: {file_name}"
        category = self._meta_category_family(meta) if meta else str(rec.category or "other")
        file_ext = str(meta.get("file_extension") or rec.file_extension or os.path.splitext(fp)[1]).strip().lower()
        if not meta:
            meta = {
                "file_path": fp,
                "file_name": file_name,
                "file_extension": file_ext,
                "doc_category": category or "other",
                "doc_category_family": category or "other",
                "doc_summary": doc_summary,
            }
        text = doc or doc_summary or file_name
        return {
            "text": text,
            "metadata": meta,
            "distance": 0.0,
            "file_name": file_name,
            "file_path": fp,
            "file_name_en": str(meta.get("file_name_en") or "").strip(),
            "folder_name_en": str(meta.get("folder_name_en") or "").strip(),
            "doc_summary": doc_summary,
            "doc_category": category or "other",
            "doc_category_leaf": self._meta_category_leaf(meta) if meta else category or "other",
            "doc_role": self._meta_doc_role(meta) if meta else "primary_source",
            "lookup_aliases": str(meta.get("lookup_aliases") or "").strip(),
            "file_extension": file_ext,
            "score": float(score or 1.0),
            "hit_chunks": 1,
            "_first_chunk_id": str(rec.chroma_id or ""),
            "_is_inventory_hit": True,
            "_keyword_index_hit": True,
        }

    def indexed_file_inventory(
        self,
        *,
        query: str = "",
        allowed_paths: Optional[List[str] | PathScopeMatcher] = None,
        category_filter: str = "",
        file_extensions: Optional[List[str]] = None,
        limit: Optional[int] = 200,
        hydrate: bool = True,
        include_documents: bool = True,
    ) -> Dict[str, Any]:
        """
        Fast file inventory/count route backed by the keyword sidecar.

        This is the DB-index replacement for query-time metadata paging used by
        requests like "list all pdf", "how many videos", and "show csv files".
        It uses in-memory keyword records and precise Chroma id hydration only.
        """
        normalized_exts = self._normalize_extension_filters(file_extensions)
        target_category = _normalize_category_en(category_filter or "", default="")
        if target_category in {"", "all", "unknown"}:
            target_category = ""

        raw_query = str(query or "").strip()
        rows: List[Tuple[KeywordIndexRecord, float]] = []
        keyword_index = getattr(self, "_keyword_index_manager", None)
        if keyword_index is None:
            return {"ready": False, "count": 0, "files": [], "raw_count": 0}
        current_records = keyword_index.records_snapshot(require_current=True)
        available_records = current_records or keyword_index.records_snapshot(require_current=False)
        if not available_records:
            return {"ready": False, "count": 0, "files": [], "raw_count": 0}
        inventory_stale = not bool(current_records)

        if raw_query:
            query_blob = build_lookup_blob(raw_query, max_terms=48)
            tokens = _tokenize_for_bm25(f"{raw_query} {query_blob}".strip())
            if tokens:
                scored = self._get_or_build_bm25(
                    query_tokens=tokens,
                    allowed_paths=allowed_paths,
                    category_filter=target_category,
                    file_extensions=normalized_exts,
                )
                records_by_id = {
                    str(rec.chroma_id or ""): rec
                    for rec in available_records
                }
                records_by_path = {
                    str(rec.file_path or ""): rec
                    for rec in records_by_id.values()
                    if str(rec.file_path or "").strip()
                }
                for hit_id, fp, bm25_score in scored:
                    rec = records_by_id.get(str(hit_id or "")) or records_by_path.get(str(fp or ""))
                    if rec is not None:
                        rows.append((rec, float(bm25_score or 0.0)))
        else:
            listed = keyword_index.list_records(
                allowed_paths=allowed_paths,  # type: ignore[arg-type]
                category_filter=target_category,
                file_extensions=normalized_exts,
                require_current=False,
            )
            rows = [(rec, 1.0) for rec in listed]

        seen_paths: set[str] = set()
        deduped: List[Tuple[KeywordIndexRecord, float]] = []
        for rec, score in rows:
            fp = str(rec.file_path or "").strip()
            if not fp or fp in seen_paths:
                continue
            seen_paths.add(fp)
            deduped.append((rec, score))

        if raw_query:
            deduped.sort(key=lambda row: (-float(row[1] or 0.0), row[0].file_name.lower(), row[0].file_path))
        else:
            deduped.sort(key=lambda row: (row[0].file_extension.lower(), row[0].file_name.lower(), row[0].file_path))

        raw_count = len(deduped)
        try:
            lim_val = None if limit is None else int(limit)
        except Exception:
            lim_val = 200
        selected = deduped if lim_val is None or lim_val <= 0 else deduped[: max(0, lim_val)]

        hydrated: Dict[str, Tuple[str, Dict[str, Any]]] = {}
        if hydrate and selected:
            hydrated = self._hydrate_keyword_records(selected, include_documents=include_documents)

        files = [
            self._inventory_item_from_record(
                rec,
                score,
                hydrated.get(str(rec.chroma_id or "")),
            )
            for rec, score in selected
        ]
        return {
            "ready": True,
            "count": raw_count,
            "files": files,
            "raw_count": raw_count,
            "limited": len(files) < raw_count,
            "stale": inventory_stale,
        }

    def _build_lightweight_placeholder_preview(self, meta: Dict[str, Any]) -> str:
        """
        Query-time retrieval is DB-only.

        Older versions tried to rescue weak spreadsheet indexes by reading a
        tiny local preview here. That made latency depend on disk/file-format
        behavior, so missing content must now be fixed by re-indexing instead.
        """
        return ""


    def _index_ignore_cache_disabled(self) -> bool:
        return os.getenv("FILEAGENT_DISABLE_INDEX_IGNORE_CACHE", "").strip().lower() in (
            "1",
            "true",
            "yes",
        )

    def _compute_ignore_rules_fingerprint(self) -> str:
        import config.settings as settings
        wh = bool(getattr(settings, "USE_WHITELIST_MODE", False))
        path_part: List[str] = []
        if wh:
            raw_paths = getattr(settings, "INCLUDE_PATHS", None) or set()
            path_part = sorted(
                os.path.normpath(os.path.abspath(os.path.expanduser(str(p))))
                for p in raw_paths
            )
        parts: List[Any] = [
            sorted(settings.ALLOWED_EXTENSIONS),
            sorted(settings.EXCLUDE_FILENAMES),
            sorted(settings.EXCLUDE_FILENAME_PREFIXES),
            sorted(settings.IGNORE_PATTERNS),
            sorted(settings.EXCLUDE_PATHS),
            sorted(getattr(settings, "IGNORE_TOP_LEVEL_DIRS", set())),
            wh,
            path_part,
        ]
        raw = json.dumps(parts, ensure_ascii=False)
        return hashlib.sha256(raw.encode("utf-8")).hexdigest()

    def _load_index_ignore_cache(self) -> None:
        if self._index_ignore_cache_disabled():
            self._ignore_cache_entries = {}
            return
        try:
            if not os.path.isfile(self._ignore_cache_path):
                self._ignore_cache_entries = {}
                return
            with open(self._ignore_cache_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            if not isinstance(data, dict) or int(data.get("version", 0)) != 1:
                self._ignore_cache_entries = {}
                return
            entries = data.get("entries") or {}
            if not isinstance(entries, dict):
                self._ignore_cache_entries = {}
                return
            pruned: Dict[str, Dict[str, Any]] = {}
            for p, meta in entries.items():
                if not isinstance(meta, dict):
                    continue
                try:
                    if os.path.isfile(p):
                        pruned[p] = meta
                except Exception:
                    continue
            self._ignore_cache_entries = pruned
            if len(pruned) != len(entries):
                self._ignore_cache_dirty = True
        except Exception as e:
            self._index_exception("加载 index_ignore_cache 失败，将忽略缓存", e)
            self._ignore_cache_entries = {}

    def _save_index_ignore_cache(self) -> None:
        if self._index_ignore_cache_disabled() or not self._ignore_cache_dirty:
            return
        try:
            payload = {"version": 1, "entries": self._ignore_cache_entries}
            d = os.path.dirname(self._ignore_cache_path)
            if d:
                os.makedirs(d, exist_ok=True)
            tmp = self._ignore_cache_path + ".tmp"
            with open(tmp, "w", encoding="utf-8") as f:
                json.dump(payload, f, ensure_ascii=False, separators=(",", ":"))
            os.replace(tmp, self._ignore_cache_path)
            self._ignore_cache_dirty = False
            self._ignore_cache_new_since_save = 0
            self._ignore_cache_last_save_time = time.time()
        except Exception as e:
            self._index_exception("保存 index_ignore_cache 失败", e)

    def _index_ignore_cache_try_skip(self, abs_path: str, rules_fp: str) -> bool:
        if self._index_ignore_cache_disabled():
            return False
        try:
            ent = self._ignore_cache_entries.get(abs_path)
            if not ent or ent.get("fp") != rules_fp:
                return False
            st = os.stat(abs_path)
            m_ns = getattr(st, "st_mtime_ns", int(st.st_mtime * 1_000_000_000))
            if int(ent.get("mtime_ns", -1)) != int(m_ns) or int(ent.get("size", -1)) != int(st.st_size):
                return False
            return True
        except Exception:
            return False

    def _index_ignore_cache_record(self, abs_path: str, rules_fp: str) -> None:
        if self._index_ignore_cache_disabled():
            return
        try:
            st = os.stat(abs_path)
            m_ns = getattr(st, "st_mtime_ns", int(st.st_mtime * 1_000_000_000))
            self._ignore_cache_entries[abs_path] = {
                "fp": rules_fp,
                "mtime_ns": int(m_ns),
                "size": int(st.st_size),
            }
            self._ignore_cache_dirty = True
            self._ignore_cache_new_since_save += 1
            
            if self._ignore_cache_new_since_save >= 20 or (time.time() - self._ignore_cache_last_save_time) > 1.0:
                self._save_index_ignore_cache()
        except Exception:
            pass

    def _load_failed_file_cache(self) -> None:
        try:
            if not os.path.isfile(self._failed_file_cache_path):
                self._failed_file_cache_entries = {}
                return
            with open(self._failed_file_cache_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            if not isinstance(data, dict) or int(data.get("version", 0)) != 1:
                self._failed_file_cache_entries = {}
                return
            entries = data.get("entries") or {}
            if not isinstance(entries, dict):
                self._failed_file_cache_entries = {}
                return

            pruned: Dict[str, Dict[str, Any]] = {}
            for p, meta in entries.items():
                if not isinstance(meta, dict):
                    continue
                try:
                    if os.path.isfile(p):
                        pruned[p] = meta
                except Exception:
                    continue
            self._failed_file_cache_entries = pruned
            if len(pruned) != len(entries):
                self._failed_file_cache_dirty = True
        except Exception as e:
            self._index_exception("加载 index_failed_files 失败，将忽略失败缓存", e)
            self._failed_file_cache_entries = {}

    def _save_failed_file_cache(self) -> None:
        if not self._failed_file_cache_dirty:
            return
        try:
            payload = {"version": 1, "entries": self._failed_file_cache_entries}
            d = os.path.dirname(self._failed_file_cache_path)
            if d:
                os.makedirs(d, exist_ok=True)
            tmp = self._failed_file_cache_path + ".tmp"
            with open(tmp, "w", encoding="utf-8") as f:
                json.dump(payload, f, ensure_ascii=False, separators=(",", ":"))
            os.replace(tmp, self._failed_file_cache_path)
            self._failed_file_cache_dirty = False
            self._failed_file_cache_new_since_save = 0
            self._failed_file_cache_last_save_time = time.time()
        except Exception as e:
            self._index_exception("保存 index_failed_files 失败", e)

    @staticmethod
    def _file_stat_fingerprint(abs_path: str) -> Optional[Dict[str, int]]:
        try:
            st = os.stat(abs_path)
            m_ns = getattr(st, "st_mtime_ns", int(st.st_mtime * 1_000_000_000))
            return {"mtime_ns": int(m_ns), "size": int(st.st_size)}
        except Exception:
            return None

    def _failed_file_cache_try_skip(self, abs_path: str) -> bool:
        """Skip an unchanged file that already failed indexing in a previous run."""
        try:
            ent = self._failed_file_cache_entries.get(abs_path)
            stat_fp = self._file_stat_fingerprint(abs_path)
            if not ent or not stat_fp:
                return False
            return (
                int(ent.get("mtime_ns", -1)) == stat_fp["mtime_ns"]
                and int(ent.get("size", -1)) == stat_fp["size"]
            )
        except Exception:
            return False

    def _failed_file_cache_record(self, abs_path: str, reason: str, stage: str = "ingest_failed") -> None:
        stat_fp = self._file_stat_fingerprint(abs_path)
        if not stat_fp:
            return
        try:
            self._failed_file_cache_entries[abs_path] = {
                **stat_fp,
                "reason": str(reason or "")[:500],
                "stage": str(stage or "ingest_failed")[:120],
                "failed_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            }
            self._failed_file_cache_dirty = True
            self._failed_file_cache_new_since_save += 1
            if (
                self._failed_file_cache_new_since_save >= 5
                or (time.time() - self._failed_file_cache_last_save_time) > 1.0
            ):
                self._save_failed_file_cache()
        except Exception as e:
            self._index_exception(f"记录失败文件缓存失败: {abs_path}", e)

    def _failed_file_cache_forget(self, abs_path: str) -> None:
        try:
            if abs_path in self._failed_file_cache_entries:
                self._failed_file_cache_entries.pop(abs_path, None)
                self._failed_file_cache_dirty = True
        except Exception:
            pass
    
    def get_indexed_files(self) -> Dict[str, str]:
        indexed: Dict[str, str] = {}
        try:
            total_count = self.collection.count()
            if total_count == 0:
                return indexed
            page_size = 5000
            for offset in range(0, total_count, page_size):
                batch = self.collection.get(include=['metadatas'], limit=page_size, offset=offset)
                for meta in (batch.get('metadatas') or []):
                    file_path = meta.get('file_path')
                    modified_time = meta.get('modified_time')
                    if file_path and modified_time:
                        indexed[file_path] = modified_time
        except Exception as e:
            self._index_exception("获取已索引文件失败", e)
        return indexed

    def get_indexed_file_paths(self) -> set:
        indexed_paths: set = set()
        try:
            keyword_index = getattr(self, "_keyword_index_manager", None)
            if keyword_index is not None:
                # This method is used by indexing write paths to decide whether a
                # file should be skipped. Stale keyword sidecars are acceptable
                # for read/query previews, but using them here can turn a deleted
                # media index into a false "already indexed" success.
                records = keyword_index.records_snapshot(require_current=True)
                if records:
                    for rec in records:
                        fp = str(rec.file_path or "").strip()
                        if fp:
                            indexed_paths.add(fp)
                    self._maybe_log_indexed_paths_stats(int(self.collection.count() or 0), len(indexed_paths))
                    return indexed_paths

            total_count = self.collection.count()
            if total_count == 0:
                return indexed_paths
            page_size = 5000
            for offset in range(0, total_count, page_size):
                batch = self.collection.get(include=['metadatas'], limit=page_size, offset=offset)
                for meta in (batch.get('metadatas') or []):
                    file_path = meta.get('file_path')
                    if file_path:
                        indexed_paths.add(file_path)
            self._maybe_log_indexed_paths_stats(total_count, len(indexed_paths))
        except Exception as e:
            self._index_exception("获取已索引文件路径失败", e)
        return indexed_paths

    @staticmethod
    def _short_name_for_log(file_name: str, max_len: int = 52) -> str:
        n = file_name or ""
        if len(n) <= max_len:
            return n
        ext = os.path.splitext(n)[1]
        base = n[: max(0, max_len - len(ext) - 3)] + "…"
        return base + ext if ext else base

    def _index_chunk_max_chars(self, ext: str = "") -> int:
        if ext.lower() in {".csv", ".tsv", ".xlsx", ".xls", ".numbers", ".json", ".jsonl"}:
            try:
                v = int(os.getenv("FILEAGENT_TABLE_CHUNK_CHARS", "3200"))
            except ValueError:
                v = 3200
            return max(800, min(v, 48000))

        try:
            v = int(os.getenv("FILEAGENT_INDEX_CHUNK_CHARS", "3800"))
        except ValueError:
            v = 3800
        return max(800, min(v, 48000))

    @staticmethod
    def _is_structured_table_ext(ext: str) -> bool:
        return str(ext or "").lower() in {".csv", ".tsv", ".xlsx", ".xls", ".numbers", ".json", ".jsonl"}

    def _table_embed_max_chars(self, ext: str = "") -> int:
        """Maximum structured-data text kept for embedding after deterministic sampling."""
        env_names = [
            "FILEAGENT_TABLE_EMBED_MAX_CHARS",
            "FILEAGENT_STRUCTURED_EMBED_MAX_CHARS",
        ]
        if str(ext or "").lower() in {".json", ".jsonl"}:
            env_names = [
                "FILEAGENT_JSON_EMBED_MAX_CHARS",
                "FILEAGENT_TABLE_EMBED_MAX_CHARS",
                "FILEAGENT_STRUCTURED_EMBED_MAX_CHARS",
            ]
        v = _env_int_first(env_names, 80000)
        return max(20000, min(v, 500000))

    def _compact_structured_text_for_embedding(
        self,
        text: str,
        metadata: Dict[str, Any],
        *,
        file_name: str,
        ext: str,
    ) -> str:
        """
        Keep structured files fast to embed without throwing away the table shape.

        Full spreadsheets can expand into hundreds of tiny embedding chunks. For
        file-level retrieval, schema + representative head/middle/tail samples
        usually beat indexing every row. The full parsed text is still available
        earlier for summary/classification; this only limits vector payload.
        """
        raw = str(text or "").strip()
        if not raw or not self._is_structured_table_ext(ext):
            return raw

        limit = self._table_embed_max_chars(ext)
        metadata["structured_embedding_original_chars"] = len(raw)
        metadata["structured_embedding_max_chars"] = limit
        if len(raw) <= limit:
            metadata["structured_embedding_compacted"] = False
            metadata["structured_embedding_chars"] = len(raw)
            return raw

        lines = [ln.rstrip() for ln in raw.splitlines() if ln.strip()]
        if not lines:
            metadata["structured_embedding_compacted"] = False
            metadata["structured_embedding_chars"] = len(raw[:limit])
            return raw[:limit]

        schema_hint = str(metadata.get("table_schema_hint") or "").strip()
        header_parts: List[str] = [
            (
                f"[Structured embedding sample] file={file_name} "
                f"original_chars={len(raw)} max_chars={limit}"
            )
        ]
        if schema_hint:
            header_parts.append(f"[Schema]\n{schema_hint}")

        header = "\n".join(header_parts).strip()
        remaining = max(4000, limit - len(header) - 300)
        head_budget = int(remaining * 0.50)
        middle_budget = int(remaining * 0.25)
        tail_budget = remaining - head_budget - middle_budget

        used_indices: set = set()

        def _collect_forward(start: int, budget: int) -> List[str]:
            out: List[str] = []
            chars = 0
            for idx in range(max(0, start), len(lines)):
                if idx in used_indices:
                    continue
                line = lines[idx]
                add_len = len(line) + (1 if out else 0)
                if out and chars + add_len > budget:
                    break
                if not out and add_len > budget:
                    out.append(line[:budget])
                    used_indices.add(idx)
                    break
                out.append(line)
                used_indices.add(idx)
                chars += add_len
            return out

        def _collect_backward(end_exclusive: int, budget: int) -> List[str]:
            out_rev: List[str] = []
            chars = 0
            for idx in range(min(len(lines), end_exclusive) - 1, -1, -1):
                if idx in used_indices:
                    continue
                line = lines[idx]
                add_len = len(line) + (1 if out_rev else 0)
                if out_rev and chars + add_len > budget:
                    break
                if not out_rev and add_len > budget:
                    out_rev.append(line[-budget:])
                    used_indices.add(idx)
                    break
                out_rev.append(line)
                used_indices.add(idx)
                chars += add_len
            out_rev.reverse()
            return out_rev

        head = _collect_forward(0, head_budget)
        middle_start = max(0, len(lines) // 2 - 20)
        middle = _collect_forward(middle_start, middle_budget)
        tail = _collect_backward(len(lines), tail_budget)

        sections = [header]
        if head:
            sections.append("[Head sample]\n" + "\n".join(head))
        if middle:
            sections.append("[Middle sample]\n" + "\n".join(middle))
        if tail:
            sections.append("[Tail sample]\n" + "\n".join(tail))
        compacted = "\n\n".join(s for s in sections if s).strip()
        if len(compacted) > limit:
            compacted = compacted[:limit].rstrip()

        metadata["structured_embedding_compacted"] = True
        metadata["structured_embedding_chars"] = len(compacted)
        metadata["structured_embedding_lines"] = len(head) + len(middle) + len(tail)
        return compacted

    def _split_text_for_max_chars(self, text: str, max_chars: int) -> List[str]:
        raw = (text or "").strip()
        if not raw:
            return []
        if len(raw) <= max_chars:
            return [raw]

        chunks: List[str] = []
        paragraphs = raw.split("\n\n")
        buf = ""
        for para in paragraphs:
            p = (para or "").strip()
            if not p:
                continue
            if len(p) > max_chars:
                if buf.strip():
                    chunks.append(buf.strip())
                    buf = ""
                lines = p.splitlines()
                piece = ""
                for ln in lines:
                    if len(piece) + len(ln) + 1 <= max_chars:
                        piece = (piece + "\n" + ln) if piece else ln
                    else:
                        if piece.strip():
                            chunks.append(piece.strip())
                        if len(ln) > max_chars:
                            for i in range(0, len(ln), max_chars):
                                chunks.append(ln[i : i + max_chars])
                            piece = ""
                        else:
                            piece = ln
                if piece.strip():
                    buf = piece
                continue
            if len(buf) + len(p) + 2 <= max_chars:
                buf = (buf + "\n\n" + p) if buf else p
            else:
                if buf.strip():
                    chunks.append(buf.strip())
                buf = p
        if buf.strip():
            chunks.append(buf.strip())
        return chunks if chunks else [raw[:max_chars]]

    def _split_merged_text_for_index_vectors(self, merged_text: str, ext: str = "") -> List[str]:
        max_chars = self._index_chunk_max_chars(ext)
        return self._split_text_for_max_chars(merged_text, max_chars)

    def _media_embed_chunk_max_chars(self, chunk_type: str = "") -> int:
        try:
            max_tokens = int(
                os.getenv(
                    "FILEAGENT_EMBED_MAX_TOKENS",
                    os.getenv("FILEAGENT_EMBED_N_BATCH", "512"),
                )
            )
        except Exception:
            max_tokens = 512
        max_tokens = max(64, min(max_tokens, 4096))
        default_chars = max(180, min(1200, int(max_tokens * 0.75)))
        env_name = "FILEAGENT_MEDIA_EMBED_CHUNK_CHARS"
        if str(chunk_type or "").strip().lower() in {
            "media_summary",
            "media_audio_summary",
            "media_visual_summary",
        }:
            env_name = "FILEAGENT_MEDIA_SUMMARY_CHUNK_CHARS"
        try:
            v = int(os.getenv(env_name, str(default_chars)))
        except Exception:
            v = default_chars
        return max(160, min(v, 2000))

    def _split_text_for_embedding_quality(
        self,
        text: str,
        *,
        reserve_text: str = "",
        fallback_max_chars: int,
    ) -> List[str]:
        raw = str(text or "").strip()
        if not raw:
            return []
        em = getattr(self, "embedding_model", None)
        split_fn = getattr(em, "split_text_for_embedding", None) if em is not None else None
        if callable(split_fn):
            try:
                parts = split_fn(raw, reserve_text=reserve_text)
                parts = [str(p or "").strip() for p in (parts or []) if str(p or "").strip()]
                if parts:
                    return parts
            except Exception as e:
                self._index_exception("embedding token-aware split 失败，回退字符切分", e)
        return self._split_text_for_max_chars(raw, max(160, int(fallback_max_chars or 160)))

    def _split_media_chunk_for_safe_embedding(
        self,
        chunk_text: str,
        chunk_meta: Dict[str, Any],
        file_name: str,
        reserve_text: str = "",
    ) -> List[Tuple[str, Dict[str, Any]]]:
        text = str(chunk_text or "").strip()
        if not text:
            return []
        chunk_type = str(chunk_meta.get("chunk_type") or "media")
        max_chars = self._media_embed_chunk_max_chars(chunk_type)
        parts = self._split_text_for_embedding_quality(
            text,
            reserve_text=reserve_text,
            fallback_max_chars=max_chars,
        )
        total = len(parts)
        if total <= 1:
            return [(text, chunk_meta)]

        logger.info(
            f"[MediaChunk] split overlong chunk file={file_name} type={chunk_type} "
            f"parts={total} max_chars={max_chars} text_len={len(text)}"
        )

        split_chunks: List[Tuple[str, Dict[str, Any]]] = []
        for idx, part in enumerate(parts, 1):
            part_meta = dict(chunk_meta)
            part_meta["chunk_part_index"] = idx
            part_meta["chunk_part_total"] = total
            part_meta["chunk_is_split"] = True
            part_meta["chunk_parent_type"] = chunk_type

            prefix = f"[{file_name} | {chunk_type} {idx}/{total}] "
            if chunk_type == "keyframe" and "keyframe_time_sec" in chunk_meta:
                try:
                    prefix = (
                        f"[{file_name} visual @ {float(chunk_meta['keyframe_time_sec']):.0f}s "
                        f"{idx}/{total}] "
                    )
                except Exception:
                    pass
            elif chunk_type in {"asr_transcript", "asr_segment"}:
                try:
                    prefix = (
                        f"[{file_name} @ {float(chunk_meta.get('asr_start_sec', 0.0)):.0f}s-"
                        f"{float(chunk_meta.get('asr_end_sec', 0.0)):.0f}s {idx}/{total}] "
                    )
                except Exception:
                    pass
            elif chunk_type in {"media_summary", "media_audio_summary", "media_visual_summary"}:
                prefix = f"[{file_name} summary {idx}/{total}] "

            split_chunks.append((prefix + part.strip(), part_meta))
        return split_chunks

    @staticmethod
    def _format_media_interval_label(start_sec: float, end_sec: float) -> str:
        def _fmt(ts: float) -> str:
            total = max(0, int(round(float(ts or 0.0))))
            minutes, seconds = divmod(total, 60)
            hours, minutes = divmod(minutes, 60)
            if hours > 0:
                return f"{hours}:{minutes:02d}:{seconds:02d}"
            return f"{minutes}:{seconds:02d}"

        return f"{_fmt(start_sec)} - {_fmt(end_sec)}"

    @staticmethod
    def _media_interval_doc_id(
        file_path: str,
        chunk_type: str,
        interval_start_sec: float,
        interval_end_sec: float,
        *,
        suffix: str = "",
    ) -> str:
        raw = "|".join(
            [
                os.path.abspath(str(file_path or "")),
                str(chunk_type or ""),
                f"{float(interval_start_sec or 0.0):.1f}",
                f"{float(interval_end_sec or 0.0):.1f}",
                str(suffix or ""),
            ]
        )
        digest = hashlib.sha256(raw.encode("utf-8")).hexdigest()[:40]
        return f"media_interval_{digest}"

    def _base_media_metadata_for_file(self, file_path: str) -> Dict[str, Any]:
        file_name = os.path.basename(file_path)
        file_name_no_ext, ext = os.path.splitext(file_name)
        parent_folder = os.path.basename(os.path.dirname(file_path))
        media_type = "video" if ext.lower() in {
            ".mp4", ".m4v", ".avi", ".mov", ".mkv", ".webm", ".flv", ".wmv", ".ts",
        } else "audio"
        default_meta: Dict[str, Any] = {
            "file_path": file_path,
            "file_name": file_name,
            "file_name_no_ext": file_name_no_ext,
            "file_name_en": file_name_no_ext,
            "folder_name_en": parent_folder,
            "file_extension": ext.lower(),
            "parent_folder": parent_folder,
            "media_type": media_type,
            "doc_category": media_type,
            "doc_category_raw": media_type,
            "doc_category_family": media_type,
            "doc_category_leaf": media_type,
            "doc_role": "primary_source",
            "doc_taxonomy_confidence": 1.0,
            "lookup_aliases": "",
        }
        try:
            existing = self.collection.get(
                where={"file_path": file_path},
                include=["metadatas"],
            )
            metadatas = list(existing.get("metadatas") or [])
        except Exception:
            metadatas = []

        preferred_order = {
            "media_summary": 0,
            "media_audio_summary": 1,
            "media_visual_summary": 2,
            "asr_transcript": 3,
            "asr_segment": 4,
            "keyframe": 5,
        }
        picked: Dict[str, Any] = {}
        best_rank = 99
        for meta in metadatas:
            chunk_type = str((meta or {}).get("chunk_type") or "")
            if chunk_type in {"interval_summary", "interval_visual"}:
                continue
            rank = preferred_order.get(chunk_type, 50)
            if rank < best_rank:
                picked = dict(meta or {})
                best_rank = rank

        merged = dict(default_meta)
        for key, value in picked.items():
            if value is not None:
                merged[key] = value

        merged["file_path"] = file_path
        merged["file_name"] = str(merged.get("file_name") or file_name)
        merged["file_name_no_ext"] = str(merged.get("file_name_no_ext") or file_name_no_ext)
        merged["file_extension"] = str(merged.get("file_extension") or ext.lower())
        merged["parent_folder"] = str(merged.get("parent_folder") or parent_folder)
        merged["media_type"] = str(merged.get("media_type") or media_type)
        merged["doc_category"] = str(merged.get("doc_category") or merged["media_type"])
        merged["doc_category_raw"] = str(merged.get("doc_category_raw") or merged["doc_category"])
        merged["doc_category_family"] = str(merged.get("doc_category_family") or merged["doc_category"])
        merged["doc_category_leaf"] = str(merged.get("doc_category_leaf") or merged["doc_category"])
        merged["lookup_aliases"] = str(merged.get("lookup_aliases") or "")
        return merged

    def get_cached_media_interval_analysis(
        self,
        file_path: str,
        interval_start_sec: float,
        interval_end_sec: float,
        *,
        language: str = "",
        tolerance_sec: float = 2.0,
    ) -> Dict[str, Any]:
        result: Dict[str, Any] = {
            "summary_text": "",
            "visual_entries": [],
            "frame_count": 0,
            "transcript_count": 0,
            "source": "none",
        }
        try:
            raw = self.collection.get(
                where={"file_path": file_path},
                include=["metadatas", "documents"],
            )
        except Exception as e:
            logger.debug(f"[media_interval_cache] lookup failed for {os.path.basename(file_path)}: {e}")
            return result

        def _matches(meta: Dict[str, Any]) -> bool:
            start = meta.get("interval_start_sec")
            end = meta.get("interval_end_sec")
            if start is None or end is None:
                return False
            return (
                abs(float(start) - float(interval_start_sec)) <= tolerance_sec
                and abs(float(end) - float(interval_end_sec)) <= tolerance_sec
            )

        requested_lang = str(language or "").strip().lower()
        summary_by_lang: Dict[str, Tuple[str, Dict[str, Any]]] = {}
        visual_entries: List[Tuple[float, str]] = []

        for meta, doc in zip(list(raw.get("metadatas") or []), list(raw.get("documents") or [])):
            meta = dict(meta or {})
            if not _matches(meta):
                continue
            chunk_type = str(meta.get("chunk_type") or "")
            if chunk_type == "interval_summary":
                answer_lang = str(meta.get("answer_language") or "").strip().lower()
                summary_text = re.sub(r"^\[.*?\]\s*", "", str(doc or "")).strip()
                if summary_text:
                    summary_by_lang[answer_lang] = (summary_text, meta)
            elif chunk_type == "interval_visual":
                ts = meta.get("interval_visual_time_sec", meta.get("keyframe_time_sec"))
                if ts is None:
                    continue
                desc = re.sub(r"^\[.*?\]\s*", "", str(doc or "")).strip()
                if desc:
                    visual_entries.append((float(ts), desc))

        summary_text = ""
        summary_meta: Dict[str, Any] = {}
        if requested_lang and requested_lang in summary_by_lang:
            summary_text, summary_meta = summary_by_lang[requested_lang]
        elif not requested_lang and summary_by_lang:
            _, (summary_text, summary_meta) = next(iter(summary_by_lang.items()))

        visual_entries.sort(key=lambda item: item[0])
        result["summary_text"] = summary_text
        result["visual_entries"] = visual_entries
        result["frame_count"] = int(summary_meta.get("interval_frame_count") or len(visual_entries) or 0)
        result["transcript_count"] = int(summary_meta.get("interval_transcript_count") or 0)
        result["source"] = "cached_interval_analysis" if (summary_text or visual_entries) else "none"
        return result

    def persist_media_interval_analysis(
        self,
        file_path: str,
        interval_start_sec: float,
        interval_end_sec: float,
        *,
        summary_text: str,
        visual_entries: Sequence[Tuple[float, str]],
        transcript_rows: Sequence[Dict[str, Any]],
        source_label: str,
        answer_language: str,
    ) -> int:
        summary_clean = re.sub(r"\s+", " ", str(summary_text or "")).strip()
        clean_visual_entries = [
            (round(float(ts), 1), re.sub(r"\s+", " ", str(desc or "")).strip())
            for ts, desc in (visual_entries or [])
            if str(desc or "").strip()
        ]
        if not summary_clean and not clean_visual_entries:
            return 0

        base_meta = self._base_media_metadata_for_file(file_path)
        file_name = str(base_meta.get("file_name") or os.path.basename(file_path))
        interval_label = self._format_media_interval_label(interval_start_sec, interval_end_sec)
        summary_doc = summary_clean[:1800]
        source = str(source_label or "db_interval_analysis")
        lang = str(answer_language or "").strip().lower() or "unknown"
        transcript_count = len(list(transcript_rows or []))
        frame_count = len(clean_visual_entries)
        cache_key = f"{float(interval_start_sec):.1f}-{float(interval_end_sec):.1f}"

        documents: List[str] = []
        metadatas: List[Dict[str, Any]] = []
        ids: List[str] = []
        embed_inputs: List[str] = []

        if summary_doc:
            summary_meta = {
                **base_meta,
                "chunk_type": "interval_summary",
                "doc_role": "summary",
                "doc_summary": summary_doc[:700],
                "interval_start_sec": round(float(interval_start_sec), 1),
                "interval_end_sec": round(float(interval_end_sec), 1),
                "interval_mid_sec": round((float(interval_start_sec) + float(interval_end_sec)) / 2.0, 1),
                "interval_frame_count": frame_count,
                "interval_transcript_count": transcript_count,
                "interval_source": source,
                "interval_cache_key": cache_key,
                "answer_language": lang,
            }
            self._set_doc_summary_provenance(summary_meta, "media_interval_analysis")
            summary_meta["lookup_aliases"] = self._build_lookup_aliases(
                file_path,
                summary_meta,
                interval_label,
                f"{float(interval_start_sec):.1f}",
                f"{float(interval_end_sec):.1f}",
                summary_doc[:400],
            )
            documents.append(summary_doc)
            metadatas.append(summary_meta)
            ids.append(
                self._media_interval_doc_id(
                    file_path,
                    "interval_summary",
                    interval_start_sec,
                    interval_end_sec,
                    suffix=f"summary:{lang}",
                )
            )
            embed_inputs.append(f"[{file_name} interval {interval_label}] {summary_doc}")

        for idx, (ts, desc) in enumerate(clean_visual_entries, 1):
            visual_doc = desc[:900]
            visual_meta = {
                **base_meta,
                "chunk_type": "interval_visual",
                "doc_role": "analysis",
                "doc_summary": visual_doc[:240],
                "interval_start_sec": round(float(interval_start_sec), 1),
                "interval_end_sec": round(float(interval_end_sec), 1),
                "interval_mid_sec": round((float(interval_start_sec) + float(interval_end_sec)) / 2.0, 1),
                "interval_visual_time_sec": round(float(ts), 1),
                "keyframe_time_sec": round(float(ts), 1),
                "interval_visual_order": idx,
                "interval_frame_count": frame_count,
                "interval_transcript_count": transcript_count,
                "interval_source": source,
                "interval_cache_key": cache_key,
                "answer_language": "en",
            }
            visual_meta["lookup_aliases"] = self._build_lookup_aliases(
                file_path,
                visual_meta,
                interval_label,
                f"{float(ts):.1f}",
                visual_doc[:220],
            )
            documents.append(visual_doc)
            metadatas.append(visual_meta)
            ids.append(
                self._media_interval_doc_id(
                    file_path,
                    "interval_visual",
                    interval_start_sec,
                    interval_end_sec,
                    suffix=f"{float(ts):.1f}",
                )
            )
            embed_inputs.append(f"[{file_name} visual @ {float(ts):.1f}s within {interval_label}] {visual_doc}")

        if not documents:
            return 0

        try:
            embeddings = self._embed_texts_for_index(embed_inputs)
        except Exception as e:
            self._index_exception("interval evidence embedding 失败", e)
            return 0

        valid_rows = [
            (doc_id, emb, doc, meta)
            for doc_id, emb, doc, meta in zip(ids, embeddings, documents, metadatas)
            if emb and any(emb)
        ]
        if not valid_rows:
            return 0

        self.collection.upsert(
            ids=[row[0] for row in valid_rows],
            embeddings=[row[1] for row in valid_rows],
            documents=[row[2] for row in valid_rows],
            metadatas=[row[3] for row in valid_rows],
        )
        self._invalidate_meta_cache()
        self._maybe_persist()
        logger.info(
            "[media_interval_cache] upserted file=%s interval=%s chunks=%s source=%s",
            file_name,
            interval_label,
            len(valid_rows),
            source,
        )
        return len(valid_rows)

    def _embed_texts_for_index(self, texts: List[str]) -> List[List[float]]:
        if not texts:
            return []
        em = self.embedding_model
        if em is None:
            raise RuntimeError("embedding model missing")
        formatted_texts = [self._format_doc_for_embedding(t) for t in texts]
        batch_fn = getattr(em, "get_text_embedding_batch", None)
        if callable(batch_fn):
            try:
                return batch_fn(formatted_texts, show_progress=False)  # type: ignore[misc]
            except TypeError:
                try:
                    return batch_fn(formatted_texts)  # type: ignore[misc]
                except Exception as e:
                    self._index_exception("get_text_embedding_batch 失败，回退逐条", e)
        batch_fn2 = getattr(em, "_get_text_embeddings", None)
        if callable(batch_fn2):
            try:
                return batch_fn2(formatted_texts)  # type: ignore[misc]
            except Exception as e:
                self._index_exception("_get_text_embeddings 失败，回退逐条", e)
        return [em.get_text_embedding(t) for t in formatted_texts]
    
    def ingest_file(
        self,
        file_path: str,
        use_smart_indexing: bool = True,
        should_cancel: Optional[Callable[[], bool]] = None,
        on_frame_progress: Optional[Callable[[int, int], None]] = None,
        on_media_progress: Optional[Callable[[Dict[str, Any]], None]] = None,
    ) -> bool:
        if not os.path.exists(file_path):
            self._index_warning(f"索引跳过：文件不存在 {file_path}")
            return False
        
        if self._should_ignore_file(file_path):
            self._index_info(f"索引跳过：命中忽略规则 {file_path}")
            return False

        self._failed_file_cache_forget(os.path.abspath(file_path))

        _start_time = time.time()
        _llm_time = 0.0
        _vl_time = 0.0
        _embed_time = 0.0
        _chunks = 0
        _doc_category = "other"

        try:
            ext = os.path.splitext(file_path)[1].lower()
            stat = os.stat(file_path)
            file_name = os.path.basename(file_path)
            self._index_info(
                f"开始索引文件: file={file_path} | ext={ext or '<none>'} | size_kb={stat.st_size / 1024:.2f} | smart={use_smart_indexing}"
            )
            self._maybe_log_startup_index_prefill_observation(
                file_path,
                use_smart_indexing=use_smart_indexing,
            )
            file_name_no_ext = os.path.splitext(file_name)[0]
            
            file_name_en = ""
            file_name_needs_translation = self._label_needs_translation(file_name_no_ext)

            _text_smart_summary_file_name_en = (
                use_smart_indexing
                and file_name_needs_translation
                and ext not in IMAGE_EXTENSIONS
                and ext not in AUDIO_EXTENSIONS
                and ext not in {".avi", ".mov", ".mkv", ".webm"}
            )

            if _text_smart_summary_file_name_en:
                self._index_info(f"文件名英文增强将由专用翻译函数单独落库: {file_name}")
            elif file_name_needs_translation:
                file_name_en = self._translate_file_name_to_en(file_name)
                self._index_info(f"文件名跨语言增强: {file_name} -> {file_name_en}")
            else:
                file_name_en = file_name_no_ext

            parent_folder = os.path.basename(os.path.dirname(file_path))
            folder_name_en = ""
            folder_needs_translation = self._label_needs_translation(parent_folder)
            
            if folder_needs_translation:
                folder_name_en = self._translate_folder_name(parent_folder)
                self._index_info(f"文件夹名跨语言增强: {parent_folder} -> {folder_name_en}")
            else:
                folder_name_en = parent_folder

            metadata = {
                "file_path": file_path,
                "file_name": file_name,
                "file_name_no_ext": file_name_no_ext,
                "file_name_en": file_name_en,
                "folder_name_en": folder_name_en,
                "file_extension": ext,
                "file_size_kb": round(stat.st_size / 1024, 2),
                "modified_time": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
                "parent_folder": parent_folder,
                "doc_category": "other",
                "doc_category_raw": "other",
                "doc_category_family": "other",
                "doc_category_leaf": "other",
                "doc_role": "primary_source",
                "doc_taxonomy_confidence": 0.0,
                "doc_taxonomy_version": 2,
                "doc_summary": "",
                "doc_summary_model_id": "",
                "doc_summary_saved_at": "",
            }
            metadata["lookup_aliases"] = self._build_lookup_aliases(file_path, metadata)
            
            metadata = {k: v for k, v in metadata.items() if v is not None}
            
            from core.media.media_expert import AUDIO_EXTENSIONS as _AUDIO_EXTS, VIDEO_EXTENSIONS as _VIDEO_EXTS, MEDIA_EXTENSIONS as _MEDIA_EXTS
            if ext.lower() in _MEDIA_EXTS:
                metadata["media_type"] = "video" if ext.lower() in _VIDEO_EXTS else "audio"
                metadata["doc_category"] = metadata["media_type"]
                metadata["doc_category_raw"] = metadata["media_type"]
                metadata["doc_category_family"] = metadata["media_type"]
                metadata["doc_category_leaf"] = metadata["media_type"]
                metadata["doc_role"] = "primary_source"
                metadata["doc_taxonomy_confidence"] = 1.0
                _doc_category = metadata["doc_category"]
                is_video = ext.lower() in _VIDEO_EXTS

                # ── Try MediaExpert processing (ASR + keyframes + summary) ────
                _media_processed = False
                _media_processing_error = ""
                try:
                    from core.media.media_expert import MediaExpert
                    # Check if pywhispercpp (primary) or faster-whisper (fallback) is available
                    import importlib
                    _asr_available = (
                        importlib.util.find_spec("pywhispercpp") is not None
                        or importlib.util.find_spec("faster_whisper") is not None
                    )
                    _llm_ready = self._test_local_llm_connection() if use_smart_indexing else False
                    _can_run_media_expert = bool(use_smart_indexing and (_asr_available or (is_video and _llm_ready)))

                    if _can_run_media_expert:
                        self._index_info(
                            f"[MediaExpert] 开始处理{'视频' if is_video else '音频'}：{file_name} "
                            f"(asr_available={_asr_available}, llm_ready={_llm_ready})"
                        )
                        _llm_client = self._get_local_llm_client() if _llm_ready else None

                        # VL describe function for video keyframes — uses a stricter
                        # prompt that forbids guessing product/brand names
                        _vl_fn = None
                        _ocr_fn = None
                        if is_video and _llm_ready:
                            _vl_fn = lambda img_path, prev_description="": self._generate_video_frame_summary(img_path, prev_description=prev_description) or ""
                            _ocr_fn = lambda img_path: self._generate_video_frame_ocr(img_path) or ""

                        expert = MediaExpert(
                            llm_client=_llm_client,
                            vl_describe_fn=_vl_fn,
                            frame_ocr_fn=_ocr_fn,
                            on_frame_progress=on_frame_progress,
                            on_media_progress=on_media_progress,
                        )

                        if is_video:
                            result = expert.process_video(file_path)
                        else:
                            result = expert.process_audio(file_path)

                        if (
                            float(getattr(result, "duration_sec", 0.0) or 0.0) <= 0.0
                            and not str(getattr(result, "transcript", "") or "").strip()
                            and not list(getattr(result, "asr_segments", []) or [])
                            and not list(getattr(result, "keyframes", []) or [])
                        ):
                            raise RuntimeError(
                                "Media parser produced no usable content "
                                "(duration=0, no ASR segments, no keyframes)"
                            )

                        # Build multi-chunk index entries
                        _visual_desc_count = sum(1 for _kf in result.keyframes if getattr(_kf, "description", ""))
                        _has_real_asr = expert._has_meaningful_asr_transcript(result.transcript)
                        if result.media_type == "audio":
                            _media_index_mode = "asr_first" if _has_real_asr else "metadata_only"
                        else:
                            if _has_real_asr:
                                _media_index_mode = "asr_first"
                            elif _visual_desc_count > 0:
                                _media_index_mode = "visual_priority"
                            else:
                                _media_index_mode = "metadata_only"

                        _base_meta = {
                            **metadata,
                            "media_index_mode": _media_index_mode,
                            "media_has_audio_track": bool(result.video_metadata.get("has_audio")),
                            "media_has_asr_transcript": _has_real_asr,
                            "media_has_visual_descriptions": bool(_visual_desc_count),
                        }

                        # Build chunks first so we can use keyframe descriptions as
                        # fallback summary when there is no ASR transcript.
                        chunks = expert.build_index_chunks(result, _base_meta)

                        # Derive doc_summary: prefer transcript summary, then
                        # keyframe visual descriptions (from chunk 0 built above),
                        # then bare filename as last resort.
                        _summary_chunk_text = chunks[0][0] if chunks else ""
                        _visual_summary_line = ""
                        for line in _summary_chunk_text.splitlines():
                            if line.startswith("Content summary:"):
                                _visual_summary_line = line[len("Content summary:"):].strip()
                                break
                            if line.startswith("内容摘要:"):
                                _visual_summary_line = line[len("内容摘要:"):].strip()
                                break
                        _transcript_summary = (
                            result.transcript_summary
                            if not expert._looks_like_low_signal_summary_text(result.transcript_summary)
                            else ""
                        )
                        _base_meta["doc_summary"] = (
                            _visual_summary_line
                            or _transcript_summary
                            or (result.transcript[:300] if _has_real_asr else "")
                            or f"Media file: {file_name}"
                        )
                        self._set_doc_summary_provenance(
                            _base_meta,
                            self._snapshot_loaded_llm_model_id() or "media_expert"
                        )

                        # ── Build en_tags for the summary ──
                        en_tags = self._build_en_tags(
                            file_name,
                            ext,
                            doc_summary=_base_meta.get("doc_summary", ""),
                            parent_folder=parent_folder,
                            folder_name_en=folder_name_en,
                        )
                        if en_tags:
                            _base_meta["en_tags"] = en_tags
                        _base_meta["lookup_aliases"] = self._build_lookup_aliases(file_path, _base_meta)

                        # Back-fill the updated doc_summary into every chunk's metadata
                        chunks = [
                            (
                                txt,
                                {
                                    **meta,
                                    "doc_summary": _base_meta["doc_summary"],
                                    "lookup_aliases": _base_meta.get("lookup_aliases", ""),
                                },
                            )
                            for txt, meta in chunks
                        ]
                        _media_folder_ctx = " ".join(
                            x for x in [
                                str(_base_meta.get("parent_folder") or "").strip(),
                                str(_base_meta.get("folder_name_en") or "").strip(),
                            ] if x
                        ).strip()
                        expanded_chunks: List[Tuple[str, Dict[str, Any]]] = []
                        for chunk_text, chunk_meta in chunks:
                            _reserve_parts: List[str] = []
                            if chunk_meta.get("chunk_type") == "keyframe":
                                _reserve_parts.append("Video Frame: ")
                            if _media_folder_ctx:
                                _reserve_parts.append(f"Folder Context: {_media_folder_ctx}\n")
                            expanded_chunks.extend(
                                self._split_media_chunk_for_safe_embedding(
                                    chunk_text,
                                    chunk_meta,
                                    file_name,
                                    reserve_text="".join(_reserve_parts),
                                )
                            )
                        chunks = expanded_chunks

                        if is_video and _llm_ready:
                            self._release_local_llm_before_media_embedding(
                                reason=f"video_media_chunks file={file_name}"
                            )

                        try:
                            _media_write_batch = int(
                                os.getenv("FILEAGENT_MEDIA_CHUNK_WRITE_BATCH", "3")
                            )
                        except Exception:
                            _media_write_batch = 3
                        _media_write_batch = max(1, min(_media_write_batch, 16))

                        _embed_start = time.time()
                        _batch_ids: List[str] = []
                        _batch_embeddings: List[List[float]] = []
                        _batch_documents: List[str] = []
                        _batch_metadatas: List[Dict[str, Any]] = []

                        def _flush_media_chunk_batch(final_flush: bool = False) -> None:
                            if not _batch_ids:
                                return
                            self.collection.add(
                                ids=list(_batch_ids),
                                embeddings=list(_batch_embeddings),
                                documents=list(_batch_documents),
                                metadatas=list(_batch_metadatas),
                            )
                            self._invalidate_meta_cache()
                            logger.info(
                                f"[MediaChunk] {'final ' if final_flush else ''}stored "
                                f"file={file_name} chunks={len(_batch_ids)}"
                            )
                            _batch_ids.clear()
                            _batch_embeddings.clear()
                            _batch_documents.clear()
                            _batch_metadatas.clear()

                        for i, (chunk_text, chunk_meta) in enumerate(chunks):
                            # Clean None values from chunk metadata
                            chunk_meta = {k: v for k, v in chunk_meta.items() if v is not None}
                            chunk_type = str(chunk_meta.get("chunk_type") or "unknown")
                            logger.info(
                                f"[MediaChunk] start file={file_name} idx={i+1}/{len(chunks)} "
                                f"type={chunk_type} text_len={len(chunk_text or '')}"
                            )

                            embed_text = chunk_text
                            if chunk_meta.get("chunk_type") == "keyframe":
                                embed_text = f"Video Frame: {chunk_text}"
                            if _media_folder_ctx:
                                embed_text = f"Folder Context: {_media_folder_ctx}\n{embed_text}"
                                
                            with self._embed_context(f"file={file_path} phase=media_chunk_{i}"):
                                embedding = self._embed_doc_text(embed_text)
                            if not any(embedding):
                                msg = f"[Embedding] zero vector | file={file_path} phase=media_chunk_{i}"
                                self._index_error(msg)
                                self._append_embedding_error_log(msg)
                                continue
                            doc_id = f"{file_path}_{uuid.uuid4().hex[:8]}"
                            _batch_ids.append(doc_id)
                            _batch_embeddings.append(embedding)
                            _batch_documents.append(chunk_text)
                            _batch_metadatas.append(chunk_meta)
                            logger.info(
                                f"[MediaChunk] embedded file={file_name} idx={i+1}/{len(chunks)} "
                                f"type={chunk_type}"
                            )
                            if len(_batch_ids) >= _media_write_batch:
                                _flush_media_chunk_batch(final_flush=False)
                        if _batch_ids:
                            _flush_media_chunk_batch(final_flush=True)
                        _embed_time = time.time() - _embed_start
                        _chunks = len(chunks)
                        self._maybe_persist()

                        self._index_info(
                            f"[MediaExpert] ✅ {file_name}: {_chunks} chunks indexed "
                            f"(ASR: {len(result.asr_segments)} segments, "
                            f"keyframes: {len(result.keyframes)}, "
                            f"visual_desc={_visual_desc_count}, "
                            f"mode={_media_index_mode}, "
                            f"duration: {result.duration_sec:.0f}s)"
                        )
                        _media_processed = True

                except ImportError as _ie:
                    _media_processing_error = str(_ie)
                    self._index_info(f"[MediaExpert] 依赖不可用 ({_ie})，将标记为索引失败")
                except Exception as _me:
                    _media_processing_error = str(_me)
                    self._index_exception("[MediaExpert] 处理失败，将标记为索引失败", _me)

                # ── Fallback: filename-only indexing (original logic) ─────────
                if not _media_processed:
                    if use_smart_indexing and _media_processing_error:
                        self._index_error(
                            f"[MediaExpert] 媒体内容解析失败，不写入文件名占位索引: "
                            f"{file_name} | error={_media_processing_error[:300]}"
                        )
                        self._log_index_file_timing(
                            file_path=file_path,
                            status="failed",
                            stage="media_parse_failed",
                            total_s=time.time() - _start_time,
                            embed_s=_embed_time,
                            llm_s=_llm_time,
                            vl_s=_vl_time,
                            chunks=0,
                            doc_category=_doc_category,
                            extra=f"error={_media_processing_error[:300]}",
                        )
                        return False

                    self._set_doc_summary_provenance(metadata, "none")
                    self._index_info(f"媒体文件快速索引（英文摘要增强）：{file_name}")

                    _AUDIO_FORMAT_DESC: Dict[str, str] = {
                        ".mp3":  "Music/audio file", ".wav":  "Audio recording",
                        ".flac": "Lossless audio",   ".aac":  "Audio file",
                        ".ogg":  "Audio file",        ".m4a":  "Audio file",
                        ".wma":  "Audio file",        ".aiff": "Lossless audio",
                        ".mp4":  "Video file",        ".m4v":  "Video file",
                        ".avi":  "Video file",        ".mov":  "Video file",
                        ".mkv":  "Video file",        ".webm": "Video file",
                    }
                    fmt_desc = _AUDIO_FORMAT_DESC.get(ext.lower(), "Media file")
                    name_no_ext = os.path.splitext(file_name)[0]

                    _has_cjk = any("\u4e00" <= ch <= "\u9fff" for ch in name_no_ext)
                    _llm_trans_desc = ""
                    if _has_cjk and use_smart_indexing and self._test_local_llm_connection():
                        try:
                            _llm_start_trans = time.time()
                            _trans_client = self._get_local_llm_client()
                            _idx_model = self._require_configured_index_model_id("媒体文件名翻译")
                            if not _idx_model:
                                raise RuntimeError("missing index model")
                            _trans_prompt = self._append_model_prompt_suffix(
                                f'Translate this Chinese audio filename into 5-8 English semantic keywords '
                                f'that describe the sound or content. Filename: "{name_no_ext}"\n'
                                f'Output ONLY English keywords separated by spaces (no punctuation, no Chinese).',
                                _idx_model
                            )
                            _trans_resp = _trans_client.chat.completions.create(
                                model=_idx_model,
                                messages=[{"role": "user", "content": _trans_prompt}],
                                max_tokens=60, temperature=0.0, stream=False,
                            )
                            _raw = (_trans_resp.choices[0].message.content or "").strip()
                            if _raw:
                                import re as _re
                                _raw = _re.sub(r"<think>.*?</think>", "", _raw, flags=_re.DOTALL).strip()
                                _llm_trans_desc = _raw[:120]
                                logger.info(f"[AudioTranslate] {file_name} → {_llm_trans_desc}")
                            _llm_time += time.time() - _llm_start_trans
                        except Exception as _e:
                            logger.debug(f"[AudioTranslate] skipped: {_e}")

                    if _llm_trans_desc:
                        rich_summary = (
                            f"{fmt_desc}: {_llm_trans_desc}. "
                            f"Format: {ext.upper().lstrip('.')}. Original filename: {file_name}"
                        )
                    else:
                        rich_summary = (
                            f"{fmt_desc}: {name_no_ext}. "
                            f"Format: {ext.upper().lstrip('.')}. Original filename: {file_name}"
                        )
                    metadata["doc_summary"] = rich_summary

                    en_tags = self._build_en_tags(
                        file_name,
                        ext,
                        doc_summary=rich_summary,
                        parent_folder=parent_folder,
                        folder_name_en=folder_name_en,
                    )
                    if en_tags:
                        metadata["en_tags"] = en_tags
                    metadata["lookup_aliases"] = self._build_lookup_aliases(file_path, metadata)

                    text = rich_summary
                    if en_tags:
                        text = f"{rich_summary}\nKeywords: {en_tags}"
                    folder_ctx = " ".join(x for x in [parent_folder, folder_name_en] if x).strip()
                    if folder_ctx:
                        text = f"Folder Context: {folder_ctx}\n{text}"
                    
                    _embed_start = time.time()
                    with self._embed_context(f"file={file_path} phase=media_name_embedding"):
                        embedding = self._embed_doc_text(text)
                    if not any(embedding):
                        msg = f"[Embedding] zero vector | file={file_path} phase=media_name_embedding"
                        self._index_error(msg)
                        self._append_embedding_error_log(msg)
                    doc_id = f"{file_path}_{uuid.uuid4().hex[:8]}"
                    self.collection.add(
                        ids=[doc_id], embeddings=[embedding],
                        documents=[text], metadatas=[metadata],
                    )
                    self._invalidate_meta_cache()
                    _embed_time = time.time() - _embed_start
                    _chunks = 1

                self._upsert_filename_lookup_for_file(file_path, metadata, persist=False)
                self._maybe_persist()
                self._log_index_file_timing(
                    file_path=file_path,
                    status="ok",
                    stage="media_expert" if _media_processed else "media_fast_path",
                    total_s=time.time() - _start_time,
                    embed_s=_embed_time,
                    llm_s=_llm_time,
                    vl_s=_vl_time,
                    chunks=_chunks,
                    doc_category=_doc_category,
                )
                return True
            
            if ext in IMAGE_EXTENSIONS:
                metadata["doc_category"] = "image"
                metadata["doc_category_raw"] = "image"
                metadata["doc_category_family"] = "image"
                metadata["doc_category_leaf"] = "image"
                metadata["doc_role"] = "primary_source"
                metadata["doc_taxonomy_confidence"] = 1.0
                _doc_category = metadata["doc_category"]

                fast_image = _env_truthy("FILEAGENT_FAST_IMAGE_INDEX")
                if use_smart_indexing and not fast_image and self._test_local_llm_connection():
                    _vl_start = time.time()
                    self._index_info(f"正在使用视觉模型分析图片: {file_name}")
                    image_summary = self._generate_image_summary(file_path)
                    if image_summary:
                        metadata["doc_summary"] = image_summary
                        text = f"Image file {file_name}. Description: {image_summary}"
                        self._set_doc_summary_provenance(
                            metadata, self._snapshot_loaded_llm_model_id() or "vl_unknown"
                        )
                    else:
                        metadata["doc_summary"] = f"Image file: {file_name}"
                        text = f"Image file {file_name}"
                        self._set_doc_summary_provenance(metadata, "none")
                    _vl_time += time.time() - _vl_start
                else:
                    if use_smart_indexing and not fast_image:
                        self._index_warning(f"图片摘要降级：本地 LLM 不可用，使用文件名占位摘要 | {file_name}")
                    metadata["doc_summary"] = f"Image file: {file_name}"
                    text = f"Image file {file_name}"
                    self._set_doc_summary_provenance(metadata, "none")

                _img_en_tags = self._build_en_tags(
                    file_name,
                    ext,
                    doc_summary=metadata.get("doc_summary", ""),
                    parent_folder=parent_folder,
                    folder_name_en=folder_name_en,
                )
                if _img_en_tags:
                    metadata["en_tags"] = _img_en_tags
                metadata["lookup_aliases"] = self._build_lookup_aliases(file_path, metadata)
                _img_folder_ctx = " ".join(x for x in [parent_folder, folder_name_en] if x).strip()
                if _img_folder_ctx:
                    text = f"Folder Context: {_img_folder_ctx}\n{text}"
                # ──────────────────────────────────────────────────────────────
                
                _embed_start = time.time()
                with self._embed_context(f"file={file_path} phase=image_summary_embedding"):
                    embedding = self._embed_doc_text(text)
                if not any(embedding):
                    msg = f"[Embedding] zero vector | file={file_path} phase=image_summary_embedding"
                    self._index_error(msg)
                    self._append_embedding_error_log(msg)
                doc_id = f"{file_path}_{uuid.uuid4().hex[:8]}"
                
                self.collection.add(
                    ids=[doc_id],
                    embeddings=[embedding],
                    documents=[text],
                    metadatas=[metadata]
                )
                self._invalidate_meta_cache()
                _embed_time = time.time() - _embed_start
                self._upsert_filename_lookup_for_file(file_path, metadata, persist=False)
                self._maybe_persist()
                self._log_index_file_timing(
                    file_path=file_path,
                    status="ok",
                    stage="image_vl_summary",
                    total_s=time.time() - _start_time,
                    embed_s=_embed_time,
                    llm_s=_llm_time,
                    vl_s=_vl_time,
                    chunks=1,
                    doc_category=_doc_category,
                )
                return True
            
            text = ""
            _tabular_exts = {".csv", ".tsv", ".xlsx", ".xls", ".numbers"}
            _tabular_parse_issue = ""
            _max_table_index_chars = max(1000, int(os.getenv("MAX_TABLE_INDEX_CHARS", "5000000")))
            _table_truncation_note = "\n\n... (Data truncated for semantic indexing)"

            def _table_cell_to_text(value: Any) -> str:
                try:
                    if value is None:
                        return ""
                    if isinstance(value, float) and math.isnan(value):
                        return ""
                    s = str(value).strip()
                    return "" if s.lower() in {"none", "nan"} else s
                except Exception:
                    return ""

            def _escape_table_cell(value: str) -> str:
                return str(value or "").replace("\n", " ").replace("|", "\\|")

            def _joined_table_len(parts: List[str]) -> int:
                return sum(len(p) for p in parts) + max(0, len(parts) - 1) * 2

            def _table_budget_remaining(parts: List[str]) -> int:
                remaining = _max_table_index_chars - _joined_table_len(parts)
                if parts:
                    remaining -= 2
                return max(0, remaining)

            def _rows_to_markdown_table(
                rows: Any,
                *,
                sheet_name: str = "",
                max_chars: Optional[int] = None,
            ) -> Tuple[str, int, bool]:
                lines: List[str] = []
                char_count = 0
                row_count = 0
                width = 0
                header_written = False
                truncated = False

                def _add_line(line: str) -> bool:
                    nonlocal char_count, truncated
                    if max_chars is not None:
                        sep_len = 1 if lines else 0
                        remaining = max_chars - char_count - sep_len
                        if remaining <= 0:
                            truncated = True
                            return False
                        if len(line) > remaining:
                            if lines:
                                char_count += 1
                            lines.append(line[:remaining])
                            char_count += remaining
                            truncated = True
                            return False
                    else:
                        sep_len = 1 if lines else 0
                    if lines:
                        char_count += sep_len
                    lines.append(line)
                    char_count += len(line)
                    return True

                for row in rows:
                    values = [_table_cell_to_text(v) for v in list(row or [])]
                    while values and values[-1] == "":
                        values.pop()
                    if not any(values):
                        continue
                    row_count += 1
                    if not header_written:
                        width = max(1, len(values))
                        if sheet_name and not _add_line(f"[Sheet: {sheet_name}]"):
                            break
                        header = values + [""] * (width - len(values))
                        if not _add_line(" | ".join(_escape_table_cell(h) for h in header)):
                            break
                        if not _add_line(" | ".join(["---"] * width)):
                            break
                        header_written = True
                        continue
                    row_width = max(width, len(values))
                    padded = values + [""] * (row_width - len(values))
                    if not _add_line(" | ".join(_escape_table_cell(v) for v in padded)):
                        break
                return "\n".join(lines), row_count, truncated

            def _xlsx_formula_value_to_text(value: Any, formula_value: Any) -> str:
                value_text = _table_cell_to_text(value)
                formula_text = _table_cell_to_text(formula_value)
                if formula_text.startswith("="):
                    if value_text and value_text != formula_text:
                        return f"{value_text} [formula: {formula_text}]"
                    return formula_text
                return value_text or formula_text

            if ext == ".xlsx":
                try:
                    from openpyxl import load_workbook
                    sheets_text = []
                    wb = load_workbook(file_path, read_only=True, data_only=True)
                    formula_wb = load_workbook(file_path, read_only=True, data_only=False)
                    _indexed_rows = 0
                    _truncated_table = False
                    try:
                        for ws, formula_ws in zip(wb.worksheets, formula_wb.worksheets):
                            _remaining = _table_budget_remaining(sheets_text)
                            if _remaining <= 0:
                                _truncated_table = True
                                break

                            def _iter_xlsx_rows() -> Any:
                                value_rows = ws.iter_rows(values_only=True)
                                formula_rows = formula_ws.iter_rows(values_only=True)
                                for value_row, formula_row in zip(value_rows, formula_rows):
                                    max_width = max(len(value_row or ()), len(formula_row or ()))
                                    values = list(value_row or ()) + [None] * max(0, max_width - len(value_row or ()))
                                    formulas = list(formula_row or ()) + [None] * max(0, max_width - len(formula_row or ()))
                                    yield [
                                        _xlsx_formula_value_to_text(values[i], formulas[i])
                                        for i in range(max_width)
                                    ]

                            part, row_count, was_truncated = _rows_to_markdown_table(
                                _iter_xlsx_rows(),
                                sheet_name=ws.title,
                                max_chars=_remaining,
                            )
                            _indexed_rows += row_count
                            if part:
                                sheets_text.append(part)
                            if was_truncated:
                                _truncated_table = True
                                break
                    finally:
                        try:
                            wb.close()
                        except Exception:
                            pass
                        try:
                            formula_wb.close()
                        except Exception:
                            pass
                    text = "\n\n".join(sheets_text)
                    if _truncated_table:
                        text = text[:_max_table_index_chars] + _table_truncation_note
                    if text.strip():
                        self._index_info(
                            f"XLSX 读取成功: {len(sheets_text)} sheet(s), "
                            f"indexed_rows={_indexed_rows}, {len(text)} chars"
                        )
                    else:
                        _tabular_parse_issue = "empty_workbook_or_no_readable_cells"
                        self._index_warning(f"XLSX 文件内容为空: {file_name}")
                except ImportError as _imp_err:
                    _tabular_parse_issue = f"missing_dependency: {_imp_err}"
                    self._index_warning(f"读取 xlsx 需要安装 openpyxl: {_imp_err}")
                    text = ""
                except Exception as e:
                    _tabular_parse_issue = str(e)
                    self._index_exception(f"XLSX 解析失败: {file_name}", e)
                    text = ""
            elif ext == ".xls":
                try:
                    import xlrd
                    sheets_text = []
                    book = xlrd.open_workbook(file_path, on_demand=True)
                    _indexed_rows = 0
                    _truncated_table = False

                    def _xlrd_cell_to_text(cell: Any) -> str:
                        try:
                            if cell.ctype == xlrd.XL_CELL_DATE:
                                try:
                                    dt = xlrd.xldate.xldate_as_datetime(cell.value, book.datemode)
                                    return dt.isoformat(sep=" ")
                                except Exception:
                                    pass
                            return _table_cell_to_text(cell.value)
                        except Exception:
                            return ""

                    try:
                        for sheet in book.sheets():
                            _remaining = _table_budget_remaining(sheets_text)
                            if _remaining <= 0:
                                _truncated_table = True
                                break
                            rows = (
                                [_xlrd_cell_to_text(sheet.cell(r, c)) for c in range(sheet.ncols)]
                                for r in range(sheet.nrows)
                            )
                            part, row_count, was_truncated = _rows_to_markdown_table(
                                rows,
                                sheet_name=sheet.name,
                                max_chars=_remaining,
                            )
                            _indexed_rows += row_count
                            if part:
                                sheets_text.append(part)
                            if was_truncated:
                                _truncated_table = True
                                break
                    finally:
                        try:
                            book.release_resources()
                        except Exception:
                            pass
                    text = "\n\n".join(sheets_text)
                    if _truncated_table:
                        text = text[:_max_table_index_chars] + _table_truncation_note
                    if text.strip():
                        self._index_info(
                            f"XLS 读取成功: {len(sheets_text)} sheet(s), "
                            f"indexed_rows={_indexed_rows}, {len(text)} chars"
                        )
                    else:
                        _tabular_parse_issue = "empty_workbook_or_no_readable_cells"
                        self._index_warning(f"XLS 文件内容为空: {file_name}")
                except ImportError as _imp_err:
                    _tabular_parse_issue = f"missing_dependency: {_imp_err}"
                    self._index_warning(f"读取 xls 需要安装 xlrd: {_imp_err}")
                    text = ""
                except Exception as e:
                    _xls_err = str(e)
                    _looks_like_xlsx_payload = False
                    try:
                        import zipfile as _zipfile
                        _looks_like_xlsx_payload = _zipfile.is_zipfile(file_path)
                    except Exception:
                        _looks_like_xlsx_payload = False
                    if ("xlsx" in _xls_err.lower() or _looks_like_xlsx_payload):
                        try:
                            from openpyxl import load_workbook
                            sheets_text = []
                            value_fh = open(file_path, "rb")
                            formula_fh = open(file_path, "rb")
                            wb = load_workbook(value_fh, read_only=True, data_only=True)
                            formula_wb = load_workbook(formula_fh, read_only=True, data_only=False)
                            _indexed_rows = 0
                            _truncated_table = False
                            try:
                                for ws, formula_ws in zip(wb.worksheets, formula_wb.worksheets):
                                    _remaining = _table_budget_remaining(sheets_text)
                                    if _remaining <= 0:
                                        _truncated_table = True
                                        break

                                    def _iter_mislabeled_xlsx_rows() -> Any:
                                        value_rows = ws.iter_rows(values_only=True)
                                        formula_rows = formula_ws.iter_rows(values_only=True)
                                        for value_row, formula_row in zip(value_rows, formula_rows):
                                            max_width = max(len(value_row or ()), len(formula_row or ()))
                                            values = list(value_row or ()) + [None] * max(0, max_width - len(value_row or ()))
                                            formulas = list(formula_row or ()) + [None] * max(0, max_width - len(formula_row or ()))
                                            yield [
                                                _xlsx_formula_value_to_text(values[i], formulas[i])
                                                for i in range(max_width)
                                            ]

                                    part, row_count, was_truncated = _rows_to_markdown_table(
                                        _iter_mislabeled_xlsx_rows(),
                                        sheet_name=ws.title,
                                        max_chars=_remaining,
                                    )
                                    _indexed_rows += row_count
                                    if part:
                                        sheets_text.append(part)
                                    if was_truncated:
                                        _truncated_table = True
                                        break
                            finally:
                                try:
                                    wb.close()
                                except Exception:
                                    pass
                                try:
                                    formula_wb.close()
                                except Exception:
                                    pass
                                try:
                                    value_fh.close()
                                except Exception:
                                    pass
                                try:
                                    formula_fh.close()
                                except Exception:
                                    pass
                            text = "\n\n".join(sheets_text)
                            if _truncated_table:
                                text = text[:_max_table_index_chars] + _table_truncation_note
                            if text.strip():
                                self._index_info(
                                    f"XLS 读取成功（实际为 xlsx 内容）: {len(sheets_text)} sheet(s), "
                                    f"indexed_rows={_indexed_rows}, {len(text)} chars"
                                )
                            else:
                                _tabular_parse_issue = "empty_mislabeled_xlsx_or_no_readable_cells"
                                self._index_warning(f"XLS 文件内容为空（实际为 xlsx 内容）: {file_name}")
                        except ImportError as _imp_err:
                            _tabular_parse_issue = f"missing_dependency: openpyxl fallback for mislabeled xls: {_imp_err}"
                            self._index_warning(f"读取错标 .xls/.xlsx 需要安装 openpyxl: {_imp_err}")
                            text = ""
                        except Exception as _fallback_e:
                            _tabular_parse_issue = f"{_xls_err}; xlsx_fallback_failed: {_fallback_e}"
                            self._index_exception(f"XLS 解析失败，xlsx 兜底也失败: {file_name}", _fallback_e)
                            text = ""
                    else:
                        _tabular_parse_issue = _xls_err
                        self._index_exception(f"XLS 解析失败: {file_name}", e)
                        text = ""
            elif ext in {".csv", ".tsv"}:
                try:
                    import csv as _csv
                    sep = "\t" if ext == ".tsv" else ","
                    _encodings = ["utf-8-sig", "utf-8", "gb18030", "latin-1"]
                    _read_ok = False
                    _indexed_rows = 0
                    _truncated_table = False
                    _read_err = None
                    try:
                        _csv.field_size_limit(max(_csv.field_size_limit(), 1024 * 1024 * 64))
                    except Exception:
                        pass
                    for _enc in _encodings:
                        try:
                            with open(file_path, "r", encoding=_enc, errors="strict", newline="") as _f:
                                reader = _csv.reader(_f, delimiter=sep)
                                text, _indexed_rows, _truncated_table = _rows_to_markdown_table(
                                    reader,
                                    max_chars=_max_table_index_chars,
                                )
                            if _truncated_table:
                                text = text[:_max_table_index_chars] + _table_truncation_note
                            self._index_info(
                                f"CSV/TSV 读取成功: encoding={_enc}, "
                                f"indexed_rows={_indexed_rows}, {len(text)} chars"
                            )
                            _read_ok = True
                            break
                        except Exception as _e:
                            _read_err = _e
                            text = ""
                            continue
                    if _read_ok and not text.strip():
                        _tabular_parse_issue = "empty_csv_or_tsv"
                        self._index_warning(f"CSV/TSV 文件为空: {file_name}")
                    elif not _read_ok:
                        _tabular_parse_issue = f"read_failed: {_read_err}"
                        self._index_warning(f"CSV/TSV 所有编码均读取失败: {file_name} | last_err={_read_err}")
                except Exception as e:
                    _tabular_parse_issue = str(e)
                    self._index_exception(f"CSV/TSV 解析失败: {file_name}", e)
                    text = ""
            elif ext == '.pdf':
                try:
                    text, page_count = extract_pdf_text(file_path)
                    if page_count:
                        metadata["page_count"] = int(page_count)
                except Exception as e:
                    self._index_exception("PDF 文本抽取失败", e)
                    text = ""
            elif ext in {".txt", ".md", ".rst", ".tex", ".org", ".yml", ".yaml"}:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
                try:
                    max_txt_chars = int(os.getenv("FILEAGENT_MAX_TXT_CHARS", "200000") or 200000)
                except Exception:
                    max_txt_chars = 200000
                max_txt_chars = max(10000, min(max_txt_chars, 5_000_000))
                if ext in {".txt", ".md", ".rst"} and len(text) > max_txt_chars:
                    original_chars = len(text)
                    text = text[:max_txt_chars] + "\n\n... (Text truncated for semantic indexing)"
                    self._index_info(
                        f"纯文本截断保护：{file_name} | original_chars={original_chars} "
                        f"kept_chars={max_txt_chars}"
                    )
            elif ext in {".html", ".htm"}:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = self._xmlish_bytes_to_text(f.read().encode("utf-8", errors="ignore"))
            elif ext == ".docx":
                try:
                    import docx2txt
                    text = docx2txt.process(file_path) or ""
                except Exception as e:
                    self._index_exception("docx 文本抽取失败", e)
                    text = self._extract_zip_text_members(file_path, ext=ext)
            elif ext == ".doc":
                text = self._extract_legacy_office_text(file_path, ext=ext)
            elif ext == ".pptx":
                try:
                    from pptx import Presentation as _Presentation
                    prs = _Presentation(file_path)
                    slide_texts = []
                    for slide_idx, slide in enumerate(prs.slides, 1):
                        parts = []
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    t = para.text.strip()
                                    if t:
                                        parts.append(t)
                            if hasattr(shape, "has_table") and shape.has_table:
                                try:
                                    for row in shape.table.rows:
                                        for cell in row.cells:
                                            t = cell.text.strip()
                                            if t:
                                                parts.append(t)
                                except Exception:
                                    pass
                        if parts:
                            slide_texts.append(f"[Slide {slide_idx}]\n" + "\n".join(parts))
                    text = "\n\n".join(slide_texts)
                    try:
                        metadata["page_count"] = int(len(prs.slides))
                    except Exception:
                        pass
                    if text.strip():
                        self._index_info(f"PPTX python-pptx extracted: {len(text)} chars, {len(prs.slides)} slides")
                    else:
                        self._index_warning("PPTX python-pptx: no text found (possibly all images)")
                except ImportError:
                    self._index_warning("python-pptx/lxml not installed, falling back to zip XML extraction")
                    text = self._extract_zip_text_members(file_path, ext=ext)
                except Exception as e:
                    self._index_exception("PPTX text extraction failed", e)
                    text = self._extract_zip_text_members(file_path, ext=ext)
            elif ext == ".ppt":
                text = self._extract_legacy_office_text(file_path, ext=ext)
            elif ext == ".numbers":
                try:
                    from numbers_parser import Document as _NumbersDoc
                    _ndoc = _NumbersDoc(file_path)
                    sheets_text = []
                    _indexed_rows = 0
                    _truncated_table = False

                    def _safe_cell_str(cell) -> str:
                        """Robustly extract a string value from a Numbers cell.
                        numbers_parser cells can be None, have None value, raise on
                        formatted_value for formula/merge cells, etc.
                        """
                        try:
                            if cell is None:
                                return ""
                            # Prefer formatted_value (keeps dates, currency, % as display strings)
                            try:
                                fv = cell.formatted_value
                                if fv is not None:
                                    s = str(fv).strip()
                                    if s and s not in ("None", "nan"):
                                        return s
                            except Exception:
                                pass
                            v = cell.value
                            if v is None:
                                return ""
                            s = str(v).strip()
                            return "" if s in ("None", "nan") else s
                        except Exception:
                            return ""

                    def _iter_numbers_rows(table: Any) -> Any:
                        try:
                            iterator = table.rows(values_only=False)
                        except TypeError:
                            iterator = table.rows()
                        for row in iterator:
                            yield [_safe_cell_str(c) for c in row]

                    for sheet in _ndoc.sheets:
                        for table in sheet.tables:
                            _remaining = _table_budget_remaining(sheets_text)
                            if _remaining <= 0:
                                _truncated_table = True
                                break
                            try:
                                table_name = str(getattr(table, "name", "") or "").strip()
                                sheet_label = f"{sheet.name}/{table_name}" if table_name else str(sheet.name)
                                part, row_count, was_truncated = _rows_to_markdown_table(
                                    _iter_numbers_rows(table),
                                    sheet_name=sheet_label,
                                    max_chars=_remaining,
                                )
                                _indexed_rows += row_count
                                if part:
                                    sheets_text.append(part)
                                if was_truncated:
                                    _truncated_table = True
                                    break
                            except Exception as _re:
                                self._index_warning(f"Numbers table read failed ({sheet.name}/{table.name}): {_re}")
                        if _truncated_table:
                            break

                    text = "\n\n".join(sheets_text)
                    if _truncated_table:
                        text = text[:_max_table_index_chars] + _table_truncation_note
                    if text.strip():
                        self._index_info(f"Numbers 读取成功: indexed_rows={_indexed_rows}, {len(text)} chars")
                    else:
                        _tabular_parse_issue = "empty_numbers_or_no_readable_cells"
                        self._index_warning(
                            f"Numbers 文件内容为空: {file_name} "
                            f"(sheets={len(list(_ndoc.sheets))}) — 可能为纯图表文件或加密文件"
                        )
                except ImportError as _imp_err:
                    _tabular_parse_issue = f"missing_dependency: {_imp_err}"
                    self._index_warning(f"解析 .numbers 文件需要安装 numbers-parser: {_imp_err}")
                    text = ""
                except Exception as e:
                    _tabular_parse_issue = str(e)
                    self._index_exception(f"Numbers 解析失败: {file_name}", e)
                    text = self._extract_zip_text_members(file_path, ext=ext)
            elif ext in {".odt", ".ods", ".odp", ".epub", ".pages", ".key"}:
                text = self._extract_zip_text_members(file_path, ext=ext)
            elif ext == ".mobi":
                text = self._extract_binary_strings_text(file_path)

            if ext in {".csv", ".tsv", ".xlsx", ".xls", ".numbers"} and text.strip():
                table_schema_hint = self._build_tabular_index_hint(text, file_name=file_name, ext=ext)
                if table_schema_hint:
                    metadata["table_schema_hint"] = table_schema_hint
            elif ext in {".json", ".jsonl"}:
                try:
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as _f:
                        raw = _f.read()
                    # Pretty-print for better readability in index
                    import json as _json
                    try:
                        parsed = _json.loads(raw)
                        text = _json.dumps(parsed, ensure_ascii=False, indent=2)
                    except Exception:
                        text = raw  # fallback to raw if not valid JSON
                    max_table_chars = int(os.getenv("MAX_TABLE_INDEX_CHARS", "100000"))
                    if len(text) > max_table_chars:
                        text = text[:max_table_chars] + "\n\n... (Data truncated for semantic indexing)"
                        self._index_info(f"JSON 文件截断：保留前 {max_table_chars} 字符")
                except Exception as e:
                    self._index_exception(f"JSON 解析失败: {file_name}", e)
                    text = ""
            elif ext == ".xml":
                try:
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as _f:
                        text = _f.read()
                    max_table_chars = int(os.getenv("MAX_TABLE_INDEX_CHARS", "100000"))
                    if len(text) > max_table_chars:
                        text = text[:max_table_chars] + "\n\n... (Data truncated for semantic indexing)"
                except Exception as e:
                    self._index_exception(f"XML 读取失败: {file_name}", e)
                    text = ""
            elif ext == ".sql":
                try:
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as _f:
                        text = _f.read()
                    max_table_chars = int(os.getenv("MAX_TABLE_INDEX_CHARS", "100000"))
                    if len(text) > max_table_chars:
                        text = text[:max_table_chars] + "\n\n... (Data truncated for semantic indexing)"
                except Exception as e:
                    self._index_exception(f"SQL 读取失败: {file_name}", e)
                    text = ""
            elif HAS_LLAMA_INDEX and not (text or "").strip():
                try:
                    reader = SimpleDirectoryReader(input_files=[file_path])
                    docs = reader.load_data()
                    text = "\n".join([doc.text for doc in docs])
                except Exception as e:
                    self._index_exception("LlamaIndex 解析失败", e)
                    text = ""
            elif not (text or "").strip():
                text = ""

            ocr_chunks: List[str] = []
            MIN_TEXT_LEN = int(os.getenv("MIN_TEXT_LEN_FOR_TEXT_MODE", "200"))
            text_chars = len((text or "").strip())
            self._index_info(f"文本抽取完成：{file_name} | text_chars={text_chars} | min_text_len={MIN_TEXT_LEN}")
            if ext in _tabular_exts and text_chars == 0:
                reason = _tabular_parse_issue or "no_readable_table_content"
                self._index_error(
                    f"表格解析无可用内容，不写入占位索引: {file_name} | reason={reason[:300]}"
                )
                self._log_index_file_timing(
                    file_path=file_path,
                    status="failed",
                    stage="tabular_parse_failed",
                    total_s=time.time() - _start_time,
                    embed_s=_embed_time,
                    llm_s=_llm_time,
                    vl_s=_vl_time,
                    chunks=0,
                    doc_category=metadata.get("doc_category", _doc_category),
                    extra=f"reason={reason[:300]}",
                )
                return False
            llm_ready = self._test_local_llm_connection() if use_smart_indexing else False
            if text_chars < MIN_TEXT_LEN and use_smart_indexing and llm_ready:
                if ext == ".pdf":
                    max_pages = int(os.getenv("PDF_OCR_MAX_PAGES", "3"))
                    dpi = int(os.getenv("PDF_OCR_DPI", "72"))
                    self._index_info(f"PDF 文字不足，启用 OCR：pages={max_pages}, dpi={dpi}")
                    _vl_start = time.time()
                    ocr_chunks = self._pdf_pages_to_ocr_chunks(file_path, max_pages=max_pages, dpi=dpi)
                    _vl_time += time.time() - _vl_start
                elif ext == ".docx":
                    max_imgs = int(os.getenv("DOCX_OCR_MAX_IMAGES", "3"))
                    self._index_info(f"DOCX 文字不足，启用图片 OCR：images={max_imgs}")
                    _vl_start = time.time()
                    ocr_chunks = self._docx_images_to_ocr_chunks(file_path, max_images=max_imgs)
                    _vl_time += time.time() - _vl_start
                elif ext in (".pptx", ".ppt"):
                    max_imgs = int(os.getenv("PPTX_OCR_MAX_IMAGES", "3"))
                    self._index_info(f"PPTX 文字不足，启用图片 OCR：images={max_imgs}")
                    _vl_start = time.time()
                    ocr_chunks = self._pptx_images_to_ocr_chunks(file_path, max_images=max_imgs)
                    _vl_time += time.time() - _vl_start
                self._index_info(f"OCR 结果：{file_name} | ocr_chunks={len(ocr_chunks)}")
            elif text_chars < MIN_TEXT_LEN and use_smart_indexing and not llm_ready:
                self._index_warning(f"文本不足但 LLM 不可用，跳过 OCR：{file_name}")

            if not (text or "").strip() and not ocr_chunks:
                placeholder = f"{ext.upper().lstrip('.')} file {file_name} (text extraction unavailable; possibly scanned/protected)"
                metadata["doc_summary"] = metadata.get("doc_summary") or "Text extraction unavailable (possibly scanned/protected)."
                self._set_doc_summary_provenance(metadata, "none")
                metadata["lookup_aliases"] = self._build_lookup_aliases(file_path, metadata)
                _embed_start = time.time()
                with self._embed_context(f"file={file_path} phase=placeholder_embedding"):
                    emb = self._embed_doc_text(placeholder)
                if not any(emb):
                    msg = f"[Embedding] zero vector | file={file_path} phase=placeholder_embedding"
                    self._index_error(msg)
                    self._append_embedding_error_log(msg)
                doc_id = f"{file_path}_{uuid.uuid4().hex[:8]}"
                self.collection.add(ids=[doc_id], embeddings=[emb], documents=[placeholder], metadatas=[metadata])
                self._invalidate_meta_cache()
                _embed_time = time.time() - _embed_start
                self._maybe_persist()
                self._upsert_folder_index_for_file(file_path)
                _doc_category = metadata.get("doc_category", _doc_category)
                self._log_index_file_timing(
                    file_path=file_path,
                    status="ok",
                    stage="placeholder_embedding",
                    total_s=time.time() - _start_time,
                    embed_s=_embed_time,
                    llm_s=_llm_time,
                    vl_s=_vl_time,
                    chunks=1,
                    doc_category=_doc_category,
                )
                return True
            
            if use_smart_indexing and self._test_local_llm_connection():
                _llm_start = time.time()
                classify_text = (text or "").strip()
                if (not classify_text) and ocr_chunks:
                    classify_text = "\n\n".join(ocr_chunks)[:4000]
                self._index_info(
                    f"智能索引（unified）：classify+summary「{file_name}」（约 {len(classify_text)} 字）"
                )
                unified = self._classify_and_summarize_unified(
                    classify_text,
                    file_name=file_name,
                    file_ext=ext,
                    page_count=metadata.get("page_count"),
                )
                taxonomy = unified.get("taxonomy") or {}
                raw_doc_category = taxonomy.get("leaf_category", "other")
                stable_doc_family = taxonomy.get("family", "other")
                metadata["doc_category_raw"] = raw_doc_category
                metadata["doc_category_family"] = stable_doc_family
                metadata["doc_category_leaf"] = raw_doc_category
                metadata["doc_role"] = taxonomy.get("role", "other")
                metadata["doc_taxonomy_confidence"] = taxonomy.get("confidence", 0.0)
                metadata["doc_category"] = stable_doc_family
                _sum_text = str(unified.get("summary") or "")
                _extracts = unified.get("extracts") if isinstance(unified.get("extracts"), list) else []
                _sum_model = str(unified.get("model") or "unknown")
                self._index_info(
                    f"智能索引（unified）：完成 family={stable_doc_family} "
                    f"leaf={raw_doc_category} role={metadata['doc_role']}"
                )
                metadata["doc_summary"] = _sum_text
                metadata["file_name_en"] = self._resolve_file_name_en_for_metadata(
                    file_name,
                    existing_translation=metadata.get("file_name_en", ""),
                )
                self._set_doc_summary_provenance(metadata, _sum_model)
                if _extracts:
                    self._store_personal_info(file_path, _extracts)
                _llm_time = time.time() - _llm_start
            else:
                metadata["file_name_en"] = self._resolve_file_name_en_for_metadata(
                    file_name,
                    existing_translation=metadata.get("file_name_en", ""),
                )
            _doc_category = metadata.get("doc_category", _doc_category)

            # Build English tags for every indexed document - this enables
            # English queries to match via the lexical fallback even for
            # Chinese-named files. Done after summary generation so Level 4
            # (summary-derived keywords) can contribute.
            _doc_en_tags = self._build_en_tags(
                file_name,
                ext,
                doc_summary=" ".join(
                    x for x in [
                        metadata.get("doc_summary", ""),
                        metadata.get("table_schema_hint", ""),
                    ] if x
                ),
                parent_folder=parent_folder,
                folder_name_en=folder_name_en,
            )
            if _doc_en_tags:
                metadata["en_tags"] = _doc_en_tags
            if ext in {".csv", ".tsv", ".xlsx", ".xls", ".numbers"}:
                table_schema_hint = str(metadata.get("table_schema_hint") or "").strip()
                if table_schema_hint:
                    if metadata.get("doc_summary"):
                        if table_schema_hint.lower() not in str(metadata.get("doc_summary") or "").lower():
                            metadata["doc_summary"] = f"{table_schema_hint}. {metadata.get('doc_summary', '')}"[:700]
                    else:
                        metadata["doc_summary"] = table_schema_hint
                        self._set_doc_summary_provenance(metadata, "table_schema_hint")
                if not str(metadata.get("doc_summary") or "").strip():
                    metadata["doc_summary"] = self._default_tabular_summary(file_name, ext)
                    self._set_doc_summary_provenance(metadata, "tabular_fallback")
            metadata["lookup_aliases"] = self._build_lookup_aliases(file_path, metadata)
            # ─────────────────────────────────────────────────────────────────
            
            merged_text = (text or "").strip()
            if ocr_chunks:
                merged_text = (merged_text + "\n\n" if merged_text else "") + "\n\n".join(ocr_chunks)
            if self._is_structured_table_ext(ext):
                before_structured_chars = len(merged_text)
                merged_text = self._compact_structured_text_for_embedding(
                    merged_text,
                    metadata,
                    file_name=file_name,
                    ext=ext,
                )
                if len(merged_text) < before_structured_chars:
                    self._index_info(
                        f"结构化文件向量文本压缩：{file_name} | "
                        f"{before_structured_chars} -> {len(merged_text)} chars"
                    )

            chunks = self._split_merged_text_for_index_vectors(merged_text, ext=ext)
            if not chunks:
                self._index_warning(f"索引失败：切分后无可用文本块 {file_name}")
                return False

            ds = metadata.get("doc_summary", "").strip()
            et = metadata.get("en_tags", "").strip()
            aliases = str(metadata.get("lookup_aliases", "") or "").strip()
            embedding_context = self._build_compact_embedding_context(file_path, metadata, ext=ext)
            embedding_prefix = f"{embedding_context}\n---\n" if embedding_context else ""
            if embedding_context:
                self._index_info(
                    f"文本向量上下文：{file_name} | context_chars={len(embedding_context)} "
                    f"| summary_chars={len(ds)} | tags_chars={len(et)} | aliases_chars={len(aliases)}"
                )

            expanded_chunks: List[str] = []
            for chunk_text in chunks:
                expanded_chunks.extend(
                    self._split_text_for_embedding_quality(
                        chunk_text,
                        reserve_text=embedding_prefix,
                        fallback_max_chars=self._index_chunk_max_chars(ext),
                    )
                )
            chunks = expanded_chunks or chunks

            log_label = self._short_name_for_log(file_name)
            n_chunks = len(chunks)
            _chunks = n_chunks
            if n_chunks > 1:
                self._index_info(
                    f"向量化：「{log_label}」共 {n_chunks} 个文本块（每块约≤{self._index_chunk_max_chars()}字），"
                    f"批量 Embedding 后写入 Chroma"
                )

            try:
                batch_sz = int(os.getenv("FILEAGENT_CHROMA_ADD_BATCH", "48"))
            except ValueError:
                batch_sz = 48
            batch_sz = max(8, min(batch_sz, 256))
            self._index_info(f"向量化参数：{file_name} | chunks={n_chunks} | chroma_add_batch={batch_sz}")

            nb = (n_chunks + batch_sz - 1) // batch_sz
            _embed_start = time.time()
            for bi, start in enumerate(range(0, n_chunks, batch_sz)):
                if should_cancel and should_cancel():
                    self._index_info(f"索引被取消：{file_path} (已处理 {start}/{n_chunks} 块)")
                    return False
                end = min(start + batch_sz, n_chunks)
                if n_chunks > 4 and (bi == 0 or bi == nb - 1 or (nb > 6 and bi == nb // 2)):
                    self._index_info(f"向量化进度：{log_label} [{end}/{n_chunks}]")
                sub = chunks[start : start + batch_sz]
                
                # Enrich each chunk with a compact anchor so the vector text keeps
                # enough context without exploding into tiny token windows.
                enriched_sub = []
                for chunk_text in sub:
                    enriched_sub.append(embedding_prefix + chunk_text if embedding_prefix else chunk_text)
                
                with self._embed_context(f"file={file_path} phase=chunk_embedding"):
                    embeddings = self._embed_texts_for_index(enriched_sub)
                zero_count = sum(1 for v in embeddings if not any(v))
                if zero_count > 0:
                    msg = (
                        f"[Embedding] zero vectors | file={file_path} phase=chunk_embedding "
                        f"batch_size={len(embeddings)} zero_count={zero_count}"
                    )
                    self._index_error(msg)
                    self._append_embedding_error_log(msg)
                ids = [f"{file_path}_{start + j}_{uuid.uuid4().hex[:8]}" for j in range(len(sub))]
                metadatas = [dict(metadata) for _ in sub]
                self.collection.add(
                    ids=ids,
                    embeddings=embeddings,
                    documents=enriched_sub,
                    metadatas=metadatas,
                )
                self._invalidate_meta_cache()
            _embed_time = time.time() - _embed_start
            self._upsert_filename_lookup_for_file(file_path, metadata, persist=False)
            self._maybe_persist()
            self._upsert_folder_index_for_file(file_path)
            self._log_index_file_timing(
                file_path=file_path,
                status="ok",
                stage="text_chunk_embedding",
                total_s=time.time() - _start_time,
                embed_s=_embed_time,
                llm_s=_llm_time,
                vl_s=_vl_time,
                chunks=_chunks,
                doc_category=_doc_category,
            )
            return True
            
        except Exception as e:
            self._index_exception(f"索引异常：{file_path}", e)
            self._log_index_file_timing(
                file_path=file_path,
                status="failed",
                stage="ingest_exception",
                total_s=time.time() - _start_time,
                embed_s=_embed_time,
                llm_s=_llm_time,
                vl_s=_vl_time,
                chunks=_chunks,
                doc_category=_doc_category,
                extra=f"error={e}",
            )
            try:
                deleted = self.delete_file(file_path, silent=True)
                self._index_info(f"索引失败后清理已执行：{file_path} | deleted={deleted}")
            except Exception as cleanup_e:
                self._index_exception(f"索引失败后清理再次失败：{file_path}", cleanup_e)
            return False

    def refresh_source(
        self,
        directories: List[str],
        progress_callback: 'Callable[[int, int, str, str], None] | None' = None,
        should_cancel: 'Callable[[], bool] | None' = None,
    ) -> Dict[str, Any]:
        """
        Wrapper that prepares settings.INCLUDE_PATHS temporarily just like scan_directory
        before running the actual increment refresh engine.
        """
        import config.settings as settings
        _orig_include_paths = None
        try:
            if getattr(settings, "USE_WHITELIST_MODE", False) and isinstance(getattr(settings, "INCLUDE_PATHS", None), set):
                _orig_include_paths = settings.INCLUDE_PATHS
                explicit = []
                for d in (directories or []):
                    try:
                        explicit.append(os.path.abspath(os.path.expanduser(str(d))))
                    except Exception:
                        continue
                settings.INCLUDE_PATHS = set(_orig_include_paths) | set(explicit)
        except Exception:
            _orig_include_paths = None
            
        try:
            return self._do_refresh_source(directories, progress_callback, should_cancel)
        finally:
            try:
                if getattr(settings, 'USE_WHITELIST_MODE', False) and _orig_include_paths is not None:
                    settings.INCLUDE_PATHS = _orig_include_paths
            except Exception:
                pass
    
    def _do_refresh_source(
        self,
        directories: List[str],
        progress_callback: 'Callable[[int, int, str, str], None] | None' = None,
        should_cancel: 'Callable[[], bool] | None' = None,
    ) -> Dict[str, Any]:
        from datetime import datetime

        dirs = [os.path.abspath(os.path.expanduser(d)) for d in (directories or [])]
        logger.info(f"[Refresh] 开始增量刷新: directories={dirs}")

        result = {"added": 0, "updated": 0, "deleted": 0, "skipped": 0, "errors": 0, "total_scanned": 0}

        logger.info("[Refresh] Phase 1 — 扫描数据库已有索引...")
        indexed: Dict[str, str] = {}  # file_path → modified_time string
        try:
            total_count = self.collection.count()
            page_size = 5000
            for offset in range(0, total_count, page_size):
                batch = self.collection.get(include=["metadatas"], limit=page_size, offset=offset)
                for meta in (batch.get("metadatas") or []):
                    fp = meta.get("file_path", "")
                    mt = meta.get("modified_time", "")
                    if fp and mt:
                        for d in dirs:
                            if fp == d or fp.startswith(d.rstrip(os.sep) + os.sep):
                                indexed[fp] = mt
                                break
        except Exception as e:
            logger.error(f"[Refresh] Phase 1 读取数据库失败: {e}")
            result["errors"] += 1

        logger.info(f"[Refresh] 数据库内该目录下共 {len(indexed)} 个已索引文件")

        import config.settings as settings
        disk_files: Dict[str, float] = {}   # file_path → disk mtime (float)
        for directory in dirs:
            if not os.path.exists(directory):
                logger.warning(f"[Refresh] 路径不存在，跳过: {directory}")
                continue
            
            if os.path.isfile(directory):
                if not self._should_ignore_file(directory):
                    try:
                        disk_files[directory] = os.stat(directory).st_mtime
                    except OSError:
                        pass
                continue

            for root, subdirs, files in os.walk(directory):
                subdirs[:] = [
                    d for d in subdirs
                    if d not in settings.IGNORE_PATTERNS and not d.startswith(".") and not d.endswith(".app")
                ]
                rel_path = os.path.relpath(root, settings.HOME_DIR)
                top_dir = rel_path.split(os.sep)[0] if os.sep in rel_path else rel_path
                if top_dir in getattr(settings, "IGNORE_TOP_LEVEL_DIRS", set()):
                    continue
                for fname in files:
                    fp = os.path.abspath(os.path.join(root, fname))
                    if self._should_ignore_file(fp):
                        continue
                    try:
                        disk_files[fp] = os.stat(fp).st_mtime
                    except OSError:
                        pass

        logger.info(f"[Refresh] 磁盘扫描到 {len(disk_files)} 个可索引文件")

        deleted_files: List[str] = []
        modified_files: List[str] = []
        new_files: List[str] = []
        unchanged_count = 0

        for fp in indexed:
            if fp not in disk_files:
                deleted_files.append(fp)

        for fp, disk_mtime in disk_files.items():
            if fp not in indexed:
                if self._failed_file_cache_try_skip(fp):
                    unchanged_count += 1
                else:
                    new_files.append(fp)
            else:
                db_mtime_str = indexed[fp]
                disk_mtime_str = datetime.fromtimestamp(disk_mtime).strftime("%Y-%m-%d %H:%M:%S")
                if disk_mtime_str != db_mtime_str:
                    modified_files.append(fp)
                else:
                    unchanged_count += 1

        result["total_scanned"] = len(disk_files)
        result["skipped"] = unchanged_count
        logger.info(
            f"[Refresh] 比对结果 — 删除:{len(deleted_files)} 修改:{len(modified_files)} "
            f"新增:{len(new_files)} 未变:{unchanged_count}"
        )

        to_process = deleted_files + modified_files + new_files
        total_ops = len(to_process)

        if total_ops == 0:
            logger.info("[Refresh] 无需更新，数据库已是最新")
            return result

        logger.info(f"[Refresh] Phase 2 — 开始执行更新，共 {total_ops} 个文件操作")

        op_idx = 0

        for fp in deleted_files:
            if should_cancel and should_cancel():
                logger.info("[Refresh] 收到取消信号，中止刷新")
                break
            op_idx += 1
            fname = os.path.basename(fp)
            if progress_callback:
                try:
                    progress_callback(op_idx, total_ops, fname, fp)
                except Exception:
                    pass
            try:
                ok = self.delete_file(fp, silent=False)
                if hasattr(self, '_metadata_index') and self._metadata_index:
                    try:
                        self._metadata_index.delete_metadata(fp)
                    except Exception:
                        pass
                if ok:
                    result["deleted"] += 1
                    logger.info(f"[Refresh] ✅ 已删除消失的文件: {fname}")
                else:
                    logger.warning(f"[Refresh] ⚠️ 删除失败（可能已不在库中）: {fp}")
            except Exception as e:
                result["errors"] += 1
                logger.error(f"[Refresh] ❌ 删除出错: {fp} | {e}")

        for fp in modified_files:
            if should_cancel and should_cancel():
                logger.info("[Refresh] 收到取消信号，中止刷新")
                break
            op_idx += 1
            fname = os.path.basename(fp)
            if progress_callback:
                try:
                    progress_callback(op_idx, total_ops, fname, fp)
                except Exception:
                    pass
            try:
                self.delete_file(fp, silent=True)
                if hasattr(self, '_metadata_index') and self._metadata_index:
                    try:
                        self._metadata_index.delete_metadata(fp)
                    except Exception:
                        pass
                ok = self.ingest_file(fp, should_cancel=should_cancel)
                if ok:
                    result["updated"] += 1
                    logger.info(f"[Refresh] ✅ 已更新: {fname}")
                else:
                    if should_cancel and should_cancel():
                        logger.info("[Refresh] 收到取消信号，中止刷新")
                        break
                    self._failed_file_cache_record(fp, "refresh 更新失败")
                    result["errors"] += 1
                    logger.warning(f"[Refresh] ⚠️ 更新失败: {fp}")
            except Exception as e:
                result["errors"] += 1
                logger.error(f"[Refresh] ❌ 更新出错: {fp} | {e}")

        for fp in new_files:
            if should_cancel and should_cancel():
                logger.info("[Refresh] 收到取消信号，中止刷新")
                break
            op_idx += 1
            fname = os.path.basename(fp)
            if progress_callback:
                try:
                    progress_callback(op_idx, total_ops, fname, fp)
                except Exception:
                    pass
            try:
                ok = self.ingest_file(fp, should_cancel=should_cancel)
                if ok:
                    result["added"] += 1
                    logger.info(f"[Refresh] ✅ 已新增: {fname}")
                else:
                    if should_cancel and should_cancel():
                        logger.info("[Refresh] 收到取消信号，中止刷新")
                        break
                    self._failed_file_cache_record(fp, "refresh 新增失败")
                    result["errors"] += 1
                    logger.warning(f"[Refresh] ⚠️ 新增失败: {fp}")
            except Exception as e:
                result["errors"] += 1
                logger.error(f"[Refresh] ❌ 新增出错: {fp} | {e}")

        self._maybe_persist()
        self._save_failed_file_cache()
        logger.info(
            f"[Refresh] 完成: added={result['added']} updated={result['updated']} "
            f"deleted={result['deleted']} errors={result['errors']} skipped={result['skipped']}"
        )
        return result

    def delete_file(self, file_path: str, silent: bool = False) -> bool:

        try:
            self._failed_file_cache_forget(os.path.abspath(file_path))
            self._save_failed_file_cache()
            results = self.collection.get(
                where={"file_path": file_path},
                include=["metadatas"]
            )
            
            if not results['ids']:
                return False
            
            doc_ids = results['ids']
            self.collection.delete(ids=doc_ids)
            self._invalidate_meta_cache()
            
            if hasattr(self, 'personal_info_db'):
                self.personal_info_db.delete_by_file(file_path)
            
            if not silent:
                logger.info(f"  ✅ 已删除索引: {os.path.basename(file_path)} ({len(doc_ids)} 条记录)")
            return True
            
        except Exception as e:
            if not silent:
                logger.error(f"删除失败 {os.path.basename(file_path)}: {e}")
            return False
    
    def delete_by_folder(self, folder_path: str) -> Dict[str, Any]:
        try:
            folder_abs = os.path.abspath(os.path.expanduser(folder_path))
            logger.info(f"开始删除索引: {folder_abs}")
            self._failed_file_cache_entries = {
                p: meta
                for p, meta in self._failed_file_cache_entries.items()
                if not (os.path.abspath(p) == folder_abs or os.path.abspath(p).startswith(folder_abs.rstrip(os.sep) + os.sep))
            }
            self._failed_file_cache_dirty = True
            
            total = self.collection.count()
            if total == 0:
                logger.info(f"数据库为空，无需删除")
                return {"ok": True, "deleted_count": 0, "folder": folder_abs}

            ids_to_delete = []
            page_size = 5000
            for offset in range(0, total, page_size):
                batch = self.collection.get(
                    include=["metadatas"],
                    limit=page_size,
                    offset=offset,
                )
                for idx, metadata in enumerate(batch.get("metadatas", [])):
                    file_path = metadata.get("file_path", "")
                    if file_path:
                        file_abs = os.path.abspath(file_path)
                        if file_abs == folder_abs or file_abs.startswith(folder_abs.rstrip(os.sep) + os.sep):
                            ids_to_delete.append(batch["ids"][idx])
            
            deleted_count = len(ids_to_delete)
            
            if deleted_count > 0:
                batch_size = 5000
                for i in range(0, deleted_count, batch_size):
                    batch_ids = ids_to_delete[i:i + batch_size]
                    self.collection.delete(ids=batch_ids)
                    logger.info(f"已删除 {i + len(batch_ids)}/{deleted_count} 条文档")
                self._invalidate_meta_cache()
                
                logger.info(f"✅ 成功删除索引: {folder_abs}，共 {deleted_count} 条文档")
                if hasattr(self, 'personal_info_db'):
                    self.personal_info_db.delete_by_folder(folder_abs)
                self._maybe_persist()
            else:
                logger.info(f"未找到相关索引文档: {folder_abs}")
            
            return {
                "ok": True,
                "deleted_count": deleted_count,
                "folder": folder_abs
            }
            
        except Exception as e:
            error_msg = f"删除索引失败: {e}"
            logger.error(f"❌ {error_msg}")
            import traceback
            traceback.print_exc()
            return {
                "ok": False,
                "error": error_msg,
                "deleted_count": 0,
                "folder": folder_path
            }
    
    def delete_by_folders(self, folder_paths: List[str]) -> Dict[str, Any]:
        try:
            folders_abs = [os.path.abspath(os.path.expanduser(p)) for p in folder_paths]
            prefixes = [p.rstrip(os.sep) + os.sep for p in folders_abs]
            folders_set = set(folders_abs)
            self._failed_file_cache_entries = {
                p: meta
                for p, meta in self._failed_file_cache_entries.items()
                if not (os.path.abspath(p) in folders_set or any(os.path.abspath(p).startswith(pfx) for pfx in prefixes))
            }
            self._failed_file_cache_dirty = True

            total = self.collection.count()
            if total == 0:
                self._failed_file_cache_entries = {}
                self._failed_file_cache_dirty = True
                self._save_failed_file_cache()
                return {"ok": True, "deleted_count": 0}

            ids_to_delete = []
            page_size = 5000
            for offset in range(0, total, page_size):
                batch = self.collection.get(
                    include=["metadatas"],
                    limit=page_size,
                    offset=offset,
                )
                for idx, metadata in enumerate(batch.get("metadatas", [])):
                    file_path = metadata.get("file_path", "")
                    if not file_path:
                        continue
                    file_abs = os.path.abspath(file_path)
                    if file_abs in folders_set or any(file_abs.startswith(pfx) for pfx in prefixes):
                        ids_to_delete.append(batch["ids"][idx])

            deleted_count = len(ids_to_delete)
            if deleted_count > 0:
                batch_size = 5000
                for i in range(0, deleted_count, batch_size):
                    self.collection.delete(ids=ids_to_delete[i:i + batch_size])
                    logger.info(f"[batch delete] {min(i + batch_size, deleted_count)}/{deleted_count} 条")
                self._invalidate_meta_cache()
                self._maybe_persist()
                if hasattr(self, 'personal_info_db'):
                    for folder in folders_abs:
                        self.personal_info_db.delete_by_folder(folder)

            logger.info(f"✅ 批量删除完成，{len(folder_paths)} 个路径，共 {deleted_count} 条文档")
            return {"ok": True, "deleted_count": deleted_count}

        except Exception as e:
            error_msg = f"批量删除索引失败: {e}"
            logger.error(f"❌ {error_msg}")
            import traceback
            traceback.print_exc()
            return {"ok": False, "error": error_msg, "deleted_count": 0}

    def clear_all(self) -> Dict[str, Any]:
        try:
            total = self.collection.count()
            if total == 0:
                return {"ok": True, "deleted_count": 0}
            batch_size = 5000
            deleted = 0
            for offset in range(0, total, batch_size):
                batch = self.collection.get(limit=batch_size, offset=0)
                ids = batch.get("ids", [])
                if not ids:
                    break
                self.collection.delete(ids=ids)
                self._invalidate_meta_cache()
                deleted += len(ids)
                logger.info(f"已清空 {deleted}/{total} 条索引")
            logger.info(f"✅ 数据库已全部清空，共删除 {deleted} 条索引")
            if hasattr(self, 'personal_info_db'):
                self.personal_info_db.clear_all()
            self._failed_file_cache_entries = {}
            self._failed_file_cache_dirty = True
            self._save_failed_file_cache()
            return {"ok": True, "deleted_count": deleted}
        except Exception as e:
            import traceback
            traceback.print_exc()
            return {"ok": False, "error": str(e), "deleted_count": 0}

    def index_file(
        self,
        file_path: str,
        should_cancel: Optional[Callable[[], bool]] = None,
        on_frame_progress: Optional[Callable[[int, int], None]] = None,
        on_media_progress: Optional[Callable[[Dict[str, Any]], None]] = None,
    ) -> bool:
        file_path = os.path.abspath(os.path.expanduser(file_path))
        
        if not os.path.exists(file_path):
            self._index_info(f"文件不存在: {file_path}")
            return False
        
        if not os.path.isfile(file_path):
            self._index_info(f"路径不是文件: {file_path}")
            return False
        
        indexed_paths = self.get_indexed_file_paths()
        if file_path in indexed_paths:
            self._index_info(f"文件已索引，跳过: {file_path}")
            return True
        
        import config.settings as settings
        
        _orig_include_paths = None
        try:
            if getattr(settings, "USE_WHITELIST_MODE", False) and isinstance(getattr(settings, "INCLUDE_PATHS", None), set):
                _orig_include_paths = settings.INCLUDE_PATHS
                settings.INCLUDE_PATHS = set(_orig_include_paths) | {file_path}
        except Exception:
            _orig_include_paths = None
            
        try:
            return self.ingest_file(
                file_path,
                should_cancel=should_cancel,
                on_frame_progress=on_frame_progress,
                on_media_progress=on_media_progress,
            )
        finally:
            if _orig_include_paths is not None:
                settings.INCLUDE_PATHS = _orig_include_paths
    
    def scan_directory(
        self, 
        directories: List[str] = None,
        progress_callback: Callable[[int, int, str], None] = None,
        should_cancel: Optional[Callable[[], bool]] = None,
        should_skip_file: Optional[Callable[[str], bool]] = None,
        on_file_indexed: Optional[Callable[[str], None]] = None,
        on_frame_progress: Optional[Callable[[int, int], None]] = None,
        on_media_progress: Optional[Callable[[Dict[str, Any]], None]] = None,
    ) -> int:
        if directories is None:
            directories = [
                os.path.join(settings.HOME_DIR, "Documents"),
                os.path.join(settings.HOME_DIR, "Downloads"),
                os.path.join(settings.HOME_DIR, "Desktop"),
            ]
        import config.settings as settings
        self._index_info(
            f"开始目录扫描：directories={directories}, whitelist_mode={getattr(settings, 'USE_WHITELIST_MODE', False)}"
        )

        _orig_include_paths = None
        try:
            if getattr(settings, "USE_WHITELIST_MODE", False) and isinstance(getattr(settings, "INCLUDE_PATHS", None), set):
                _orig_include_paths = settings.INCLUDE_PATHS
                explicit = []
                for d in (directories or []):
                    try:
                        explicit.append(os.path.abspath(os.path.expanduser(str(d))))
                    except Exception:
                        continue
                settings.INCLUDE_PATHS = set(_orig_include_paths) | set(explicit)
        except Exception:
            _orig_include_paths = None

        scan_ignore_fp = self._compute_ignore_rules_fingerprint()


        try:
            indexed_paths = self.get_indexed_file_paths()
            self._index_info(f"已有 {len(indexed_paths)} 个已索引文件")

            all_files = []
            failed_files = []
            skipped_missing_dirs = 0
            skipped_by_ignore_rules = 0
            skipped_by_cache = 0
            skipped_by_failed_cache = 0
            for directory in directories:
                if not os.path.exists(directory):
                    skipped_missing_dirs += 1
                    self._index_warning(f"扫描跳过：目录不存在 {directory}")
                    continue

                for root, dirs, files in os.walk(directory):
                    dirs[:] = [
                        d
                        for d in dirs
                        if d not in settings.IGNORE_PATTERNS and not d.startswith(".") and not d.endswith(".app")
                    ]

                    rel_path = os.path.relpath(root, settings.HOME_DIR)
                    top_dir = rel_path.split(os.sep)[0] if os.sep in rel_path else rel_path
                    if top_dir in getattr(settings, "IGNORE_TOP_LEVEL_DIRS", set()):
                        continue

                    for file in files:
                        file_path = os.path.join(root, file)
                        abs_fp = os.path.abspath(file_path)
                        if self._index_ignore_cache_try_skip(abs_fp, scan_ignore_fp):
                            skipped_by_cache += 1
                            continue
                        if self._should_ignore_file(file_path):
                            self._index_ignore_cache_record(abs_fp, scan_ignore_fp)
                            skipped_by_ignore_rules += 1
                            failed_files.append((file_path, "配置规则忽略(前缀/后缀/黑白名单)"))
                            continue
                        if self._failed_file_cache_try_skip(abs_fp):
                            skipped_by_failed_cache += 1
                            failed_files.append((file_path, "此前索引失败且文件未变化，跳过重试"))
                            continue
                        all_files.append(file_path)

            files_to_process = [f for f in all_files if f not in indexed_paths]

            self._index_info(
                f"扫描到 {len(all_files)} 个文件，其中 {len(all_files) - len(files_to_process)} 个已索引，{len(files_to_process)} 个待索引"
            )

            success_count = 0
            total = len(files_to_process)

            if total == 0:
                self._index_info(f"没有新文件需要索引")
                return 0
        
            import time

            start_time = time.time()
            last_report_time = start_time
            report_interval = 30

            def format_time(seconds):
                if seconds < 60:
                    return f"{seconds:.0f}秒"
                elif seconds < 3600:
                    return f"{seconds/60:.1f}分钟"
                else:
                    return f"{seconds/3600:.1f}小时"

            for i, file_path in enumerate(files_to_process, 1):
                if should_cancel and should_cancel():
                    self._index_info(f"收到取消信号：停止本次目录索引")
                    break
                if should_skip_file:
                    try:
                        if should_skip_file(file_path):
                            self._index_info(f"⏭️ 用户跳过: {file_path}")
                            continue
                    except Exception:
                        pass
                file_name = os.path.basename(file_path)

                if progress_callback:
                    try:
                        progress_callback(i, total, file_name, file_path)
                    except TypeError:
                        progress_callback(i, total, file_name)

                file_start_time = time.time()
                if self.ingest_file(
                    file_path,
                    should_cancel=should_cancel,
                    on_frame_progress=on_frame_progress,
                    on_media_progress=on_media_progress,
                ):
                    success_count += 1
                    file_elapsed = time.time() - file_start_time
                    self._index_info(
                        f"✅ [{i}/{total}] 成功建立索引: {file_path} | 耗时: {file_elapsed:.2f}秒"
                    )
                    if on_file_indexed:
                        try:
                            on_file_indexed(file_path)
                        except Exception:
                            pass
                else:
                    if should_cancel and should_cancel():
                        self._index_info(f"收到取消信号：停止本次目录索引")
                        break
                    file_elapsed = time.time() - file_start_time
                    self._index_error(
                        f"❌ [{i}/{total}] 索引失败: {file_path}（详见上一条异常日志） | 耗时: {file_elapsed:.2f}秒"
                    )
                    self._failed_file_cache_record(
                        os.path.abspath(file_path),
                        "ingest_file 失败(可能文件为空或不支持)",
                    )
                    failed_files.append((file_path, "ingest_file 失败(可能文件为空或不支持)"))

                _indexing_cooperative_yield()

                current_time = time.time()
                if current_time - last_report_time >= report_interval or i == total:
                    elapsed = current_time - start_time
                    avg_time = elapsed / i
                    remaining = total - i
                    eta_seconds = remaining * avg_time

                    self._index_info(
                        f"\n[进度] {i}/{total} ({i*100/total:.1f}%) | "
                        f"已用: {format_time(elapsed)} | "
                        f"速度: {avg_time:.1f}秒/文件 | "
                        f"预计剩余: {format_time(eta_seconds)}\n"
                    )
                    last_report_time = current_time

            total_time = time.time() - start_time

            from datetime import datetime

            log_msg = (
                f"[{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] "
                f"索引完成 - 成功建立索引: {success_count}/{total}, "
                f"总耗时: {format_time(total_time)}, "
                f"平均: {total_time/total:.2f}秒/文件\n"
            )

            data_dir = (os.getenv("FILEAGENT_DATA_DIR") or "").strip()
            if data_dir.startswith("~"):
                data_dir = os.path.expanduser(data_dir)
            base_dir = os.path.abspath(data_dir) if data_dir else os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
            logs_dir = os.path.join(base_dir, "logs")
            try:
                os.makedirs(logs_dir, exist_ok=True)
            except Exception:
                pass
            
            log_path = os.path.join(logs_dir, "index_log.txt")
            with open(log_path, "a", encoding="utf-8") as f:
                f.write(log_msg)
                
            detailed_msg = f"[{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] 文件夹扫描任务完成\n"
            detailed_msg += f"扫描目录: {directories}\n"
            detailed_msg += f"共发现文件: {len(all_files)}\n"
            detailed_msg += f"成功索引数: {success_count}\n"
            detailed_msg += f"失败/忽略数: {len(all_files) - success_count}\n"
            
            if failed_files:
                detailed_msg += "失败/被忽略的文件列表 (部分):\n"
                for fp, reason in failed_files[:100]:
                    detailed_msg += f"  - {fp}: {reason}\n"
                if len(failed_files) > 100:
                    detailed_msg += f"  ... 以及其他 {len(failed_files) - 100} 个文件\n"
            detailed_msg += "-" * 50 + "\n"
            
            self._append_index_details_log(detailed_msg.rstrip("\n"))

            self._index_info(
                f"索引完成，成功建立索引: {success_count}/{total}，总耗时: {format_time(total_time)}"
            )
            self._index_info(
                "扫描统计："
                f"missing_dirs={skipped_missing_dirs},"
                f" skipped_by_cache={skipped_by_cache},"
                f" skipped_by_failed_cache={skipped_by_failed_cache},"
                f" skipped_by_ignore_rules={skipped_by_ignore_rules},"
                f" failed_or_ignored={len(failed_files)}"
            )
            self._index_info(f"日志已写入: {self._index_details_log_path or os.path.abspath(os.path.join(logs_dir, 'index_details.log'))}")
            return success_count
        finally:
            try:
                self._save_index_ignore_cache()
            except Exception:
                pass
            try:
                self._save_failed_file_cache()
            except Exception:
                pass

            try:
                if getattr(settings, 'USE_WHITELIST_MODE', False) and _orig_include_paths is not None:
                    settings.INCLUDE_PATHS = _orig_include_paths
            except Exception:
                pass
